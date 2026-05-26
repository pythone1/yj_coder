# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import build_dma_ai_presentation_v2 as b


OUT = b.OUT / "AI供水管网DMA漏损检测_模型应用教学汇报_高信息密度版.pptx"
AUDIT = b.OUT / "AI供水管网DMA漏损检测_高信息密度版_逐页检查.txt"


def label(slide, text, x, y, w, color=b.BLUE):
    b.tb(slide, text, x, y, w, Inches(0.25), 14, color, True)


def info_panel(slide, title, items, x, y, w, h, color=b.BLUE):
    b.rect(slide, x, y, w, h, b.WHITE, b.GRAY_LINE, False)
    b.tb(slide, title, x + Inches(0.16), y + Inches(0.14), w - Inches(0.3), Inches(0.26), 16, color, True)
    b.bullets(slide, items, x + Inches(0.16), y + Inches(0.52), w - Inches(0.32), h - Inches(0.62), 14.5, b.TEXT)


def bottom_teach(slide, title, cols):
    y = Inches(5.05)
    b.rect(slide, Inches(0.72), y, Inches(11.85), Inches(1.55), b.PALE2, b.GRAY_LINE, False)
    b.tb(slide, title, Inches(0.92), y + Inches(0.15), Inches(2.0), Inches(0.25), 16, b.BLUE, True)
    x0 = Inches(2.7)
    for i, (h, body) in enumerate(cols):
        x = x0 + Inches(i * 3.15)
        b.tb(slide, h, x, y + Inches(0.18), Inches(2.9), Inches(0.25), 14.5, b.BLUE, True)
        b.tb(slide, body, x, y + Inches(0.55), Inches(2.9), Inches(0.75), 13.2, b.MUTED)


def table(slide, headers, rows, x, y, widths, row_h=0.48, head_color=b.BLUE, font=12.5):
    for i, head in enumerate(headers):
        xx = x + sum(Inches(w) for w in widths[:i])
        b.rect(slide, xx, y, Inches(widths[i]), Inches(row_h), head_color, head_color, False)
        b.tb(slide, head, xx + Inches(0.04), y + Inches(0.13), Inches(widths[i] - 0.08), Inches(0.18), font, b.WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y + Inches(row_h + 0.08 + r * (row_h + 0.08))
        for c, val in enumerate(row):
            xx = x + sum(Inches(w) for w in widths[:c])
            b.rect(slide, xx, yy, Inches(widths[c]), Inches(row_h), b.PALE2 if r % 2 == 0 else b.WHITE, b.GRAY_LINE, False)
            b.tb(slide, val, xx + Inches(0.07), yy + Inches(0.12), Inches(widths[c] - 0.14), Inches(0.20), font, b.BLUE if c == 0 else b.TEXT, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)


def dense_header_slide(prs, title, sub, no):
    return b.content(prs, title, sub, no)


def cover(prs):
    s = b.blank(prs, b.NAVY)
    b.fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.25), Inches(7.5)), b.NAVY)
    b.fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(1.05), Inches(0.13), Inches(5.15)), b.CYAN)
    b.pic_fit(s, b.R("00_algorithm_overall_architecture.png"), Inches(6.35), Inches(0.55), Inches(6.55), Inches(6.15), crop=False)
    b.tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.25), Inches(5.05), Inches(1.45), 35, b.WHITE, True)
    b.tb(s, "模型机理 · 算法选型 · 场景落地 · 持续运营", Inches(1.0), Inches(3.15), Inches(5.0), Inches(0.34), 18, b.CYAN, True)
    b.tb(s, "教学汇报版", Inches(1.0), Inches(5.62), Inches(2.1), Inches(0.30), 16, b.CYAN, True)
    b.tb(s, "基于DMA分区计量、时序预测、异常检测与水力模型融合", Inches(1.0), Inches(6.05), Inches(5.05), Inches(0.32), 12.5, b.WHITE)
    bottom_teach(s, "本讲目标", [("理解对象", "DMA如何把管网漏损问题转化为可建模对象"), ("掌握方法", "AI模型如何完成预警、定位、排序和复盘"), ("落地路径", "数据治理、模型验证和运营机制如何支撑规模化")])


def toc(prs):
    s = b.blank(prs)
    b.header(s, "目录", "CONTENTS", 2)
    sections = [
        ("01", "DMA漏损检测业务背景", "对象、数据、传统方法局限"),
        ("02", "AI模型应用发展逻辑", "阈值、基线、多源、组合模型"),
        ("03", "核心算法体系", "LSTM、异常检测、树模型、机理融合"),
        ("04", "典型应用场景", "预警、定位、复盘、压力优化、规划"),
        ("05", "模型建设与落地路径", "治理、训练、验证、运营、推广"),
    ]
    for i, (n, title, desc) in enumerate(sections):
        y = Inches(1.18 + i * 0.97)
        b.tb(s, n, Inches(0.95), y, Inches(0.55), Inches(0.28), 17, b.CYAN, True)
        b.tb(s, title, Inches(1.75), y - Inches(0.04), Inches(4.0), Inches(0.32), 21, b.TEXT, True)
        b.tb(s, desc, Inches(6.05), y, Inches(5.3), Inches(0.25), 15.5, b.MUTED)
        line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), y + Inches(0.50), Inches(10.85), Inches(0.01))
        b.fill(line, b.PALE)
    b.pic_fit(s, b.A("04_small_infographic_plugins", "02_water_ai_droplet.png"), Inches(11.15), Inches(5.6), Inches(0.85), Inches(0.85))
    b.footer(s)


def build():
    prs = Presentation(str(b.TEMPLATE))
    b.clear_slides(prs)
    cover(prs)
    toc(prs)

    b.section(prs, "PART 01", "DMA漏损检测业务背景", "先讲清检测对象，再进入模型方法。", b.A("01_water_equipment_plugins", "04_dma_boundary_map.png"), 3)

    s = dense_header_slide(prs, "DMA系统在漏损管控中的位置", "PART 01 业务背景", 4)
    b.pic_fit(s, b.A("01_water_equipment_plugins", "04_dma_boundary_map.png"), Inches(0.82), Inches(1.25), Inches(2.25), Inches(2.25))
    table(s, ["对象", "工程含义", "对应模型任务"], [
        ("入口计量", "记录DMA总进水量和夜间最小流量", "建立流量基线"),
        ("压力监测", "反映供水服务压力和异常响应", "识别压力联动"),
        ("边界阀门", "决定DMA是否封闭、是否串水", "判断边界异常"),
        ("管网资产", "管龄、材质、口径影响漏损概率", "管段风险排序"),
    ], Inches(3.35), Inches(1.22), [1.45, 3.25, 3.05], 0.48, font=12.5)
    bottom_teach(s, "讲课展开", [("一句话定义", "DMA是AI建模的基本业务单元。"), ("讲解重点", "没有清晰边界和可靠计量，模型无法区分漏损与正常用水。"), ("过渡", "因此模型建设首先要看数据基础。")])

    s = dense_header_slide(prs, "传统漏损检测方法与局限", "PART 01 业务背景", 5)
    table(s, ["方法", "适用场景", "主要问题"], [
        ("MNF判断", "夜间低用水窗口", "需扣除合法夜间用水和计量误差"),
        ("固定阈值", "监测启动阶段", "难以适配不同DMA差异"),
        ("人工巡检", "现场确认", "效率受经验和排查范围影响"),
        ("投诉/报修", "明显爆管或用户感知事件", "滞后，难以发现隐性漏损"),
    ], Inches(0.8), Inches(1.22), [1.55, 3.55, 5.55], 0.52, font=13)
    bottom_teach(s, "引出AI", [("业务痛点", "误报、滞后、定位范围大。"), ("模型机会", "学习分区正常状态，识别持续偏离。"), ("讲法建议", "不要贬低传统方法，强调其是AI模型的基础。")])

    s = dense_header_slide(prs, "DMA数据基础：模型能否有效取决于数据是否可用", "PART 01 业务背景", 6)
    table(s, ["数据类型", "关键字段", "用途"], [
        ("时序数据", "流量、压力、采样时间、缺失标记", "预测基线、残差、异常时段"),
        ("空间数据", "管段、阀门、拓扑、DMA归属", "候选区域和管段定位"),
        ("资产数据", "管龄、材质、口径、维修历史", "风险排序和改造计划"),
        ("工单数据", "报警、核查、维修、误报原因", "标签回填和模型迭代"),
        ("外部数据", "天气、节假日、施工、用户结构", "解释正常扰动"),
    ], Inches(0.72), Inches(1.16), [1.55, 4.35, 4.25], 0.44, font=12.5)
    bottom_teach(s, "落地提醒", [("先治理", "统一编码、时序对齐、缺失处理。"), ("再建模", "先做可解释基线，再做复杂模型。"), ("常见问题", "工单没有坐标或结果，会直接影响监督模型。")])

    b.section(prs, "PART 02", "AI模型应用发展逻辑", "从经验阈值走向动态基线、多源证据和模型组合。", b.A("04_small_infographic_plugins", "10_timeseries_mini_chart.png"), 7)

    s = dense_header_slide(prs, "从阈值报警到动态基线", "PART 02 发展逻辑", 8)
    b.pic_fit(s, b.A("02_algorithm_plugins", "13_prediction_residual_chart.png"), Inches(0.82), Inches(1.18), Inches(2.3), Inches(2.3))
    table(s, ["阶段", "判断依据", "升级价值"], [
        ("固定阈值", "流量或压力超过人工设定上限", "部署简单，适合启动阶段"),
        ("统计基线", "历史均值、移动平均、同比环比", "能识别周期性偏离"),
        ("动态基线", "按DMA学习预测区间和残差", "适应分区差异，降低误报"),
        ("风险预警", "残差持续性+压力响应+历史事件", "输出可复核证据链"),
    ], Inches(3.45), Inches(1.18), [1.45, 4.0, 3.0], 0.50, font=12.5)
    bottom_teach(s, "讲课展开", [("核心概念", "动态基线就是每个DMA自己的正常曲线。"), ("判定逻辑", "持续偏离比单点超限更重要。"), ("输出方式", "风险等级、异常证据、复核建议。")])

    s = dense_header_slide(prs, "从单点报警到多源联动识别", "PART 02 发展逻辑", 9)
    table(s, ["证据源", "异常表现", "解释价值"], [
        ("入口流量", "夜间底流抬升、日内曲线偏移", "判断是否存在持续异常"),
        ("压力点", "末梢压力下降、波动加剧", "辅助判断位置和影响范围"),
        ("GIS资产", "老旧管段、材质脆弱、阀门异常", "解释风险来源"),
        ("历史工单", "同一区域多次维修或投诉", "提高候选区域优先级"),
        ("外部上下文", "节假日、施工、商业用水", "排除正常扰动"),
    ], Inches(0.72), Inches(1.16), [1.55, 4.1, 4.35], 0.44, font=12.3)
    bottom_teach(s, "讲课展开", [("关键转变", "从超限报警转为证据链判断。"), ("模型输出", "异常时段、主证据、疑似原因、候选管段。"), ("管理价值", "让调度和巡检人员知道下一步做什么。")])

    s = dense_header_slide(prs, "从单一算法到模型组合", "PART 02 发展逻辑", 10)
    table(s, ["业务环节", "推荐模型", "业务输出"], [
        ("正常状态预测", "LSTM、GRU、CNN-LSTM", "预测区间、残差、趋势偏离"),
        ("异常候选发现", "孤立森林、DBSCAN、自编码器", "异常分数、异常时段"),
        ("候选区域收敛", "水力仿真+机器学习、GNN", "候选片区、管段TopN"),
        ("风险排序", "随机森林、GBDT、HGB", "风险分、检漏优先级"),
        ("解释协同", "知识图谱、大模型智能体", "报告、问答、派单建议"),
    ], Inches(0.72), Inches(1.16), [1.9, 4.0, 4.05], 0.44, font=12.2)
    bottom_teach(s, "讲课展开", [("讲法", "漏损检测不是一个算法问题，是一条业务链。"), ("重点", "不同模型解决不同环节。"), ("风险", "单一算法很难同时完成预警、定位和解释。")])

    s = dense_header_slide(prs, "供水漏损AI模型技术演进路线图", "PART 02 发展逻辑", 11)
    table(s, ["阶段", "方法", "适用边界"], [
        ("经验阈值", "固定上下限、MNF经验值", "基础监测，可快速启动"),
        ("统计分析", "移动平均、控制图、季节分解", "识别周期偏离，解释性强"),
        ("机器学习", "孤立森林、DBSCAN、GBDT", "少标签异常和结构化风险排序"),
        ("深度学习", "LSTM、GRU、自编码器、GNN", "高频时序、多变量耦合、拓扑学习"),
        ("智能协同", "知识图谱、智能体、工具调用", "解释、报告、问答和派单辅助"),
    ], Inches(0.72), Inches(1.16), [1.55, 4.35, 4.35], 0.44, font=12.2)
    bottom_teach(s, "汇报口径", [("主线", "规则是基础，AI是增强。"), ("判断", "模型成熟度取决于数据基础和闭环程度。"), ("落脚", "最终目标是漏损治理效率提升。")])

    b.section(prs, "PART 03", "核心算法体系", "算法讲解围绕业务任务展开。", b.A("02_algorithm_plugins", "14_model_training_pipeline.png"), 12)

    algorithm_pages = [
        ("LSTM / GRU：时序预测与动态基线", b.A("02_algorithm_plugins", "01_lstm_sequence_prediction.png"), [
            ("输入", "过去窗口流量、压力、天气、日期、节假日"),
            ("处理", "学习长短期时间依赖，输出未来正常曲线"),
            ("输出", "预测值、预测区间、残差、趋势偏离"),
            ("适用", "日周期、周周期明显，监测数据连续的DMA"),
        ], [("讲解重点", "它先学习正常运行边界，再通过残差判断是否需要复核。"), ("适配场景", "入口流量连续、压力点稳定、日周周期明显的DMA。"), ("注意事项", "训练验证要按时间切分，避免未来信息泄漏。")]),
        ("CNN-LSTM：多传感器特征融合", b.A("02_algorithm_plugins", "03_cnn_lstm_hybrid.png"), [
            ("输入", "多个压力点、入口流量、局部时序片段"),
            ("处理", "CNN提取局部特征，LSTM学习时间依赖"),
            ("输出", "综合残差、异常趋势、压力响应模式"),
            ("适用", "多点监测、压力联动明显的分区"),
        ], [("讲解重点", "先看多个传感器的局部变化，再看变化是否持续。"), ("业务价值", "适合解释多个压力点与入口流量同步变化的异常。"), ("风险", "压力点布局不足或拓扑错误时，模型解释会变弱。")]),
        ("孤立森林 / DBSCAN：少标签异常发现", b.A("02_algorithm_plugins", "05_isolation_forest_outlier.png"), [
            ("孤立森林", "快速筛查多维异常，输出异常分数"),
            ("DBSCAN", "识别低密度离群点和异常簇"),
            ("输入", "MNF、残差、压力波动、多维统计特征"),
            ("定位", "作为人工复核清单，不直接确认漏损"),
        ], [("讲解重点", "少标签阶段先发现可疑对象，再由人工复核形成标签。"), ("输出方式", "异常时段、异常分数、主要特征和复核建议。"), ("注意", "输出不能直接等同漏损结论，必须进入工单闭环。")]),
        ("自编码器：重构误差识别异常", b.A("02_algorithm_plugins", "04_autoencoder_reconstruction.png"), [
            ("训练对象", "正常状态下的多变量组合关系"),
            ("判断依据", "重构误差升高代表偏离正常模式"),
            ("输入变量", "流量、压力、MNF、残差、上下文特征"),
            ("应用", "复杂多变量异常检测和早期预警"),
        ], [("讲解重点", "模型记住正常状态的变量关系，异常状态会产生较高重构误差。"), ("优势", "适合漏损标签不足但时序监测较连续的早期阶段。"), ("注意", "训练集要尽量排除已知异常，否则异常会被学成正常。")]),
        ("随机森林 / GBDT / HGB：管段风险排序", b.A("02_algorithm_plugins", "07_random_forest_ensemble.png"), [
            ("输入", "管龄、材质、口径、压力、维修、投诉、道路"),
            ("处理", "学习历史漏损与资产特征之间的关系"),
            ("输出", "管段风险分、风险等级、检漏优先级"),
            ("适用", "资产台账和历史工单较完整的区域"),
        ], [("讲解重点", "它回答哪里更容易漏，不直接回答当前是否正在漏。"), ("管理价值", "支撑年度巡检计划、检漏优先级和管网改造排序。"), ("前提", "资产台账完整，工单结构化，漏损标签口径一致。")]),
        ("水力模型 + AI：机理与数据融合", b.A("03_system_architecture_modules", "07_hydraulic_simulation_module.png"), [
            ("机理侧", "拓扑、供水路径、压力敏感性、仿真样本"),
            ("数据侧", "真实流量压力、异常模式、历史工单"),
            ("融合方式", "仿真样本辅助训练，AI结果受水力约束"),
            ("输出", "候选片区、候选管段TopN、定位置信度"),
        ], [("讲解重点", "机理模型提供物理边界，AI模型提供数据识别能力。"), ("价值", "将异常识别结果限制在水力学合理范围内，提高定位可信度。"), ("前提", "水力模型需要校准，阀门状态和拓扑关系要可靠。")]),
        ("GNN与知识图谱：拓扑关系与解释能力", b.A("02_algorithm_plugins", "10_gnn_pipe_topology.png"), [
            ("GNN", "把管网拓扑和压力传播关系纳入学习"),
            ("知识图谱", "连接DMA、管段、设备、工单、案例"),
            ("智能体", "基于模型结果生成报告和问答解释"),
            ("适用", "拓扑质量较高、业务系统数据可关联的场景"),
        ], [("讲解重点", "图模型解决管网连接关系，知识图谱解决业务解释关系。"), ("业务价值", "辅助异常原因分析、历史案例检索和经验沉淀。"), ("控制", "专业模型判断、权限控制和人工确认必须保留。")]),
    ]
    no = 13
    for title, img, rows, teach in algorithm_pages:
        s = dense_header_slide(prs, title, "PART 03 核心算法", no)
        b.pic_fit(s, img, Inches(0.82), Inches(1.18), Inches(2.25), Inches(2.25))
        table(s, ["维度", "内容"], rows, Inches(3.45), Inches(1.18), [1.35, 6.85], 0.48, font=13)
        bottom_teach(s, "讲课展开", teach)
        no += 1

    s = dense_header_slide(prs, "算法选型矩阵", "PART 03 核心算法", 20)
    table(s, ["业务任务", "推荐模型", "输入数据", "输出结果"], [
        ("动态基线", "LSTM/GRU/CNN-LSTM", "流量、压力、天气、日历", "预测区间、残差"),
        ("异常发现", "孤立森林/DBSCAN/自编码器", "MNF、残差、多维时序", "异常分数、异常时段"),
        ("风险排序", "随机森林/GBDT/HGB", "管龄、材质、维修、投诉", "风险分、优先级"),
        ("定位收敛", "水力仿真+ML/GNN", "拓扑、多点压力、仿真样本", "候选管段TopN"),
        ("业务协同", "知识图谱/智能体", "模型输出、GIS、工单", "报告、问答、派单建议"),
    ], Inches(0.62), Inches(1.12), [1.55, 3.05, 3.4, 2.65], 0.43, font=12.0)
    bottom_teach(s, "讲课展开", [("选型原则", "先看业务任务，再看数据条件。"), ("落地原则", "先可解释模型，再复杂模型。"), ("风险控制", "所有模型输出都要进入复核闭环。")])

    b.section(prs, "PART 04", "典型应用场景", "把模型输出转化为调度、巡检和管理动作。", b.A("01_water_equipment_plugins", "05_pipe_leak_alert.png"), 21)

    scenario_pages = [
        ("事前预警：DMA异常识别", b.A("04_small_infographic_plugins", "01_leak_warning_badge.png"), [
            ("模型输入", "入口流量、压力点、MNF、天气、节假日"),
            ("识别逻辑", "实际曲线持续偏离动态基线，且压力响应同步异常"),
            ("输出内容", "异常时段、风险等级、主要证据、建议复核动作"),
            ("业务价值", "提前发现隐性漏损，减少被动报修"),
        ]),
        ("事中定位：候选管段TopN收敛", b.A("03_system_architecture_modules", "08_leak_localization_module.png"), [
            ("模型输入", "DMA异常、压力响应、GIS拓扑、水力仿真"),
            ("识别逻辑", "从全分区收敛到候选片区和候选管段"),
            ("输出内容", "管段编号、位置、置信度、证据、影响用户数"),
            ("业务价值", "缩短巡检路线，提高现场核查效率"),
        ]),
        ("事后复盘：工单回填与持续学习", b.A("03_system_architecture_modules", "12_model_retraining_loop.png"), [
            ("记录字段", "报警时间、异常特征、现场结果、漏点位置、误报原因"),
            ("复盘动作", "修正标签、更新特征、调整阈值、再训练"),
            ("输出内容", "误报类型、真实漏损样本、模型版本表现"),
            ("业务价值", "把每次处置沉淀为训练数据"),
        ]),
        ("压力优化：降漏、稳压、节能平衡", b.A("01_water_equipment_plugins", "14_pump_station.png"), [
            ("模型输入", "需水预测、压力监测、泵阀状态、水力模型"),
            ("约束条件", "最不利点压力、消防保障、二次供水补水、用户体验"),
            ("输出内容", "分时压力策略、泵站启停、阀门调节建议"),
            ("业务价值", "在保障供水安全前提下降漏节能"),
        ]),
        ("DMA规划：分区与监测点优化", b.A("04_small_infographic_plugins", "05_dma_boundary_icon.png"), [
            ("评价维度", "拓扑、地形、入口数量、边界阀数量、施工成本"),
            ("监测点优化", "仿真不同漏点情景，比较压力点组合贡献"),
            ("输出内容", "分区方案评分、压力点优先级、建设建议"),
            ("业务价值", "提高有限预算下的监测覆盖效率"),
        ]),
    ]
    no = 22
    for title, img, rows in scenario_pages:
        s = dense_header_slide(prs, title, "PART 04 应用场景", no)
        b.pic_fit(s, img, Inches(0.82), Inches(1.18), Inches(2.25), Inches(2.25))
        table(s, ["环节", "讲解内容"], rows, Inches(3.45), Inches(1.18), [1.45, 6.75], 0.48, font=13)
        bottom_teach(s, "现场讲法", [("先讲输入", "说明模型依赖哪些数据，以及这些数据来自SCADA、GIS还是工单系统。"), ("再讲逻辑", "解释模型如何从异常曲线、压力响应和历史事件中形成判断。"), ("最后讲输出", "强调输出必须进入调度、巡检、派单或复盘动作。")])
        no += 1

    b.section(prs, "PART 05", "模型建设与落地路径", "把模型做成可运行、可验证、可迭代的工程系统。", b.A("03_system_architecture_modules", "03_data_cleaning_filter.png"), 27)

    s = dense_header_slide(prs, "数据治理：五张基础表", "PART 05 落地路径", 28)
    table(s, ["数据表", "关键字段", "建模用途"], [
        ("设备表", "流量计、压力计、RTU、边缘网关", "识别采集来源和设备状态"),
        ("管网表", "管段、阀门、DMA归属、资产属性", "支撑空间定位和风险排序"),
        ("时序表", "采样时间、流量压力、缺失标记", "训练预测和异常检测模型"),
        ("工单表", "报警、核查、维修、处置结果", "形成业务闭环和标签"),
        ("标签表", "真实漏损、误报原因、模型版本", "监督学习和模型评估"),
    ], Inches(0.72), Inches(1.16), [1.4, 4.35, 4.25], 0.44, font=12.3)
    bottom_teach(s, "检查清单", [("统一编码", "设备、管段、DMA编码必须一致。"), ("时序对齐", "流量压力采样时间要可比。"), ("工单结构化", "现场结果要能回填为标签。")])

    s = dense_header_slide(prs, "模型训练路径：基线 - 异常 - 定位 - 排序", "PART 05 落地路径", 29)
    table(s, ["阶段", "核心任务", "模型方法", "检查点"], [
        ("1 基线", "建立MNF和移动基线", "统计基线", "能解释正常波动"),
        ("2 预测", "学习动态运行曲线", "LSTM/GRU", "时间切分验证"),
        ("3 异常", "识别候选异常时段", "孤立森林/自编码器", "人工复核闭环"),
        ("4 定位", "收敛候选片区管段", "水力仿真+ML", "候选TopN命中率"),
        ("5 排序", "生成检漏优先级", "GBDT/HGB", "工单标签质量"),
    ], Inches(0.62), Inches(1.12), [1.25, 3.0, 3.0, 2.8], 0.43, font=12.0)
    bottom_teach(s, "讲课展开", [("顺序", "先基础能力，再复杂模型。"), ("验证", "避免随机切分导致未来信息泄漏。"), ("输出", "每一步都要对应业务动作。")])

    s = dense_header_slide(prs, "模型验证：算法指标与业务指标并重", "PART 05 落地路径", 30)
    table(s, ["指标类型", "代表指标", "说明"], [
        ("预测误差", "MAE、RMSE", "衡量动态基线预测稳定性"),
        ("分类识别", "精确率、召回率、F1、AUC", "衡量异常识别效果"),
        ("定位效果", "TopN命中率、平均排查范围", "衡量现场定位价值"),
        ("业务闭环", "提前预警时间、误报工单比例、复盘回填率", "衡量系统落地价值"),
        ("维修效果", "维修后MNF回落、压力恢复", "衡量治理结果"),
    ], Inches(0.72), Inches(1.16), [1.45, 3.55, 5.0], 0.44, font=12.4)
    bottom_teach(s, "汇报建议", [("少讲单一准确率", "准确率容易掩盖漏报和误报。"), ("多讲业务指标", "预警提前量、排查范围、MNF回落更有说服力。"), ("保留口径", "指标必须与工单闭环一致。")])

    s = dense_header_slide(prs, "模型运营：漂移监控与再训练", "PART 05 落地路径", 31)
    table(s, ["触发因素", "影响", "运营动作"], [
        ("季节变化", "用水模式和夜间底流发生改变", "季度检查基线偏移，必要时更新分区模型"),
        ("管网改造", "拓扑、阀门状态和压力制度变化", "重新校准水力模型，并复核DMA边界"),
        ("传感器更换", "采样误差、量程和数据分布改变", "复核采集质量，建立新旧设备对照"),
        ("用户结构变化", "商业、工业或居民用水占比变化", "更新动态基线和上下文特征"),
        ("新增工单", "真实漏损和误报标签样本增加", "半年再训练，并比较模型版本效果"),
    ], Inches(0.72), Inches(1.16), [1.6, 3.55, 4.85], 0.44, font=12.4)
    bottom_teach(s, "运营节奏", [("每月", "复盘误报警报。"), ("每季度", "检查模型漂移。"), ("每半年", "结合新增工单再训练。")])

    s = dense_header_slide(prs, "试点到规模化推广路线", "PART 05 落地路径", 32)
    table(s, ["阶段", "工作重点", "交付物"], [
        ("试点选区", "选择边界清晰、数据稳定、治理需求明确的DMA", "试点DMA清单"),
        ("数据治理", "编码统一、时序对齐、工单结构化", "建模数据集"),
        ("模型试运行", "预警复核、误报分析、阈值校准", "试运行报告"),
        ("效果评估", "业务指标与算法指标同步评估", "评估表和优化建议"),
        ("规模推广", "沉淀接口、流程、版本和培训体系", "推广标准包"),
    ], Inches(0.72), Inches(1.16), [1.45, 5.35, 3.25], 0.44, font=12.3)
    bottom_teach(s, "落地判断", [("能不能复制", "看数据接口和工单闭环。"), ("能不能持续", "看模型运营机制。"), ("能不能讲清", "看业务指标是否改善。")])

    s = dense_header_slide(prs, "总结：以业务闭环定义AI模型价值", "SUMMARY", 33)
    table(s, ["核心结论", "讲解口径"], [
        ("先数据，再模型", "没有清晰边界、可靠计量和结构化工单，复杂算法价值有限。"),
        ("先预警，再定位", "动态基线解决发现问题，机理融合和空间数据解决排查范围。"),
        ("先复核，再自动化", "模型输出必须进入工单闭环，持续回填才能形成长期能力。"),
        ("先业务，再算法", "算法指标必须转化为提前预警、减少误报、缩小排查范围和MNF回落。"),
    ], Inches(0.72), Inches(1.18), [2.0, 7.95], 0.50, font=12.8)
    bottom_teach(s, "收束语", [("一句话", "AI不是单个模型，而是漏损治理闭环的一组能力。"), ("价值", "能发现、能解释、能派单、能复盘、能持续优化。"), ("落地", "从试点DMA开始，逐步标准化推广。")])

    s = dense_header_slide(prs, "附：素材与讲解页对应关系", "APPENDIX", 34)
    table(s, ["素材类型", "使用页面", "使用原则"], [
        ("大图", "封面、章节页、总体架构", "只做主视觉，不硬裁切、不塞满正文页"),
        ("算法插件", "LSTM、异常检测、树模型、GNN页", "配合输入-处理-输出讲解"),
        ("系统模块", "机理融合、定位、治理、运营页", "说明业务流程和平台模块"),
        ("小图标", "目录、总结、提示区", "用于辅助识别，不替代内容"),
    ], Inches(0.72), Inches(1.18), [1.7, 3.25, 5.1], 0.50, font=12.7)
    bottom_teach(s, "使用规范", [("正文页", "必须有表格或流程支撑。"), ("字体", "正文尽量18pt以上，表格不低于12pt。"), ("检查", "每页必须有足够讲解信息。")])

    prs.save(OUT)
    audit(prs)
    return OUT


def audit(prs):
    lines = []
    fail = False
    for idx, slide in enumerate(prs.slides, 1):
        chars = 0
        pics = 0
        min_font = 99
        max_bottom = 0
        off = 0
        for sh in slide.shapes:
            max_bottom = max(max_bottom, sh.top + sh.height)
            if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs.slide_width + 50000 or sh.top + sh.height > prs.slide_height + 50000:
                off += 1
            if sh.shape_type == 13:
                pics += 1
            if hasattr(sh, "text_frame") and sh.text.strip():
                chars += len(sh.text.strip())
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            min_font = min(min_font, r.font.size.pt)
        role = "section" if idx in [3, 7, 12, 21, 27] else "content"
        flags = []
        if off:
            flags.append(f"越界{off}")
        if min_font < 10:
            flags.append(f"小字{min_font}")
        if role == "content" and idx not in [1, 2] and chars < 210:
            flags.append(f"文字量不足{chars}")
        if role == "content" and max_bottom < Inches(6.1):
            flags.append("下半页利用不足")
        shape_count = len(slide.shapes)
        table_like = shape_count >= 32
        if role == "content" and idx not in [2] and pics == 0 and not table_like:
            flags.append("缺少图示或表格支撑")
        status = "FAIL" if flags else "OK"
        if flags:
            fail = True
        lines.append(f"{idx:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(max_bottom/914400,2)}\t{'；'.join(flags)}")
    AUDIT.write_text("\n".join(lines), encoding="utf-8")
    if fail:
        print("\n".join(lines))
        raise SystemExit("audit failed")


if __name__ == "__main__":
    print(build())
