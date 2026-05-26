# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = max([p for p in ROOT.glob("*.pptx") if not p.name.startswith("~$")], key=lambda p: p.stat().st_size)
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_完整混合可编辑版.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_完整混合可编辑版_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_完整混合可编辑版_预览联系表.png"


NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
FONT = "微软雅黑"


ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "dma": ASSET_DIR / "dma_clean.png",
    "data": ASSET_DIR / "data_clean.png",
    "ai": ASSET_DIR / "ai_clean.png",
    "locate": ASSET_DIR / "locate_clean.png",
    "loop": ASSET_DIR / "loop_clean.png",
    "algorithm": ASSET_DIR / "algorithm_clean.png",
}


def ensure_assets():
    missing = [str(v) for v in ASSETS.values() if not v.exists()]
    if missing:
        raise FileNotFoundError("missing assets:\n" + "\n".join(missing))


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, trans=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = trans
    shape.line.fill.background()


def stroke(shape, color=LINE, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def tb(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(typ, x, y, w, h)
    fill(shp, color, trans)
    stroke(shp, line, 1.0)
    return shp


def connector(slide, x1, y1, x2, y2, color=BLUE, width=1.8):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def pic(slide, path, x, y, w, h, crop=False):
    p = slide.shapes.add_picture(str(path), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.25), Inches(0.12), Inches(0.48))
    fill(bar, BLUE)
    tb(slide, title, Inches(0.62), Inches(0.14), Inches(9.1), Inches(0.40), 25, TEXT, True)
    tb(slide, sub, Inches(0.64), Inches(0.62), Inches(8.0), Inches(0.22), 11.5, MUTED)
    tb(slide, f"{no:02d}", Inches(12.05), Inches(0.22), Inches(0.7), Inches(0.25), 12, BLUE, True, PP_ALIGN.RIGHT)
    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.90), Inches(11.7), Inches(0.016))
    fill(divider, PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA漏损检测 · 完整混合可编辑版", Inches(0.62), Inches(7.08), Inches(6.2), Inches(0.2), 10, MUTED)


def content(prs, title, sub, no):
    slide = blank(prs)
    header(slide, title, sub, no)
    footer(slide)
    return slide


def pill(slide, text, x, y, w, color=PALE, text_color=BLUE, size=12):
    rect(slide, x, y, w, Inches(0.32), color, color, True)
    tb(slide, text, x + Inches(0.06), y + Inches(0.075), w - Inches(0.12), Inches(0.14), size, text_color, True, PP_ALIGN.CENTER)


def note_band(slide, text, y=6.00, color=BLUE):
    rect(slide, Inches(0.72), Inches(y), Inches(11.85), Inches(0.55), RGBColor(247, 252, 255), LINE, True)
    tb(slide, text, Inches(0.98), Inches(y + 0.15), Inches(11.3), Inches(0.22), 15.0, color, True, PP_ALIGN.CENTER)


def table(slide, headers, rows, x, y, widths, row_h=0.45, font=11.6):
    for i, head in enumerate(headers):
        xx = x + sum(Inches(v) for v in widths[:i])
        rect(slide, xx, y, Inches(widths[i]), Inches(row_h), BLUE, BLUE, False)
        tb(slide, head, xx + Inches(0.04), y + Inches(0.115), Inches(widths[i] - 0.08), Inches(0.16), font + 1, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y + Inches(row_h + 0.07 + r * (row_h + 0.07))
        for c, val in enumerate(row):
            xx = x + sum(Inches(v) for v in widths[:c])
            rect(slide, xx, yy, Inches(widths[c]), Inches(row_h), PALE2 if r % 2 == 0 else WHITE, LINE, False)
            tb(slide, val, xx + Inches(0.06), yy + Inches(0.11), Inches(widths[c] - 0.12), Inches(0.18),
               font, BLUE if c == 0 else TEXT, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)


def bullets(slide, title, items, x, y, w, h, color=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE, True)
    tb(slide, title, x + Inches(0.18), y + Inches(0.16), w - Inches(0.36), Inches(0.24), 16, color, True)
    top = y + Inches(0.56)
    for i, item in enumerate(items):
        tb(slide, "· " + item, x + Inches(0.20), top + Inches(i * 0.35), w - Inches(0.40), Inches(0.22), 12.8, TEXT)


def visual_text_slide(prs, no, title, sub, image_key, rows, note, img_x=0.75, img_w=4.70):
    slide = content(prs, title, sub, no)
    rect(slide, Inches(img_x), Inches(1.18), Inches(img_w), Inches(4.62), WHITE, RGBColor(205, 229, 243), True)
    pic(slide, ASSETS[image_key], Inches(img_x + 0.18), Inches(1.38), Inches(img_w - 0.36), Inches(4.20), crop=False)
    table(slide, ["环节", "讲解内容"], rows, Inches(5.85), Inches(1.22), [1.55, 4.95], 0.48, 12.0)
    note_band(slide, note, 6.05)
    audit_slide(slide, no)
    return slide


def module(slide, img_path, title, desc, x, y, w, h, color):
    rect(slide, x, y, w, h, WHITE, RGBColor(201, 228, 242), True)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.43))
    fill(top, color)
    tb(slide, title, x, y + Inches(0.095), w, Inches(0.20), 14.2, WHITE, True, PP_ALIGN.CENTER)
    pic(slide, img_path, x + Inches(0.14), y + Inches(0.58), w - Inches(0.28), h - Inches(1.28), crop=False)
    tb(slide, desc, x + Inches(0.14), y + h - Inches(0.43), w - Inches(0.28), Inches(0.32), 11.8, MUTED, False, PP_ALIGN.CENTER)


def section_slide(prs, no, part, title, subtitle, image_key):
    slide = blank(prs, NAVY)
    fill(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5)), CYAN)
    tb(slide, part, Inches(0.85), Inches(1.22), Inches(1.4), Inches(0.3), 16, CYAN, True)
    tb(slide, title, Inches(0.85), Inches(1.95), Inches(6.55), Inches(0.65), 36, WHITE, True)
    tb(slide, subtitle, Inches(0.90), Inches(3.0), Inches(6.55), Inches(0.72), 17, RGBColor(210, 232, 245))
    rect(slide, Inches(8.00), Inches(1.20), Inches(4.35), Inches(4.35), WHITE, RGBColor(205, 229, 243), True)
    pic(slide, ASSETS[image_key], Inches(8.25), Inches(1.48), Inches(3.85), Inches(3.75), crop=False)
    tb(slide, f"{no:02d}", Inches(11.75), Inches(6.58), Inches(0.65), Inches(0.30), 14, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)
    audit_slide(slide, no, section=True)
    return slide


def audit_slide(slide, no, section=False):
    prs_w, prs_h = 12192000, 6858000
    chars = 0
    pics = 0
    min_font = 99
    off = []
    bottom = 0
    for sh in slide.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs_w + 50000 or sh.top + sh.height > prs_h + 50000:
            off.append(str(sh.shape_type))
        if sh.shape_type == 13:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off:
        flags.append(f"越界{len(off)}")
    if min_font < 10:
        flags.append(f"小字{min_font}")
    if pics < 1:
        flags.append("缺少主视觉")
    if not section and chars < 150:
        flags.append(f"文字不足{chars}")
    if not section and bottom < Inches(6.0):
        flags.append("下半页利用不足")
    status = "FAIL" if flags else "OK"
    AUDIT_LINES.append(f"{no:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise SystemExit(f"slide {no} audit failed: {'; '.join(flags)}")


AUDIT_LINES = []


def slide_cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.30), Inches(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(1.05), Inches(0.13), Inches(5.10)), CYAN)
    tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.28), Inches(5.1), Inches(1.55), 35, WHITE, True)
    tb(s, "教学汇报版 · 混合可编辑PPT", Inches(1.0), Inches(3.20), Inches(4.9), Inches(0.32), 18, CYAN, True)
    tb(s, "从DMA分区异常识别，讲到AI溯源定位、候选管段收敛与工单闭环处置。", Inches(1.0), Inches(5.58), Inches(5.05), Inches(0.30), 12.5, RGBColor(210, 232, 245))
    tb(s, "面向教学演讲，重点说明模型如何从“发现异常”走向“定位管段”和“闭环复盘”。", Inches(1.0), Inches(5.98), Inches(5.05), Inches(0.34), 12.0, RGBColor(210, 232, 245))
    tb(s, "适用于智慧水务、管网运维、漏损治理与AI模型应用培训。", Inches(1.0), Inches(6.34), Inches(5.05), Inches(0.28), 11.5, RGBColor(210, 232, 245))
    rect(s, Inches(6.42), Inches(0.74), Inches(6.45), Inches(5.32), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS["cover"], Inches(6.58), Inches(0.92), Inches(6.12), Inches(4.95), crop=False)
    pill(s, "DMA锁定区域", Inches(6.95), Inches(6.18), Inches(1.70), PALE, BLUE, 12)
    pill(s, "AI溯源定位", Inches(8.90), Inches(6.18), Inches(1.70), PALE, GREEN, 12)
    pill(s, "工单闭环", Inches(10.85), Inches(6.18), Inches(1.35), PALE, ORANGE, 12)
    audit_slide(s, 1)


def slide_toc(prs):
    s = content(prs, "目录", "CONTENTS", 2)
    rect(s, Inches(0.78), Inches(1.20), Inches(4.00), Inches(4.65), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS["cover"], Inches(0.98), Inches(1.45), Inches(3.58), Inches(4.10), crop=False)
    sections = [
        ("01", "DMA漏损检测业务背景", "讲清对象、数据和传统方法局限"),
        ("02", "AI模型应用发展逻辑", "从阈值报警到动态基线和模型组合"),
        ("03", "核心算法体系", "时序预测、异常检测、风险排序、机理融合"),
        ("04", "典型应用场景", "预警、定位、复盘、调度与规划"),
        ("05", "模型建设与落地路径", "治理、训练、验证、运营和推广"),
    ]
    for i, (n, title, desc) in enumerate(sections):
        y = Inches(1.25 + i * 0.78)
        tb(s, n, Inches(5.25), y, Inches(0.50), Inches(0.20), 14, CYAN, True)
        tb(s, title, Inches(5.95), y - Inches(0.04), Inches(3.0), Inches(0.26), 18, TEXT, True)
        tb(s, desc, Inches(9.05), y, Inches(3.10), Inches(0.22), 12.3, MUTED)
    note_band(s, "本套PPT按“一个主视觉 + 可编辑结构 + 讲课要点”组织，适合现场教学演讲。", 6.05)
    audit_slide(s, 2)


def slide_architecture(prs, no):
    s = content(prs, "DMA宏观锁定 + AI微观溯源总体架构", "总体架构", no)
    y = Inches(1.18)
    w = Inches(2.62)
    module(s, ASSETS["dma"], "1 DMA锁定区域", "先判断哪个片区异常", Inches(0.62), y, w, Inches(4.18), BLUE)
    module(s, ASSETS["data"], "2 汇集关键数据", "流量、压力、管网、工单", Inches(3.55), y, w, Inches(4.18), CYAN)
    module(s, ASSETS["ai"], "3 AI分析原因", "识别压力和流量变化", Inches(6.48), y, w, Inches(4.18), GREEN)
    module(s, ASSETS["locate"], "4 定位并派单", "给出管段和坐标", Inches(9.41), y, w, Inches(4.18), ORANGE)
    for x in [Inches(3.28), Inches(6.21), Inches(9.14)]:
        connector(s, x, Inches(3.02), x + Inches(0.23), Inches(3.02), BLUE, 2.0)
    note_band(s, "讲课主线：先圈定片区，再分析原因，最后给出管段坐标并派单处置到现场。", 5.82)
    audit_slide(s, no)


def build():
    ensure_assets()
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    slide_cover(prs)
    slide_toc(prs)

    section_slide(prs, 3, "PART 01", "DMA漏损检测业务背景", "先讲清检测对象，再进入模型方法。", "dma")
    visual_text_slide(prs, 4, "DMA系统在漏损管控中的位置", "PART 01 业务背景", "dma", [
        ("DMA分区", "把复杂管网切成可计量、可分析、可考核的区域"),
        ("边界阀门", "控制分区边界，避免串水影响判断"),
        ("入口计量", "记录进水量和夜间最小流量，发现持续异常"),
        ("压力监测", "观察末梢压力和波动，辅助判断影响范围"),
        ("业务定位", "DMA能锁定异常片区，但不能直接指出具体漏点"),
    ], "结论：DMA是AI建模的基本单元，边界清晰和计量可靠决定模型上限。")

    visual_text_slide(prs, 5, "传统漏损检测方法与局限", "PART 01 业务背景", "dma", [
        ("MNF判断", "适合发现夜间底流异常，但需扣除合法夜间用水"),
        ("固定阈值", "部署简单，但难适配不同DMA的用水差异"),
        ("人工巡检", "能够确认现场，但排查范围大、效率受经验影响"),
        ("投诉报修", "适合明显爆管，难发现隐性漏损"),
        ("AI切入点", "从单点超限转为动态基线和多源证据判断"),
    ], "讲课重点：传统方法不是被替代，而是成为AI模型的基础输入和复核依据。")

    visual_text_slide(prs, 6, "DMA数据基础：模型能否有效取决于数据是否可用", "PART 01 业务背景", "data", [
        ("时序数据", "流量、压力、采样时间、缺失标记，用于预测基线"),
        ("空间数据", "管段、阀门、拓扑、DMA归属，用于定位收敛"),
        ("资产数据", "管龄、材质、口径、维修历史，用于风险排序"),
        ("工单数据", "报警、核查、维修、误报原因，用于标签回填"),
        ("上下文数据", "天气、节假日、施工、用户结构，用于解释扰动"),
    ], "落地提醒：统一编码、时序对齐、工单结构化，是建模前必须完成的基础工作。")

    section_slide(prs, 7, "PART 02", "AI模型应用发展逻辑", "从经验阈值走向动态基线、多源证据和模型组合。", "ai")
    visual_text_slide(prs, 8, "从阈值报警到动态基线", "PART 02 发展逻辑", "algorithm", [
        ("固定阈值", "按人工设定上下限报警，适合系统启动阶段"),
        ("统计基线", "用移动平均、同比环比刻画历史波动"),
        ("动态基线", "按每个DMA学习正常运行曲线和预测区间"),
        ("残差预警", "关注持续偏离，而不是单个时间点超限"),
        ("闭环更新", "用工单反馈持续修正阈值和模型"),
    ], "讲课重点：动态基线就是每个DMA自己的正常曲线。")

    visual_text_slide(prs, 9, "从单点报警到多源联动识别", "PART 02 发展逻辑", "data", [
        ("入口流量", "夜间底流抬升、日内曲线偏移，提示持续异常"),
        ("压力响应", "末梢压力下降或波动加剧，辅助判断位置范围"),
        ("GIS资产", "管龄、材质、阀门状态帮助解释风险来源"),
        ("历史工单", "同一区域多次维修或投诉，提高候选优先级"),
        ("外部扰动", "节假日、施工、商业夜间用水用于排除误报"),
    ], "结论：模型输出要带证据链，不能只给一个异常分数。")

    visual_text_slide(prs, 10, "从单一算法到模型组合", "PART 02 发展逻辑", "algorithm", [
        ("预测", "LSTM、GRU学习正常运行曲线"),
        ("发现", "孤立森林、DBSCAN、自编码器发现候选异常"),
        ("定位", "水力仿真和机器学习收敛候选管段"),
        ("排序", "随机森林、GBDT输出管段风险优先级"),
        ("解释", "知识图谱和智能体生成报告、问答和派单建议"),
    ], "讲法：漏损检测是一条业务链，因此模型也应是组合体系。")

    slide_architecture(prs, 11)

    section_slide(prs, 12, "PART 03", "核心算法体系", "算法讲解围绕业务任务展开，而不是罗列模型名。", "algorithm")
    visual_text_slide(prs, 13, "LSTM / GRU：时序预测与动态基线", "PART 03 核心算法", "algorithm", [
        ("输入", "过去窗口的流量、压力、天气、日期和节假日"),
        ("处理", "学习日周期、周周期和季节性变化"),
        ("输出", "未来短时窗口预测值、预测区间和残差"),
        ("适用", "入口流量连续、压力点稳定、用水周期明显的DMA"),
        ("注意", "按时间切分训练验证，避免未来信息泄漏"),
    ], "讲课重点：时序模型先学正常边界，再用残差判断是否需要复核。")

    visual_text_slide(prs, 14, "CNN-LSTM：多传感器特征融合", "PART 03 核心算法", "ai", [
        ("输入", "多个压力点、入口流量和局部时序片段"),
        ("卷积层", "提取多个传感器之间的局部变化特征"),
        ("循环层", "判断这些变化是否持续、是否符合异常模式"),
        ("输出", "综合残差、压力响应模式和预警建议"),
        ("限制", "压力点布局不足或拓扑错误时，解释能力会变弱"),
    ], "讲课重点：先看多个点怎么一起变，再看这种变化是否持续。")

    visual_text_slide(prs, 15, "孤立森林 / DBSCAN：少标签异常发现", "PART 03 核心算法", "algorithm", [
        ("孤立森林", "快速筛查多维异常，输出异常分数"),
        ("DBSCAN", "识别低密度离群点和异常簇"),
        ("输入特征", "MNF、残差、压力波动、多维统计特征"),
        ("业务输出", "异常时段、异常分数、主要特征和复核建议"),
        ("注意", "无监督结果是复核清单，不等同真实漏损结论"),
    ], "落地方式：少标签阶段先发现可疑对象，再由人工复核形成标签。")

    visual_text_slide(prs, 16, "自编码器：用重构误差识别异常状态", "PART 03 核心算法", "algorithm", [
        ("训练对象", "正常状态下的流量、压力和上下文组合关系"),
        ("判断依据", "重构误差升高，代表当前状态偏离正常模式"),
        ("适用场景", "漏损标签不足但监测数据连续的早期阶段"),
        ("业务输出", "异常时段、异常变量和复核优先级"),
        ("注意", "训练集要尽量排除已知异常，避免把异常学成正常"),
    ], "讲法：模型记住正常状态，异常状态会重构失败。")

    visual_text_slide(prs, 17, "随机森林 / GBDT / HGB：管段风险排序", "PART 03 核心算法", "algorithm", [
        ("输入", "管龄、材质、口径、压力、维修、投诉、道路等级"),
        ("处理", "学习历史漏损与资产特征之间的关系"),
        ("输出", "管段风险分、风险等级、检漏优先级"),
        ("用途", "支撑年度巡检计划和管网改造排序"),
        ("前提", "资产台账完整，工单结构化，标签口径一致"),
    ], "讲课重点：它回答哪里更容易漏，不直接回答当前是否正在漏。")

    visual_text_slide(prs, 18, "水力模型 + AI：机理与数据融合", "PART 03 核心算法", "ai", [
        ("机理侧", "提供拓扑、供水路径、压力敏感性和仿真样本"),
        ("数据侧", "提供真实流量压力、异常模式和历史工单"),
        ("融合方式", "仿真样本辅助训练，AI结果受水力约束"),
        ("业务输出", "候选片区、候选管段TopN、定位置信度"),
        ("前提", "水力模型需要校准，阀门状态和拓扑关系要可靠"),
    ], "价值：机理给边界，AI做识别，两者结合提高定位可信度。")

    visual_text_slide(prs, 19, "GNN与知识图谱：拓扑关系与解释能力", "PART 03 核心算法", "ai", [
        ("GNN", "把管网拓扑和压力传播关系纳入学习"),
        ("知识图谱", "连接DMA、管段、设备、工单和历史案例"),
        ("智能体", "基于模型结果生成报告、问答和派单建议"),
        ("应用", "异常原因解释、相似案例检索、经验沉淀"),
        ("控制", "保留专业判断、权限控制和人工确认"),
    ], "讲法：图模型解决连接关系，知识图谱解决解释关系。")

    visual_text_slide(prs, 20, "算法选型矩阵：按业务任务配置模型", "PART 03 核心算法", "algorithm", [
        ("动态基线", "LSTM、GRU、CNN-LSTM；输出预测区间和残差"),
        ("异常发现", "孤立森林、DBSCAN、自编码器；输出异常时段"),
        ("风险排序", "随机森林、GBDT、HGB；输出风险分和优先级"),
        ("定位收敛", "水力仿真+ML、GNN；输出候选管段TopN"),
        ("业务协同", "知识图谱、智能体；输出报告、问答和派单建议"),
    ], "选型原则：先看业务任务，再看数据条件，最后选择模型组合。")

    section_slide(prs, 21, "PART 04", "典型应用场景", "把模型输出转化为调度、巡检和管理动作。", "locate")
    visual_text_slide(prs, 22, "事前预警：DMA异常识别", "PART 04 应用场景", "loop", [
        ("模型输入", "入口流量、压力点、MNF、天气和节假日"),
        ("判断逻辑", "实际曲线持续偏离动态基线，同时出现压力响应"),
        ("输出内容", "异常时段、风险等级、主要证据、复核建议"),
        ("业务动作", "调度确认、现场复核、生成预警工单"),
        ("价值", "提前发现隐性漏损，减少被动报修"),
    ], "讲法：预警不是直接判定漏点，而是告诉人员哪里需要优先复核。")

    visual_text_slide(prs, 23, "事中定位：候选管段TopN收敛", "PART 04 应用场景", "locate", [
        ("模型输入", "DMA异常、压力响应、GIS拓扑、水力仿真"),
        ("判断逻辑", "从全分区逐步收敛到候选片区和候选管段"),
        ("输出内容", "管段编号、空间位置、置信度、证据和影响用户数"),
        ("业务动作", "按TopN顺序安排检漏和现场核查"),
        ("价值", "缩短巡检路线，提高现场定位效率"),
    ], "核心价值：把“哪个区可能漏”变成“哪段管优先查”。")

    visual_text_slide(prs, 24, "事后复盘：工单回填与持续学习", "PART 04 应用场景", "loop", [
        ("记录字段", "报警时间、异常特征、现场结果、漏点位置、误报原因"),
        ("复盘动作", "修正标签、更新特征、调整阈值、再训练"),
        ("输出内容", "误报类型、真实漏损样本、模型版本表现"),
        ("业务动作", "形成样本库和复盘报告"),
        ("价值", "把每次处置沉淀为训练数据"),
    ], "模型越用越准的前提，是每一次现场结果都能回填。")

    visual_text_slide(prs, 25, "压力优化：降漏、稳压、节能平衡", "PART 04 应用场景", "ai", [
        ("模型输入", "需水预测、压力监测、泵阀状态、水力模型"),
        ("约束条件", "最不利点压力、消防保障、二次供水补水、用户体验"),
        ("输出内容", "分时压力策略、泵站启停、阀门调节建议"),
        ("业务动作", "调度人员复核后执行压力优化方案"),
        ("价值", "在保障供水安全前提下降漏节能"),
    ], "讲课重点：压力优化必须先满足服务压力，再谈降漏和节能。")

    visual_text_slide(prs, 26, "DMA规划：分区与监测点优化", "PART 04 应用场景", "dma", [
        ("分区评价", "拓扑结构、地形高差、入口数量、边界阀数量"),
        ("监测点优化", "仿真不同漏点情景，比较压力点组合贡献"),
        ("输出内容", "分区方案评分、压力点优先级、建设建议"),
        ("业务动作", "指导DMA建设、计量点改造和压力点布设"),
        ("价值", "提高有限预算下的监测覆盖效率"),
    ], "规划类模型解决的是怎么建得更合理，而不是单次漏点定位。")

    section_slide(prs, 27, "PART 05", "模型建设与落地路径", "把模型做成可运行、可验证、可迭代的工程系统。", "data")
    visual_text_slide(prs, 28, "数据治理：五张基础表", "PART 05 落地路径", "data", [
        ("设备表", "流量计、压力计、RTU、边缘网关；识别采集来源"),
        ("管网表", "管段、阀门、DMA归属、资产属性；支撑空间定位"),
        ("时序表", "采样时间、流量压力、缺失标记；训练预测模型"),
        ("工单表", "报警、核查、维修、处置结果；形成业务闭环"),
        ("标签表", "真实漏损、误报原因、模型版本；支撑监督学习"),
    ], "检查清单：统一编码、时序对齐、缺失处理、工单结构化、标签口径统一。")

    visual_text_slide(prs, 29, "模型训练路径：基线 - 异常 - 定位 - 排序", "PART 05 落地路径", "algorithm", [
        ("1 基线", "建立MNF和移动基线，形成基础预警能力"),
        ("2 预测", "训练LSTM/GRU动态基线，识别持续残差"),
        ("3 异常", "引入无监督模型，形成候选异常清单"),
        ("4 定位", "结合水力仿真，输出候选管段TopN"),
        ("5 排序", "使用监督学习生成检漏优先级"),
    ], "顺序：先做可解释基础能力，再逐步引入复杂模型。")

    visual_text_slide(prs, 30, "模型验证：算法指标与业务指标并重", "PART 05 落地路径", "algorithm", [
        ("预测误差", "MAE、RMSE；衡量动态基线预测稳定性"),
        ("识别效果", "精确率、召回率、F1、AUC；衡量异常识别"),
        ("定位效果", "TopN命中率、平均排查范围；衡量现场价值"),
        ("业务闭环", "提前预警时间、误报比例、复盘回填率"),
        ("治理效果", "维修后MNF回落、压力恢复、漏损持续时间缩短"),
    ], "对外汇报中，业务指标通常比单一准确率更有说服力。")

    visual_text_slide(prs, 31, "模型运营：漂移监控与再训练", "PART 05 落地路径", "ai", [
        ("季节变化", "用水模式改变；季度检查基线偏移"),
        ("管网改造", "拓扑和压力制度变化；重新校准水力模型"),
        ("传感器更换", "采样误差变化；复核采集质量"),
        ("用户结构变化", "商业、工业或居民用水占比变化；更新基线"),
        ("新增工单", "标签样本增加；半年再训练并比较版本效果"),
    ], "模型上线不是结束，而是运营开始。")

    visual_text_slide(prs, 32, "试点到规模化推广路线", "PART 05 落地路径", "loop", [
        ("试点选区", "选择边界清晰、数据稳定、治理需求明确的DMA"),
        ("数据治理", "统一编码、时序对齐、工单结构化"),
        ("模型试运行", "预警复核、误报分析、阈值校准"),
        ("效果评估", "算法指标和业务指标同步评估"),
        ("规模推广", "沉淀接口、流程、版本和培训体系"),
    ], "能否复制，取决于数据接口、工单闭环和模型运营机制。")

    visual_text_slide(prs, 33, "总结：以业务闭环定义AI模型价值", "SUMMARY", "cover", [
        ("先数据，再模型", "没有清晰边界、可靠计量和结构化工单，复杂算法价值有限"),
        ("先预警，再定位", "动态基线发现问题，机理融合缩小排查范围"),
        ("先复核，再自动化", "模型输出进入工单闭环，才能持续改进"),
        ("先业务，再算法", "算法指标必须转化为提前预警、减少误报和MNF回落"),
        ("最终目标", "能发现、能解释、能派单、能复盘、能持续优化"),
    ], "AI模型的价值，不在于展示算法名词，而在于提升漏损治理效率。")

    visual_text_slide(prs, 34, "附：PPT素材与版式使用规范", "APPENDIX", "cover", [
        ("主视觉", "每页保留一个高质量插图，不使用被裁切的小图"),
        ("文字层", "标题、标签、表格和讲解要点全部可编辑"),
        ("结构层", "箭头、流程、模块用PPT形状绘制，便于修改"),
        ("检查项", "逐页检查字号、越界、空白和信息密度"),
        ("扩展原则", "主图讲场景，表格讲逻辑，底部一句话讲结论"),
    ], "这套规则可继续用于后续整套教学PPT扩展。")

    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    return PPTX


def build_preview():
    prs = Presentation(str(PPTX))
    W, H = 320, 180
    sx, sy = W / prs.slide_width, H / prs.slide_height
    try:
        font_title = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 8)
        font_body = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 6)
    except Exception:
        font_title = font_body = ImageFont.load_default()
    thumbs = []
    for slide in prs.slides:
        im = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(im)
        for sh in slide.shapes:
            x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
            if w <= 0 or h <= 0:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                    p.thumbnail((w, h))
                    im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
                except Exception:
                    pass
            else:
                fill_color = None
                line_color = None
                try:
                    rgb = sh.fill.fore_color.rgb
                    if rgb:
                        fill_color = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
                except Exception:
                    pass
                try:
                    rgb = sh.line.color.rgb
                    if rgb:
                        line_color = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
                except Exception:
                    pass
                if fill_color:
                    d.rounded_rectangle([x, y, x+w, y+h], radius=4, fill=fill_color, outline=line_color)
                elif line_color:
                    d.rectangle([x, y, x+w, y+h], outline=line_color)
                if hasattr(sh, "text") and sh.text.strip() and w > 20 and h > 8:
                    font = font_title if any(r.font.bold for pgh in sh.text_frame.paragraphs for r in pgh.runs) else font_body
                    txt = sh.text.strip().split("\n")[0][:24]
                    d.text((x+2, y+2), txt, fill=(20, 40, 60), font=font)
        thumbs.append(im)
    cols = 4
    rows = (len(thumbs) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * W, rows * (H + 18)), (230, 236, 242))
    d = ImageDraw.Draw(sheet)
    for i, t in enumerate(thumbs):
        x = (i % cols) * W
        y = (i // cols) * (H + 18)
        sheet.paste(t, (x, y))
        d.text((x + 4, y + H + 2), f"{i+1:02d}", fill=(0, 70, 148), font=font_title)
    sheet.save(PREVIEW)
    return PREVIEW


if __name__ == "__main__":
    print(build())
    print(build_preview())
