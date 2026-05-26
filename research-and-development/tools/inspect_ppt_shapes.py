from pathlib import Path

from pptx import Presentation

p = Path(r"D:\Users\Downloads\一厂一网_智慧水务系统全景图_文字版_editable.pptx")
prs = Presentation(p)
for si, slide in enumerate(prs.slides, 1):
    print("slide", si)
    for i, shape in enumerate(slide.shapes):
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip().replace("\n", "|")
            if text:
                print(i, text, shape.left, shape.top, shape.width, shape.height)
