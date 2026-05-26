# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO
from collections import defaultdict
import hashlib

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "机理、算法与实践：人工智能行业应用实证分析.pptx"
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
MOD = ASSET_DIR / "v4_modules"
ARCH = ASSET_DIR / "v4_arch"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_对外讲解重构版_v4.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_对外讲解重构版_v4_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_对外讲解重构版_v4_预览联系表.png"

NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
FONT = "微软雅黑"

ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "dma": MOD / "01_dma_boundary.png",
    "hydraulic": MOD / "02_hydraulic_model.png",
    "simulation": MOD / "03_simulation_dashboard.png",
    "butterfly": MOD / "04_error_butterfly.png",
    "night": MOD / "05_night_demand.png",
    "ai_evo": MOD / "06_ai_evolution.png",
    "lstm": MOD / "07_lstm_gru.png",
    "iforest": MOD / "08_isolation_forest.png",
    "dbscan": MOD / "09_dbscan.png",
    "autoencoder": MOD / "10_autoencoder.png",
    "fusion": MOD / "11_ai_hydraulic_fusion.png",
    "risk": MOD / "12_risk_ranking.png",
    "gnn": MOD / "13_gnn_topology.png",
    "governance": MOD / "14_data_governance.png",
    "edge": MOD / "15_edge_cloud.png",
    "arch": ARCH / "overall_architecture.png",
    "pain": ARCH / "pain_ai_path.png",
    "alg_combo": ARCH / "algorithm_combo.png",
    "impl": ARCH / "implementation_path.png",
    "workorder": ARCH / "workorder_no_person.png",
}

AUDIT_LINES = []


def emu(v): return Inches(v)


def clear_slides(prs):
    ids = prs.slides._sldIdLst
    for sid in list(ids):
        prs.part.drop_rel(sid.rId)
        ids.remove(sid)


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def fill(shp, color, trans=0):
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = trans
    shp.line.fill.background()


def stroke(shp, color=LINE, width=1):
    shp.line.color.rgb = color
    shp.line.width = Pt(width)


def rect(s, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = s.shapes.add_shape(typ, x, y, w, h)
    fill(shp, color, trans)
    stroke(shp, line, 1)
    return shp


def tb(s, text, x, y, w, h, size=15, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pic(s, key, x, y, w, h, crop=False):
    path = ASSETS[key] if isinstance(key, str) else key
    p = s.shapes.add_picture(str(path), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def connector(s, x1, y1, x2, y2, color=BLUE, width=1.6, arrow=True):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(s, title, sub, no):
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.38), emu(0.23), emu(0.12), emu(0.50)), BLUE)
    tb(s, title, emu(0.62), emu(0.12), emu(9.5), emu(0.44), 23.5, TEXT, True)
    tb(s, sub, emu(0.64), emu(0.62), emu(8.8), emu(0.24), 11.2, MUTED)
    tb(s, f"{no:02d}", emu(12.05), emu(0.20), emu(0.70), emu(0.24), 12, BLUE, True, PP_ALIGN.RIGHT)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.62), emu(0.90), emu(11.70), emu(0.016)), PALE)


def footer(s):
    tb(s, "AI供水管网DMA系统漏损检测 · 教学演讲重构版", emu(0.62), emu(7.08), emu(6.3), emu(0.18), 9.5, MUTED)


def content(prs, title, sub, no):
    s = blank(prs)
    header(s, title, sub, no)
    footer(s)
    return s


def label(s, title, body, x, y, w, h, color=BLUE, fill_color=PALE2, size=11.4):
    rect(s, x, y, w, h, fill_color, LINE, True)
    tb(s, title, x + emu(0.14), y + emu(0.10), w - emu(0.28), emu(0.22), 13.5, color, True)
    tb(s, body, x + emu(0.14), y + emu(0.43), w - emu(0.28), h - emu(0.50), size, TEXT)


def note(s, text, y=6.08, color=BLUE):
    rect(s, emu(0.72), emu(y), emu(11.85), emu(0.56), PALE2, LINE, True)
    tb(s, text, emu(0.98), emu(y + 0.14), emu(11.30), emu(0.27), 14.0, color, True, PP_ALIGN.CENTER)


def section(prs, no, part, title, subtitle, img_key):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(0.24), emu(7.5)), CYAN)
    tb(s, part, emu(0.88), emu(1.14), emu(1.4), emu(0.30), 16, CYAN, True)
    tb(s, title, emu(0.88), emu(1.82), emu(6.65), emu(0.72), 33, WHITE, True)
    tb(s, subtitle, emu(0.92), emu(3.00), emu(6.55), emu(0.75), 16, RGBColor(210, 232, 245))
    rect(s, emu(8.05), emu(1.18), emu(4.25), emu(4.25), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img_key, emu(8.27), emu(1.43), emu(3.80), emu(3.75), crop=False)
    tb(s, f"{no:02d}", emu(11.75), emu(6.60), emu(0.65), emu(0.24), 14, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)
    audit(s, no, allow_short=True)


def image_left(prs, no, title, sub, img_key, points, bottom, color=BLUE):
    s = content(prs, title, sub, no)
    rect(s, emu(0.72), emu(1.18), emu(5.00), emu(4.62), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img_key, emu(0.95), emu(1.42), emu(4.54), emu(4.12), crop=False)
    for i, (head, body, c) in enumerate(points):
        label(s, head, body, emu(6.05), emu(1.18 + i * 1.19), emu(5.85), emu(0.88), c, PALE2, 11.2)
    note(s, bottom, 6.08, color)
    audit(s, no)


def image_right(prs, no, title, sub, img_key, points, bottom, color=BLUE):
    s = content(prs, title, sub, no)
    rect(s, emu(7.48), emu(1.18), emu(4.80), emu(4.62), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img_key, emu(7.70), emu(1.42), emu(4.36), emu(4.12), crop=False)
    for i, (head, body, c) in enumerate(points):
        label(s, head, body, emu(0.78), emu(1.18 + i * 1.19), emu(6.25), emu(0.88), c, PALE2, 11.2)
    note(s, bottom, 6.08, color)
    audit(s, no)


def full_visual(prs, no, title, sub, img_key, takeaway, top_labels=None):
    s = content(prs, title, sub, no)
    rect(s, emu(0.55), emu(1.06), emu(12.25), emu(5.58), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img_key, emu(0.72), emu(1.18), emu(11.90), emu(5.20), crop=False)
    if top_labels:
        for x, text, color in top_labels:
            rect(s, emu(x), emu(1.18), emu(2.35), emu(0.40), WHITE, color, True)
            tb(s, text, emu(x + 0.10), emu(1.29), emu(2.15), emu(0.15), 11.5, color, True, PP_ALIGN.CENTER)
    note(s, takeaway, 6.66)
    audit(s, no, allow_short=True)


def two_visual(prs, no, title, sub, left_key, right_key, points, bottom):
    s = content(prs, title, sub, no)
    rect(s, emu(0.72), emu(1.18), emu(3.90), emu(3.15), WHITE, RGBColor(205, 229, 243), True)
    pic(s, left_key, emu(0.92), emu(1.38), emu(3.50), emu(2.75), crop=False)
    rect(s, emu(4.95), emu(1.18), emu(3.90), emu(3.15), WHITE, RGBColor(205, 229, 243), True)
    pic(s, right_key, emu(5.15), emu(1.38), emu(3.50), emu(2.75), crop=False)
    for i, (head, body, c) in enumerate(points):
        label(s, head, body, emu(9.18), emu(1.18 + i * 1.06), emu(3.25), emu(0.78), c, PALE2, 10.7)
    note(s, bottom, 5.98)
    audit(s, no)


def flow_cards(prs, no, title, sub, cards, bottom):
    s = content(prs, title, sub, no)
    for i, (head, body, img, color) in enumerate(cards):
        x = 0.70 + (i % 4) * 3.05
        y = 1.18 + (i // 4) * 2.35
        rect(s, emu(x), emu(y), emu(2.68), emu(1.98), WHITE, RGBColor(205, 229, 243), True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(2.68), emu(0.34)), color)
        tb(s, head, emu(x + 0.08), emu(y + 0.08), emu(2.52), emu(0.14), 11.5, WHITE, True, PP_ALIGN.CENTER)
        pic(s, img, emu(x + 0.18), emu(y + 0.50), emu(2.32), emu(0.88), crop=False)
        tb(s, body, emu(x + 0.16), emu(y + 1.43), emu(2.36), emu(0.36), 10.4, TEXT, False, PP_ALIGN.CENTER)
        if i % 4 != 3:
            connector(s, emu(x + 2.72), emu(y + 1.02), emu(x + 2.95), emu(y + 1.02), BLUE, 1.4)
    note(s, bottom, 6.12)
    audit(s, no)


def table_slide(prs, no, title, sub, img_key, rows, bottom):
    s = content(prs, title, sub, no)
    pic(s, img_key, emu(0.78), emu(1.18), emu(3.70), emu(4.72), crop=False)
    headers = ["主题", "讲解重点", "落地输出"]
    widths = [1.55, 3.25, 3.25]
    x = emu(4.70); y = emu(1.18)
    for i, h in enumerate(headers):
        xx = x + sum(emu(w) for w in widths[:i])
        rect(s, xx, y, emu(widths[i]), emu(0.42), BLUE, BLUE, False)
        tb(s, h, xx, y + emu(0.11), emu(widths[i]), emu(0.15), 11.2, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y + emu(0.54 + r * 0.82)
        for c, val in enumerate(row):
            xx = x + sum(emu(w) for w in widths[:c])
            rect(s, xx, yy, emu(widths[c]), emu(0.68), PALE2 if r % 2 == 0 else WHITE, LINE, False)
            tb(s, val, xx + emu(0.07), yy + emu(0.12), emu(widths[c] - 0.14), emu(0.38), 10.2, BLUE if c == 0 else TEXT, c == 0)
    note(s, bottom, 6.10)
    audit(s, no)


def audit(s, no, allow_short=False):
    prs_w, prs_h = 12192000, 6858000
    chars = pics = off = 0
    bottom = 0
    min_font = 99
    for sh in s.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs_w + 50000 or sh.top + sh.height > prs_h + 50000:
            off += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off: flags.append(f"越界{off}")
    if min_font < 9.5: flags.append(f"小字{min_font}")
    if pics < 1: flags.append("缺主视觉")
    if not allow_short and chars < 150: flags.append(f"文字偏少{chars}")
    if bottom < emu(6.1): flags.append("页面偏空")
    status = "FAIL" if flags else "OK"
    AUDIT_LINES.append(f"{no:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise RuntimeError(f"slide {no}: {'; '.join(flags)}")


def cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(5.45), emu(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.58), emu(0.86), emu(0.13), emu(5.85)), CYAN)
    tb(s, "AI模型在供水管网\nDMA系统漏损检测\n中的应用", emu(0.98), emu(1.10), emu(4.20), emu(1.78), 30, WHITE, True)
    tb(s, "对外教学演讲 · 重构版", emu(1.02), emu(3.17), emu(3.30), emu(0.30), 17, CYAN, True)
    tb(s, "基于原始Word提纲扩展：机理模型、AI发展、痛点分析、核心算法、应用场景、实施路径与展望。", emu(1.02), emu(5.25), emu(4.10), emu(0.62), 12.0, RGBColor(218, 237, 248))
    tb(s, "主视觉采用生图架构图，正文保留PPT可编辑标题、标签和讲解要点。", emu(1.02), emu(6.05), emu(4.10), emu(0.36), 11.5, RGBColor(218, 237, 248))
    pic(s, "arch", emu(5.72), emu(0.52), emu(7.25), emu(6.15), crop=False)
    audit(s, 1, allow_short=True)


def agenda(prs):
    cards = [
        ("01 模型应用背景", "政策趋势、机理模型、AI发展、痛点与价值", "ai_evo", BLUE),
        ("02 技术基础", "AI适配逻辑、数据部署、算法与业务价值", "governance", GREEN),
        ("03 核心应用", "事前预警、事中定位、事后闭环、调度规划", "fusion", ORANGE),
        ("04 实施路径", "规划、治理、训练、集成、运营与保障", "impl", PURPLE),
        ("05 典型案例", "用闭环工单讲清从报警到复盘", "workorder", CYAN),
        ("06 总结展望", "从漏损检测走向全生命周期智能管控", "cover", RED),
    ]
    flow_cards(prs, 2, "目录：按原始Word六部分重构", "每页围绕一个可讲问题展开，避免只摆图或只放表格", cards, "这版不再压缩成少量概念页，而是按原始提纲逐章展开为教学讲稿。")


def build():
    missing = [str(p) for p in ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("\n".join(missing))
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    cover(prs)
    agenda(prs)
    section(prs, 3, "PART 01", "模型应用背景", "从行业刚需、机理模型演进、AI发展和传统痛点建立问题背景。", "butterfly")
    image_right(prs, 4, "1.1 行业刚需：漏损管控从成本问题上升为韧性问题", "供水漏损影响水资源、能耗、水质安全和城市运行韧性", "night", [
        ("国家与行业导向", "漏损治理与水资源集约利用、城市更新、智慧水务建设同步推进。", BLUE),
        ("管理对象变化", "从单点抢修转向分区计量、连续监测和全过程绩效考核。", GREEN),
        ("讲课落点", "引出DMA：只有先把管网拆成可度量单元，AI才有稳定建模对象。", ORANGE),
    ], "本页用于开场：漏损检测不是单一检修问题，而是城市供水系统的长期治理问题。")
    image_left(prs, 5, "1.2 机理模型演进：从平差到瞬变分析", "传统水力模型是供水管网分析的基础，不应被简单否定", "hydraulic", [
        ("哈代克罗斯", "基于环路流量平差，适合早期管网计算与工程设计教学。", BLUE),
        ("牛顿拉夫逊", "基于节点水头和非线性方程迭代，提高复杂网络求解能力。", GREEN),
        ("特征线法", "用于瞬变流和水锤分析，解释爆管、阀门动作等短时动态。", ORANGE),
        ("教学表达", "机理模型提供物理约束，AI模型提供动态学习和误差容忍。", PURPLE),
    ], "这一页要讲清：AI不是取代机理模型，而是在运行场景中增强它。")
    table_slide(prs, 6, "1.2 模型软件生态：不同工具适合不同任务", "原始提纲中的EPANET、SWMM、WaterGEMS、Hammer等应作为工程背景讲清", "simulation", [
        ("开源模型", "EPANET偏供水水力分析；SWMM偏排水和雨洪过程。", "适合教学、验证和轻量试点。"),
        ("商业软件", "WaterGEMS、Hammer等强调建模、校核、水锤和工程应用。", "适合复杂管网建模与专业分析。"),
        ("GIS衍生", "Innovyze、MIKE等更强调空间数据、水环境和综合平台能力。", "适合多系统融合与规划管理。"),
        ("AI关系", "AI可接收仿真输出，也可反向辅助参数校准和候选漏点收敛。", "形成机理+数据双驱动。"),
    ], "软件生态页不是罗列名称，而是说明每类工具在AI漏损检测中的位置。")
    image_left(prs, 7, "1.3 AI模型应用发展：从规则阈值到智能协同", "AI发展脉络应服务于供水漏损检测场景，而不是泛泛讲技术史", "ai_evo", [
        ("规则与统计阶段", "固定阈值、MNF经验判断、统计控制图，优势是透明，弱点是适配性有限。", BLUE),
        ("机器学习阶段", "孤立森林、DBSCAN、随机森林、GBDT用于异常识别与风险排序。", GREEN),
        ("深度学习阶段", "LSTM、GRU、自编码器、CNN-LSTM处理多变量时序和复杂非线性。", ORANGE),
        ("智能体阶段", "结合知识图谱和业务系统，实现解释、报告、派单和复盘辅助。", PURPLE),
    ], "AI发展史只讲与本课相关的主线：从阈值判断走向模型组合和业务协同。")
    full_visual(prs, 8, "1.3 供水漏损可适配算法全景", "把算法放回预警、定位、排序、解释和闭环任务中讲", "alg_combo", "算法页可以作为本部分的总览：每个算法必须对应明确业务输出。")
    full_visual(prs, 9, "1.4 传统机理模型痛点：为什么DMA运行场景需要AI增强", "基础数据、拓扑、需水模式和压力漏损耦合都会影响定位结果", "pain", "痛点不是为了否定机理模型，而是说明运行阶段需要数据驱动能力补位。")
    image_right(prs, 10, "基础数据误差的蝴蝶效应", "微小台账或仪表误差会沿仿真链路逐级放大", "butterfly", [
        ("误差入口", "管径、粗糙系数、高程、阀门状态、计量时间戳都可能成为偏差源。", RED),
        ("传播链条", "阻抗计算偏差会传导到节点压力、管段流量和MNF拟合。", ORANGE),
        ("业务后果", "理论漏损量虚高或虚低，候选位置跑偏，现场排查成本上升。", BLUE),
    ], "这页对应原始Word中的“微小基础数据偏差→维修决策误判”链条。", RED)
    two_visual(prs, 11, "拓扑失真与用水模式时变", "运行管网每年改造、换管、新增用户，模型维护难度持续上升", "butterfly", "night", [
        ("拓扑变化", "隐性连通、边界阀状态不准，会直接破坏DMA水量平衡。", BLUE),
        ("用水变化", "居民、商业、工业和学校医院片区夜间用水规律差异明显。", GREEN),
        ("AI价值", "动态基线可学习本DMA的运行规律，降低单一阈值误报。", ORANGE),
    ], "讲这里时要把“管网在变、用户在变、数据在变”作为AI建模必要性讲清。")
    image_left(prs, 12, "压力-漏损耦合不足与模型维护成本", "传统模型依赖人工调参，对短时动态和低压场景响应有限", "hydraulic", [
        ("耦合难点", "漏损与压力、流量、阀门状态、用户需求存在复杂非线性关系。", BLUE),
        ("实时性不足", "突发爆管或短时异常需要更快的数据识别和分级响应。", RED),
        ("维护成本", "管网改造、新增用户、设备更换都会带来模型参数重校。", ORANGE),
        ("改进方向", "用AI识别异常，用水力模型约束物理合理性。", GREEN),
    ], "这一页把传统模型痛点自然引到“机理+AI双驱动”。")
    flow_cards(prs, 13, "AI破解DMA漏损管控问题的四项核心价值", "对应原始提纲：数据集搭建、自主学习、自动决策、迭代更新", [
        ("数据集搭建", "抗误差、容错性高、时序可对齐", "governance", BLUE),
        ("自主学习", "学习复杂规律和非线性耦合", "lstm", GREEN),
        ("自动决策", "参数寻优、风险排序、辅助派单", "autoencoder", ORANGE),
        ("迭代更新", "工单回填、漂移监控、再训练", "iforest", PURPLE),
    ], "AI价值要落到工程动作：更稳预警、更小范围、更少无效排查、更快闭环。")

    section(prs, 14, "PART 02", "技术基础", "讲清AI在水务场景的边界、架构、数据端部署和核心模型能力。", "iforest")
    full_visual(prs, 15, "2.1 AI在DMA系统中的总体架构", "数据层、模型层、业务层和反馈层共同支撑漏损检测闭环", "arch", "这张架构图作为全课核心页：后续算法、场景和实施路径都围绕它展开。")
    image_left(prs, 16, "DMA作为AI建模基本单元", "模型不应只看全网均值，而要学习不同分区的运行规律", "dma", [
        ("分区边界", "边界阀、入口流量计和压力点决定DMA是否具备建模条件。", BLUE),
        ("分区差异", "居民区、商业区、工业园区和老旧片区的用水规律不同。", GREEN),
        ("模型策略", "分区模型学习本DMA基线，全局模型学习跨DMA共性风险。", ORANGE),
        ("部署要求", "数据端需保证采样频率、时间同步、设备编码和边界状态一致。", PURPLE),
    ], "这一页回应原始Word中“DMA分区如何设计、数据端怎么部署”的要求。")
    image_right(prs, 17, "时序数据分析：动态基线与异常检测", "流量、压力数据本质上是高频时间序列", "lstm", [
        ("输入特征", "历史流量、压力、日期、节假日、天气和事件信息。", BLUE),
        ("预测对象", "下一时段正常流量或压力区间，而非直接给出漏点。", GREEN),
        ("异常证据", "残差幅度、持续时间、夜间窗口和压力响应共同判断。", RED),
    ], "时序模型的核心作用是建立“正常应当是什么样”的动态参照。")
    table_slide(prs, 18, "LSTM/GRU/CNN-LSTM：适配不同数据条件", "用户本地有LSTM环境，本页保留可用于演示的算法主线", "lstm", [
        ("LSTM", "适合长周期依赖，学习日周期、周周期、季节性和短期扰动。", "动态基线、残差预警。"),
        ("GRU", "结构更简化，参数更少，适合轻量部署和快速迭代。", "边缘端或小样本试点。"),
        ("CNN-LSTM", "先提取多压力点局部特征，再处理时间依赖。", "多监测点联合判断。"),
        ("讲解建议", "现场演示可用LSTM预测正常曲线，再展示实际曲线偏离。", "形成可理解证据。"),
    ], "这页不要讲成公式课，重点是“预测正常曲线—计算残差—触发复核”。")
    image_left(prs, 19, "孤立森林：快速筛查少而不同的异常点", "适合标签不足阶段的多维异常识别", "iforest", [
        ("基本思想", "异常点通常更容易被随机切分隔离，平均路径更短。", BLUE),
        ("适用数据", "流量残差、压力波动、MNF增幅、设备状态等多维特征。", GREEN),
        ("业务输出", "异常时段、异常分数和需要人工复核的片区清单。", ORANGE),
        ("注意事项", "无监督结果不能直接等同漏损，需要工单回填验证。", RED),
    ], "孤立森林适合早期建设，重点价值是快速把疑似异常筛出来。")
    two_visual(prs, 20, "DBSCAN与自编码器：异常识别的两类补充思路", "一个看密度离群，一个看重构误差", "dbscan", "autoencoder", [
        ("DBSCAN", "不预设聚类数量，可发现低密度离群点和异常簇。", BLUE),
        ("自编码器", "学习正常状态组合关系，重构误差高时提示异常。", GREEN),
        ("组合使用", "与LSTM残差、孤立森林交叉验证，减少单模型误判。", ORANGE),
    ], "无监督模型要服务复核清单，而不是直接替代现场判断。")
    image_right(prs, 21, "机器学习与水力模型融合", "利用仿真样本和水力约束提升漏点定位可信度", "fusion", [
        ("水力仿真", "模拟不同漏点位置和漏量下的压力、流量响应。", BLUE),
        ("AI学习", "把仿真输出、监测残差和资产特征一起输入模型。", GREEN),
        ("定位收敛", "将候选范围从整个DMA收敛到若干管段或节点。", ORANGE),
        ("前提条件", "拓扑、阀门状态、粗糙系数和需水模式需具备基本可信度。", RED),
    ], "融合模型是本课定位部分的关键：既用数据，也用物理合理性约束。")
    image_left(prs, 22, "随机森林、GBDT与HGB：管段风险排序", "结构化资产和工单数据适合树模型建模", "risk", [
        ("输入字段", "管龄、材质、口径、压力等级、维修历史、投诉频次、道路等级。", BLUE),
        ("模型优势", "对非线性特征友好，可输出特征重要性和风险排序。", GREEN),
        ("业务输出", "管段风险分、风险等级、改造优先级和巡检计划。", ORANGE),
        ("应用边界", "依赖历史工单质量，标签不准会直接影响排序可信度。", RED),
    ], "树模型适合对外讲清楚，因为它能把“为什么这段管风险高”解释出来。")
    two_visual(prs, 23, "知识图谱、GNN与端边云协同", "管网天然是图结构，业务系统天然是关系网络", "gnn", "edge", [
        ("图结构", "节点、管段、阀门、泵站、压力点和用户形成网络关系。", BLUE),
        ("知识组织", "连接DMA、设备、工单、维修、投诉、材质和管龄。", GREEN),
        ("部署协同", "边缘端做采集与轻量识别，云端做训练、分析和全局决策。", ORANGE),
    ], "这页用于扩展到智能体和智慧水务平台，但仍要围绕漏损检测闭环。")

    section(prs, 24, "PART 03", "核心应用", "按事前防控、事中处置、事后复盘，并补充调度和DMA规划场景。", "alg_combo")
    image_right(prs, 25, "3.1 事前防控：AI漏损智能预警体系", "目标是提前识别风险，而不是一报警就开挖", "night", [
        ("DMA分区管控", "自动识别边界异常、MNF异常和高风险片区。", BLUE),
        ("运行状态预警", "基于时序预测识别流量、压力异常并分级。", RED),
        ("健康度评估", "结合资产属性评估老化、腐蚀、材质和历史维修风险。", GREEN),
    ], "事前防控讲的是风险识别和准备动作，避免把预警系统讲成简单报警器。")
    image_left(prs, 26, "DMA分区智能管控：从区域锁定开始", "DMA先回答哪个区异常，AI再继续解释原因和位置", "dma", [
        ("边界异常", "识别边界阀误开、隐性连通和计量不平衡。", BLUE),
        ("MNF异常", "发现夜间最小流量持续抬升或波动模式异常。", RED),
        ("风险分级", "根据持续时间、残差幅度、压力响应和影响范围分级。", ORANGE),
        ("输出动作", "提示调度复核、现场巡检或进入定位模型。", GREEN),
    ], "这里要强调“DMA宏观锁定”是后续定位的前置条件。")
    image_left(prs, 27, "管网健康度评估：把漏损检测前移到风险治理", "模型不仅发现已发生漏损，也服务前置改造排序", "risk", [
        ("资产风险", "管龄、材质、口径、埋深、历史维修和腐蚀环境。", BLUE),
        ("运行风险", "压力波动、低压区、高流速和供水方向变化。", GREEN),
        ("外部风险", "道路施工、地面荷载、温度变化和投诉频次。", ORANGE),
        ("管理输出", "高风险管段清单、年度改造建议和巡检优先级。", PURPLE),
    ], "这一页把AI漏损检测从“发现问题”扩展到“降低问题发生概率”。")
    full_visual(prs, 28, "3.2 事中处置：漏点候选区域收敛", "融合SCADA实时数据、水力模型和AI算法，减少无效排查", "fusion", "从片区锁定到候选管段，是AI在现场检漏中最容易讲出价值的环节。")
    image_right(prs, 29, "水力模型 + AI：从片区到点位的定位逻辑", "候选位置必须同时满足数据异常和水力合理性", "hydraulic", [
        ("监测证据", "入口流量异常、压力点响应和MNF偏离构成第一层证据。", BLUE),
        ("仿真比对", "不同假想漏点的压力响应与实际曲线进行匹配。", GREEN),
        ("候选输出", "形成管段编号、空间位置、置信度和建议复核方式。", ORANGE),
        ("工程表达", "输出“优先查哪里”，不是直接宣布“这里一定漏”。", RED),
    ], "这页可与前面用户给出的DMA宏观锁定+AI微观溯源素材对应。")
    full_visual(prs, 30, "3.3 事后管控：从维修结果回到模型训练", "报警、定位、派单、维修、标签回填和再训练形成闭环", "workorder", "这张图不含人物，用系统界面讲清工单闭环和数据回流。")
    image_left(prs, 31, "调度类应用：压力优化与需水预测", "漏损治理可进一步延伸到稳压、降漏和节能协同", "edge", [
        ("需水预测", "预测分时用水需求，为泵站和阀门调度提供依据。", BLUE),
        ("压力优化", "在满足最不利点服务压力的前提下降低背景漏损风险。", GREEN),
        ("约束条件", "消防保障、二次供水补水、用户体验和安全冗余必须保留。", RED),
        ("协同价值", "把漏损检测与能耗优化、调度运行和韧性保障连接起来。", ORANGE),
    ], "调度类应用能让演讲从漏损检测扩展到智慧水务运行优化。")
    two_visual(prs, 32, "DMA规划类应用：分区与监测点优化", "AI可辅助评价分区方案和压力点布设效果", "dma", "gnn", [
        ("分区评价", "结合拓扑、地形、用户数量、压力制度和入口数量进行多目标评价。", BLUE),
        ("监测点优化", "仿真不同漏点场景，比较压力点组合对定位效果的贡献。", GREEN),
        ("预算约束", "在有限设备预算下优先布设对定位贡献最大的监测点。", ORANGE),
    ], "这一页补齐原始Word中提出的DMA规划类场景。")

    section(prs, 33, "PART 04", "实施路径", "从前期规划到长效运营，形成可落地的建设路线。", "impl")
    full_visual(prs, 34, "4.1-4.5 AI系统实施路径总览", "前期规划、数据治理、模型建设、工程落地、长效运营", "impl", "实施路径必须讲成项目推进路线，而不是技术清单。")
    table_slide(prs, 35, "4.2 数据治理：模型上线前的基础工程", "数据治理决定模型效果上限，也决定现场是否信任模型", "governance", [
        ("设备表", "流量计、压力计、RTU、网关、采样频率、设备状态。", "设备统一编码。"),
        ("管网表", "管段、阀门、节点、DMA归属、管材、管龄、口径。", "拓扑与资产底座。"),
        ("时序表", "采样时间、流量压力、缺失标记、清洗版本。", "模型训练输入。"),
        ("工单表", "报警、核查、维修、照片、处置时长和结果。", "闭环过程记录。"),
        ("标签表", "真实漏点、误报原因、模型版本和置信度。", "再训练样本库。"),
    ], "数据治理不是附属工作，而是AI漏损检测能不能落地的前置条件。")
    full_visual(prs, 36, "4.3 模型建设：选型、训练、验证与调优", "从基础基线到融合模型，逐步提高复杂度", "alg_combo", "模型建设要按数据成熟度推进，不能一开始就追求最复杂算法。")
    image_right(prs, 37, "模型验证：算法指标与业务指标并重", "对外汇报中，业务指标通常更能体现应用价值", "risk", [
        ("算法指标", "MAE、RMSE、精确率、召回率、F1、AUC等用于评估模型本身。", BLUE),
        ("业务指标", "提前预警时间、误报比例、排查范围、TopN命中率、闭环时长。", GREEN),
        ("验证方式", "按DMA、季节、用水类型和事件类型分组验证，避免平均值掩盖问题。", ORANGE),
        ("表达建议", "把准确率转换成少跑多少路、少挖多少次、少漏多少小时。", RED),
    ], "验证页要让管理层听懂模型价值，也让技术人员知道如何评估。")
    two_visual(prs, 38, "4.4 工程落地：系统集成与端边云协同", "避免重复建设，接入既有SCADA、GIS和运维系统", "edge", "simulation", [
        ("硬件适配", "兼容在线监测、巡检、内窥检测、压力计和边缘网关。", BLUE),
        ("系统集成", "与SCADA、GIS、运维管理系统对接，形成统一工作台。", GREEN),
        ("端边云协同", "边缘端快速识别，云端训练分析，业务端派单复盘。", ORANGE),
    ], "工程落地页要讲清系统接入和职责分工，避免PPT停留在模型实验。")
    image_left(prs, 39, "4.5 长效运营：人员能力与算法运维保障", "模型上线后需要有人维护、有人复盘、有人负责指标", "simulation", [
        ("人员培训", "管理人员看指标，技术人员管模型，运维人员用工单反馈结果。", BLUE),
        ("模型运营", "监控漂移、复盘误报、维护版本、定期再训练。", GREEN),
        ("制度保障", "明确报警响应、现场核查、数据回填和模型更新责任。", ORANGE),
        ("长期目标", "把漏损检测能力沉淀为可复制的城市管网治理体系。", PURPLE),
    ], "长效运营不是收尾页，而是决定AI系统能否持续产生价值的核心。")

    section(prs, 40, "PART 05", "典型案例与总结展望", "用闭环案例讲清应用过程，并回到未来技术、架构和行业升级。", "cover")
    full_visual(prs, 41, "典型案例讲法：从一次DMA异常到闭环复盘", "预警、定位、派单、维修、回填、再训练的完整链条", "workorder", "案例页建议按时间线讲，不只讲模型分数，而要讲工单如何闭环。")
    image_left(prs, 42, "6.1 总结：AI赋能漏损管控的三条主线", "对应原始Word的核心认知、核心逻辑和实施路径", "arch", [
        ("核心认知", "AI推动漏损管控从被动处置向主动防控转型。", BLUE),
        ("核心逻辑", "以数据为基础，以算法为核心，机理模型与AI模型融合。", GREEN),
        ("实施路径", "先规划后建设、先试点后推广、兼顾先进性与实用性。", ORANGE),
        ("讲课落点", "高质量漏损检测不是单个算法，而是系统工程。", RED),
    ], "总结页要把整套演讲重新收束到“数据、模型、场景、闭环”。")
    image_right(prs, 43, "6.2 展望：从漏损检测走向全生命周期智能管控", "技术升级、架构升级、体系升级和行业升级", "ai_evo", [
        ("技术升级", "水务大模型和智能体用于自然语言交互、报告生成和案例沉淀。", BLUE),
        ("架构升级", "端边云一体化与智能传感终端协同，实现更广域的实时感知。", GREEN),
        ("体系升级", "从漏损检测延伸到规划、建设、运行、维护和改造全链条。", ORANGE),
        ("行业升级", "标准化、模块化、轻量化降低应用门槛，推动行业普惠。", PURPLE),
    ], "最后一页不夸大AI，而是强调标准化落地和持续运营。")

    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    build_preview(prs)
    check_repetition(prs)
    return PPTX


def check_repetition(prs):
    by = defaultdict(list)
    for si, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                by[hashlib.md5(sh.image.blob).hexdigest()[:10]].append(si)
    lines = ["\n图片重复检查："]
    for h, slides in sorted(by.items(), key=lambda kv: len(kv[1]), reverse=True):
        if len(slides) > 2:
            lines.append(f"{h}\t{len(slides)}\t{slides}")
    with AUDIT.open("a", encoding="utf-8") as f:
        f.write("\n".join(lines))


def build_preview(prs):
    W, H = 330, 186
    cols = 5
    rows = (len(prs.slides) + cols - 1) // cols
    sx, sy = W / prs.slide_width, H / prs.slide_height
    try:
        font_b = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 8)
        font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 7)
    except Exception:
        font_b = font = ImageFont.load_default()
    sheet = Image.new("RGB", (cols * W, rows * (H + 20)), (230, 236, 242))
    d_sheet = ImageDraw.Draw(sheet)
    for idx, slide in enumerate(prs.slides):
        im = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(im)
        try:
            bg = slide.background.fill.fore_color.rgb
            if bg:
                im.paste(tuple(int(str(bg)[i:i+2], 16) for i in (0, 2, 4)), [0, 0, W, H])
        except Exception:
            pass
        for sh in slide.shapes:
            x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
            if w <= 0 or h <= 0:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                    p.thumbnail((w, h))
                    im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
                except Exception:
                    pass
                continue
            fc = lc = None
            try:
                rgb = sh.fill.fore_color.rgb
                if rgb: fc = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            try:
                rgb = sh.line.color.rgb
                if rgb: lc = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            if fc:
                d.rounded_rectangle([x, y, x+w, y+h], radius=4, fill=fc, outline=lc)
            elif lc:
                d.rectangle([x, y, x+w, y+h], outline=lc)
            txt = getattr(sh, "text", "").strip()
            if txt and w > 28 and h > 8:
                bold = False
                try:
                    bold = any(r.font.bold for p in sh.text_frame.paragraphs for r in p.runs)
                except Exception:
                    pass
                d.text((x+2, y+2), txt.split("\n")[0][:23], fill=(18, 38, 58), font=font_b if bold else font)
        cx = (idx % cols) * W
        cy = (idx // cols) * (H + 20)
        sheet.paste(im, (cx, cy))
        d_sheet.text((cx + 5, cy + H + 3), f"{idx+1:02d}", fill=(0, 70, 148), font=font_b)
    sheet.save(PREVIEW)


if __name__ == "__main__":
    print(build())
    print(AUDIT)
    print(PREVIEW)
