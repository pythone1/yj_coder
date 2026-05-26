# -*- coding: utf-8 -*-
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
GEN_ROOT = Path(r"C:\Users\Administrator\.codex\generated_images\019e1a29-183a-70c1-bc80-166e84fb9582")
TEMPLATE = max([p for p in ROOT.glob("*.pptx") if not p.name.startswith("~$")], key=lambda p: p.stat().st_size)
OUT = ROOT / "output" / "ppt"
ASSET_OUT = ROOT / "output" / "mixed_editable_assets"
OUT.mkdir(parents=True, exist_ok=True)
ASSET_OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_混合可编辑视觉样张版.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_混合可编辑视觉样张版_逐页检查.txt"


NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
SKY = RGBColor(72, 191, 235)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
FONT = "微软雅黑"


def latest_generated():
    files = sorted(GEN_ROOT.glob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not files:
        raise FileNotFoundError(f"no generated png in {GEN_ROOT}")
    return files[0]


def crop_assets():
    src = latest_generated()
    im = Image.open(src).convert("RGB")
    w, h = im.size
    crops = {
        "cover": (0.012, 0.02, 0.61, 0.545),
        "dma": (0.622, 0.02, 0.988, 0.545),
        "ai": (0.012, 0.57, 0.335, 0.985),
        "locate": (0.345, 0.57, 0.65, 0.985),
        "algorithm": (0.66, 0.57, 0.988, 0.985),
    }
    out = {}
    for name, box in crops.items():
        px = tuple(int(v * (w if i % 2 == 0 else h)) for i, v in enumerate(box))
        cropped = im.crop(px)
        path = ASSET_OUT / f"{name}.png"
        cropped.save(path, quality=95)
        out[name] = path
        if name == "cover":
            cw, ch = cropped.size
            # Remove the large blank rounded-card bands from the sheet crop.
            focus = cropped.crop((int(cw * 0.04), int(ch * 0.26), int(cw * 0.96), int(ch * 0.78)))
            focus_path = ASSET_OUT / "cover_focus.png"
            focus.save(focus_path, quality=95)
            out["cover_focus"] = focus_path
    for clean_name in ["cover_clean", "loop_clean", "dma_clean", "data_clean", "ai_clean", "locate_clean", "algorithm_clean"]:
        clean_path = ASSET_OUT / f"{clean_name}.png"
        if clean_path.exists():
            out[clean_name] = clean_path
    return src, out


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, trans=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = trans
    shape.line.fill.background()


def stroke(shape, color=LINE, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def tb(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=14.5, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.line_spacing = 0.92
        r = p.add_run()
        r.text = "· " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(typ, x, y, w, h)
    fill(s, color, trans)
    stroke(s, line, 1.0)
    return s


def connector(slide, x1, y1, x2, y2, color=BLUE, width=1.8):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    c.line.end_arrowhead = True
    return c


def pic(slide, path, x, y, w, h, crop=False):
    p = slide.shapes.add_picture(str(path), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.25), Inches(0.12), Inches(0.48))
    fill(bar, BLUE)
    tb(slide, title, Inches(0.62), Inches(0.16), Inches(8.8), Inches(0.38), 26, TEXT, True)
    tb(slide, sub, Inches(0.64), Inches(0.62), Inches(7.8), Inches(0.22), 11.5, MUTED)
    tb(slide, f"{no:02d}", Inches(12.05), Inches(0.22), Inches(0.7), Inches(0.25), 12, BLUE, True, PP_ALIGN.RIGHT)
    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.90), Inches(11.7), Inches(0.016))
    fill(divider, PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA漏损检测 · 混合可编辑样张", Inches(0.62), Inches(7.08), Inches(5.8), Inches(0.2), 10, MUTED)


def content(prs, title, sub, no):
    s = blank(prs)
    header(s, title, sub, no)
    footer(s)
    return s


def pill(slide, text, x, y, w, color=PALE, text_color=BLUE, size=12):
    rect(slide, x, y, w, Inches(0.32), color, color, True)
    tb(slide, text, x + Inches(0.06), y + Inches(0.075), w - Inches(0.12), Inches(0.14), size, text_color, True, PP_ALIGN.CENTER)


def note_band(slide, text, y=5.95):
    rect(slide, Inches(0.72), Inches(y), Inches(11.85), Inches(0.55), RGBColor(247, 252, 255), LINE, True)
    tb(slide, text, Inches(0.98), Inches(y + 0.16), Inches(11.3), Inches(0.20), 15.5, BLUE, True, PP_ALIGN.CENTER)


def module(slide, img_path, title, desc, x, y, w, h, color):
    rect(slide, x, y, w, h, WHITE, RGBColor(201, 228, 242), True)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.43))
    fill(top, color)
    tb(slide, title, x, y + Inches(0.095), w, Inches(0.20), 14.5, WHITE, True, PP_ALIGN.CENTER)
    pic(slide, img_path, x + Inches(0.14), y + Inches(0.58), w - Inches(0.28), h - Inches(1.28), crop=False)
    tb(slide, desc, x + Inches(0.14), y + h - Inches(0.43), w - Inches(0.28), Inches(0.32), 12.0, MUTED, False, PP_ALIGN.CENTER)


def slide_cover(prs, assets):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.30), Inches(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(1.05), Inches(0.13), Inches(5.10)), CYAN)
    tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.28), Inches(5.1), Inches(1.55), 35, WHITE, True)
    tb(s, "混合可编辑样张 · 风格确认版", Inches(1.0), Inches(3.20), Inches(4.9), Inches(0.32), 18, CYAN, True)
    tb(s, "生图负责高级插图，PPT负责可编辑文字、箭头、结构和讲解", Inches(1.0), Inches(5.78), Inches(5.0), Inches(0.45), 12.5, RGBColor(210, 232, 245))
    tb(s, "用于教学演讲：从DMA分区异常识别，讲到AI溯源定位、候选管段收敛与工单闭环处置。", Inches(1.0), Inches(6.28), Inches(5.05), Inches(0.38), 11.5, RGBColor(210, 232, 245))
    rect(s, Inches(6.42), Inches(0.74), Inches(6.45), Inches(5.32), WHITE, RGBColor(205, 229, 243), True)
    pic(s, assets.get("cover_clean", assets["cover_focus"]), Inches(6.58), Inches(0.92), Inches(6.12), Inches(4.95), crop=False)
    pill(s, "DMA锁定区域", Inches(6.95), Inches(6.18), Inches(1.70), PALE, BLUE, 12)
    pill(s, "AI溯源定位", Inches(8.90), Inches(6.18), Inches(1.70), PALE, GREEN, 12)
    pill(s, "工单闭环", Inches(10.85), Inches(6.18), Inches(1.35), PALE, ORANGE, 12)


def slide_arch(prs, assets):
    s = content(prs, "DMA宏观锁定 + AI微观溯源总体架构", "图片主视觉 + 可编辑文字/箭头/讲解区", 2)
    y = Inches(1.18)
    w = Inches(2.62)
    module(s, assets.get("dma_clean", assets["dma"]), "1 DMA锁定区域", "先判断哪个片区异常", Inches(0.62), y, w, Inches(4.18), BLUE)
    module(s, assets.get("data_clean", assets["cover_focus"]), "2 汇集关键数据", "流量、压力、管网、工单", Inches(3.55), y, w, Inches(4.18), CYAN)
    module(s, assets.get("ai_clean", assets["ai"]), "3 AI分析原因", "识别压力和流量变化", Inches(6.48), y, w, Inches(4.18), GREEN)
    module(s, assets.get("locate_clean", assets["locate"]), "4 定位并派单", "给出管段和坐标", Inches(9.41), y, w, Inches(4.18), ORANGE)
    for x in [Inches(3.28), Inches(6.21), Inches(9.14)]:
        connector(s, x, Inches(3.02), x + Inches(0.23), Inches(3.02), BLUE, 2.0)
    note_band(s, "讲课主线：先圈定片区，再分析原因，最后给出管段坐标并派单。", 5.82)


def slide_algorithm(prs, assets):
    s = content(prs, "AI算法体系：按业务任务组合模型", "高级算法插图 + 可编辑选型矩阵", 3)
    rect(s, Inches(0.70), Inches(1.18), Inches(5.45), Inches(4.78), WHITE, RGBColor(205, 229, 243), True)
    pic(s, assets.get("algorithm_clean", assets["algorithm"]), Inches(0.92), Inches(1.40), Inches(5.02), Inches(4.30), crop=False)
    headers = ["任务", "推荐模型", "输出"]
    rows = [
        ("动态基线", "LSTM / GRU", "预测区间、残差"),
        ("异常发现", "孤立森林 / 自编码器", "异常分数、异常时段"),
        ("风险排序", "GBDT / HGB", "管段风险分"),
        ("定位收敛", "水力仿真 + ML", "候选管段TopN"),
        ("解释协同", "知识图谱 / 智能体", "报告、问答、派单"),
    ]
    x0, y0 = Inches(6.35), Inches(1.30)
    widths = [1.45, 2.18, 2.03]
    for i, head in enumerate(headers):
        xx = x0 + sum(Inches(v) for v in widths[:i])
        rect(s, xx, y0, Inches(widths[i]), Inches(0.42), BLUE, BLUE, False)
        tb(s, head, xx, y0 + Inches(0.11), Inches(widths[i]), Inches(0.18), 13.5, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y0 + Inches(0.52 + r * 0.58)
        for c, val in enumerate(row):
            xx = x0 + sum(Inches(v) for v in widths[:c])
            rect(s, xx, yy, Inches(widths[c]), Inches(0.44), PALE2 if r % 2 == 0 else WHITE, LINE, False)
            tb(s, val, xx + Inches(0.06), yy + Inches(0.11), Inches(widths[c] - 0.12), Inches(0.18), 12.4, BLUE if c == 0 else TEXT, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)
    rect(s, Inches(6.35), Inches(4.58), Inches(5.66), Inches(0.95), RGBColor(247, 252, 255), LINE, True)
    tb(s, "讲法：预测模型负责判断正常边界；异常模型负责发现可疑信号；机理模型负责缩小位置范围。", Inches(6.55), Inches(4.82), Inches(5.25), Inches(0.34), 14.2, BLUE, True, PP_ALIGN.CENTER)
    note_band(s, "不要把算法讲成名词堆砌，要讲清楚每个模型解决哪一步业务问题。", 6.02)


def slide_loop(prs, assets):
    s = content(prs, "典型应用闭环：事前预警、事中定位、事后复盘", "场景插图 + 可编辑流程层", 4)
    rect(s, Inches(0.72), Inches(1.18), Inches(4.55), Inches(4.55), WHITE, RGBColor(205, 229, 243), True)
    pic(s, assets.get("loop_clean", assets["locate"]), Inches(0.92), Inches(1.38), Inches(4.15), Inches(4.12), crop=False)
    stages = [
        ("事前预警", "输入：流量、压力、MNF\n判断：持续偏离动态基线\n输出：异常等级和复核建议", BLUE),
        ("事中定位", "输入：压力响应、拓扑、工单\n判断：候选区域逐步收敛\n输出：管段TopN和节点坐标", CYAN),
        ("现场处置", "输入：定位坐标和证据链\n动作：派单、核查、维修\n输出：维修记录和现场结果", ORANGE),
        ("事后复盘", "输入：真实漏损或误报原因\n动作：标签回填、阈值校准\n输出：模型再训练样本", GREEN),
    ]
    for i, (title, body, color) in enumerate(stages):
        x = Inches(5.75 + (i % 2) * 3.18)
        y = Inches(1.18 + (i // 2) * 2.18)
        rect(s, x, y, Inches(2.82), Inches(1.63), WHITE, RGBColor(205, 229, 243), True)
        tb(s, title, x + Inches(0.18), y + Inches(0.14), Inches(2.45), Inches(0.22), 16, color, True)
        tb(s, body, x + Inches(0.18), y + Inches(0.50), Inches(2.45), Inches(0.88), 11.8, TEXT, True)
    connector(s, Inches(8.56), Inches(1.98), Inches(8.91), Inches(1.98), BLUE, 1.7)
    connector(s, Inches(7.18), Inches(2.82), Inches(7.18), Inches(3.34), BLUE, 1.7)
    connector(s, Inches(8.56), Inches(4.15), Inches(8.91), Inches(4.15), BLUE, 1.7)
    note_band(s, "核心价值：把“哪个区可能漏”转成“哪段管优先查、谁去查、结果怎么回填”。", 5.82)


def slide_spec(prs, assets):
    s = content(prs, "版式规范：混合可编辑PPT怎么扩展整套", "确认后按此方式重做全套PPT", 5)
    module(s, assets.get("dma_clean", assets["dma"]), "插图层", "高清图片只负责视觉", Inches(0.75), Inches(1.25), Inches(2.65), Inches(2.85), BLUE)
    module(s, assets.get("ai_clean", assets["ai"]), "结构层", "箭头、模块、流程可改", Inches(3.65), Inches(1.25), Inches(2.65), Inches(2.85), GREEN)
    module(s, assets.get("algorithm_clean", assets["algorithm"]), "文字层", "标题、标签、表格可改", Inches(6.55), Inches(1.25), Inches(2.65), Inches(2.85), ORANGE)
    module(s, assets.get("locate_clean", assets["locate"]), "检查层", "逐页检查空白和字号", Inches(9.45), Inches(1.25), Inches(2.65), Inches(2.85), PURPLE)
    rect(s, Inches(0.82), Inches(4.55), Inches(11.15), Inches(1.05), RGBColor(247, 252, 255), LINE, True)
    tb(s, "扩展规则", Inches(1.05), Inches(4.75), Inches(1.3), Inches(0.24), 16, BLUE, True)
    tb(s, "每页只保留一个主视觉；文字全部放在PPT文本框里；箭头和模块用PPT形状；正文不低于16pt；避免把复杂管网线条和小字堆在一起。", Inches(2.25), Inches(4.72), Inches(9.35), Inches(0.36), 14.2, TEXT, True)
    note_band(s, "确认这5页后，整套PPT按“主视觉 + 可编辑结构 + 讲解信息区”的规则重做。", 5.88)


def audit(prs):
    lines = []
    fail = False
    for idx, slide in enumerate(prs.slides, 1):
        chars = 0
        pics = 0
        min_font = 99
        off = 0
        for sh in slide.shapes:
            if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs.slide_width + 50000 or sh.top + sh.height > prs.slide_height + 50000:
                off += 1
            if sh.shape_type == 13:
                pics += 1
            if hasattr(sh, "text_frame") and sh.text.strip():
                chars += len(sh.text.strip())
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            min_font = min(min_font, r.font.size.pt)
        flags = []
        if off:
            flags.append(f"越界{off}")
        if min_font < 10:
            flags.append(f"小字{min_font}")
        if pics < 1:
            flags.append("缺少主视觉")
        if chars < 120:
            flags.append(f"文字不足{chars}")
        status = "FAIL" if flags else "OK"
        if flags:
            fail = True
        lines.append(f"{idx:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\t{'；'.join(flags)}")
    AUDIT.write_text("\n".join(lines), encoding="utf-8")
    if fail:
        print("\n".join(lines))
        raise SystemExit("audit failed")


def build():
    src, assets = crop_assets()
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    slide_cover(prs, assets)
    slide_arch(prs, assets)
    slide_algorithm(prs, assets)
    slide_loop(prs, assets)
    slide_spec(prs, assets)
    prs.save(PPTX)
    audit(prs)
    print(f"source={src}")
    return PPTX


if __name__ == "__main__":
    print(build())
