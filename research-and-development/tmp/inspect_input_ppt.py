from pptx import Presentation
from pathlib import Path
p=Path(r'D:\Users\Downloads\Jinshiyuan_AI_Production_Blueprint.pptx')
prs=Presentation(str(p))
print('slides', len(prs.slides))
print('size', prs.slide_width, prs.slide_height)
for i,s in enumerate(prs.slides,1):
    types={}
    texts=[]
    pics=0
    for sh in s.shapes:
        types[sh.shape_type]=types.get(sh.shape_type,0)+1
        if getattr(sh,'has_text_frame',False) and sh.text.strip():
            texts.append(sh.text.strip())
        if sh.shape_type == 13:
            pics += 1
    print(i, 'shapes', len(s.shapes), 'pics', pics, 'texts', len(texts), 'types', types)
    if texts:
        print('TEXT:', texts[:3])
