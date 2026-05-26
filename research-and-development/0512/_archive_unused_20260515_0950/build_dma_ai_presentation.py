# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
ASSET_ROOT = ROOT / "output" / "generated_ppt_assets_v6"
SPLIT = ASSET_ROOT / "split_plugins"
RAW = ASSET_ROOT / "raw_sheets"
OUTPUT = OUT / "AI供水管网DMA漏损检测_模型应用教学汇报_模板风格版.pptx"


W, H = Inches(13.333333), Inches(7.5)
NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(36, 157, 213)
PALE = RGBColor(232, 244, 252)
PALE2 = RGBColor(244, 249, 252)
TEXT = RGBColor(28, 45, 65)
MUTED = RGBColor(92, 109, 128)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(45, 150, 116)
ORANGE = RGBColor(242, 143, 45)
RED = RGBColor(208, 75, 75)

FONT = "微软雅黑"
FONT_BOLD = "微软雅黑"


def asset(folder, name):
    return str(SPLIT / folder / name)


def raw(name):
    return str(RAW / name)


def rgb(hexstr):
    hexstr = hexstr.strip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if transparency:
        shape.fill.transparency = transparency
    shape.line.fill.background()


def set_line(shape, color=BLUE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def text_box(slide, text, x, y, w, h, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.TOP, font=FONT, line_spacing=1.0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def add_multiline(slide, lines, x, y, w, h, size=14, color=TEXT, bullet=False, gap=0.88):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = gap
        if bullet:
            p.level = 0
        r = p.add_run()
        r.text = ("· " if bullet else "") + line
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def add_picture_fit(slide, path, x, y, w, h, crop=False):
    pic = slide.shapes.add_picture(path, x, y, width=w)
    scale = min(w / pic.width, h / pic.height)
    if crop:
        scale = max(w / pic.width, h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)
    return pic


def header(slide, title, section, num):
    text_box(slide, title, Inches(0.55), Inches(0.22), Inches(8.6), Inches(0.35), 22, TEXT, True)
    text_box(slide, "AI供水管网漏损检测研究", Inches(1.35), Inches(0.62), Inches(6.2), Inches(0.18), 8.5, MUTED)
    mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.36), Inches(0.26), Inches(0.11), Inches(0.45))
    set_fill(mark, BLUE)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.36), Inches(0.78), Inches(11.0), Inches(0.015))
    set_fill(line, PALE)
    text_box(slide, section, Inches(9.55), Inches(0.30), Inches(2.55), Inches(0.25), 9.5, BLUE, True, PP_ALIGN.RIGHT)
    text_box(slide, f"{num:02d}", Inches(12.15), Inches(0.25), Inches(0.62), Inches(0.32), 12, BLUE, True, PP_ALIGN.RIGHT)


def footer(slide):
    text_box(slide, "DMA · AI模型 · 漏损检测 · 工程落地", Inches(0.55), Inches(7.08), Inches(4.2), Inches(0.18), 8.5, MUTED)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.01))
    set_fill(line, PALE)


def content_slide(prs, title, section, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    header(slide, title, section, num)
    footer(slide)
    return slide


def section_slide(prs, part, title, subtitle, icon_path, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    set_fill(rect, NAVY)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    set_fill(accent, CYAN)
    text_box(slide, part, Inches(0.75), Inches(1.25), Inches(1.4), Inches(0.35), 15, CYAN, True)
    text_box(slide, title, Inches(0.75), Inches(2.0), Inches(6.8), Inches(0.75), 36, WHITE, True)
    text_box(slide, subtitle, Inches(0.8), Inches(3.05), Inches(6.8), Inches(0.75), 15, RGBColor(205, 226, 240))
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.15), Inches(1.45), Inches(3.55), Inches(3.55))
    set_fill(circ, RGBColor(10, 58, 116), 8)
    set_line(circ, CYAN, 1.2)
    add_picture_fit(slide, icon_path, Inches(8.55), Inches(1.88), Inches(2.75), Inches(2.75))
    text_box(slide, f"{num:02d}", Inches(11.8), Inches(6.55), Inches(0.75), Inches(0.35), 14, RGBColor(157, 198, 230), True, PP_ALIGN.RIGHT)
    return slide


def pill(slide, text, x, y, w, h, fill=PALE, color=BLUE, size=12, bold=True):
    sh = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(sh, fill)
    set_line(sh, RGBColor(204, 228, 241), 0.8)
    tf = sh.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return sh


def icon_label(slide, img, label, x, y, w=1.2, label_size=10.5):
    add_picture_fit(slide, img, x, y, Inches(w), Inches(w))
    text_box(slide, label, x - Inches(0.1), y + Inches(w) + Inches(0.05), Inches(w + 0.2), Inches(0.28),
             label_size, TEXT, True, PP_ALIGN.CENTER)


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.3):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    ln.line.end_arrowhead = True
    return ln


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_picture_fit(slide, raw("00_algorithm_overall_architecture.png"), Inches(6.0), Inches(0.25), Inches(7.3), Inches(6.9), crop=True)
    overlay = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(7.3), Inches(7.5))
    set_fill(overlay, NAVY)
    overlay.fill.transparency = 3
    side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(0.92), Inches(0.14), Inches(5.45))
    set_fill(side, CYAN)
    text_box(slide, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.15), Inches(5.8), Inches(1.55), 33, WHITE, True)
    text_box(slide, "模型机理 · 算法选型 · 场景落地 · 持续运营", Inches(1.0), Inches(3.02), Inches(5.5), Inches(0.32), 15, RGBColor(198, 226, 242))
    text_box(slide, "教学汇报版", Inches(1.0), Inches(5.75), Inches(2.2), Inches(0.36), 15, CYAN, True)
    text_box(slide, "基于DMA分区计量、时序预测、异常检测与水力模型融合", Inches(1.0), Inches(6.18), Inches(5.4), Inches(0.32), 10.5, RGBColor(214, 230, 242))

    # 2 contents
    slide = content_slide(prs, "目录", "CONTENTS", 2)
    items = [
        ("PART 01", "发展逻辑", "从阈值报警到动态基线，再到模型协同"),
        ("PART 02", "模型与算法", "时序预测、异常检测、监督排序、机理融合"),
        ("PART 03", "应用场景", "事前预警、事中定位、事后复盘与调度拓展"),
        ("PART 04", "落地路径", "数据治理、模型验证、漂移监控和持续迭代"),
    ]
    for i, (part, title, desc) in enumerate(items):
        y = Inches(1.35 + i * 1.25)
        text_box(slide, part, Inches(0.95), y, Inches(1.2), Inches(0.32), 11, CYAN, True)
        text_box(slide, title, Inches(2.25), y - Inches(0.05), Inches(2.2), Inches(0.42), 22, TEXT, True)
        text_box(slide, desc, Inches(4.65), y + Inches(0.02), Inches(6.0), Inches(0.28), 13, MUTED)
        line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.95), y + Inches(0.62), Inches(10.6), Inches(0.01))
        set_fill(line, PALE)
    add_picture_fit(slide, asset("04_small_infographic_plugins", "02_water_ai_droplet.png"), Inches(10.85), Inches(4.85), Inches(1.25), Inches(1.25))

    section_slide(prs, "PART 01", "发展逻辑", "供水管网AI模型应用从经验判断走向动态、联动和闭环。", asset("04_small_infographic_plugins", "10_timeseries_mini_chart.png"), 3)

    # 4 evolution
    slide = content_slide(prs, "AI模型应用演进路线", "PART 01 发展逻辑", 4)
    phases = [
        ("经验阈值", "MNF经验值\n压力低限报警", BLUE),
        ("统计分析", "移动平均\n同比环比", CYAN),
        ("机器学习", "孤立森林\n随机森林", GREEN),
        ("深度学习", "LSTM/GRU\n自编码器/GNN", ORANGE),
        ("智能协同", "知识图谱\n智能体闭环", RED),
    ]
    x0 = Inches(0.9)
    y = Inches(2.05)
    for i, (name, desc, color) in enumerate(phases):
        x = x0 + Inches(i * 2.25)
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(1.02), Inches(1.02))
        set_fill(circ, color)
        text_box(slide, str(i + 1), x, y + Inches(0.19), Inches(1.02), Inches(0.32), 20, WHITE, True, PP_ALIGN.CENTER)
        text_box(slide, name, x - Inches(0.38), y + Inches(1.25), Inches(1.8), Inches(0.3), 16, TEXT, True, PP_ALIGN.CENTER)
        text_box(slide, desc, x - Inches(0.45), y + Inches(1.65), Inches(1.95), Inches(0.58), 11, MUTED, False, PP_ALIGN.CENTER)
        if i < len(phases)-1:
            arrow(slide, x + Inches(1.18), y + Inches(0.51), x + Inches(2.02), y + Inches(0.51), RGBColor(154, 191, 214), 1.4)
    text_box(slide, "核心变化：从单点超限报警，升级为按DMA学习正常运行基线，并通过多源证据形成可解释风险判断。", Inches(1.0), Inches(5.45), Inches(10.8), Inches(0.45), 15, BLUE, True, PP_ALIGN.CENTER)

    # 5 dynamic baseline
    slide = content_slide(prs, "动态基线：把每个DMA的正常波动学出来", "PART 01 发展逻辑", 5)
    add_picture_fit(slide, asset("02_algorithm_plugins", "13_prediction_residual_chart.png"), Inches(0.95), Inches(1.25), Inches(4.0), Inches(4.0))
    text_box(slide, "固定阈值的限制", Inches(5.25), Inches(1.30), Inches(2.8), Inches(0.32), 17, TEXT, True)
    add_multiline(slide, ["分区差异大：用户结构、地形、压力制度不同", "正常扰动多：节假日、气温、商业夜间用水", "统一阈值容易带来误报或漏报"], Inches(5.25), Inches(1.80), Inches(6.2), Inches(1.2), 13.5, MUTED, True)
    text_box(slide, "AI模型的作用", Inches(5.25), Inches(3.35), Inches(2.8), Inches(0.32), 17, TEXT, True)
    add_multiline(slide, ["按DMA建立动态基线和预测区间", "持续跟踪残差幅度、持续时间和压力响应", "将异常转化为风险等级与复核建议"], Inches(5.25), Inches(3.85), Inches(6.2), Inches(1.25), 13.5, MUTED, True)
    pill(slide, "实际值持续偏离预测区间 + 漏损特征吻合 = 预警信号", Inches(5.25), Inches(5.55), Inches(6.1), Inches(0.55), PALE, BLUE, 13.5, True)

    # 6 multi-source
    slide = content_slide(prs, "多源联动：形成可解释的异常证据链", "PART 01 发展逻辑", 6)
    labels = [
        (asset("01_water_equipment_plugins", "01_smart_inlet_flow_meter.png"), "入口流量"),
        (asset("01_water_equipment_plugins", "02_pressure_sensor_node.png"), "压力点"),
        (asset("01_water_equipment_plugins", "07_gis_network_map.png"), "GIS资产"),
        (asset("01_water_equipment_plugins", "08_work_order_clipboard.png"), "历史工单"),
        (asset("01_water_equipment_plugins", "10_ai_model_cube.png"), "AI模型"),
    ]
    for i, (img, lab) in enumerate(labels):
        x = Inches(0.85 + i*2.25)
        icon_label(slide, img, lab, x, Inches(1.45), 1.05, 11)
        if i < 4:
            arrow(slide, x + Inches(1.25), Inches(1.95), x + Inches(1.85), Inches(1.95), RGBColor(165, 200, 220), 1.2)
    add_multiline(slide, [
        "运行时序：流量、压力、MNF、残差",
        "时间上下文：日期、节假日、天气、季节",
        "空间资产：管龄、材质、口径、阀门状态",
        "历史事件：维修、投诉、巡检、误报原因",
    ], Inches(1.0), Inches(4.05), Inches(5.0), Inches(1.55), 13, TEXT, True)
    add_multiline(slide, [
        "输出应面向处置：异常时段、主要证据、疑似原因、候选管段、复核方式和优先级。",
        "汇报重点应强调证据链，而不是只展示异常分数。"
    ], Inches(6.65), Inches(4.05), Inches(5.2), Inches(1.25), 13.5, BLUE, False)

    section_slide(prs, "PART 02", "模型与算法", "以业务任务选择模型组合，避免单一算法覆盖全部场景。", asset("02_algorithm_plugins", "14_model_training_pipeline.png"), 7)

    # 8 architecture
    slide = content_slide(prs, "DMA漏损检测AI总体架构", "PART 02 模型与算法", 8)
    add_picture_fit(slide, raw("00_algorithm_overall_architecture.png"), Inches(0.7), Inches(1.0), Inches(11.95), Inches(5.55))
    pill(slide, "数据层", Inches(0.8), Inches(6.3), Inches(1.35), Inches(0.35), PALE, BLUE, 11, True)
    pill(slide, "模型层", Inches(2.35), Inches(6.3), Inches(1.35), Inches(0.35), PALE, BLUE, 11, True)
    pill(slide, "业务层", Inches(3.9), Inches(6.3), Inches(1.35), Inches(0.35), PALE, BLUE, 11, True)
    pill(slide, "闭环层", Inches(5.45), Inches(6.3), Inches(1.35), Inches(0.35), PALE, BLUE, 11, True)

    # 9 time-series
    slide = content_slide(prs, "时序预测模型：建立正常运行曲线", "PART 02 模型与算法", 9)
    for i, (img, title, desc) in enumerate([
        (asset("02_algorithm_plugins", "01_lstm_sequence_prediction.png"), "LSTM", "学习长周期依赖，适合日/周周期明显的流量压力预测"),
        (asset("02_algorithm_plugins", "02_gru_gate_motif.png"), "GRU", "结构更轻，适合快速迭代和算力有限场景"),
        (asset("02_algorithm_plugins", "03_cnn_lstm_hybrid.png"), "CNN-LSTM", "先提取多传感器局部特征，再处理时间依赖"),
    ]):
        x = Inches(0.85 + i*3.95)
        add_picture_fit(slide, img, x, Inches(1.35), Inches(1.35), Inches(1.35))
        text_box(slide, title, x + Inches(1.55), Inches(1.50), Inches(1.45), Inches(0.3), 18, TEXT, True)
        text_box(slide, desc, x + Inches(1.55), Inches(1.92), Inches(2.05), Inches(0.7), 11.2, MUTED)
    pill(slide, "输入：历史流量 + 压力 + 天气 + 日历", Inches(1.25), Inches(4.1), Inches(3.2), Inches(0.48), PALE2, TEXT, 12)
    arrow(slide, Inches(4.55), Inches(4.34), Inches(5.55), Inches(4.34))
    pill(slide, "输出：预测值 + 预测区间 + 残差", Inches(5.7), Inches(4.1), Inches(3.2), Inches(0.48), PALE2, TEXT, 12)
    arrow(slide, Inches(9.05), Inches(4.34), Inches(10.05), Inches(4.34))
    pill(slide, "应用：异常预警和复核建议", Inches(10.2), Inches(4.1), Inches(2.35), Inches(0.48), PALE2, TEXT, 12)
    text_box(slide, "讲解要点：时序模型不是直接判定漏点，而是先建立每个DMA的正常运行边界。", Inches(1.05), Inches(5.55), Inches(10.9), Inches(0.35), 14, BLUE, True, PP_ALIGN.CENTER)

    # 10 unsupervised
    slide = content_slide(prs, "无监督异常检测：适合漏损标签不足阶段", "PART 02 模型与算法", 10)
    models = [
        (asset("02_algorithm_plugins", "05_isolation_forest_outlier.png"), "孤立森林", "快速筛查多维异常"),
        (asset("02_algorithm_plugins", "06_dbscan_cluster_outlier.png"), "DBSCAN", "识别低密度离群点"),
        (asset("02_algorithm_plugins", "04_autoencoder_reconstruction.png"), "自编码器", "用重构误差发现异常"),
    ]
    for i, (img, title, desc) in enumerate(models):
        x = Inches(1.05 + i*3.75)
        add_picture_fit(slide, img, x, Inches(1.25), Inches(1.45), Inches(1.45))
        text_box(slide, title, x + Inches(1.65), Inches(1.47), Inches(1.6), Inches(0.3), 17, TEXT, True)
        text_box(slide, desc, x + Inches(1.65), Inches(1.88), Inches(1.85), Inches(0.5), 11.5, MUTED)
    add_picture_fit(slide, asset("02_algorithm_plugins", "12_anomaly_score_gauge.png"), Inches(0.95), Inches(4.15), Inches(1.2), Inches(1.2))
    add_multiline(slide, [
        "输出作为人工复核清单使用，不直接等同于真实漏损。",
        "复核结果回填为真实漏损、合法用水、设备故障、边界异常和施工扰动。",
        "标签积累后，再逐步引入监督学习模型。"
    ], Inches(2.45), Inches(4.05), Inches(8.8), Inches(1.35), 13.3, TEXT, True)

    # 11 supervised risk
    slide = content_slide(prs, "监督学习：从报警识别走向管段风险排序", "PART 02 模型与算法", 11)
    add_picture_fit(slide, asset("02_algorithm_plugins", "07_random_forest_ensemble.png"), Inches(0.8), Inches(1.25), Inches(1.4), Inches(1.4))
    add_picture_fit(slide, asset("02_algorithm_plugins", "08_gradient_boosting_trees.png"), Inches(0.8), Inches(3.1), Inches(1.4), Inches(1.4))
    text_box(slide, "典型模型", Inches(2.45), Inches(1.28), Inches(2.2), Inches(0.32), 18, TEXT, True)
    add_multiline(slide, ["随机森林：稳定、抗噪声，适合解释特征贡献", "GBDT/HGB：对结构化资产和工单数据表现稳定"], Inches(2.45), Inches(1.82), Inches(4.2), Inches(1.2), 13, MUTED, True)
    text_box(slide, "输入变量", Inches(7.05), Inches(1.28), Inches(2.2), Inches(0.32), 18, TEXT, True)
    add_multiline(slide, ["管龄、材质、口径、压力等级", "历史维修、投诉频次、道路等级", "爆管记录、阀门状态、施工影响"], Inches(7.05), Inches(1.82), Inches(4.4), Inches(1.35), 13, MUTED, True)
    add_picture_fit(slide, asset("04_small_infographic_plugins", "12_risk_heatmap_tile.png"), Inches(7.25), Inches(4.15), Inches(1.2), Inches(1.2))
    pill(slide, "输出：管段风险分、风险等级、检漏优先级和改造排序", Inches(2.45), Inches(4.55), Inches(5.0), Inches(0.52), PALE, BLUE, 12.5, True)

    # 12 hydraulic
    slide = content_slide(prs, "机理模型与AI融合：提高漏点定位可信度", "PART 02 模型与算法", 12)
    icon_label(slide, asset("03_system_architecture_modules", "07_hydraulic_simulation_module.png"), "水力仿真", Inches(0.9), Inches(1.35), 1.15)
    arrow(slide, Inches(2.35), Inches(1.92), Inches(3.15), Inches(1.92))
    icon_label(slide, asset("04_small_infographic_plugins", "03_pressure_wave_pulse.png"), "压力响应", Inches(3.25), Inches(1.35), 1.15)
    arrow(slide, Inches(4.70), Inches(1.92), Inches(5.50), Inches(1.92))
    icon_label(slide, asset("03_system_architecture_modules", "08_leak_localization_module.png"), "候选定位", Inches(5.60), Inches(1.35), 1.15)
    arrow(slide, Inches(7.05), Inches(1.92), Inches(7.85), Inches(1.92))
    icon_label(slide, asset("01_water_equipment_plugins", "13_field_inspection_tablet.png"), "现场复核", Inches(7.95), Inches(1.35), 1.15)
    arrow(slide, Inches(9.40), Inches(1.92), Inches(10.20), Inches(1.92))
    icon_label(slide, asset("01_water_equipment_plugins", "16_maintenance_feedback_loop.png"), "结果回填", Inches(10.30), Inches(1.35), 1.15)
    add_multiline(slide, [
        "水力模型提供拓扑、供水路径、压力敏感性和仿真样本。",
        "AI模型负责从真实监测数据中识别异常模式和候选区域。",
        "融合后可把识别结果限制在水力学合理范围内，提升定位可信度。"
    ], Inches(1.15), Inches(4.2), Inches(10.7), Inches(1.2), 14, TEXT, True)

    # 13 graph/intelligent
    slide = content_slide(prs, "图模型、知识图谱与智能体：支撑解释和协同", "PART 02 模型与算法", 13)
    blocks = [
        (asset("02_algorithm_plugins", "10_gnn_pipe_topology.png"), "图神经网络", "把管网拓扑和压力传播关系纳入学习"),
        (asset("02_algorithm_plugins", "11_knowledge_graph.png"), "知识图谱", "连接DMA、管段、设备、工单和案例"),
        (asset("04_small_infographic_plugins", "08_ai_chip_icon.png"), "智能体协同", "生成报告、问答解释和派单建议"),
    ]
    for i, (img, title, desc) in enumerate(blocks):
        x = Inches(1.0 + i*3.85)
        add_picture_fit(slide, img, x, Inches(1.35), Inches(1.55), Inches(1.55))
        text_box(slide, title, x, Inches(3.1), Inches(2.6), Inches(0.35), 17, TEXT, True, PP_ALIGN.CENTER)
        text_box(slide, desc, x - Inches(0.15), Inches(3.62), Inches(2.9), Inches(0.72), 12.2, MUTED, False, PP_ALIGN.CENTER)
    pill(slide, "定位：专业模型负责判断，知识系统负责解释，人工确认负责闭环。", Inches(1.65), Inches(5.55), Inches(10.0), Inches(0.52), PALE, BLUE, 13, True)

    # 14 matrix
    slide = content_slide(prs, "算法选型矩阵：按业务任务配置模型组合", "PART 02 模型与算法", 14)
    rows = [
        ("动态基线", "LSTM / GRU / CNN-LSTM", "流量、压力、天气、日历", "预测区间、残差"),
        ("异常发现", "孤立森林 / DBSCAN / 自编码器", "MNF、残差、多维时序", "异常分数、异常时段"),
        ("风险排序", "随机森林 / GBDT / HGB", "管龄、材质、维修、投诉", "风险分、优先级"),
        ("定位收敛", "水力仿真 + ML / GNN", "拓扑、多点压力、仿真样本", "候选管段TopN"),
        ("业务协同", "知识图谱 / 智能体", "模型输出、GIS、工单", "报告、问答、派单建议"),
    ]
    x, y = Inches(0.65), Inches(1.15)
    widths = [1.45, 3.1, 3.0, 2.75]
    headers = ["业务任务", "推荐模型", "输入数据", "输出结果"]
    for i, head in enumerate(headers):
        pill(slide, head, x + sum(Inches(w) for w in widths[:i]), y, Inches(widths[i]), Inches(0.42), BLUE, WHITE, 11, True)
    for r, row in enumerate(rows):
        yy = y + Inches(0.55 + r*0.78)
        for c, val in enumerate(row):
            xx = x + sum(Inches(w) for w in widths[:c])
            box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, xx, yy, Inches(widths[c]), Inches(0.58))
            set_fill(box, PALE2 if r % 2 == 0 else WHITE)
            set_line(box, RGBColor(220, 235, 244), 0.4)
            text_box(slide, val, xx + Inches(0.08), yy + Inches(0.14), Inches(widths[c]-0.16), Inches(0.2),
                     9.6 if c != 1 else 9.2, TEXT if c else BLUE, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)
    text_box(slide, "选型原则：先明确业务目标，再选择算法；复杂模型必须建立在可靠数据和稳定边界之上。", Inches(0.9), Inches(6.15), Inches(11.3), Inches(0.32), 13.5, BLUE, True, PP_ALIGN.CENTER)

    section_slide(prs, "PART 03", "应用场景", "模型价值体现在预警、定位、复盘和调度等业务动作中。", asset("01_water_equipment_plugins", "05_pipe_leak_alert.png"), 15)

    # 16 before
    slide = content_slide(prs, "事前防控：基于AI的DMA异常预警", "PART 03 应用场景", 16)
    icon_label(slide, asset("04_small_infographic_plugins", "10_timeseries_mini_chart.png"), "动态基线", Inches(0.9), Inches(1.35), 1.1)
    arrow(slide, Inches(2.25), Inches(1.90), Inches(3.1), Inches(1.90))
    icon_label(slide, asset("04_small_infographic_plugins", "11_prediction_confidence_band.png"), "预测区间", Inches(3.25), Inches(1.35), 1.1)
    arrow(slide, Inches(4.6), Inches(1.90), Inches(5.45), Inches(1.90))
    icon_label(slide, asset("02_algorithm_plugins", "12_anomaly_score_gauge.png"), "异常评分", Inches(5.6), Inches(1.35), 1.1)
    arrow(slide, Inches(6.95), Inches(1.90), Inches(7.8), Inches(1.90))
    icon_label(slide, asset("04_small_infographic_plugins", "01_leak_warning_badge.png"), "风险预警", Inches(7.95), Inches(1.35), 1.1)
    arrow(slide, Inches(9.3), Inches(1.90), Inches(10.15), Inches(1.90))
    icon_label(slide, asset("01_water_equipment_plugins", "08_work_order_clipboard.png"), "复核建议", Inches(10.3), Inches(1.35), 1.1)
    add_multiline(slide, [
        "重点看残差持续时间、夜间低用水窗口、压力点响应和历史相似事件。",
        "输出应包括异常时间、持续时长、残差幅度、疑似原因和建议动作。",
    ], Inches(1.05), Inches(4.2), Inches(10.8), Inches(1.05), 14, TEXT, True)

    # 17 during
    slide = content_slide(prs, "事中处置：漏点候选区域收敛", "PART 03 应用场景", 17)
    add_picture_fit(slide, asset("01_water_equipment_plugins", "04_dma_boundary_map.png"), Inches(0.85), Inches(1.12), Inches(3.1), Inches(3.1))
    add_picture_fit(slide, asset("03_system_architecture_modules", "08_leak_localization_module.png"), Inches(4.75), Inches(1.12), Inches(2.4), Inches(2.4))
    add_picture_fit(slide, asset("01_water_equipment_plugins", "07_gis_network_map.png"), Inches(8.35), Inches(1.12), Inches(2.8), Inches(2.8))
    arrow(slide, Inches(3.85), Inches(2.55), Inches(4.55), Inches(2.55))
    arrow(slide, Inches(7.15), Inches(2.55), Inches(8.05), Inches(2.55))
    text_box(slide, "从整个DMA", Inches(1.15), Inches(4.33), Inches(2.5), Inches(0.3), 16, TEXT, True, PP_ALIGN.CENTER)
    text_box(slide, "收敛为候选片区", Inches(4.55), Inches(4.33), Inches(2.8), Inches(0.3), 16, TEXT, True, PP_ALIGN.CENTER)
    text_box(slide, "输出候选管段TopN", Inches(8.35), Inches(4.33), Inches(3.0), Inches(0.3), 16, TEXT, True, PP_ALIGN.CENTER)
    add_multiline(slide, [
        "候选项包含：管段编号、空间位置、置信度、主要证据、影响用户数、建议复核方式。",
        "该输出比单纯报警更符合现场检漏和派单实际。"
    ], Inches(1.0), Inches(5.35), Inches(10.8), Inches(0.85), 13.5, BLUE, True)

    # 18 after
    slide = content_slide(prs, "事后复盘：把每一次处置变成训练样本", "PART 03 应用场景", 18)
    steps = [
        (asset("01_water_equipment_plugins", "08_work_order_clipboard.png"), "报警/派单"),
        (asset("01_water_equipment_plugins", "13_field_inspection_tablet.png"), "现场核查"),
        (asset("04_small_infographic_plugins", "14_repair_wrench_pipe.png"), "维修处置"),
        (asset("03_system_architecture_modules", "11_repair_verification_module.png"), "效果验证"),
        (asset("03_system_architecture_modules", "12_model_retraining_loop.png"), "再训练"),
    ]
    for i, (img, lab) in enumerate(steps):
        x = Inches(0.85 + i*2.25)
        icon_label(slide, img, lab, x, Inches(1.35), 1.05, 10.5)
        if i < len(steps)-1:
            arrow(slide, x + Inches(1.22), Inches(1.88), x + Inches(1.80), Inches(1.88))
    add_multiline(slide, [
        "结构化记录：DMA编号、报警时间、异常特征、现场结果、漏点位置、漏损类型、误报原因、模型版本。",
        "复盘目的：修正标签、更新特征、调整阈值、降低同类误报。"
    ], Inches(1.05), Inches(4.25), Inches(10.8), Inches(1.0), 13.8, TEXT, True)

    # 19 dispatch and planning
    slide = content_slide(prs, "拓展应用：压力优化与DMA规划", "PART 03 应用场景", 19)
    icon_label(slide, asset("01_water_equipment_plugins", "14_pump_station.png"), "泵站运行", Inches(1.0), Inches(1.45), 1.2)
    icon_label(slide, asset("02_algorithm_plugins", "09_genetic_algorithm_optimization.png"), "优化算法", Inches(3.3), Inches(1.45), 1.2)
    icon_label(slide, asset("04_small_infographic_plugins", "05_dma_boundary_icon.png"), "DMA边界", Inches(5.6), Inches(1.45), 1.2)
    icon_label(slide, asset("01_water_equipment_plugins", "02_pressure_sensor_node.png"), "压力点", Inches(7.9), Inches(1.45), 1.2)
    icon_label(slide, asset("03_system_architecture_modules", "16_executive_dashboard.png"), "管理看板", Inches(10.2), Inches(1.45), 1.2)
    add_multiline(slide, [
        "压力优化：在保障最不利点服务压力、消防保障和用户体验前提下，平衡降漏、节能与稳压。",
        "DMA规划：结合拓扑结构、地形高差、入口数量、边界阀数量和施工成本进行多目标评价。",
        "监测点优化：通过仿真比较不同压力点组合对定位效果的贡献。"
    ], Inches(1.05), Inches(4.25), Inches(10.8), Inches(1.35), 13.5, TEXT, True)

    section_slide(prs, "PART 04", "落地路径", "数据治理、训练验证和模型运营共同决定系统长期价值。", asset("03_system_architecture_modules", "03_data_cleaning_filter.png"), 20)

    # 21 data governance
    slide = content_slide(prs, "数据治理：模型上线前的基础工程", "PART 04 落地路径", 21)
    columns = [
        ("设备表", "流量计、压力计、RTU、边缘网关", asset("01_water_equipment_plugins", "12_edge_gateway.png")),
        ("管网表", "管段、阀门、DMA归属、资产属性", asset("01_water_equipment_plugins", "07_gis_network_map.png")),
        ("时序表", "采样时间、流量压力、缺失标记", asset("04_small_infographic_plugins", "10_timeseries_mini_chart.png")),
        ("工单表", "报警、核查、维修、处置结果", asset("01_water_equipment_plugins", "08_work_order_clipboard.png")),
        ("标签表", "真实漏损、误报原因、模型版本", asset("02_algorithm_plugins", "16_human_in_loop_validation.png")),
    ]
    for i, (title, desc, img) in enumerate(columns):
        x = Inches(0.75 + i*2.42)
        add_picture_fit(slide, img, x + Inches(0.4), Inches(1.32), Inches(1.0), Inches(1.0))
        text_box(slide, title, x, Inches(2.60), Inches(1.85), Inches(0.28), 15, BLUE, True, PP_ALIGN.CENTER)
        text_box(slide, desc, x - Inches(0.05), Inches(3.05), Inches(1.95), Inches(0.65), 10.7, MUTED, False, PP_ALIGN.CENTER)
    text_box(slide, "治理重点：统一编码、时序对齐、缺失处理、毛刺剔除、工单结构化、标签口径统一。", Inches(1.1), Inches(5.55), Inches(10.8), Inches(0.35), 14, BLUE, True, PP_ALIGN.CENTER)

    # 22 train verify operate
    slide = content_slide(prs, "训练、验证与运营：从模型可用到长期可信", "PART 04 落地路径", 22)
    stages = [
        ("1", "基础基线", "MNF + 移动基线\n形成可用预警"),
        ("2", "动态预测", "LSTM / GRU\n识别持续残差"),
        ("3", "异常检测", "孤立森林 / 自编码器\n交叉验证"),
        ("4", "定位排序", "水力模型 + 工单标签\n候选管段TopN"),
        ("5", "持续运营", "漂移监控 + 再训练\n闭环迭代"),
    ]
    for i, (n, title, desc) in enumerate(stages):
        x = Inches(0.75 + i*2.42)
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.52), Inches(1.35), Inches(0.72), Inches(0.72))
        set_fill(circ, [BLUE, CYAN, GREEN, ORANGE, RED][i])
        text_box(slide, n, x + Inches(0.52), Inches(1.50), Inches(0.72), Inches(0.25), 16, WHITE, True, PP_ALIGN.CENTER)
        text_box(slide, title, x, Inches(2.33), Inches(1.75), Inches(0.3), 15, TEXT, True, PP_ALIGN.CENTER)
        text_box(slide, desc, x - Inches(0.12), Inches(2.82), Inches(2.0), Inches(0.65), 10.8, MUTED, False, PP_ALIGN.CENTER)
        if i < 4:
            arrow(slide, x + Inches(1.55), Inches(1.70), x + Inches(2.10), Inches(1.70), RGBColor(170, 200, 218), 1.1)
    pill(slide, "验证指标：算法指标看准确性，业务指标看提前预警、误报比例、TopN命中率和维修后MNF回落。", Inches(1.1), Inches(5.35), Inches(11.05), Inches(0.55), PALE, BLUE, 12.5, True)

    # 23 summary
    slide = content_slide(prs, "总结：以业务闭环定义AI模型价值", "SUMMARY", 23)
    add_picture_fit(slide, asset("04_small_infographic_plugins", "16_closed_loop_arrows.png"), Inches(0.95), Inches(1.35), Inches(2.0), Inches(2.0))
    text_box(slide, "建设顺序", Inches(3.45), Inches(1.25), Inches(2.0), Inches(0.32), 18, TEXT, True)
    add_multiline(slide, ["先数据，再基线", "先预警，再定位", "先复核，再自动化"], Inches(3.45), Inches(1.75), Inches(2.8), Inches(1.1), 14, MUTED, True)
    text_box(slide, "模型组合", Inches(6.55), Inches(1.25), Inches(2.0), Inches(0.32), 18, TEXT, True)
    add_multiline(slide, ["时序预测建立边界", "异常检测发现候选", "监督学习排序风险", "水力模型约束定位"], Inches(6.55), Inches(1.75), Inches(3.0), Inches(1.35), 14, MUTED, True)
    text_box(slide, "落地关键", Inches(9.85), Inches(1.25), Inches(2.0), Inches(0.32), 18, TEXT, True)
    add_multiline(slide, ["证据链可解释", "工单可回填", "模型可监控", "流程可闭环"], Inches(9.85), Inches(1.75), Inches(2.4), Inches(1.35), 14, MUTED, True)
    text_box(slide, "供水管网DMA漏损检测的AI应用，应从算法展示转向业务闭环：能发现、能解释、能派单、能复盘、能持续优化。", Inches(1.25), Inches(5.15), Inches(10.8), Inches(0.65), 18, BLUE, True, PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
