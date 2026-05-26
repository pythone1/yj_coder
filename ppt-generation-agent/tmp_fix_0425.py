from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu


BASE = Path(r"E:\PY\research\0425")
INPUT = BASE / "生产流程.pptx"
OUTPUT = BASE / "生产流程_文本修正_补底部模板.pptx"

REPLACEMENTS = {
    "输 迭线": "输送线",
    "锡护": "锅炉",
    "异崇检到": "异常检测",
    "预湖": "预测",
    "遇终功能效益": "最终功能效益",
    "停 机凤险": "停机风险",
    "橙惨计划": "检修计划",
    "摄矢": "损失",
    "赐标": "贴标",
    "波位": "液位",
    "酒体思浮异物": "酒体悬浮异物",
    "主霎算法任务": "主要算法任务",
    "图像分劃": "图像分割",
    "轨迹识刖": "轨迹识别",
    "复较压力": "复核压力",
    "缺船追": "缺陷追",
    "潮和质量": "溯和质量",
    "任努": "任务",
    "最 功 能 效 益": "最终功能效益",
    "上瓶": "上甑",
    "昇常": "异常",
    "现场传燃数据": "现场传感数据",
    "温温度": "温度",
    "酒酷温度": "酒醅温度",
    "异案检测": "异常检测",
    "控制围": "控制图",
    "等缀": "等级",
    "掠入": "输入",
    "专案舰则": "专家规则",
    "检素": "检索",
    "规钏约束": "规则约束",
    "牧据": "数据",
    "措荐抉行": "推荐执行",
    "反懊学习": "反馈学习",
    "横型童训": "模型重训",
    "辖出": "输出",
    "建议璃参": "建议微调",
    "知玖沉渡": "知识沉淀",
    "产续": "产线",
    "发观": "发现",
    "该杀率": "误杀率",
    "采集：NG图片、复核络": "采集：NG图片、复核结论",
    "罄线": "基线",
    "健康指数横型": "健康指数模型",
    "尖时时序": "实时序列",
    "数据 | 汇粲": "数据 | 汇聚",
    "逮度": "速度",
    "裁重": "载重",
    "起焚点": "起终点",
    "紫行区": "禁行区",
    "A\"": "A*",
    "多车落径": "多车路径",
    "规刺": "规划",
    "重规刘": "重规划",
    "逼过仿真": "通过仿真",
    "交忖时间": "交付时间",
    "温装": "混装",
    "多仓睾求": "多仓要求",
    "SKU 期间": "SKU 周转",
    "关肤规则": "关联规则",
    "库位保分": "库位评分",
    "密令": "订单",
    "启寂丸排序": "启发式排序",
    "液灰计刻": "波次计划",
    "捺选顺序": "拣选顺序",
    "月台支排": "月台安排",
    "辅入": "输入",
    "装魉舰物": "装载规划",
    "约柬伏化": "约束优化",
    "籍出": "输出",
    "多合节赤": "多仓节点",
    "韩出": "输出",
    "肠逆路线": "配送路线",
    "庠区": "库区",
    "月台、又车": "月台、叉车",
    "拣 遮路径": "拣选路径",
    "疲儋": "疲倦",
    "淀粉辜": "淀粉率",
    "Al": "AI",
    "AI质检": "AI质检",
}


def fix_text(value):
    for old, new in REPLACEMENTS.items():
        value = value.replace(old, new)
    return value.replace("ANotebookLM", "").replace("NotebookLM", "")


def tune_shape(shape):
    if hasattr(shape, "text_frame") and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = fix_text(run.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = fix_text(cell.text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            tune_shape(child)


def put_after_background(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(3, el)


def add_footer(slide, slide_width, slide_height):
    footer_h = Emu(213333)
    footer_y = Emu(int(slide_height) - int(footer_h))
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, footer_y, slide_width, footer_h)
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(71, 105, 158)
    footer.line.fill.background()
    put_after_background(slide, footer)


prs = Presentation(str(INPUT))
for slide in prs.slides:
    for shape in slide.shapes:
        tune_shape(shape)
    add_footer(slide, prs.slide_width, prs.slide_height)

prs.save(str(OUTPUT))
print(OUTPUT)
