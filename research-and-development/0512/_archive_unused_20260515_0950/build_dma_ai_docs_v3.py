from pathlib import Path
import math
import random

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output"
DOC_DIR = OUT / "doc"
ASSET_DIR = OUT / "assets_v3"
PPT_DIR = OUT / "ppt"
for d in (DOC_DIR, ASSET_DIR, PPT_DIR):
    d.mkdir(parents=True, exist_ok=True)

FONT_PATH = Path(r"C:\Windows\Fonts\simhei.ttf")
F_XL = ImageFont.truetype(str(FONT_PATH), 48)
F_H = ImageFont.truetype(str(FONT_PATH), 36)
F_M = ImageFont.truetype(str(FONT_PATH), 27)
F_S = ImageFont.truetype(str(FONT_PATH), 22)
F_XS = ImageFont.truetype(str(FONT_PATH), 18)

INK = (24, 35, 52)
MUTED = (92, 105, 124)
GRID = (217, 226, 237)
BLUE = (37, 99, 168)
CYAN = (14, 148, 168)
GREEN = (52, 141, 93)
ORANGE = (222, 139, 33)
REDP = (203, 73, 65)
PURPLE = (105, 93, 190)
YELLOW = (247, 190, 70)

DOC_RED = RGBColor(192, 0, 0)
DOC_BLACK = RGBColor(0, 0, 0)


def wrap(draw, text, font, max_width):
    lines, cur = [], ""
    for ch in text:
        test = cur + ch
        if draw.textlength(test, font=font) <= max_width:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = ch
    if cur:
        lines.append(cur)
    return lines


def tc(draw, box, text, font, fill=INK, spacing=5):
    x1, y1, x2, y2 = box
    lines = []
    for raw in text.split("\n"):
        lines.extend(wrap(draw, raw, font, x2 - x1 - 18))
    hs = [draw.textbbox((0, 0), line, font=font)[3] - draw.textbbox((0, 0), line, font=font)[1] for line in lines]
    total = sum(hs) + spacing * max(0, len(lines) - 1)
    y = y1 + (y2 - y1 - total) / 2
    for line, h in zip(lines, hs):
        bb = draw.textbbox((0, 0), line, font=font)
        draw.text((x1 + (x2 - x1 - (bb[2] - bb[0])) / 2, y), line, font=font, fill=fill)
        y += h + spacing


def save(img, name):
    p = ASSET_DIR / name
    img.save(p)
    return p


def canvas():
    return Image.new("RGBA", (1600, 900), (255, 255, 255, 0))


def rounded(draw, box, color, width=3, fill=(255, 255, 255, 235), radius=24):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=color + (255,), width=width)


def soft_label(draw, box, text, color, font=F_XS):
    draw.rounded_rectangle(box, radius=18, fill=color + (35,), outline=color + (220,), width=2)
    tc(draw, box, text, font, INK)


def fig_arch():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "DMA + AI漏损检测总体架构", font=F_XL, fill=INK)
    d.text((72, 105), "按层拆解，避免把AI误解为单个算法或单个软件功能", font=F_M, fill=MUTED)
    layers = [
        ("工程分区", ["边界阀核验", "入口/出口计量", "压力分区", "用户结构"], BLUE),
        ("感知采集", ["入口流量", "关键压力点", "远传水表", "设备状态"], CYAN),
        ("数据治理", ["时序对齐", "缺失补齐", "异常剔除", "工单标签"], GREEN),
        ("模型识别", ["MNF基线", "LSTM/GRU", "无监督异常", "水力融合"], ORANGE),
        ("业务闭环", ["风险分级", "候选管段", "巡检派单", "复盘再训练"], REDP),
    ]
    x0, y0, w, h, gap = 70, 210, 265, 470, 35
    for i, (name, items, color) in enumerate(layers):
        x = x0 + i * (w + gap)
        rounded(d, (x, y0, x + w, y0 + h), color, 4, radius=26)
        d.ellipse((x + 28, y0 + 28, x + 78, y0 + 78), fill=color + (255,))
        d.text((x + 95, y0 + 35), name, font=F_M, fill=INK)
        for j, item in enumerate(items):
            soft_label(d, (x + 28, y0 + 125 + j * 78, x + w - 28, y0 + 178 + j * 78), item, color, F_S)
    d.rounded_rectangle((250, 742, 1350, 810), radius=22, fill=(245, 248, 252, 235), outline=GRID + (255,), width=2)
    tc(d, (270, 742, 1330, 810), "PPT讲法：先讲DMA工程边界，再讲数据质量，最后讲AI模型，否则算法无法落地", F_S, INK)
    return save(img, "v3_01_architecture_no_arrow.png")


def fig_dma_blueprint():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "DMA分区建设蓝图", font=F_XL, fill=INK)
    d.text((72, 105), "一个可建模DMA必须同时满足边界、计量、压力、台账、工单五类条件", font=F_M, fill=MUTED)
    # central zone
    d.rounded_rectangle((455, 205, 1145, 705), radius=60, fill=(239, 246, 255, 210), outline=BLUE + (255,), width=5)
    d.text((735, 236), "DMA-A", font=F_H, fill=BLUE)
    nodes = [(560, 365), (720, 325), (900, 365), (1025, 500), (810, 570), (610, 555)]
    edges = [(0,1), (1,2), (2,3), (3,4), (4,5), (5,0), (1,4), (0,4), (2,4)]
    for a, b in edges:
        d.line((nodes[a], nodes[b]), fill=(132, 174, 219, 255), width=14)
        d.line((nodes[a], nodes[b]), fill=(223, 237, 253, 255), width=5)
    for i, (x, y) in enumerate(nodes, 1):
        d.ellipse((x - 22, y - 22, x + 22, y + 22), fill=(255, 255, 255, 255), outline=BLUE + (255,), width=3)
        tc(d, (x - 20, y - 18, x + 20, y + 18), f"J{i}", F_XS, INK)
    # devices
    d.rounded_rectangle((325, 330, 445, 400), radius=14, fill=CYAN + (255,))
    tc(d, (325, 330, 445, 400), "入口\n流量", F_XS, (255, 255, 255))
    d.line((445, 365, 538, 365), fill=CYAN + (255,), width=7)
    for box, text, color in [
        ((690, 270, 820, 320), "P1高点", ORANGE),
        ((938, 535, 1068, 585), "P2末梢", ORANGE),
        ((1060, 445, 1185, 495), "边界阀", GREEN),
        ((760, 460, 885, 525), "疑似\n漏点", REDP),
    ]:
        d.rounded_rectangle(box, radius=14, fill=color + (255,))
        tc(d, box, text, F_XS, (255, 255, 255))
    # side conditions, no arrows
    checks = [
        ("边界封闭", "阀门状态可核验\n无隐性连通", BLUE),
        ("计量闭合", "入口/出口可计量\n低流量段可靠", CYAN),
        ("压力可观测", "高点、低点、末梢\n至少有代表点", ORANGE),
        ("台账一致", "GIS、SCADA、工单\n编码能关联", GREEN),
        ("反馈闭环", "核查结果回填\n形成训练标签", REDP),
    ]
    for i, (a, b, color) in enumerate(checks):
        x = 70 if i < 3 else 1220
        y = 220 + (i % 3) * 150 if i < 3 else 300 + (i - 3) * 180
        rounded(d, (x, y, x + 260, y + 105), color, 3, radius=22)
        d.text((x + 24, y + 18), a, font=F_S, fill=color)
        tc(d, (x + 22, y + 48, x + 238, y + 100), b, F_XS, MUTED)
    return save(img, "v3_02_dma_blueprint.png")


def fig_mnf_formula():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "MNF不是漏损量：需要拆分与校核", font=F_XL, fill=INK)
    d.text((72, 105), "夜间最小流量是漏损判断入口，但必须扣除合法夜间用水并校验压力和边界", font=F_M, fill=MUTED)
    # chart
    left, top, right, bottom = 120, 230, 920, 690
    d.rectangle((left, top, right, bottom), fill=(255, 255, 255, 230), outline=GRID + (255,), width=2)
    for i in range(5):
        y = top + i * (bottom - top) / 4
        d.line((left, y, right, y), fill=GRID + (180,), width=1)
    pts = []
    random.seed(4)
    for i in range(96):
        hour = i / 4
        v = 0.45 + 0.25 * math.sin((hour - 7) / 24 * 2 * math.pi) + 0.07 * math.sin(hour / 24 * 6 * math.pi)
        if 2 <= hour <= 4:
            v = 0.24 + 0.01 * math.sin(i)
        v += random.uniform(-0.015, 0.015)
        x = left + i * (right - left) / 95
        y = bottom - v * (bottom - top)
        pts.append((x, y))
    d.line(pts, fill=BLUE + (255,), width=5)
    x1 = left + 8 / 96 * (right - left)
    x2 = left + 16 / 96 * (right - left)
    d.rectangle((x1, top, x2, bottom), fill=(252, 211, 77, 95))
    d.text((x1 + 14, top + 18), "MNF\n窗口", font=F_XS, fill=ORANGE)
    mnfy = bottom - 0.24 * (bottom - top)
    d.line((left, mnfy, right, mnfy), fill=REDP + (255,), width=3)
    d.text((right + 16, mnfy - 15), "MNF", font=F_S, fill=REDP)
    # formula blocks
    formula = [
        ("入口MNF", BLUE),
        ("合法夜间用水", CYAN),
        ("背景漏损", REDP),
        ("突发漏损", ORANGE),
        ("仪表/边界误差", MUTED),
    ]
    x = 1010
    for i, (name, color) in enumerate(formula):
        y = 250 + i * 82
        d.rounded_rectangle((x, y, x + 370, y + 54), radius=18, fill=color + (45,), outline=color + (230,), width=2)
        tc(d, (x, y, x + 370, y + 54), name, F_S, INK)
        if i in [0, 1]:
            sign = "=" if i == 0 else "-"
        elif i < len(formula) - 1:
            sign = "+"
        else:
            sign = "±"
        d.text((x - 54, y + 10), sign, font=F_H, fill=MUTED)
    d.rounded_rectangle((1010, 700, 1470, 790), radius=18, fill=(245, 248, 252, 240), outline=GRID + (255,), width=2)
    tc(d, (1030, 700, 1450, 790), "PPT讲法：MNF用于发现“异常水量”，AI用于判断异常是否持续、是否符合漏损模式", F_XS, INK)
    return save(img, "v3_03_mnf_formula.png")


def fig_lstm_residual():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "LSTM/GRU：学习正常曲线，再看残差", font=F_XL, fill=INK)
    d.text((72, 105), "更适合做动态基线和早期预警，不宜单独承诺直接给出漏点坐标", font=F_M, fill=MUTED)
    left, top, right, bottom = 120, 230, 1120, 700
    d.rectangle((left, top, right, bottom), fill=(255, 255, 255, 230), outline=GRID + (255,), width=2)
    for i in range(6):
        y = top + i * (bottom - top) / 5
        d.line((left, y, right, y), fill=GRID + (170,), width=1)
    actual, pred = [], []
    random.seed(8)
    for i in range(140):
        base = 0.42 + 0.16 * math.sin(i / 12) + 0.05 * math.sin(i / 3.5)
        leak = 0 if i < 88 else 0.16 + 0.0015 * (i - 88)
        x = left + i * (right - left) / 139
        pred.append((x, bottom - base * (bottom - top)))
        actual.append((x, bottom - (base + leak + random.uniform(-0.015, 0.015)) * (bottom - top)))
    d.line(pred, fill=GREEN + (255,), width=5)
    d.line(actual, fill=REDP + (255,), width=5)
    sx = left + 88 * (right - left) / 139
    d.rectangle((sx, top, right, bottom), fill=(254, 226, 226, 80))
    d.text((sx + 20, top + 22), "持续正残差区", font=F_S, fill=REDP)
    # legend and residual cards
    d.rounded_rectangle((1185, 260, 1490, 430), radius=20, fill=(255, 255, 255, 240), outline=GRID + (255,), width=2)
    d.line((1215, 305, 1285, 305), fill=GREEN + (255,), width=5)
    d.text((1302, 288), "预测正常值", font=F_XS, fill=INK)
    d.line((1215, 355, 1285, 355), fill=REDP + (255,), width=5)
    d.text((1302, 338), "实际监测值", font=F_XS, fill=INK)
    cards = [("输入", "过去N小时\n流量+压力+日历", BLUE), ("输出", "未来正常值\n预测区间", GREEN), ("判断", "残差持续超阈\n触发复核", REDP)]
    for i, (a, b, c) in enumerate(cards):
        x = 1185
        y = 475 + i * 90
        rounded(d, (x, y, x + 305, y + 66), c, 2, radius=18)
        d.text((x + 22, y + 13), a, font=F_S, fill=c)
        tc(d, (x + 105, y + 8, x + 285, y + 58), b, F_XS, INK)
    return save(img, "v3_04_lstm_residual.png")


def fig_algorithm_matrix():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "漏损检测算法选型矩阵", font=F_XL, fill=INK)
    d.text((72, 105), "不要按“先进程度”选模型，按数据条件、标签质量和业务目标选", font=F_M, fill=MUTED)
    headers = ["任务", "优先算法", "输入要求", "输出给业务"]
    colw = [230, 380, 420, 420]
    x0, y0, rowh = 75, 190, 76
    colors = [BLUE, CYAN, GREEN, ORANGE]
    x = x0
    for h, cw, color in zip(headers, colw, colors):
        d.rectangle((x, y0, x + cw, y0 + rowh), fill=color + (240,))
        tc(d, (x, y0, x + cw, y0 + rowh), h, F_S, (255, 255, 255))
        x += cw
    rows = [
        ["正常基线", "移动基线 / Prophet / LSTM / GRU", "连续入口流量、关键压力、日历特征", "预测值、预测区间、残差"],
        ["少标签异常", "孤立森林 / DBSCAN / 自编码器", "正常样本多，漏损标签少", "异常分数、异常时段、复核清单"],
        ["管段风险", "随机森林 / GBDT / 逻辑回归", "管龄、材质、压力、维修工单、投诉", "风险排序、改造优先级"],
        ["定位收敛", "RF/GBDT + 水力仿真 / GNN", "拓扑可靠，多点压力，可生成仿真样本", "候选片区、候选管段TopN"],
        ["协同解释", "知识库 + 大模型智能体", "模型结果、GIS、工单接口、权限控制", "日报、问答、派单建议、复盘材料"],
    ]
    for r, row in enumerate(rows):
        y = y0 + rowh + r * rowh
        fill = (255, 255, 255, 235) if r % 2 == 0 else (247, 250, 253, 235)
        x = x0
        for txt, cw in zip(row, colw):
            d.rectangle((x, y, x + cw, y + rowh), fill=fill, outline=GRID + (255,), width=1)
            tc(d, (x + 10, y + 5, x + cw - 10, y + rowh - 5), txt, F_XS, INK)
            x += cw
    d.rounded_rectangle((180, 775, 1420, 835), radius=18, fill=(245, 248, 252, 235), outline=GRID + (255,), width=2)
    tc(d, (200, 775, 1400, 835), "对外口径：文献准确率只能说明研究条件下的表现，项目承诺必须基于本地DMA数据验证", F_XS, REDP)
    return save(img, "v3_05_algorithm_matrix.png")


def fig_fusion():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "机理-数据融合：不是替代，而是互补", font=F_XL, fill=INK)
    d.text((72, 105), "机理模型提供物理约束和仿真样本，AI模型提供识别、排序和泛化能力", font=F_M, fill=MUTED)
    areas = [
        ((105, 245, 505, 665), "水力机理模型", ["拓扑/管径/高程", "泵阀边界", "压力-流量响应", "漏点情景仿真"], BLUE),
        ((600, 245, 1000, 665), "融合层", ["仿真样本库", "特征工程", "参数校准", "置信度校验"], GREEN),
        ((1095, 245, 1495, 665), "AI数据模型", ["LSTM残差", "异常检测", "RF/GBDT排序", "GNN拓扑学习"], ORANGE),
    ]
    for box, title, items, color in areas:
        rounded(d, box, color, 4, radius=30)
        d.text((box[0] + 35, box[1] + 28), title, font=F_M, fill=color)
        for i, item in enumerate(items):
            soft_label(d, (box[0] + 38, box[1] + 105 + i * 72, box[2] - 38, box[1] + 155 + i * 72), item, color, F_S)
    # no arrows, use bridge labels
    d.rounded_rectangle((505, 415, 600, 495), radius=12, fill=(245, 248, 252, 240), outline=GRID + (255,), width=2)
    tc(d, (505, 415, 600, 495), "校准\n样本", F_XS, MUTED)
    d.rounded_rectangle((1000, 415, 1095, 495), radius=12, fill=(245, 248, 252, 240), outline=GRID + (255,), width=2)
    tc(d, (1000, 415, 1095, 495), "识别\n反馈", F_XS, MUTED)
    d.rounded_rectangle((300, 735, 1300, 810), radius=20, fill=(245, 248, 252, 235), outline=GRID + (255,), width=2)
    tc(d, (320, 735, 1280, 810), "工程输出：异常等级 + 可能原因 + 候选管段TopN + 现场复核建议", F_S, INK)
    return save(img, "v3_06_hydraulic_ai_fusion.png")


def fig_data_governance():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "模型上线前的数据治理清单", font=F_XL, fill=INK)
    d.text((72, 105), "数据不稳，模型结论一定不稳；治理对象要落到表和字段", font=F_M, fill=MUTED)
    groups = [
        ("设备表", ["设备编码", "安装位置", "量程精度", "在线状态"], BLUE),
        ("管网表", ["管段编号", "管径材质", "阀门状态", "DMA归属"], GREEN),
        ("时序表", ["采样时间", "流量压力", "缺失标记", "清洗版本"], CYAN),
        ("工单表", ["报警编号", "漏点位置", "原因分类", "处置结果"], ORANGE),
        ("标签表", ["真漏/误报", "异常类型", "影响范围", "模型版本"], REDP),
    ]
    for i, (title, fields, color) in enumerate(groups):
        x = 90 + (i % 3) * 500
        y = 220 + (i // 3) * 270
        rounded(d, (x, y, x + 390, y + 205), color, 4, radius=25)
        d.text((x + 28, y + 25), title, font=F_M, fill=color)
        for j, f in enumerate(fields):
            d.text((x + 45, y + 78 + j * 31), f"· {f}", font=F_S, fill=INK)
    d.rounded_rectangle((610, 590, 1510, 765), radius=25, fill=(245, 248, 252, 235), outline=GRID + (255,), width=2)
    tc(d, (640, 605, 1480, 750), "关键原则：SCADA、GIS、营收、工单里的同一对象必须能关联；否则AI只能发现异常，无法形成可执行定位和复盘", F_S, INK)
    return save(img, "v3_07_data_governance.png")


def fig_implementation():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "实施路径：从试点到规模化", font=F_XL, fill=INK)
    d.text((72, 105), "先跑通一个可复盘的DMA，再复制到全网", font=F_M, fill=MUTED)
    phases = [
        ("01", "诊断", "DMA边界、仪表、GIS、工单\n现状体检"),
        ("02", "治理", "统一编码、时间对齐、缺失补齐\n标签口径"),
        ("03", "建模", "MNF基线、LSTM残差、异常检测\n初版规则"),
        ("04", "试点", "灰度报警、人工复核、误报归因\n阈值校准"),
        ("05", "闭环", "工单派发、维修回填、模型再训练\n运营机制"),
        ("06", "推广", "跨DMA迁移、分级运营、月度评估\n标准化复制"),
    ]
    for i, (num, title, desc) in enumerate(phases):
        x = 85 + i * 250
        y = 300 if i % 2 == 0 else 420
        color = [BLUE, CYAN, GREEN, ORANGE, REDP, PURPLE][i]
        d.ellipse((x, y, x + 86, y + 86), fill=color + (255,))
        tc(d, (x, y, x + 86, y + 86), num, F_M, (255, 255, 255))
        d.text((x - 5, y + 112), title, font=F_M, fill=color)
        tc(d, (x - 60, y + 154, x + 160, y + 245), desc, F_XS, INK)
        if i < len(phases) - 1:
            d.line((x + 95, y + 43, x + 240, 360 + ((i + 1) % 2) * 120 + 43), fill=GRID + (255,), width=4)
    d.rounded_rectangle((305, 735, 1295, 810), radius=20, fill=(245, 248, 252, 235), outline=GRID + (255,), width=2)
    tc(d, (325, 735, 1275, 810), "试点选择：边界清晰、仪表稳定、历史事件较多、运维配合度高的DMA", F_S, INK)
    return save(img, "v3_08_implementation_roadmap.png")


def fig_case_story():
    img = canvas()
    d = ImageDraw.Draw(img)
    d.text((70, 45), "教学案例页素材：一次漏损事件如何讲", font=F_XL, fill=INK)
    d.text((72, 105), "用于PPT案例章节，不虚构城市和效果数字，只展示可复用叙事框架", font=F_M, fill=MUTED)
    cards = [
        ("1 异常出现", "MNF连续抬升\n实际值偏离预测区间", BLUE),
        ("2 模型识别", "LSTM残差持续为正\n无监督模型给高分", ORANGE),
        ("3 定位收敛", "压力点响应+拓扑\n输出候选管段TopN", GREEN),
        ("4 现场复核", "听漏/阀门/压力试验\n确认真实原因", CYAN),
        ("5 复盘学习", "工单结果回填\n模型阈值再校准", REDP),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 95 + i * 295
        y = 280
        rounded(d, (x, y, x + 245, y + 330), color, 4, radius=28)
        d.text((x + 28, y + 28), title, font=F_M, fill=color)
        tc(d, (x + 25, y + 105, x + 220, y + 220), body, F_S, INK)
        d.rounded_rectangle((x + 58, y + 250, x + 187, y + 292), radius=14, fill=color + (45,), outline=color + (255,), width=2)
        tc(d, (x + 58, y + 250, x + 187, y + 292), "PPT一页", F_XS, color)
    return save(img, "v3_09_case_storyboard.png")


def make_figs():
    return [
        ("DMA + AI漏损检测总体架构", fig_arch(), "总览页，解释DMA、数据、模型、业务闭环的层次关系。"),
        ("DMA分区建设蓝图", fig_dma_blueprint(), "说明边界、计量、压力点、阀门和疑似漏点的关系。"),
        ("MNF夜间最小流量拆分逻辑", fig_mnf_formula(), "讲清MNF不是漏损量，需扣除合法夜间用水和边界误差。"),
        ("LSTM/GRU残差预警逻辑", fig_lstm_residual(), "展示预测值、实际值、残差和异常区间。"),
        ("漏损检测算法选型矩阵", fig_algorithm_matrix(), "用于算法章节，按任务和数据条件选择模型。"),
        ("机理-数据融合路线", fig_fusion(), "说明水力模型与AI模型互补。"),
        ("数据治理清单", fig_data_governance(), "用于实施章节，强调表结构和字段治理。"),
        ("试点到规模化实施路径", fig_implementation(), "用于工程落地章节。"),
        ("教学案例叙事框架", fig_case_story(), "用于后续制作案例PPT。"),
    ]


def set_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def doc_style(doc):
    sec = doc.sections[0]
    sec.top_margin = Inches(0.72)
    sec.bottom_margin = Inches(0.72)
    sec.left_margin = Inches(0.82)
    sec.right_margin = Inches(0.82)
    for name in ["Normal", "List Bullet"]:
        st = doc.styles[name]
        st.font.name = "宋体"
        st._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
        st.font.size = Pt(10.5)
    for name in ["Heading 1", "Heading 2", "Heading 3", "Heading 4"]:
        st = doc.styles[name]
        st.font.name = "黑体"
        st._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        st.font.color.rgb = DOC_BLACK


def rs(run, red=True, bold=False, size=10.5, font="宋体"):
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = DOC_RED if red else DOC_BLACK


def add_title(doc, title, subtitle):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(title)
    rs(r, False, True, 22, "黑体")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(subtitle)
    rs(r, True, False, 11)
    doc.add_paragraph()


def h(doc, text, level):
    p = doc.add_heading(text, level=level)
    for r in p.runs:
        rs(r, False, True, 16 if level == 1 else 13 if level == 2 else 11.5, "黑体")


def para(doc, text, red=True):
    p = doc.add_paragraph()
    p.paragraph_format.first_line_indent = Pt(21)
    p.paragraph_format.line_spacing = 1.25
    r = p.add_run(text)
    rs(r, red)
    return p


def bullets(doc, items, red=True):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.line_spacing = 1.15
        r = p.add_run(item)
        rs(r, red)


def table(doc, headers, rows):
    t = doc.add_table(rows=1, cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    for i, head in enumerate(headers):
        c = t.rows[0].cells[i]
        c.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        set_shading(c, "D9EAF7")
        p = c.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(head)
        rs(r, False, True, 9.2, "黑体")
    for row in rows:
        cells = t.add_row().cells
        for i, text in enumerate(row):
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            p = cells[i].paragraphs[0]
            p.paragraph_format.line_spacing = 1.05
            r = p.add_run(text)
            rs(r, True, False, 8.8)
    doc.add_paragraph()


def picture(doc, fig_path, caption):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(fig_path), width=Inches(6.7))
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(caption)
    rs(r, True, False, 9)


def make_doc(figs):
    doc = Document()
    doc_style(doc)
    add_title(doc, "AI模型在供水管网DMA系统漏损检测中的应用", "v3完善版：以《AI供水管网漏损检测研究》为素材库，结合公开搜调重写；红色文字为补充内容")

    para(doc, "写作说明：本版吸收了《AI供水管网漏损检测研究》中关于LSTM/GRU、孤立森林、DBSCAN、自编码器、遗传算法、随机森林、梯度提升树、GNN、边缘计算和数据治理的素材，但对其中未经核验的百分比、具体城市效果和营销化表达进行了降级处理。对外演讲中不建议把文献条件下的准确率直接说成本项目承诺。", True)
    para(doc, "本版以原提纲为主线进行扩写：第一部分讲背景和AI发展，第二部分讲DMA架构与核心模型，第三部分讲应用场景，第四部分讲实施路径，第五部分讲案例和PPT转化。内容写成讲稿型，便于后续拆分为PPT页面、讲者备注和宣传材料。", True)

    h(doc, "第一部分 模型应用背景", 1)
    h(doc, "1.1 供水管网漏损管控的国家战略与行业刚需", 2)
    para(doc, "供水管网漏损控制已经从单纯的运维降本事项，升级为水资源节约、韧性城市建设、供水安全保障和水务企业精细化经营的共同议题。住建部、国家发展改革委关于加强公共供水管网漏损控制的通知提出，全国城市公共供水管网漏损率力争控制在9%以内，并强调分区计量、压力调控、老旧管网改造和智能化管理。这个政策导向决定了漏损控制不能只依赖阶段性突击检漏，而需要形成长期在线、可量化、可闭环的管理体系。", True)
    para(doc, "在实际运营中，漏损带来的影响不是单一水量损失。漏损会造成制水药耗、电耗和输配能耗浪费，抬高无收益水量；会改变局部压力状态，诱发爆管、道路沉陷和二次污染风险；会增加抢修、投诉和开挖协调压力；也会让新增供水能力被无效水量占用。对老城区而言，漏损治理还与管网更新、道路施工、用户体验和营商环境相关。", True)
    para(doc, "DMA系统是漏损精细化管控的基础。它通过边界阀、入口流量计、必要出口计量和压力监测，把大管网拆成可计量、可比较、可定位的管理单元。AI模型的价值，是在DMA数据基础上建立动态基线、识别异常模式、收敛候选区域、辅助派单复核，并将每一次处置结果沉淀为下一轮模型迭代的数据资产。", True)

    h(doc, "1.2 供水管网机理模型的发展与局限", 2)
    para(doc, "供水管网机理模型以水力学方程为基础，核心包括质量守恒、能量守恒、管道阻力、节点需水、泵阀边界和水池水位等。工程算法经历了哈代-克罗斯平差、牛顿-拉夫逊迭代、特征线法瞬变分析等阶段。EPANET等软件可对有压供水管网进行延时段水力和水质模拟，输出管段流量、节点压力、水池水位、水龄和水质组分等结果。", True)
    para(doc, "机理模型在规划设计、压力分区、泵阀调度、供水安全校核和水龄分析中仍然不可替代。但在DMA漏损检测中，它依赖高质量基础数据：管径、材质、粗糙系数、阀门状态、高程、用户需水曲线、仪表同步性、边界连通关系都必须可靠。任何基础误差都可能沿水力计算链路放大，导致夜间最小流量拟合失真、漏损量估算偏差和定位误判。", True)
    para(doc, "因此，AI与机理模型不是替代关系。更合理的框架是：机理模型提供物理约束、拓扑关系和仿真样本；AI模型负责从高频监测数据中学习正常模式、识别异常、排序候选管段；业务系统负责派单、复核和回填。只有三者结合，才可能从“模型演示”走向“漏损治理效果”。", True)

    h(doc, "1.3 供水管网AI模型应用的发展（详细）", 2)
    picture(doc, figs[0][1], "图1 DMA + AI漏损检测总体架构")
    para(doc, "人工智能的发展可以概括为四个阶段：符号智能、统计学习、深度学习和大模型/智能体。符号智能依靠专家规则和逻辑推理，适合表达明确规则；统计学习从样本中学习分类、回归和聚类规律；深度学习通过神经网络学习复杂非线性特征，LSTM、GRU、CNN、自编码器等模型推动了时序和多源信号分析；大模型和智能体则增强了知识问答、报告生成、工具调用和跨系统协同能力。", True)
    para(doc, "供水管网AI应用并不是从大模型开始，而是从分区计量和运行数据积累开始。早期系统主要依靠固定阈值，例如入口流量超过上限、夜间最小流量超过经验值、压力低于设定值即报警。随后进入统计分析阶段，使用同比、环比、移动平均、季节分解和控制图识别异常。再往后进入机器学习阶段，模型开始综合流量、压力、天气、节假日、用户类型、管龄、材质、维修工单等多源变量。当前趋势是时序深度学习、拓扑图模型、机理-数据融合和业务智能体协同。", True)
    para(doc, "在DMA漏损检测中，AI模型主要处理四类问题。第一类是预测：学习每个DMA在不同日期、时段、季节、天气和用户结构下的正常流量或压力。第二类是异常识别：判断实际曲线是否持续偏离正常基线，偏离是否符合漏损特征。第三类是定位收敛：结合压力点响应、管网拓扑、水力仿真和历史事件，把异常从DMA级收敛到片区或候选管段。第四类是业务决策：输出预警等级、巡检优先级、复核方式和复盘建议。", True)
    para(doc, "算法的发展也呈现从单点阈值到多模型组合的趋势。LSTM/GRU适合做动态基线预测；孤立森林、DBSCAN、自编码器适合少标签异常检测；随机森林、GBDT适合管段风险排序和漏点区域分类；遗传算法适合水力模型参数校准和DMA分区优化；GNN适合在高质量拓扑和多点压力条件下学习空间传播关系；大模型智能体适合把模型结果转化为管理语言、报告和派单建议。", True)

    h(doc, "1.4 传统机理模型在DMA系统漏损管控中的痛点", 2)
    bullets(doc, [
        "基础数据误差的传导放大：管径、粗糙系数、阀门开度、节点高程、需水模式等小偏差，会导致压力和流量仿真偏移。",
        "拓扑结构失真：实际管网存在历史资料缺失、阀门误开误关、临时连通、旁通管未入库等问题，导致模型边界与现场边界不一致。",
        "夜间用水复杂：MNF中包含合法夜间用水、二供补水、商业用水、消防或市政用水，不能直接等同于漏损量。",
        "压力-漏损耦合难：压力变化会影响背景漏损，但压力控制又受服务水压、消防保障和用户体验约束。",
        "模型维护成本高：管网改造、换管、新用户接入、设备更换和阀门调整都会改变模型参数，若缺少持续维护机制，模型很快失真。"
    ], True)

    h(doc, "1.5 人工智能模型破解问题的核心价值", 2)
    para(doc, "AI模型的第一项价值是建立动态基线。传统阈值通常是静态的，而DMA运行曲线受季节、日期、天气、节假日、用户结构和压力制度影响。LSTM/GRU等模型可以学习每个DMA自己的正常曲线，从而减少统一阈值带来的误报。", True)
    para(doc, "第二项价值是处理多源非线性关系。漏损并不是单变量问题，流量、压力、管龄、材质、泵阀运行、历史工单、道路施工和用户行为共同影响结果。机器学习模型可以在一定程度上学习这些变量之间的耦合关系，辅助识别“像漏损”的异常模式。", True)
    para(doc, "第三项价值是沉淀经验。传统漏损排查依赖个人经验，AI系统通过报警、派单、现场核查、维修结果和误报原因回填，把经验变成可复用的数据标签和模型规则。", True)
    para(doc, "第四项价值是形成闭环。一个成熟系统不只是报警，还应输出异常等级、异常证据、可能原因、候选管段、建议复核方式和复盘动作。只有模型输出能进入工单，工单结果能回到模型，AI才真正参与漏损治理。", True)

    h(doc, "第二部分 技术基础", 1)
    h(doc, "2.1 DMA架构与数据端部署", 2)
    picture(doc, figs[1][1], "图2 DMA分区建设蓝图")
    para(doc, "DMA的核心是边界和水量闭合。一个DMA通常通过边界阀或边界控制形成相对独立区域，并对进入和离开的水量进行计量。公开DMA研究通常把DMA规模、入口数量、边界管段、压力稳定性、施工成本和运维便利性作为重要评价因素。对于AI模型而言，DMA不是简单的地图区域，而是模型训练和业务处置的最小管理单元。", True)
    para(doc, "一个适合建模的DMA至少需要五类条件：边界封闭可核验、入口出口可计量、压力点可观测、GIS和SCADA编码能关联、工单结果能回填。若边界不清或入口计量不准，AI模型会把边界误差误判为漏损；若压力点布设不足，模型只能判断DMA异常，难以定位候选区域；若工单不能回填，模型无法持续学习。", True)
    para(doc, "压力点布设建议覆盖高点、低点、末梢和水力敏感区域。入口流量回答“是否异常”，压力点回答“异常可能影响哪里”。在预算有限时，应优先布设能最大区分漏点位置的压力点，而不是简单追求传感器数量。可以利用水力模型仿真不同漏点情景，比较不同布点方案对定位的贡献。", True)

    h(doc, "2.2 MNF夜间最小流量与AI动态基线", 2)
    picture(doc, figs[2][1], "图3 MNF夜间最小流量拆分逻辑")
    para(doc, "MNF是DMA漏损分析的基础指标。凌晨低用水时段，合法用水较少，漏损在总流量中的比例更高，因此更容易观察异常。但MNF不是漏损量本身，它包含合法夜间用水、背景漏损、突发漏损、仪表误差和边界误差。对外演讲时必须讲清这一点，否则容易把传统指标讲得过于绝对。", True)
    para(doc, "传统MNF分析通常先确定夜间窗口，再估计合法夜间用水，结合压力和历史基线判断异常。AI的改进在于：它不只看一个夜间最小点，而是学习全天和多天的动态基线；不只看入口流量，还能引入压力点、天气、节假日和用户类型；不只输出是否超阈值，还能输出异常持续性、残差幅度和相似历史事件。", True)

    h(doc, "2.3 时序预测模型：LSTM/GRU/CNN-LSTM", 2)
    picture(doc, figs[3][1], "图4 LSTM/GRU残差预警逻辑")
    para(doc, "LSTM通过输入门、遗忘门和输出门处理长序列依赖，适合学习供水流量和压力的日周期、周周期、季节性和短期扰动。GRU结构更轻，参数更少，适合算力受限或需要快速迭代的场景。CNN-LSTM或CNN-GRU可先用一维卷积提取多传感器局部特征，再用循环网络处理时间依赖，适用于多压力点、多流量点联合建模。", True)
    para(doc, "在DMA漏损预警中，LSTM/GRU更适合作为“正常曲线预测器”，而不是直接作为“漏点定位器”。典型流程是：输入过去N小时或N天的流量、压力、天气、节假日和日历特征，输出未来短时窗口的正常值或预测区间；实际值持续偏离预测区间时形成残差信号；残差信号再与MNF、压力响应和工单历史结合，形成漏损风险判断。", True)
    para(doc, "模型训练应注意时间切分，不能随机切分训练集和测试集，否则容易高估效果。更稳妥的做法是用较早时段训练，用后续时段验证，并保留真实漏损事件做事件级检验。评价指标不应只看RMSE，还应看提前预警时间、误报率、漏报率、异常持续识别能力和业务可解释性。", True)

    h(doc, "2.4 无监督异常检测：孤立森林、DBSCAN、自编码器", 2)
    para(doc, "水务项目常见问题是漏损标签不足。历史工单可能只有位置，没有准确发生时间；可能记录了维修原因，但没有关联到DMA和传感器；也可能存在大量误报和非漏损事件。此时，孤立森林、DBSCAN和自编码器等无监督或半监督方法具有实际价值。", True)
    para(doc, "孤立森林通过随机切分特征空间识别“少且不同”的样本，适合多维运行特征的快速异常筛查。DBSCAN通过密度聚类识别低密度离群点，适合发现夜间流量和压力模式中的异常簇。自编码器用正常样本训练重构能力，当输入异常模式时会产生较大重构误差，适合多变量联合异常识别。", True)
    para(doc, "这类模型的工程使用方式应是“发现候选异常”，而不是“自动确认漏损”。建议把异常结果推送给运维人员复核，并让复核结果反向标注为真实漏损、合法用水、设备故障、边界异常或未知原因。随着标签积累，再逐步训练监督模型。", True)

    h(doc, "2.5 监督学习与风险排序：RF、GBDT、HGB、逻辑回归", 2)
    para(doc, "当工单、管网资产和历史事件标签较完整时，可以引入监督学习模型进行风险排序和候选管段识别。随机森林具有抗过拟合和解释性较好的特点，适合处理管龄、材质、历史维修、压力等级和道路环境等结构化变量。GBDT和直方图梯度提升树适合处理非线性特征组合，在管段风险排序和区域分类中常见。逻辑回归虽然简单，但在样本量有限、需要解释系数时仍有价值。", True)
    para(doc, "监督学习模型的输出应面向业务，而不是只给算法分数。建议输出管段风险分、主要贡献因素、历史相似事件、建议巡检优先级和复核方式。对于对外教学，可以把树模型讲成“把老师傅经验结构化”：哪些管龄、材质、压力、维修历史和道路条件组合更容易出现漏损。", True)

    h(doc, "2.6 水力模型与AI融合", 2)
    picture(doc, figs[5][1], "图5 机理-数据融合路线")
    para(doc, "融合模型通常有三种方式。第一是特征融合，把水力模型输出的压力、流量、压力敏感性、供水路径、上下游关系作为AI输入特征。第二是样本融合，用EPANET等工具模拟不同漏点、漏量和起始时间下的压力响应，生成训练样本库。第三是约束融合，AI输出候选管段后，再用水力常识、现场可达性和阀门边界进行过滤。", True)
    para(doc, "融合方法能缓解真实漏损样本不足的问题，但前提是水力模型本身要经过校准。如果拓扑、阀门状态、粗糙系数和需水模式不准，仿真样本会把错误规律传递给AI模型。因此建议先用已知事件或压力试验校验仿真响应，再把仿真数据用于训练。", True)

    h(doc, "2.7 算法选型总表", 2)
    picture(doc, figs[4][1], "图6 漏损检测算法选型矩阵")
    table(doc, ["任务", "可选算法", "输入数据", "输出", "落地注意点"], [
        ["正常基线预测", "移动基线、Prophet、LSTM、GRU、CNN-LSTM", "入口流量、压力点、天气、日历、节假日", "预测值、预测区间、残差", "按DMA建模；时间切分验证；不要只看RMSE"],
        ["少标签异常识别", "孤立森林、DBSCAN、自编码器、LSTM-AE", "多维时序特征、MNF特征、压力残差", "异常分数、异常时段", "结果需人工复核；误报原因必须回填"],
        ["管段风险排序", "随机森林、GBDT、HGB、逻辑回归", "管龄、材质、维修、压力、道路、投诉", "风险分、优先级", "标签偏差会影响模型；需给出贡献因素"],
        ["漏点定位收敛", "水力仿真+RF/GBDT、GNN、图信号处理", "拓扑、多点压力、仿真样本、工单坐标", "候选片区、候选管段TopN", "依赖拓扑和压力点质量；不宜输出唯一坐标"],
        ["调度与压力优化", "遗传算法、贝叶斯优化、强化学习", "泵阀状态、压力约束、需水预测", "调压策略、泵阀建议", "必须满足服务水压和安全约束"],
        ["业务协同", "知识图谱、大模型智能体", "模型结果、GIS、工单、知识库", "问答、报告、派单建议", "必须有权限控制和人工确认"],
    ])

    h(doc, "第三部分 核心应用", 1)
    h(doc, "3.1 事前防控：基于AI的DMA智能预警", 2)
    para(doc, "事前防控的目标是让漏损在早期被发现。系统可为每个DMA建立动态基线，持续监测入口流量、夜间最小流量、关键压力点、残差趋势和历史同类日期差异。当异常持续出现且无法由合法用水、节假日或设备故障解释时，系统输出风险等级。", True)
    para(doc, "预警分级建议采用业务口径：关注、核查、派单建议、重点处置。每一级都应对应动作，而不是只改变颜色。关注级进入观察清单；核查级要求调度人员确认近期阀门、施工和设备状态；派单建议级进入巡检任务；重点处置级需要现场复核和管理层关注。", True)

    h(doc, "3.2 事中处置：漏点识别、定位与派单", 2)
    para(doc, "事中处置的关键是缩小排查范围。入口流量异常只能说明DMA整体异常，多点压力响应和管网拓扑才能帮助定位。AI系统应输出候选管段TopN，而不是单一绝对坐标。每个候选管段应包含证据链：异常时段、压力响应、历史工单、管龄材质、模型置信度和建议复核方式。", True)
    para(doc, "派单应与GIS和移动端联动。巡检人员到现场后，可通过听漏、阀门操作、压力试验、用户走访和开挖验证确认原因。现场结果必须结构化回填，包括真实漏损、误报、仪表故障、边界异常、合法用水、施工扰动等分类。", True)

    h(doc, "3.3 事后复盘：漏损事件全生命周期管理", 2)
    para(doc, "事后复盘决定AI系统能否越用越准。每个异常事件都应记录报警时间、DMA编号、异常特征、派单时间、现场核查结果、漏点位置、漏损类型、维修完成时间、估算漏量、误报原因和模型版本。没有复盘标签，模型就无法从错误中学习。", True)
    para(doc, "复盘还应沉淀案例库。案例库不仅服务算法训练，也服务对外汇报和内部培训。典型案例应包含异常曲线、模型判断、候选区域、现场处置、维修前后MNF变化和经验总结。", True)

    h(doc, "3.4 调度类应用：压力优化与泵阀策略", 2)
    para(doc, "压力管理是漏损控制的重要抓手。漏损流量通常与压力相关，高压会增加背景漏损和爆管风险。AI可结合需水预测和压力监测，辅助识别高压区、低压区和压力波动时段，给出泵站启停、阀门调节、分时压力控制建议。", True)
    para(doc, "但压力优化必须受服务水压、消防保障、二供补水和用户体验约束。对外演讲中应避免把调压说成简单降压，更准确的表达是“在满足服务压力前提下，寻找降漏和节能空间”。", True)

    h(doc, "3.5 DMA规划类应用：分区与监测点优化", 2)
    para(doc, "AI可辅助DMA规划和监测点优化。分区规划可综合管网拓扑、地形高差、道路施工条件、用户数量、入口数量、边界阀数量、压力制度和运维半径，生成候选方案并进行多目标评价。监测点优化可通过仿真不同漏点情景，判断哪些压力点对定位贡献最大。", True)

    h(doc, "第四部分 实施路径", 1)
    h(doc, "4.1 数据治理与模型上线前准备", 2)
    picture(doc, figs[6][1], "图7 模型上线前的数据治理清单")
    para(doc, "实施AI漏损检测前，必须先解决数据治理。至少应形成设备表、管网表、时序表、工单表和标签表五类基础表。SCADA、GIS、营收、工单里的同一对象必须能通过统一编码关联。否则模型最多只能发现异常，无法定位、派单和复盘。", True)
    para(doc, "时间对齐尤其重要。流量和压力的采样粒度可能不同，设备上报可能延迟，工单记录可能只有日期没有具体时间。建模前必须统一时间粒度，标记缺失、漂移和异常毛刺。对于LSTM/GRU等时序模型，时间错位会直接造成错误学习。", True)

    h(doc, "4.2 试点建设与灰度运行", 2)
    picture(doc, figs[7][1], "图8 试点到规模化实施路径")
    para(doc, "建议选择边界清晰、仪表稳定、历史漏损事件较多、运维队伍配合度高的DMA作为试点。试点阶段不宜追求全自动派单，而应采用灰度运行：模型报警先进入观察清单，由调度和巡检人员复核；复核结果回填后，再校准阈值和模型。", True)
    para(doc, "试点评价应同时看算法指标和业务指标。算法指标包括MAE、RMSE、精确率、召回率、F1、异常检测延迟等；业务指标包括提前预警时间、误报工单比例、排查范围缩小程度、平均处置时间、维修后MNF回落情况和复盘回填率。", True)

    h(doc, "4.3 长效运营机制", 2)
    para(doc, "AI漏损系统不是交付后结束，而是进入运营期。建议每月复盘误报、漏报和典型事件；每季度检查模型漂移和阈值有效性；每半年根据新增工单和管网改造情况进行再训练；重大管网改造、压力制度调整或传感器更换后，应重新评估模型基线。", True)
    para(doc, "人员能力建设也需要分层。管理层关注指标、投入产出和治理闭环；调度人员关注报警解释、阀泵策略和派单规则；巡检人员关注候选区域、现场复核和标签回填；算法人员关注数据质量、模型漂移、版本管理和评估报告。", True)

    h(doc, "第五部分 典型案例与PPT转化", 1)
    h(doc, "5.1 教学型案例组织方式", 2)
    picture(doc, figs[8][1], "图9 教学案例叙事框架")
    para(doc, "在没有真实案例数据时，不建议虚构具体城市和降漏数字。可采用教学型案例：某DMA连续三天凌晨MNF抬升，LSTM预测残差持续为正，P1与P2压力点出现轻微响应，模型输出候选管段TopN，巡检人员现场复核后确认漏点，维修后MNF回落，工单结果回填为训练标签。", True)
    para(doc, "最终PPT建议按五页讲案例：第一页展示异常曲线，第二页展示LSTM残差，第三页展示DMA拓扑和压力响应，第四页展示候选管段和派单， 第五页展示维修后回落和复盘学习。这样既能体现技术深度，也能让非算法听众理解业务价值。", True)

    h(doc, "5.2 PPT图片材料使用建议", 2)
    para(doc, "本版图片均按PPT素材思路制作：透明底、16:9、减少网页式页眉、尽量避免复杂箭头。建议在正式PPT中把图片作为主体，占页面60%以上，正文只放一句观点和少量标签。算法表格不要塞进一页讲完，可拆成“时序预测”“异常检测”“机理融合”“智能体协同”四页。", True)

    h(doc, "第六部分 总结与展望", 1)
    para(doc, "AI模型在DMA漏损检测中的核心逻辑可以概括为：DMA提供边界，MNF提供低用水窗口，流量压力提供实时信号，LSTM/GRU提供动态基线，无监督模型提供异常筛查，水力模型提供物理约束，工单系统提供闭环标签。任何一个环节缺失，AI都难以真正落地。", True)
    para(doc, "未来方向不是单一大模型替代专业模型，而是专业模型、机理模型、知识库和智能体协同。底层用时序模型、异常检测、树模型、GNN和水力仿真完成专业判断；中层用数据治理和工单系统完成闭环；上层用大模型智能体完成解释、问答、报告和协同。", True)

    h(doc, "参考资料", 1)
    refs = [
        "住房和城乡建设部办公厅、国家发展改革委办公厅. 关于加强公共供水管网漏损控制的通知（建办城〔2022〕2号）.",
        "US EPA. EPANET: Application for Modeling Drinking Water Distribution Systems.",
        "Di Nardo A. et al. Water Network Partitioning into District Metered Areas: A State-Of-The-Art Review. Water, 2020.",
        "Zuniga-Uribe M. et al. Artificial Intelligence in Water Distribution Networks: A Systematic Review of Models, Input Variables, Databases, and Output Strategies for Leak Detection. Smart Cities, 2026.",
        "Hochreiter S., Schmidhuber J. Long Short-Term Memory. Neural Computation, 1997.",
        "Vaswani A. et al. Attention Is All You Need. arXiv:1706.03762, 2017.",
        "Adedeji K. B. et al. Leak detection in water distribution networks: an introductory overview. Smart Water, 2019.",
        "《AI供水管网漏损检测研究.docx》作为算法素材库使用，未将其中未经核验的性能百分比作为项目承诺。",
    ]
    for ref in refs:
        para(doc, ref, True)

    out = DOC_DIR / "AI模型在供水管网DMA系统漏损检测中的应用_完善扩充红色标注版_v3.docx"
    doc.save(out)
    return out


def make_asset_doc(figs):
    doc = Document()
    doc_style(doc)
    add_title(doc, "AI供水管网漏损检测PPT图片材料_v3", "透明底PNG + 使用建议；用于后续正式PPT排版")
    para(doc, "本版图片重做为PPT素材组件，尽量减少复杂箭头，避免网页式长页眉。建议后续在PPT中搭配少量标题和讲者备注使用。", True)
    for i, (name, path, usage) in enumerate(figs, 1):
        h(doc, f"{i}. {name}", 2)
        picture(doc, path, name)
        para(doc, f"文件路径：{path}", False)
        para(doc, f"建议用途：{usage}", False)
    out = DOC_DIR / "AI供水管网漏损检测_PPT图片材料说明_v3.docx"
    doc.save(out)
    return out


def make_ppt(figs):
    prs = Presentation()
    prs.slide_width = PInches(13.333333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for name, path, usage in figs:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PRGBColor(255, 255, 255)
        slide.shapes.add_picture(str(path), PInches(0), PInches(0), width=PInches(13.333333))
    out = PPT_DIR / "AI供水管网漏损检测_PPT图片素材包_v3.pptx"
    prs.save(out)
    return out


def main():
    figs = make_figs()
    doc = make_doc(figs)
    asset_doc = make_asset_doc(figs)
    ppt = make_ppt(figs)
    print(doc)
    print(asset_doc)
    print(ppt)
    for _, p, _ in figs:
        print(p)


if __name__ == "__main__":
    main()
