from pathlib import Path

from pptx import Presentation


PPT = Path(r"E:\PY\research\0425\生产流程.pptx")
prs = Presentation(str(PPT))

print("slides", len(prs.slides), "size", prs.slide_width, prs.slide_height)
for si, slide in enumerate(prs.slides, 1):
    print(f"\n--- slide {si} ---")
    bottom_shapes = 0
    for i, shape in enumerate(slide.shapes):
        top = int(shape.top)
        left = int(shape.left)
        width = int(shape.width)
        height = int(shape.height)
        if top > prs.slide_height * 0.86:
            bottom_shapes += 1
        text = ""
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.replace("\n", " | ").strip()
        if text:
            print(i, shape.shape_type, left, top, width, height, repr(text[:240]))
    print("bottom-shape-count", bottom_shapes)
