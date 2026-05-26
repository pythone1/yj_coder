# -*- coding: utf-8 -*-
from pathlib import Path

import build_teaching_deck_v6 as b
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

b.PPTX = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_内容强化版_v8.pptx"
b.AUDIT = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_内容强化版_v8_逐页检查.txt"
b.PREVIEW = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_内容强化版_v8_预览联系表.png"
b.PREVIEW_DIR = OUT / "v8_page_previews"
b.PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
b.AUDIT_LINES.clear()


def clear_and_new():
    prs = Presentation(str(b.TEMPLATE))
    b.clear_slides(prs)
    return prs


def dense_card(s, x, y, w, h, title, lines, color=b.BLUE, size=12.5):
    b.rect(s, x, y, w, h, b.WHITE, color, True)
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x), b.emu(y), b.emu(w), b.emu(0.36)), color)
    b.tb(s, title, x + 0.10, y + 0.11, w - 0.20, 0.14, 13.5, b.WHITE, True, b.PP_ALIGN.CENTER)
    if not lines:
        return
    content_top = y + 0.62
    content_h = h - 0.82
    step = content_h / len(lines)
    item_h = min(0.48, max(0.32, step * 0.72))
    for i, line in enumerate(lines):
        iy = content_top + i * step + max(0, (step - item_h) / 2)
        if h >= 1.25:
            b.rect(s, x + 0.18, iy, w - 0.36, item_h, b.RGBColor(248, 252, 255), b.RGBColor(224, 239, 248), True)
            b.tb(s, line, x + 0.30, iy + 0.10, w - 0.60, item_h - 0.13, size, b.TEXT, False, b.PP_ALIGN.LEFT)
        else:
            b.tb(s, line, x + 0.18, iy + 0.08, w - 0.36, item_h - 0.10, size, b.TEXT)


def cover(prs):
    s = b.blank(prs, b.NAVY)
    b.tb(s, "杨佳负责部分", 0.90, 0.88, 3.20, 0.35, 22, b.CYAN, True)
    b.tb(s, "AI模型、核心算法\n与实施路径", 0.90, 1.45, 4.65, 1.20, 37, b.WHITE, True)
    b.tb(s, "对应提纲：1.3、2.2、4.1-4.5", 0.92, 3.05, 4.20, 0.24, 15.5, b.RGBColor(218, 237, 248))
    for i, (t, c) in enumerate([("AI发展", b.CYAN), ("算法技术栈", b.TEAL), ("工程落地", b.ORANGE)]):
        b.chip(s, t, 0.92 + i * 1.32, 6.25, 1.12, c)
    b.pic(s, "positioning", 5.20, 0.52, 7.85, 6.55, crop=False)
    border = s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, b.emu(5.10), b.emu(0.42), b.emu(8.05), b.emu(6.75))
    border.fill.background()
    b.stroke(border, b.RGBColor(30, 120, 190), 1.2)
    b.audit(s, 1, allow_short=True)


def scope(prs):
    s = b.content_slide(prs, 2, "本部分内容边界：只讲杨佳负责的模型与落地", "按提纲范围组织，不展开国家政策和传统机理模型背景")
    dense_card(s, 0.60, 1.25, 3.80, 4.95, "1.3 AI模型应用发展", [
        "讲AI模型与智能体演进",
        "讲漏损检测可适配算法谱系",
        "讲AI从报警到诊断、派单、复盘的角色变化",
        "输出：AI发展逻辑 + 算法地图"
    ], b.BLUE, 14)
    dense_card(s, 4.75, 1.25, 3.80, 4.95, "2.2 核心AI技术及价值", [
        "时序数据分析与异常检测",
        "机器学习与水力模型融合",
        "大数据融合与智能决策",
        "输出：模型输入、处理逻辑、业务价值"
    ], b.TEAL, 14)
    dense_card(s, 8.90, 1.25, 3.80, 4.95, "4.1-4.5 实施路径", [
        "前期规划与目标边界",
        "数据治理与数据中台",
        "模型选型、训练、验证、优化",
        "系统集成、试点验证、长效运营"
    ], b.ORANGE, 14)
    b.audit(s, 2)


def section(prs, no, part, title, subtitle, color):
    b.section(prs, no, part, title, subtitle, color)


def ai_evolution(prs):
    s = b.content_slide(prs, 4, "1.3 AI模型应用发展：从经验阈值到智能体协同", "这一页用于讲清AI能力为什么会逐步进入DMA漏损检测")
    stages = [
        ("规则阈值", "MNF固定阈值\n人工巡检复核\n问题：误报多、定位粗", b.BLUE),
        ("机器学习", "孤立森林/DBSCAN\n随机森林/提升树\n价值：少标签也能筛查", b.TEAL),
        ("深度学习", "LSTM/GRU\n自编码器\n价值：学习动态基线", b.ORANGE),
        ("机理融合", "水力模型约束\n仿真样本训练\n价值：片区收敛到管段", b.PURPLE),
        ("智能体协同", "知识库+工具调用\n工单复盘+调度建议\n价值：形成闭环工作流", b.RED),
    ]
    for i, (h, body, col) in enumerate(stages):
        x = 0.45 + i * 2.55
        dense_card(s, x, 1.35, 2.20, 3.55, h, body.split("\n"), col, 12.2)
        if i < 4:
            b.line(s, x + 2.23, 3.12, x + 2.47, 3.12, b.BLUE, 2)
    b.rect(s, 0.85, 5.62, 11.72, 0.70, b.PALE2, b.BLUE, True)
    b.tb(s, "讲解口径：AI的发展不是“替代人工”，而是把异常发现、定位证据、派单处置和结果回填串成连续流程。", 1.20, 5.86, 11.05, 0.18, 17, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 4)


def agent_arch(prs):
    s = b.content_slide(prs, 5, "AI模型与智能体：模型负责判断，智能体负责任务组织", "智能体的价值在于把模型能力接入知识库、工具和业务流程")
    b.pic(s, "overall", 0.60, 1.18, 5.15, 4.55, crop=False)
    dense_card(s, 6.05, 1.12, 3.00, 1.45, "模型层", ["LSTM动态基线", "异常检测模型", "定位排序模型"], b.BLUE, 13.2)
    dense_card(s, 9.35, 1.12, 3.00, 1.45, "知识层", ["DMA边界", "管网资产", "历史工单", "维修案例"], b.TEAL, 13.2)
    dense_card(s, 6.05, 3.05, 3.00, 1.45, "工具层", ["GIS查询", "水力仿真", "SCADA取数", "工单生成"], b.ORANGE, 13.2)
    dense_card(s, 9.35, 3.05, 3.00, 1.45, "业务层", ["预警解释", "核查路线", "派单建议", "复盘报告"], b.PURPLE, 13.2)
    b.tb(s, "落地边界：智能体输出建议，不直接替代调度和现场确认；关键动作仍需人工审核。", 1.00, 6.25, 11.40, 0.25, 17.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 5)


def algorithm_map(prs):
    s = b.content_slide(prs, 7, "供水管网漏损检测适配算法：按任务选择，不按名词堆叠", "把算法放回业务任务中讲，听众更容易理解")
    headers = ["任务", "算法/技术", "输入", "输出", "业务价值"]
    widths = [1.35, 2.25, 3.00, 2.20, 2.70]
    rows = [
        ("动态基线", "LSTM/GRU", "历史流量、压力、日历、天气", "预测区间、残差", "识别正常范围"),
        ("异常筛查", "孤立森林", "残差、波动、压力响应", "异常分数", "快速分级预警"),
        ("状态聚类", "DBSCAN", "压力-流量状态点", "离群簇、异常片段", "识别异常运行状态"),
        ("隐性异常", "自编码器", "多变量时序窗口", "重构误差", "发现复杂异常模式"),
        ("模型校核", "遗传算法", "压力实测、粗糙度、需水量", "校核参数", "提升水力模型精度"),
        ("定位排序", "RF/GBT", "拓扑、仿真样本、资产特征", "候选管段TopN", "减少无效巡检"),
    ]
    x0, y0 = 0.50, 1.18
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.50), b.emu(0.38)), b.BLUE)
    for i, h in enumerate(headers):
        b.tb(s, h, x0 + sum(widths[:i]) + 0.04, y0 + 0.12, widths[i] - 0.08, 0.12, 11.5, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.38 + r * 0.68
        b.rect(s, x0, y, 11.50, 0.62, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(222, 238, 248), False)
        for i, txt in enumerate(row):
            col = b.BLUE if i in [0, 1] else b.TEXT
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.05, y + 0.18, widths[i] - 0.10, 0.18, 11.2, col, i < 2, b.PP_ALIGN.CENTER)
    b.tb(s, "课堂重点：先让听众知道“每个算法解决什么问题”，再讲模型结构和工程实现。", 1.00, 6.18, 11.30, 0.24, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 7)


def full_visual(prs, no, key, note=None):
    s = b.blank(prs)
    b.pic(s, key, 0, 0, 13.333, 7.5, crop=False)
    if note:
        b.rect(s, 0.78, 6.78, 11.80, 0.42, b.WHITE, b.RGBColor(210, 232, 245), True, 8)
        b.tb(s, note, 1.05, 6.91, 11.25, 0.14, 13.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no, allow_short=True)


def lstm_detail(prs):
    s = b.content_slide(prs, 10, "2.2.1 LSTM/GRU时序预测：动态基线怎么服务漏损预警", "本地演示可只使用 LSTM 环境，重点讲输入、训练、残差和业务输出")
    b.pic(s, "baseline", 0.52, 1.08, 5.70, 4.85, crop=False)
    dense_card(s, 6.55, 1.10, 2.85, 1.18, "输入特征", ["历史流量/压力", "小时、星期、节假日", "天气、施工、阀门状态"], b.BLUE, 12.3)
    dense_card(s, 9.70, 1.10, 2.85, 1.18, "模型训练", ["按时间切分训练集/验证集", "学习正常用水节律", "输出预测区间"], b.TEAL, 12.3)
    dense_card(s, 6.55, 2.75, 2.85, 1.18, "异常判断", ["实际值 - 预测值 = 残差", "连续偏离触发预警", "结合MNF时段降低误报"], b.ORANGE, 12.3)
    dense_card(s, 9.70, 2.75, 2.85, 1.18, "业务输出", ["异常DMA", "异常等级", "复核建议", "工单触发条件"], b.PURPLE, 12.3)
    dense_card(s, 6.55, 4.40, 6.00, 1.15, "讲解要点", ["动态基线不是固定阈值；它会随季节、日期、用户规律自动变化。", "漏损预警看的是持续偏离，而不是单个点的波动。"], b.RED, 12.5)
    b.audit(s, 10)


def anomaly_detail(prs):
    s = b.content_slide(prs, 11, "2.2.1 异常检测技术：LSTM残差之后怎么继续筛查", "孤立森林、DBSCAN、自编码器适合少标签或弱标签阶段")
    dense_card(s, 0.60, 1.15, 3.85, 4.50, "孤立森林", [
        "逻辑：异常点更容易被随机切分隔离",
        "输入：残差、流量波动、压力响应",
        "输出：异常分数、报警等级",
        "优点：速度快，适合在线初筛",
        "注意：需要结合业务阈值做复核"
    ], b.BLUE, 12.5)
    dense_card(s, 4.75, 1.15, 3.85, 4.50, "DBSCAN聚类", [
        "逻辑：正常状态形成密度簇",
        "输入：压力-流量状态点、多维特征",
        "输出：离群点、异常片段",
        "优点：不需要预设类别数",
        "注意：参数eps需按DMA特性调节"
    ], b.TEAL, 12.5)
    dense_card(s, 8.90, 1.15, 3.85, 4.50, "自编码器", [
        "逻辑：只学习正常模式",
        "输入：多变量时序窗口",
        "输出：重构误差、异常阈值",
        "优点：适合复杂隐性异常",
        "注意：需要稳定的正常样本"
    ], b.ORANGE, 12.5)
    b.tb(s, "组合方式：LSTM先形成动态残差，孤立森林做快速筛查，DBSCAN识别状态簇，自编码器补充隐性模式。", 0.95, 6.28, 11.50, 0.22, 17, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 11)


def hydraulic_ml(prs):
    s = b.content_slide(prs, 13, "2.2.2 机器学习与水力模型融合：从片区锁定到管段定位", "定位问题必须引入管网拓扑、水力响应和现场约束")
    b.pic(s, "positioning", 0.45, 1.15, 5.85, 4.85, crop=False)
    dense_card(s, 6.60, 1.15, 2.75, 1.25, "1 建立样本", ["不同漏点位置", "不同漏量和时段", "仿真压力/流量响应"], b.BLUE, 12.2)
    dense_card(s, 9.70, 1.15, 2.75, 1.25, "2 提取特征", ["压力降幅", "相邻测点响应", "拓扑距离、管龄管材"], b.TEAL, 12.2)
    dense_card(s, 6.60, 2.90, 2.75, 1.25, "3 训练模型", ["随机森林/RF", "梯度提升树/GBT", "学习响应-管段关系"], b.ORANGE, 12.2)
    dense_card(s, 9.70, 2.90, 2.75, 1.25, "4 输出结果", ["候选管段TopN", "置信度", "核查顺序"], b.PURPLE, 12.2)
    dense_card(s, 6.60, 4.65, 5.85, 1.00, "业务价值", ["把“哪个DMA可能漏”进一步收敛为“哪几根管优先查”，减少无效巡检和开挖。"], b.RED, 12.8)
    b.audit(s, 13)


def ga_rf_gbt(prs):
    s = b.content_slide(prs, 14, "2.2.2 GA、RF、GBT在漏损定位中的分工", "这页用于把机器学习与水力模型融合讲实")
    rows = [
        ("遗传算法 GA", "水力模型参数校核", "粗糙度、需水量、背景漏失、压力实测", "校核参数，降低仿真误差", "适合模型校准"),
        ("随机森林 RF", "漏点位置分类", "仿真样本、压力响应、管段属性", "候选管段类别/概率", "抗噪声、解释性较好"),
        ("梯度提升树 GBT", "风险评分与排序", "残差、拓扑、管龄、维修历史", "TopN风险管段", "适合精细排序"),
    ]
    widths = [2.00, 2.05, 3.35, 2.35, 2.10]
    headers = ["算法", "任务", "输入", "输出", "适用特点"]
    x0, y0 = 0.55, 1.20
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.85), b.emu(0.42)), b.ORANGE)
    for i, h in enumerate(headers):
        b.tb(s, h, x0 + sum(widths[:i]), y0 + 0.13, widths[i], 0.13, 12.3, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.48 + r * 1.18
        b.rect(s, x0, y, 11.85, 1.04, b.RGBColor(255, 250, 244) if r % 2 == 0 else b.WHITE, b.RGBColor(238, 219, 198), False)
        for i, txt in enumerate(row):
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.05, y + 0.28, widths[i] - 0.10, 0.32, 12.5, b.ORANGE if i < 2 else b.TEXT, i < 2, b.PP_ALIGN.CENTER)
    b.tb(s, "讲解重点：GA偏参数反演，RF/GBT偏候选定位和风险排序，三者与水力模型配合使用。", 1.05, 6.15, 11.25, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 14)


def bigdata_edge(prs):
    s = b.content_slide(prs, 16, "2.2.3 大数据融合、知识图谱与边云协同", "把模型结果接入资产、工单和调度系统，才是真正可执行的AI")
    b.pic(s, "training", 0.45, 1.08, 5.80, 4.95, crop=False)
    dense_card(s, 6.55, 1.10, 2.80, 1.20, "多源融合", ["SCADA、GIS、DMA", "管网资产、工单", "天气、施工信息"], b.BLUE, 12.2)
    dense_card(s, 9.65, 1.10, 2.80, 1.20, "知识图谱", ["设备-管段-事件关联", "原因追踪", "案例知识沉淀"], b.TEAL, 12.2)
    dense_card(s, 6.55, 2.80, 2.80, 1.20, "边缘计算", ["校时、缓存、初筛", "断网续传", "轻量化推理"], b.ORANGE, 12.2)
    dense_card(s, 9.65, 2.80, 2.80, 1.20, "云端AI", ["模型训练", "全局分析", "跨DMA策略优化"], b.PURPLE, 12.2)
    dense_card(s, 6.55, 4.50, 5.90, 1.00, "业务价值", ["支撑漏损决策、巡检路径优化、改造优先级排序和全生命周期管理。"], b.GREEN, 12.8)
    b.audit(s, 16)


def planning(prs):
    s = b.content_slide(prs, 20, "4.1 前期规划：先定目标和边界，再谈模型", "实施路径不是从算法开始，而是从现状、目标和试点范围开始")
    dense_card(s, 0.65, 1.22, 3.75, 4.55, "现状诊断", [
        "数据资产：SCADA、GIS、DMA、工单",
        "硬件设备：流量计、压力计、水质仪",
        "系统基础：SCADA、GIS、运维平台",
        "漏损现状：NRW、MNF、历史爆管",
        "DMA现状：边界阀、入口计量、监测点"
    ], b.BLUE, 12.4)
    dense_card(s, 4.78, 1.22, 3.75, 4.55, "目标设定", [
        "短期：降低误报、提升发现效率",
        "中期：预警-定位-派单闭环",
        "长期：资产健康评估与改造排序",
        "指标：TopN命中率、派单有效率、处置时长"
    ], b.TEAL, 12.4)
    dense_card(s, 8.90, 1.22, 3.75, 4.55, "实施路径", [
        "选择典型DMA试点",
        "先接入核心数据和核心传感器",
        "先跑通LSTM动态基线和异常复核",
        "再扩展水力融合、风险排序和智能体"
    ], b.ORANGE, 12.4)
    b.audit(s, 20)


def data_governance(prs):
    s = b.content_slide(prs, 21, "4.2 数据治理：AI模型上线前必须打好的底座", "没有可信数据，模型只能输出不可复核的结果")
    headers = ["治理任务", "具体内容", "对模型的作用"]
    widths = [2.20, 5.40, 4.00]
    rows = [
        ("多源归集", "流量、压力、水质、GIS、SCADA、管网属性、工单、天气", "形成完整特征空间"),
        ("统一编码", "设备ID、管段ID、DMA编号、阀门编号统一", "保证跨系统可关联"),
        ("时间对齐", "统一采样粒度、时间戳、补传和断点处理", "避免时序错位"),
        ("质量校验", "缺失、漂移、异常值、重复点、仪表故障识别", "减少误报与误判"),
        ("标签体系", "误报、漏损、施工、阀门调整、维修结果回填", "支撑监督学习和再训练"),
    ]
    x0, y0 = 0.65, 1.18
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.60), b.emu(0.40)), b.TEAL)
    for i, h in enumerate(headers):
        b.tb(s, h, x0 + sum(widths[:i]), y0 + 0.12, widths[i], 0.13, 12.5, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.42 + r * 0.82
        b.rect(s, x0, y, 11.60, 0.75, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(224, 239, 248), False)
        for i, txt in enumerate(row):
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.08, y + 0.22, widths[i] - 0.16, 0.22, 12.2, b.TEAL if i == 0 else b.TEXT, i == 0, b.PP_ALIGN.CENTER)
    b.tb(s, "输出物：设备表、管网表、时序表、工单表、标签表、规则表。", 1.10, 6.18, 11.10, 0.24, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 21)


def model_build(prs):
    s = b.content_slide(prs, 22, "4.3 模型建设：选型、训练、验证、优化", "这一页把模型建设讲成可执行流程")
    steps = [
        ("场景拆解", ["预警、定位、排序、复盘", "明确模型输出要服务哪个动作"], b.BLUE),
        ("模型选型", ["轻量优先、可解释优先", "LSTM + 异常检测 + 水力融合"], b.TEAL),
        ("训练验证", ["按时间切分训练/验证/测试", "历史回放 + 在线灰度 + 现场复核"], b.ORANGE),
        ("版本迭代", ["工单结果回填", "漂移监测、阈值调整、模型发布"], b.PURPLE),
    ]
    for i, (h, lines, col) in enumerate(steps):
        x = 0.70 + i * 3.08
        dense_card(s, x, 1.35, 2.55, 3.95, h, lines, col, 12.6)
        if i < 3:
            b.line(s, x + 2.58, 3.35, x + 2.98, 3.35, b.BLUE, 2)
    b.tb(s, "验收口径：不仅看准确率，还要看TopN命中率、误报率、派单有效率、平均处置时长和闭环回填率。", 0.95, 6.20, 11.50, 0.24, 17.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 22)


def engineering(prs):
    s = b.content_slide(prs, 24, "4.4 工程落地：系统集成、端边云协同与试点验证", "AI模型必须进入业务系统，才能从预测结果变成处置动作")
    b.pic(s, "impl", 0.45, 1.08, 5.80, 4.95, crop=False)
    dense_card(s, 6.55, 1.10, 2.80, 1.20, "硬件适配", ["兼容现有流量计/压力计", "处理离线、补传、校时"], b.BLUE, 12.2)
    dense_card(s, 9.65, 1.10, 2.80, 1.20, "系统集成", ["SCADA、GIS、工单系统", "统一看板与移动端"], b.TEAL, 12.2)
    dense_card(s, 6.55, 2.80, 2.80, 1.20, "端边云协同", ["边缘初筛", "云端训练", "业务端闭环"], b.ORANGE, 12.2)
    dense_card(s, 9.65, 2.80, 2.80, 1.20, "试点验证", ["典型DMA先行", "指标复盘", "形成复制模板"], b.PURPLE, 12.2)
    dense_card(s, 6.55, 4.50, 5.90, 1.00, "落地判断", ["不是证明算法能跑，而是证明预警能复核、定位能派单、结果能回填。"], b.RED, 12.8)
    b.audit(s, 24)


def operation(prs):
    s = b.content_slide(prs, 25, "4.5 长效运营：人员能力与运维保障体系", "系统上线只是开始，长期有效依赖组织机制和模型版本管理")
    rows = [
        ("管理层", "关注漏损率、投资回报、试点推广节奏", "确定目标、预算和考核口径"),
        ("调度人员", "关注预警解释、定位证据、派单优先级", "确认业务动作与处置策略"),
        ("运维人员", "关注现场核查、维修记录、结果回填", "提供高质量标签"),
        ("算法/IT", "关注数据漂移、接口稳定、模型性能", "维护模型和系统版本"),
    ]
    headers = ["角色", "关注内容", "责任重点"]
    widths = [1.70, 5.10, 4.55]
    x0, y0 = 0.80, 1.22
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.35), b.emu(0.42)), b.BLUE)
    for i, h in enumerate(headers):
        b.tb(s, h, x0 + sum(widths[:i]), y0 + 0.13, widths[i], 0.13, 12.8, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.42 + r * 0.88
        b.rect(s, x0, y, 11.35, 0.80, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(224, 239, 248), False)
        for i, txt in enumerate(row):
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.08, y + 0.24, widths[i] - 0.16, 0.20, 12.5, b.BLUE if i == 0 else b.TEXT, i == 0, b.PP_ALIGN.CENTER)
    b.rect(s, 1.00, 5.45, 11.25, 0.82, b.PALE2, b.BLUE, True)
    b.tb(s, "保障机制：传感器在线率监测、接口巡检、数据漂移监控、模型复盘会、版本发布审批。", 1.30, 5.75, 10.60, 0.20, 17, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 25)


def summary(prs):
    s = b.blank(prs, b.NAVY)
    b.tb(s, "杨佳部分最终讲法", 0.90, 0.85, 4.80, 0.45, 34, b.WHITE, True)
    dense_card(s, 0.95, 1.75, 3.70, 3.35, "1.3", ["AI模型应用发展", "从阈值报警到智能体协同", "重点讲算法演进和业务角色变化"], b.BLUE, 13.2)
    dense_card(s, 4.85, 1.75, 3.70, 3.35, "2.2", ["核心AI技术与业务价值", "时序异常、水力融合、大数据决策", "重点讲输入、输出、价值"], b.TEAL, 13.2)
    dense_card(s, 8.75, 1.75, 3.70, 3.35, "4.1-4.5", ["实施路径", "规划、治理、建模、集成、运营", "重点讲项目如何落地"], b.ORANGE, 13.2)
    b.tb(s, "一句话收束：DMA负责缩小范围，AI负责形成定位证据，工单闭环负责让模型持续变准。", 1.20, 6.18, 11.00, 0.24, 18, b.RGBColor(218, 237, 248), True, b.PP_ALIGN.CENTER)
    b.audit(s, 26, allow_short=True)


def build():
    prs = clear_and_new()
    cover(prs)
    scope(prs)
    section(prs, 3, "PART 01", "1.3 AI模型应用发展", "用发展逻辑和算法地图讲清AI为什么适配DMA漏损检测。", b.CYAN)
    ai_evolution(prs)
    agent_arch(prs)
    full_visual(prs, 6, "positioning", "用于讲“DMA宏观锁区 + AI微观定位 + 工单闭环”的整体逻辑。")
    algorithm_map(prs)
    section(prs, 8, "PART 02", "2.2 核心AI技术及业务价值", "按时序异常、水力融合、大数据决策三条技术线展开。", b.TEAL)
    full_visual(prs, 9, "baseline", "用于讲动态基线、残差预警和业务动作之间的关系。")
    lstm_detail(prs)
    anomaly_detail(prs)
    full_visual(prs, 12, "overall", "用于讲数据层、模型层、业务层、反馈层的总体架构。")
    hydraulic_ml(prs)
    ga_rf_gbt(prs)
    full_visual(prs, 15, "training", "用于讲数据治理、训练路径和模型迭代闭环。")
    bigdata_edge(prs)
    section(prs, 17, "PART 03", "第四部分 实施路径", "把AI模型从演示能力变成可运维系统。", b.ORANGE)
    full_visual(prs, 18, "impl", "用于讲从规划、治理、模型建设到长期运营的实施路径。")
    full_visual(prs, 19, "workorder", "用于讲预警、定位、派单、维修、回填的业务闭环。")
    planning(prs)
    data_governance(prs)
    model_build(prs)
    full_visual(prs, 23, "training", "用于讲训练验证、质量提升、定位精度提升和持续迭代。")
    engineering(prs)
    operation(prs)
    summary(prs)
    prs.save(b.PPTX)
    b.AUDIT.write_text("\n".join(b.AUDIT_LINES), encoding="utf-8")
    b.build_previews(prs)
    b.append_repetition_check(prs)
    return b.PPTX


if __name__ == "__main__":
    print(build())
    print(b.AUDIT)
    print(b.PREVIEW)
