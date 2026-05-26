from pathlib import Path
from pptx import Presentation

prs = Presentation(str(Path(r"E:\PY\research\0425\生产流程_最终版_v2.pptx")))


def walk(shapes, prefix=""):
    for i, shape in enumerate(shapes):
        text = ""
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.replace("\n", " | ").strip()
        if text:
            print(prefix + str(i), int(shape.left), int(shape.top), int(shape.width), int(shape.height), repr(text[:300]))
        if hasattr(shape, "shapes"):
            walk(shape.shapes, prefix + str(i) + ".")


for idx in [4]:
    print("SLIDE", idx + 1)
    walk(prs.slides[idx].shapes)
