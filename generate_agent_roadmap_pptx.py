# -*- coding: utf-8 -*-
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_agent_presentation():
    prs = Presentation()
    # Set to 16:9 widescreen layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ------------------ THEME COLORS (Modern Tech Cyberpunk/Dark Mode) ------------------
    BG_COLOR = RGBColor(10, 15, 29)       # Very dark tech blue #0A0F1D
    CARD_BG = RGBColor(21, 30, 51)        # Slate tech card background #151E33
    CARD_BORDER = RGBColor(0, 180, 216)   # Cyan border #00B4D8
    
    TEXT_TITLE = RGBColor(255, 255, 255)  # White
    TEXT_BODY = RGBColor(226, 232, 240)   # Light slate #E2E8F0
    TEXT_MUTED = RGBColor(148, 163, 184)  # Slate gray #94A3B8
    
    ACCENT_CYAN = RGBColor(0, 245, 212)   # Neon Cyan #00F5D4
    ACCENT_BLUE = RGBColor(0, 180, 216)   # Electric Blue #00B4D8
    ACCENT_PURPLE = RGBColor(114, 9, 183) # Deep Purple #7209B7
    
    blank_layout = prs.slide_layouts[6]
    
    # Helper: Set solid background for a slide
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR
        
    # Helper: Add Slide Header (Title + Subtitle + Accent Line)
    def add_slide_header(slide, title, subtitle):
        # Decorative left accent bar
        dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.4), Inches(0.08), Inches(0.65))
        dec.fill.solid()
        dec.fill.fore_color.rgb = ACCENT_CYAN
        dec.line.fill.background()
        
        # Textbox for Title + Subtitle
        tb = slide.shapes.add_textbox(Inches(0.85), Inches(0.3), Inches(11.8), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        # Title
        p_t = tf.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Microsoft YaHei'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_TITLE
        p_t.space_after = Pt(2)
        
        # Subtitle
        p_s = tf.add_paragraph()
        p_s.text = subtitle
        p_s.font.name = 'Microsoft YaHei'
        p_s.font.size = Pt(11)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_BLUE
        
    # Helper: Create a styled "Card" with title and bullets
    def draw_card(slide, left, top, width, height, title, items, badge_text=None, accent_color=ACCENT_CYAN):
        # 1. Card base shape
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        # 2. Add badge if present
        if badge_text:
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left + Inches(0.2), top - Inches(0.15), Inches(1.5), Inches(0.3))
            badge.fill.solid()
            badge.fill.fore_color.rgb = accent_color
            badge.line.fill.background()
            
            btf = badge.text_frame
            btf.word_wrap = True
            btf.margin_left = btf.margin_top = btf.margin_right = btf.margin_bottom = 0
            bp = btf.paragraphs[0]
            bp.text = badge_text
            bp.alignment = PP_ALIGN.CENTER
            bp.font.name = 'Microsoft YaHei'
            bp.font.size = Pt(9)
            bp.font.bold = True
            bp.font.color.rgb = BG_COLOR
            
        # 3. Card Title
        tb_title = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.25), width - Inches(0.5), Inches(0.45))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_ct = tf_title.paragraphs[0]
        p_ct.text = title
        p_ct.font.name = 'Microsoft YaHei'
        p_ct.font.size = Pt(14)
        p_ct.font.bold = True
        p_ct.font.color.rgb = accent_color
        
        # 4. Bullet Items
        tb_body = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.75), width - Inches(0.5), height - Inches(0.9))
        tf_body = tb_body.text_frame
        tf_body.word_wrap = True
        tf_body.margin_left = tf_body.margin_top = tf_body.margin_right = tf_body.margin_bottom = 0
        
        for idx, item in enumerate(items):
            p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
            p.text = item
            p.font.name = 'Microsoft YaHei'
            p.font.size = Pt(10.5)
            p.font.color.rgb = TEXT_BODY
            p.space_after = Pt(8)
            p.level = 0
            
    # ====================================================================================
    # ------------------ SLIDE 1: COVER (封面页) ------------------
    # ====================================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide1)
    
    # Left decorative vertical bars
    dec1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.25), Inches(7.5))
    dec1.fill.solid()
    dec1.fill.fore_color.rgb = ACCENT_PURPLE
    dec1.line.fill.background()
    
    dec2 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.25), Inches(0), Inches(0.1), Inches(7.5))
    dec2.fill.solid()
    dec2.fill.fore_color.rgb = ACCENT_CYAN
    dec2.line.fill.background()
    
    # Cover Card (Large Center-Left Card)
    cover_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(4.5))
    tf1 = cover_box.text_frame
    tf1.word_wrap = True
    
    # Category Tag
    p_tag = tf1.paragraphs[0]
    p_tag.text = "TECHNICAL ROADMAP & CORE ARCHITECTURE"
    p_tag.font.name = 'Courier New'
    p_tag.font.size = Pt(13)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_CYAN
    p_tag.space_after = Pt(12)
    
    # Main Title
    p_title = tf1.add_paragraph()
    p_title.text = "AI Agent 技术路线与核心技术深度解析"
    p_title.font.name = 'Microsoft YaHei'
    p_title.font.size = Pt(36)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_TITLE
    p_title.space_after = Pt(8)
    
    # English Subtitle
    p_sub_en = tf1.add_paragraph()
    p_sub_en.text = "Deep Dive into Autonomous AI Agents: Planning, Memory, Action and Multi-Agent Orchestration"
    p_sub_en.font.name = 'Arial'
    p_sub_en.font.size = Pt(14)
    p_sub_en.font.italic = True
    p_sub_en.font.color.rgb = TEXT_MUTED
    p_sub_en.space_after = Pt(28)
    
    # Decorative line
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(4.3), Inches(5.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()
    
    # Metadata info
    p_meta = tf1.add_paragraph()
    p_meta.text = "\n主讲人：全栈 AI 算法 & 智能体系统架构专家\n报告定位：前沿技术分享 · 架构备战答辩 · 核心机制拆解\n导出方式：Python 原生矢量渲染幻灯片 (100% 可编辑)"
    p_meta.font.name = 'Microsoft YaHei'
    p_meta.font.size = Pt(12.5)
    p_meta.font.color.rgb = TEXT_BODY
    p_meta.space_after = Pt(10)
    
    # ====================================================================================
    # ------------------ SLIDE 2: ARCHITECTURE (架构总览) ------------------
    # ====================================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide2)
    add_slide_header(slide2, "一、AI Agent 经典架构总览", "面向闭环控制的“感知 - 思考 - 行动”核心三要素体系")
    
    # Draw 3 column cards (Perception, Brain, Action)
    col_w = Inches(3.7)
    gap = Inches(0.3)
    start_x = Inches(0.6)
    top_y = Inches(1.8)
    height_y = Inches(4.8)
    
    draw_card(
        slide2, start_x, top_y, col_w, height_y,
        "01 / 感知层 (Perception)",
        [
            "■ 多模态输入解析：不仅支持传统的文本 Prompt 输入，还支持图像、语音、视频及多模态混合输入的特征提取与语义理解。",
            "■ 环境感知与触觉：在具身智能或机器人控制场景中，包含各类传感器数据输入（激光雷达、IMU、触觉传感器）。",
            "■ 主动环境扫描：定期主动拉取数据流（如监控摄像头帧、系统日志），而非被动等待用户下发指令。"
        ],
        "INPUT SYSTEM", ACCENT_PURPLE
    )
    
    draw_card(
        slide2, start_x + col_w + gap, top_y, col_w, height_y,
        "02 / 核心大脑 (Brain)",
        [
            "■ 目标分解 (Planning)：将复杂的大目标拆解为可操作的微小步骤，并动态调整计划路线。",
            "■ 知识检索 (Memory)：从短期对话上下文和长期外挂向量数据库中快速检索、召回相关事实知识。",
            "■ 任务分配与决策：依据 Prompt 约束与 Fine-tuned 规则，匹配最佳执行策略，决定调用什么工具。"
        ],
        "DECISION CORE", ACCENT_CYAN
    )
    
    draw_card(
        slide2, start_x + (col_w + gap) * 2, top_y, col_w, height_y,
        "03 / 执行层 (Action)",
        [
            "■ 结构化工具调用：利用大模型 Function Calling 输出标准 JSON 格式参数，驱动外部 API 接口。",
            "■ 自主代码执行：自动编写 Python/Bash 脚本并在隔离的安全沙盒中运行，处理数学计算与数据可视化。",
            "■ 环境交互响应：将执行结果反向输入给大脑，形成持续修正、滚动的自适应反馈控制环路。"
        ],
        "EXECUTION UNIT", ACCENT_BLUE
    )
    
    # ====================================================================================
    # ------------------ SLIDE 3: PLANNING (规划与推理核心) ------------------
    # ====================================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide3)
    add_slide_header(slide3, "二、大脑核心之：Planning 规划与推理技术", "从一步直觉式回复走向深度搜索与反思反馈控制")
    
    col_w2 = Inches(5.8)
    start_x2 = Inches(0.6)
    gap2 = Inches(0.533)
    
    draw_card(
        slide3, start_x2, top_y, col_w2, height_y,
        "经典思维路径拓扑 (From Chain to Tree)",
        [
            "■ 思位链 (Chain of Thought, CoT):",
            "  让模型写出思考步骤：y = f(x) -> y = f(f(f(x)))。强制大模型逐步推理，激活多步注意特征，大幅降低逻辑推导错误。",
            "■ 思维树 (Tree of Thoughts, ToT):",
            "  在多步规划中引入搜索树架构，对每个子步骤（思考节点）生成多个可行分支，结合深度优先(DFS)或广度优先(BFS)评估各分支价值。",
            "■ 蒙特卡洛树搜索 (MCTS) 智能体:",
            "  结合大模型打分作为启发式评估函数，在庞大的行动空间中进行多轮 Rollout 模拟，从而求解出极长任务链路的最优动作路径。"
        ],
        "THOUGHT PATTERNS"
    )
    
    draw_card(
        slide3, start_x2 + col_w2 + gap2, top_y, col_w2, height_y,
        "动作与反思闭环机制 (Act & Reflexion)",
        [
            "■ ReAct (Reason + Act) 协同框架:",
            "  交替执行『Thought -> Action -> Observation』三步循环。思考驱动动作，动作获得环境反馈（观察），观察修正下一次思考。完全打破了纯静态生成的黑盒模式。",
            "■ Reflexion 自我修正网络:",
            "  Agent 执行任务失败或报错时，将其失败日志和执行路径作为上下文输入给“反思模型”（Evaluator）。",
            "  反思模型生成具体的失败归因和避坑提示词（Reflection Vector），并将其注入下一轮迭代的短期记忆中，从而实现无需微调的自主进化。"
        ],
        "REFLEXION LOOP"
    )
    
    # ====================================================================================
    # ------------------ SLIDE 4: MEMORY (记忆体系设计) ------------------
    # ====================================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide4)
    add_slide_header(slide4, "三、大脑核心之：Memory 记忆机制与知识存取", "融合短期序列关联与长期外挂多维检索的双轨记忆架构")
    
    draw_card(
        slide4, start_x2, top_y, col_w2, height_y,
        "短期记忆：上下文与会话状态管理",
        [
            "■ 会话上下文缓冲区 (Session Buffer):",
            "  存储原生对话、工具调用入参和执行结果。为应对大模型上下文长度极限与显存吞吐压力，通常采用滑动窗口或摘要算法进行遗忘压缩。",
            "■ 结构化系统状态机 (State Machine Memory):",
            "  使用严格的状态变量记录 Agent 当前所处的业务步骤（如：已支付、待核销、已退款），确保即使大模型在多轮对话中发生偏置，业务主骨架依旧绝对可控。",
            "■ KV Cache 读写性能优化:",
            "  在智能体快速迭代的极速推理模式中，通过前置缓存 System Prompt 的 Prefill 特征，大幅度降低首Token延迟。"
        ],
        "SHORT-TERM MEMORY"
    )
    
    draw_card(
        slide4, start_x2 + col_w2 + gap2, top_y, col_w2, height_y,
        "长期记忆：基于向量与RAG的非结构化检索",
        [
            "■ 混合检索机制 (Hybrid Search):",
            "  结合基于稠密向量（Dense Vector）的语义内积检索与基于稀疏向量（Sparse Vector, 如 BM25）的关键词碰撞检索，兼顾同义词理解与专业词汇匹配。",
            "■ 高性能向量底座 (Vector DB):",
            "  引入 Milvus、Qdrant 等分布式向量数据库，实现海量嵌入向量（Embedding）的毫秒级 ANN（近似最近邻）搜索。",
            "■ 重排模型过滤 (Reranker):",
            "  针对初筛获得的 Top-K 文本块，使用 Cross-Encoder 双向交互重排模型计算二次相关度得分，深度剔除噪声上下文，大幅减缓大模型生成幻觉。"
        ],
        "LONG-TERM MEMORY"
    )
    
    # ====================================================================================
    # ------------------ SLIDE 5: ACTIONS (工具调用与执行安全) ------------------
    # ====================================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide5)
    add_slide_header(slide5, "四、工具调用与沙盒安全 (Actions & Tool Integration)", "打通大模型与外界真实操作系统的矢量互通与安全防护隔离")
    
    draw_card(
        slide5, start_x, top_y, col_w, height_y,
        "工具描述与函数调用",
        [
            "■ JSON Schema 规范化定义：",
            "  大模型无法直接点击按钮。开发者必须提供细致的 API 函数描述，包括参数名称、类型及每个参数的业务含义说明，以便模型进行语义模式匹配。",
            "■ 动态意图分类匹配：",
            "  大模型提取用户 Prompt 中的实体与动作，计算其与注册工具列表中描述向量的夹角余弦，自动路由到最切合的函数。"
        ],
        "SCHEMA DESIGN", ACCENT_CYAN
    )
    
    draw_card(
        slide5, start_x + col_w + gap, top_y, col_w, height_y,
        "代码解释器 (Code Interpreter)",
        [
            "■ 动态脚本编写：",
            "  面对高维数学题、大规模 CSV 数据分析或生成统计图表，传统 LLM 的文本推理极易出错。Agent 则是生成完整的 Python 代码，借由写出代码来间接解决问题。",
            "■ 代码回传与反馈：",
            "  如果代码报错，报错信息回传给 Brain，触发 CoT 自动修复，直至代码成功运行并得出结果。"
        ],
        "DYNAMIC RUNTIME", ACCENT_BLUE
    )
    
    draw_card(
        slide5, start_x + (col_w + gap) * 2, top_y, col_w, height_y,
        "隔离沙盒与系统安全防护",
        [
            "■ 敏感系统级防御：",
            "  绝不允许 Agent 生成的任意代码直接在宿主机裸奔。必须构建极致的沙盒隔离防线，防范勒索软件和拒绝服务攻击。",
            "■ 行业级沙盒实现：",
            "  使用 Docker 轻量级容器隔离或 E2B 安全沙盒运行时。严格限制 CPU/内存限额，限制出网流量，并开启系统调用级别审计监控。"
        ],
        "SANDBOX SECURITY", ACCENT_PURPLE
    )
    
    # ====================================================================================
    # ------------------ SLIDE 6: MULTI-AGENT (多智能体协同与编排) ------------------
    # ====================================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide6)
    add_slide_header(slide6, "五、多智能体协同与编排架构 (Multi-Agent Systems)", "从单兵作战演进为有组织、有状态的多角色生产线网络")
    
    draw_card(
        slide6, start_x2, top_y, col_w2, height_y,
        "多智能体协同范式分类",
        [
            "■ 主从架构 (Manager-Worker Model):",
            "  Manager Agent 统一负责接收用户任务、制定主骨架大纲、进行任务细分指派，并对各 Worker Agent 的产出进行严格质检与汇总。",
            "■ 对等协同与辩论机制 (Peer-to-Peer & Debate):",
            "  多个同等级 Agent 各司其职。例如『研发 Agent』编写代码，『测试 Agent』运行测试，『安全 Agent』审计漏洞。多角色在消息总线上异步通信，通过辩论博弈提升输出健壮性。",
            "■ 分布式共识驱动:",
            "  多 Agent 通过特定协议（如类似于 Raft 或投票表决）对某个全局决策达成最终共识，适合集群分布式寻优场景。"
        ],
        "COLLABORATION MODELS"
    )
    
    draw_card(
        slide6, start_x2 + col_w2 + gap2, top_y, col_w2, height_y,
        "主流 Agent 开发框架对比与剖析",
        [
            "■ LangGraph (基于状态图编排 - 强烈推荐):",
            "  将任务流程抽象为图结构（Nodes & Edges）。图的边支持带条件的逻辑分支判断。核心是以全局一致的状态（State）为核心进行流转，完美支持任意循环迭代逻辑，极其可控。",
            "■ CrewAI (面向角色扮演与团队管理):",
            "  以角色（Role）、任务（Task）和团队（Crew）为核心概念。极其易用，善于模拟企业内部多部门配合的工作流逻辑。",
            "■ AutoGen (面向对话式自主协作):",
            "  由微软主导开发，基于 Agent 间的多轮对话来协同解决问题，支持极高自由度的自主交互，适合开放式的创意与研究性探索。"
        ],
        "FRAMEWORK COMPARISON"
    )
    
    # ====================================================================================
    # ------------------ SLIDE 7: ROADMAP (落地路线与未来挑战) ------------------
    # ====================================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide7)
    add_slide_header(slide7, "六、智能体落地演进路线与行业挑战", "从低阶提示词包装，向深度微调、高并发状态编排迈进")
    
    # 3 horizontal flow cards
    flow_w = Inches(3.7)
    flow_gap = Inches(0.3)
    flow_start_x = Inches(0.6)
    flow_top_y = Inches(2.2)
    flow_height_y = Inches(4.3)
    
    draw_card(
        slide7, flow_start_x, flow_top_y, flow_w, flow_height_y,
        "L1 - RAG 与工具浅层包装",
        [
            "■ 主攻方向：基于 LangChain、LlamaIndex 快速搭建本地知识库检索与单 API 接口调用。",
            "■ 特点：构建成本低廉。大模型仅作为单次语义解析与静态回复的外壳，缺乏自我修正和多步循环规划能力。",
            "■ 适用场景：标准文档检索问答、自动化邮件通知草拟。"
        ],
        "STAGE 1: STANDARD", ACCENT_PURPLE
    )
    
    draw_card(
        slide7, flow_start_x + flow_w + flow_gap, flow_top_y, flow_w, flow_height_y,
        "L2 - 复杂图状态工作流编排",
        [
            "■ 主攻方向：引入 LangGraph 进行多角色协作图构建，引入多模态感知与安全代码执行沙盒。",
            "■ 特点：可控性高。通过把控全局状态，加入人类审批机制（Human-in-the-loop），有效杜绝无限调用死循环。",
            "■ 适用场景：复杂合同分析评审、自动化软件 Bug 修复与集成测试。"
        ],
        "STAGE 2: ADVANCED", ACCENT_CYAN
    )
    
    draw_card(
        slide7, flow_start_x + (flow_w + flow_gap) * 2, flow_top_y, flow_w, flow_height_y,
        "L3 - 垂直微调与完全自主进化",
        [
            "■ 主攻方向：通过大模型微调（SFT）将 Agent 推理轨迹、工具调用规范（Function Schema）强固化到模型权重中。",
            "■ 特点：推理延迟减半，模型摆脱对巨量 Prompt 的依赖。利用长时反思日志，实现离线策略梯度更新与持续自进化。",
            "■ 适用场景：特定业务场景的完全自主数字员工、具身智能控制中枢。"
        ],
        "STAGE 3: AUTONOMOUS", ACCENT_BLUE
    )
    
    # Draw simple connecting arrows between flow cards
    # Step 1 -> Step 2
    arrow1 = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, flow_start_x + flow_w + Inches(0.05), flow_top_y + Inches(2.0), Inches(0.2), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = ACCENT_CYAN
    arrow1.line.fill.background()
    
    # Step 2 -> Step 3
    arrow2 = slide7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, flow_start_x + (flow_w + flow_gap) + flow_w + Inches(0.05), flow_top_y + Inches(2.0), Inches(0.2), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = ACCENT_BLUE
    arrow2.line.fill.background()
    
    # ====================================================================================
    # ------------------ SLIDE 8: END (谢谢观看) ------------------
    # ====================================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide8)
    
    # Geometric accent shapes in the background for a modern look
    shape_dec1 = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.0), Inches(-1.0), Inches(5.0), Inches(5.0))
    shape_dec1.fill.solid()
    shape_dec1.fill.fore_color.rgb = CARD_BG
    shape_dec1.line.color.rgb = ACCENT_PURPLE
    shape_dec1.line.width = Pt(2.0)
    
    shape_dec2 = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(-1.5), Inches(4.5), Inches(4.0), Inches(4.0))
    shape_dec2.fill.solid()
    shape_dec2.fill.fore_color.rgb = CARD_BG
    shape_dec2.line.color.rgb = ACCENT_CYAN
    shape_dec2.line.width = Pt(1.5)
    
    # Thanks text box
    tb_thanks = slide8.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.0))
    t_tf = tb_thanks.text_frame
    t_tf.word_wrap = True
    
    p_th = t_tf.paragraphs[0]
    p_th.text = "谢谢您的观看与专业评审"
    p_th.font.name = 'Microsoft YaHei'
    p_th.font.size = Pt(36)
    p_th.font.bold = True
    p_th.font.color.rgb = ACCENT_CYAN
    p_th.alignment = PP_ALIGN.CENTER
    p_th.space_after = Pt(12)
    
    p_sub_th = t_tf.add_paragraph()
    p_sub_th.text = "AI Agent 技术路线与核心技术深度解析  |  演示文稿制作完毕\nLet's Build the Future of Agentic AI Together."
    p_sub_th.font.name = 'Microsoft YaHei'
    p_sub_th.font.size = Pt(16)
    p_sub_th.font.color.rgb = TEXT_TITLE
    p_sub_th.alignment = PP_ALIGN.CENTER
    p_sub_th.space_after = Pt(24)
    
    p_details = t_tf.add_paragraph()
    p_details.text = "您可以随时双击编辑任何幻灯片中的卡片内容或调整色调。\n期待与您进一步交流智能体系统在各行各业的实际落地方法论！"
    p_details.font.name = 'Microsoft YaHei'
    p_details.font.size = Pt(12.5)
    p_details.font.color.rgb = TEXT_MUTED
    p_details.alignment = PP_ALIGN.CENTER
    
    # Save the presentation
    filename = "AI_Agent_技术路线与核心技术深度解析.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as '{filename}'!")

if __name__ == '__main__':
    create_agent_presentation()
