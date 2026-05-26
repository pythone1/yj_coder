# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO
import hashlib
from collections import defaultdict

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "机理、算法与实践：人工智能行业应用实证分析.pptx"
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
MOD = ASSET_DIR / "v4_modules"
ARCH = ASSET_DIR / "v4_arch"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_教学清爽版_v5.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_教学清爽版_v5_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_教学清爽版_v5_预览联系表.png"

NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
FONT = "微软雅黑"

ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "dma": MOD / "01_dma_boundary.png",
    "hydraulic": MOD / "02_hydraulic_model.png",
    "simulation": MOD / "03_simulation_dashboard.png",
    "butterfly": MOD / "04_error_butterfly.png",
    "night": MOD / "05_night_demand.png",
    "ai_evo": MOD / "06_ai_evolution.png",
    "lstm": MOD / "07_lstm_gru.png",
    "iforest": MOD / "08_isolation_forest.png",
    "dbscan": MOD / "09_dbscan.png",
    "autoencoder": MOD / "10_autoencoder.png",
    "fusion": MOD / "11_ai_hydraulic_fusion.png",
    "risk": MOD / "12_risk_ranking.png",
    "gnn": MOD / "13_gnn_topology.png",
    "governance": MOD / "14_data_governance.png",
    "edge": MOD / "15_edge_cloud.png",
    "arch": ARCH / "overall_architecture.png",
    "pain": ARCH / "pain_ai_path.png",
    "alg_combo": ARCH / "algorithm_combo.png",
    "impl": ARCH / "implementation_path.png",
    "workorder": ARCH / "workorder_no_person.png",
}

AUDIT_LINES = []


def emu(v): return Inches(v)


def clear_slides(prs):
    ids = prs.slides._sldIdLst
    for sid in list(ids):
        prs.part.drop_rel(sid.rId)
        ids.remove(sid)


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def fill(shp, color, trans=0):
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = trans
    shp.line.fill.background()


def stroke(shp, color=LINE, width=1):
    shp.line.color.rgb = color
    shp.line.width = Pt(width)


def rect(s, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = s.shapes.add_shape(typ, x, y, w, h)
    fill(shp, color, trans)
    stroke(shp, line, 1)
    return shp


def tb(s, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pic(s, key, x, y, w, h, crop=False):
    p = s.shapes.add_picture(str(ASSETS[key]), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def connector(s, x1, y1, x2, y2, color=BLUE, width=2, arrow=True):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(s, title, no):
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.40), emu(0.25), emu(0.12), emu(0.50)), BLUE)
    tb(s, title, emu(0.65), emu(0.16), emu(10.15), emu(0.43), 25.5, TEXT, True)
    tb(s, f"{no:02d}", emu(12.0), emu(0.24), emu(0.72), emu(0.22), 13.5, BLUE, True, PP_ALIGN.RIGHT)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.65), emu(0.88), emu(11.65), emu(0.018)), PALE)


def content(prs, title, no):
    s = blank(prs)
    header(s, title, no)
    return s


def pill(s, text, x, y, w, color=BLUE, size=15):
    rect(s, emu(x), emu(y), emu(w), emu(0.46), PALE2, color, True)
    tb(s, text, emu(x + 0.08), emu(y + 0.10), emu(w - 0.16), emu(0.26), size, color, True, PP_ALIGN.CENTER)


def big_label(s, title, body, x, y, w, h, color=BLUE):
    rect(s, emu(x), emu(y), emu(w), emu(h), PALE2, LINE, True)
    tb(s, title, emu(x + 0.16), emu(y + 0.14), emu(w - 0.32), emu(0.25), 17, color, True)
    tb(s, body, emu(x + 0.16), emu(y + 0.52), emu(w - 0.32), emu(h - 0.62), 14.5, TEXT)


def section(prs, no, part, title, subtitle, img):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(0.24), emu(7.5)), CYAN)
    tb(s, part, emu(0.88), emu(1.18), emu(1.5), emu(0.30), 17, CYAN, True)
    tb(s, title, emu(0.88), emu(1.92), emu(6.4), emu(0.70), 36, WHITE, True)
    tb(s, subtitle, emu(0.92), emu(3.10), emu(6.35), emu(0.70), 18, RGBColor(210, 232, 245))
    rect(s, emu(8.05), emu(1.18), emu(4.25), emu(4.25), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img, emu(8.30), emu(1.44), emu(3.75), emu(3.70), crop=False)
    tb(s, f"{no:02d}", emu(11.75), emu(6.58), emu(0.65), emu(0.25), 15, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)
    audit(s, no, allow_short=True)


def visual_full(prs, no, title, img, caption=None):
    s = content(prs, title, no)
    rect(s, emu(0.55), emu(1.02), emu(12.25), emu(5.95), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img, emu(0.72), emu(1.16), emu(11.90), emu(5.55), crop=False)
    if caption:
        rect(s, emu(3.75), emu(6.20), emu(5.85), emu(0.52), WHITE, RGBColor(205, 229, 243), True)
        tb(s, caption, emu(3.92), emu(6.35), emu(5.50), emu(0.18), 15, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no, allow_short=True)


def visual_split(prs, no, title, img, labels):
    s = content(prs, title, no)
    rect(s, emu(0.72), emu(1.18), emu(6.10), emu(5.20), WHITE, RGBColor(205, 229, 243), True)
    pic(s, img, emu(0.95), emu(1.45), emu(5.64), emu(4.66), crop=False)
    for i, (a, b, c) in enumerate(labels):
        big_label(s, a, b, 7.25, 1.22 + i * 1.55, 4.72, 1.10, c)
    audit(s, no)


def two_images(prs, no, title, left, right, labels):
    s = content(prs, title, no)
    rect(s, emu(0.70), emu(1.18), emu(3.95), emu(5.15), WHITE, RGBColor(205, 229, 243), True)
    pic(s, left, emu(0.92), emu(1.38), emu(3.50), emu(4.75), crop=False)
    rect(s, emu(4.92), emu(1.18), emu(3.95), emu(5.15), WHITE, RGBColor(205, 229, 243), True)
    pic(s, right, emu(5.14), emu(1.38), emu(3.50), emu(4.75), crop=False)
    for i, (a, b, c) in enumerate(labels):
        big_label(s, a, b, 9.18, 1.18 + i * 1.46, 3.10, 1.05, c)
    audit(s, no)


def cards(prs, no, title, items, cols=4):
    s = content(prs, title, no)
    single_row = len(items) <= cols
    card_h = 5.20 if single_row else 2.20
    img_h = 2.00 if single_row else 0.90
    body_y = 3.02 if single_row else 1.63
    for i, (head, body, img, color) in enumerate(items):
        col = i % cols
        row = i // cols
        x = 0.72 + col * (11.95 / cols)
        y = 1.20 + row * 2.60
        w = 2.65 if cols == 4 else 3.72
        rect(s, emu(x), emu(y), emu(w), emu(card_h), WHITE, RGBColor(205, 229, 243), True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(w), emu(0.42)), color)
        tb(s, head, emu(x + 0.08), emu(y + 0.11), emu(w - 0.16), emu(0.17), 14, WHITE, True, PP_ALIGN.CENTER)
        pic(s, img, emu(x + 0.22), emu(y + 0.68), emu(w - 0.44), emu(img_h), crop=False)
        tb(s, body, emu(x + 0.16), emu(y + body_y), emu(w - 0.32), emu(0.46), 14.5 if single_row else 13.5, TEXT, False, PP_ALIGN.CENTER)
    audit(s, no)


def timeline(prs, no, title, items, img):
    s = content(prs, title, no)
    pic(s, img, emu(0.80), emu(1.12), emu(4.35), emu(5.35), crop=False)
    y = 2.10
    for i, (head, body, color) in enumerate(items):
        x = 5.35 + i * 1.75
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(x), emu(y), emu(0.58), emu(0.58)), color)
        tb(s, str(i + 1), emu(x), emu(y + 0.15), emu(0.58), emu(0.18), 15, WHITE, True, PP_ALIGN.CENTER)
        if i < len(items) - 1:
            connector(s, emu(x + 0.62), emu(y + 0.29), emu(x + 1.63), emu(y + 0.29), BLUE, 2)
        tb(s, head, emu(x - 0.30), emu(y + 0.85), emu(1.25), emu(0.28), 15, color, True, PP_ALIGN.CENTER)
        tb(s, body, emu(x - 0.42), emu(y + 1.24), emu(1.50), emu(0.70), 13.2, TEXT, False, PP_ALIGN.CENTER)
    pill(s, "机理模型 = 物理约束", 5.60, 5.75, 2.10, BLUE, 15)
    pill(s, "AI模型 = 动态学习", 8.10, 5.75, 2.10, GREEN, 15)
    audit(s, no)


def native_flow(prs, no, title, steps, img=None):
    s = content(prs, title, no)
    if img:
        pic(s, img, emu(0.78), emu(1.18), emu(3.75), emu(4.90), crop=False)
        x0 = 4.85
    else:
        x0 = 0.90
    for i, (head, body, color) in enumerate(steps):
        x = x0 + (i % 3) * 2.55
        y = 1.20 + (i // 3) * 2.85
        rect(s, emu(x), emu(y), emu(2.22), emu(1.95), WHITE, RGBColor(205, 229, 243), True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(x + 0.16), emu(y + 0.16), emu(0.42), emu(0.42)), color)
        tb(s, str(i + 1), emu(x + 0.16), emu(y + 0.24), emu(0.42), emu(0.16), 13, WHITE, True, PP_ALIGN.CENTER)
        tb(s, head, emu(x + 0.72), emu(y + 0.18), emu(1.30), emu(0.22), 15, color, True)
        tb(s, body, emu(x + 0.18), emu(y + 0.72), emu(1.86), emu(0.75), 13.5, TEXT)
        if i < len(steps) - 1 and i % 3 != 2:
            connector(s, emu(x + 2.25), emu(y + 0.80), emu(x + 2.50), emu(y + 0.80), BLUE, 1.8)
    audit(s, no)


def cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(5.15), emu(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.62), emu(0.88), emu(0.13), emu(5.82)), CYAN)
    tb(s, "AI模型在供水管网\nDMA系统漏损检测\n中的应用", emu(1.02), emu(1.10), emu(3.95), emu(1.85), 31, WHITE, True)
    tb(s, "教学演讲版", emu(1.05), emu(3.25), emu(2.30), emu(0.33), 18, CYAN, True)
    tb(s, "从机理模型到AI溯源，从异常预警到工单闭环", emu(1.05), emu(5.68), emu(3.95), emu(0.35), 14.5, RGBColor(218, 237, 248))
    tb(s, "大图讲解 · 少字呈现 · 可编辑结构", emu(1.05), emu(6.18), emu(3.95), emu(0.28), 14, RGBColor(218, 237, 248))
    pic(s, "arch", emu(5.45), emu(0.62), emu(7.45), emu(6.20), crop=False)
    audit(s, 1, allow_short=True)


def agenda(prs):
    items = [
        ("背景", "为什么需要AI", "hydraulic", BLUE),
        ("基础", "数据与模型怎么建", "lstm", GREEN),
        ("算法", "模型组合怎么用", "alg_combo", ORANGE),
        ("应用", "预警定位闭环", "workorder", CYAN),
        ("实施", "项目怎么落地", "impl", PURPLE),
        ("展望", "长期如何演进", "ai_evo", RED),
    ]
    cards(prs, 2, "课程地图：六个教学模块", items, cols=3)


def audit(s, no, allow_short=False):
    prs_w, prs_h = 12192000, 6858000
    chars = pics = off = 0
    min_font = 99
    bottom = 0
    for sh in s.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs_w + 50000 or sh.top + sh.height > prs_h + 50000:
            off += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off: flags.append(f"越界{off}")
    if min_font < 13: flags.append(f"小字{min_font}")
    if pics < 1: flags.append("缺主视觉")
    if not allow_short and chars < 50: flags.append(f"可见内容偏少{chars}")
    if bottom < emu(6.0): flags.append("页面下半部偏空")
    status = "FAIL" if flags else "OK"
    AUDIT_LINES.append(f"{no:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise RuntimeError(f"slide {no}: {'; '.join(flags)}")


def build():
    missing = [str(p) for p in ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("\n".join(missing))
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    cover(prs)
    agenda(prs)

    section(prs, 3, "PART 01", "模型应用背景", "先讲清问题，再讲模型。", "butterfly")
    visual_split(prs, 4, "DMA的定位：先锁区，再找点", "dma", [
        ("DMA能解决", "把复杂管网划成可计量、可考核的片区。", BLUE),
        ("DMA不能直接解决", "只能指出哪个区异常，不能直接指出哪根管。", RED),
        ("AI接着做", "融合压力、流量、拓扑和工单，输出候选管段。", GREEN),
    ])
    timeline(prs, 5, "机理模型演进：供水建模的基础", [
        ("平差", "哈代克罗斯", BLUE),
        ("迭代", "牛顿拉夫逊", GREEN),
        ("瞬变", "特征线法", ORANGE),
        ("仿真", "EPANET等", PURPLE),
    ], "hydraulic")
    visual_split(prs, 6, "模型软件生态：工具各有边界", "simulation", [
        ("供水仿真", "EPANET、WaterGEMS用于水力分析。", BLUE),
        ("瞬变分析", "Hammer等用于水锤和短时动态。", ORANGE),
        ("GIS融合", "偏空间管理和资产联动。", GREEN),
    ])
    timeline(prs, 7, "AI发展路线：从阈值到智能协同", [
        ("阈值", "MNF经验判断", BLUE),
        ("机器学习", "异常筛查", GREEN),
        ("深度学习", "动态基线", ORANGE),
        ("智能体", "解释与协同", PURPLE),
    ], "ai_evo")
    visual_full(prs, 8, "传统机理模型痛点与AI增强路径", "pain")
    two_images(prs, 9, "痛点拆解：数据误差与用水变化", "butterfly", "night", [
        ("数据误差", "台账、仪表、阀门状态会放大定位偏差。", RED),
        ("时变用水", "夜间合法用水与真实漏损需要区分。", ORANGE),
        ("AI价值", "动态学习、容错识别、持续更新。", GREEN),
    ])
    cards(prs, 10, "AI破解问题的四项价值", [
        ("数据集", "容错、对齐、融合", "governance", BLUE),
        ("自主学习", "学习复杂规律", "lstm", GREEN),
        ("辅助决策", "排序与派单", "risk", ORANGE),
        ("迭代更新", "反馈再训练", "autoencoder", PURPLE),
    ], cols=4)

    section(prs, 11, "PART 02", "技术基础", "把架构和数据讲清楚。", "governance")
    visual_full(prs, 12, "DMA漏损检测AI总体架构", "arch")
    visual_split(prs, 13, "数据层：模型效果的上限", "governance", [
        ("时序数据", "流量、压力、水质、设备状态。", BLUE),
        ("空间数据", "GIS拓扑、管段、阀门、DMA边界。", GREEN),
        ("业务数据", "工单、维修、投诉、施工、天气。", ORANGE),
    ])
    native_flow(prs, 14, "数据端部署：从采集到模型输入", [
        ("采集", "流量与压力", BLUE),
        ("对齐", "统一时间戳", GREEN),
        ("清洗", "缺失与毛刺", ORANGE),
        ("编码", "设备管段统一", PURPLE),
        ("入库", "时序与标签", CYAN),
        ("训练", "形成样本", RED),
    ], "edge")
    visual_split(prs, 15, "动态基线：每个DMA有自己的正常曲线", "lstm", [
        ("输入窗口", "历史流量、压力、日期与节假日。", BLUE),
        ("预测区间", "输出未来短时正常范围。", GREEN),
        ("残差报警", "持续偏离才进入复核。", RED),
    ])

    section(prs, 16, "PART 03", "核心算法", "算法按任务分工，不按名词堆叠。", "alg_combo")
    visual_full(prs, 17, "多模型组合：预测、识别、定位、排序、解释", "alg_combo")
    cards(prs, 18, "无监督异常检测：少标签阶段先可用", [
        ("孤立森林", "快速筛出离群异常", "iforest", BLUE),
        ("DBSCAN", "识别异常簇", "dbscan", GREEN),
        ("自编码器", "重构误差识别异常", "autoencoder", ORANGE),
        ("人工复核", "把结果回填成标签", "simulation", PURPLE),
    ], cols=4)
    visual_split(prs, 19, "水力模型 + AI：定位可信度的关键", "fusion", [
        ("AI发现异常", "从实时数据中提取异常证据。", BLUE),
        ("水力约束", "判断候选漏点是否符合压力响应。", GREEN),
        ("管段输出", "给出TopN候选管段和核查顺序。", ORANGE),
    ])
    two_images(prs, 20, "风险排序与图结构建模", "risk", "gnn", [
        ("风险排序", "管龄、材质、压力、维修历史。", ORANGE),
        ("拓扑传播", "异常影响沿管网关系扩散。", BLUE),
        ("解释输出", "让现场知道为什么优先查。", GREEN),
    ])

    section(prs, 21, "PART 04", "核心应用", "按事前、事中、事后讲现场价值。", "fusion")
    visual_split(prs, 22, "事前预警：把异常提前识别出来", "night", [
        ("MNF异常", "夜间底流持续抬升。", RED),
        ("压力响应", "局部压力出现异常变化。", BLUE),
        ("分级预警", "先复核，再定位，再派单。", GREEN),
    ])
    visual_full(prs, 23, "事中定位：从片区收敛到候选管段", "fusion")
    visual_full(prs, 24, "事后闭环：工单回填与再训练", "workorder")
    two_images(prs, 25, "拓展场景：压力优化与DMA规划", "edge", "gnn", [
        ("压力优化", "稳压、降漏、节能协同。", GREEN),
        ("监测点优化", "有限预算下提升定位贡献。", BLUE),
        ("分区优化", "让DMA更适合计量和建模。", ORANGE),
    ])

    section(prs, 26, "PART 05", "实施路径", "先试点，再推广，持续迭代。", "impl")
    visual_full(prs, 27, "AI系统实施路径", "impl")
    native_flow(prs, 28, "模型训练路径：由浅入深", [
        ("MNF", "基础预警", BLUE),
        ("LSTM", "动态基线", GREEN),
        ("异常检测", "交叉验证", ORANGE),
        ("水力融合", "候选收敛", PURPLE),
        ("风险排序", "管段优先级", CYAN),
        ("闭环更新", "再训练", RED),
    ], "iforest")
    cards(prs, 29, "数据治理五张表", [
        ("设备表", "流量计、压力计、网关", "edge", BLUE),
        ("管网表", "管段、阀门、DMA", "dma", GREEN),
        ("时序表", "采样时间、流量压力", "night", ORANGE),
        ("工单表", "报警、维修、结果", "simulation", PURPLE),
        ("标签表", "漏点、误报、版本", "iforest", RED),
        ("规则表", "阀门、施工、阈值", "butterfly", CYAN),
    ], cols=3)
    two_images(prs, 30, "模型验证：算法指标 + 业务指标", "risk", "autoencoder", [
        ("算法指标", "MAE、RMSE、召回率、F1。", BLUE),
        ("业务指标", "TopN命中、误报比例、闭环时长。", GREEN),
        ("教学重点", "把准确率翻译成现场效率。", ORANGE),
    ])
    native_flow(prs, 31, "长效运营：上线后怎么保持有效", [
        ("监控", "数据漂移", BLUE),
        ("复盘", "误报原因", ORANGE),
        ("标注", "工单回填", GREEN),
        ("再训练", "版本更新", PURPLE),
        ("灰度", "试点验证", CYAN),
        ("推广", "规模复制", RED),
    ], "workorder")

    section(prs, 32, "PART 06", "总结与展望", "高质量应用是系统工程。", "cover")
    visual_split(prs, 33, "最后三句话", "arch", [
        ("第一句", "DMA负责把问题锁定到片区。", BLUE),
        ("第二句", "AI负责把片区异常收敛到候选管段。", GREEN),
        ("第三句", "工单闭环负责让模型持续变准。", ORANGE),
    ])
    timeline(prs, 34, "未来演进：从漏损检测到全生命周期管控", [
        ("技术", "大模型与智能体", BLUE),
        ("架构", "端边云一体化", GREEN),
        ("体系", "规划建设运维", ORANGE),
        ("行业", "标准化轻量化", PURPLE),
    ], "ai_evo")

    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    build_preview(prs)
    append_repetition_check(prs)
    return PPTX


def build_preview(prs):
    W, H = 420, 236
    cols = 4
    rows = (len(prs.slides) + cols - 1) // cols
    sx, sy = W / prs.slide_width, H / prs.slide_height
    try:
        font_b = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 10)
        font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 8)
    except Exception:
        font_b = font = ImageFont.load_default()
    sheet = Image.new("RGB", (cols * W, rows * (H + 22)), (230, 236, 242))
    sd = ImageDraw.Draw(sheet)
    for idx, slide in enumerate(prs.slides):
        im = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(im)
        try:
            bg = slide.background.fill.fore_color.rgb
            if bg:
                im.paste(tuple(int(str(bg)[i:i + 2], 16) for i in (0, 2, 4)), [0, 0, W, H])
        except Exception:
            pass
        for sh in slide.shapes:
            x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
            if w <= 0 or h <= 0:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                    p.thumbnail((w, h))
                    im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
                except Exception:
                    pass
                continue
            fc = lc = None
            try:
                rgb = sh.fill.fore_color.rgb
                if rgb: fc = tuple(int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            try:
                rgb = sh.line.color.rgb
                if rgb: lc = tuple(int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            if fc:
                d.rounded_rectangle([x, y, x + w, y + h], radius=5, fill=fc, outline=lc)
            elif lc:
                d.rectangle([x, y, x + w, y + h], outline=lc)
            txt = getattr(sh, "text", "").strip()
            if txt and w > 35 and h > 10:
                bold = False
                try:
                    bold = any(r.font.bold for p in sh.text_frame.paragraphs for r in p.runs)
                except Exception:
                    pass
                d.text((x + 2, y + 2), txt.split("\n")[0][:26], fill=(18, 38, 58), font=font_b if bold else font)
        cx = (idx % cols) * W
        cy = (idx // cols) * (H + 22)
        sheet.paste(im, (cx, cy))
        sd.text((cx + 5, cy + H + 4), f"{idx + 1:02d}", fill=(0, 70, 148), font=font_b)
    sheet.save(PREVIEW)


def append_repetition_check(prs):
    by = defaultdict(list)
    for si, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                by[hashlib.md5(sh.image.blob).hexdigest()[:10]].append(si)
    with AUDIT.open("a", encoding="utf-8") as f:
        f.write("\n\n图片重复检查：\n")
        f.write(f"唯一图片={len(by)}，总图片={sum(len(v) for v in by.values())}\n")
        for h, slides in sorted(by.items(), key=lambda kv: len(kv[1]), reverse=True):
            if len(slides) > 3:
                f.write(f"{h}\t{len(slides)}\t{slides}\n")


if __name__ == "__main__":
    print(build())
    print(AUDIT)
    print(PREVIEW)
