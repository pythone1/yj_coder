from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Pt


BASE = Path(r"E:\PY\research\0425")
INPUT = BASE / "生产流程.pptx"
OUTPUT = BASE / "生产流程_最终修正版.pptx"

REPLACEMENTS = {
    "输 迭线": "输送线",
    "输\n迭线": "输送线",
    "锡护": "锅炉",
    "异崇检到": "异常检测",
    "异崇检测": "异常检测",
    "预湖": "预测",
    "遇终功能效益": "最终功能效益",
    "停 机凤险": "停机风险",
    "停\n机凤险": "停机风险",
    "凤险": "风险",
    "橙惨计划": "检修计划",
    "摄矢": "损失",
    "赐标": "贴标",
    "波位": "液位",
    "酒体思浮异物": "酒体悬浮异物",
    "主霎算法任务": "主要算法任务",
    "图像分劃": "图像分割",
    "识刖": "识别",
    "复较": "复核",
    "缺船": "缺陷",
    "追\n潮": "追溯",
    "追潮": "追溯",
    "任努": "任务",
    "最 功 能 效 益": "最终功能效益",
    "上瓶": "上甑",
    "昇常": "异常",
    "传燃数据": "传感数据",
    "温温度": "温度",
    "酒酷温度": "酒醅温度",
    "异案检测": "异常检测",
    "控制围": "控制图",
    "等缀": "等级",
    "掠入": "输入",
    "专案舰则": "专家规则",
    "检素": "检索",
    "规钏": "规则",
    "牧据": "数据",
    "措荐抉行": "推荐执行",
    "反懊学习": "反馈学习",
    "横型童训": "模型重训",
    "辖出": "输出",
    "璃参": "微调",
    "知玖沉渡": "知识沉淀",
    "产续": "产线",
    "发观": "发现",
    "该杀率": "误杀率",
    "复核络": "复核结论",
    "计罄线": "统计基线",
    "健康指数横型": "健康指数模型",
    "尖时时序": "实时序列",
    "数据\n汇粲": "数据\n汇聚",
    "汇粲": "汇聚",
    "逮度": "速度",
    "裁重": "载重",
    "起焚点": "起终点",
    "紫行区": "禁行区",
    "A\"": "A*",
    "多车落径": "多车路径",
    "规刺": "规划",
    "重规刘": "重规划",
    "逼过仿真": "通过仿真",
    "交忖时间": "交付时间",
    "温装": "混装",
    "多仓睾求": "多仓要求",
    "SKU 期间": "SKU 周转",
    "关肤规则": "关联规则",
    "库位保分": "库位评分",
    "密令": "订单",
    "启寂丸排序": "启发式排序",
    "液灰计刻": "波次计划",
    "捺选顺序": "拣选顺序",
    "月台支排": "月台安排",
    "辅入": "输入",
    "装魉舰物": "装载规划",
    "约柬伏化": "约束优化",
    "籍出": "输出",
    "多合节赤": "多仓节点",
    "韩出": "输出",
    "肠逆路线": "配送路线",
    "庠区": "库区",
    "又车": "叉车",
    "拣 遮": "拣选",
    "疲儋": "疲倦",
    "淀粉辜": "淀粉率",
    "搠眼": "摊晾",
    "发   醇": "发酵",
    "出客": "出窖",
    "缺船追": "缺陷追",
    "Al": "AI",
}


def fix_text(value):
    value = value.replace("ANotebookLM", "").replace("NotebookLM", "")
    for old, new in REPLACEMENTS.items():
        value = value.replace(old, new)
    return value


def recolor_text(shape, slide_idx):
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return
    text = shape.text_frame.text.strip()
    is_header = shape.top < 900000 and text
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            run.font.name = "Microsoft YaHei"
            if is_header:
                run.font.color.rgb = RGBColor(31, 78, 121)
            elif slide_idx > 0:
                run.font.color.rgb = RGBColor(245, 250, 255)


def tune_shape(shape, slide_idx):
    if hasattr(shape, "text_frame") and shape.text_frame:
        original = shape.text_frame.text
        fixed = fix_text(original)
        if fixed != original:
            shape.text_frame.text = fixed
        recolor_text(shape, slide_idx)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = fix_text(cell.text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            tune_shape(child, slide_idx)


def put_after_background(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(3, el)


def add_footer(slide, slide_width, slide_height):
    footer_h = Emu(213333)
    footer_y = Emu(int(slide_height) - int(footer_h))
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, footer_y, slide_width, footer_h)
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(71, 105, 158)
    footer.line.fill.background()
    put_after_background(slide, footer)


def add_body_card(slide, x, y, w, h, text):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(4, 28, 43)
    rect.fill.transparency = 8
    rect.line.fill.background()

    box = slide.shapes.add_textbox(x + 120000, y + 90000, w - 240000, h - 180000)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for pi, paragraph in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        para.text = paragraph
        para.line_spacing = 1.08
        for run in para.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(245, 250, 255)


prs = Presentation(str(INPUT))
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        tune_shape(shape, idx)
    if idx == 1:
        add_body_card(
            slide,
            10120000,
            5830000,
            2550000,
            3050000,
            "核心工艺对象：厂区物流节点、AGV任务、车辆位置、充电、避障、任务优先级。\n"
            "主要算法任务：路径规划、任务分配、多车冲突消解、预防、仿真评估。\n"
            "最终功能效益：减少等待和绕路，提高厂内物流准时率和设备利用率。",
        )
        add_body_card(
            slide,
            13220000,
            5830000,
            2550000,
            3050000,
            "核心工艺对象：库位、订单、库存、装车、车辆、客户路线、多仓发货。\n"
            "主要算法任务：库位优化、订单波次、装车排程、车辆路径优化。\n"
            "最终功能效益：缩短装车发货时间，提高库位利用率、满载率和配送效率。",
        )
    add_footer(slide, prs.slide_width, prs.slide_height)

prs.save(str(OUTPUT))
print(OUTPUT)
