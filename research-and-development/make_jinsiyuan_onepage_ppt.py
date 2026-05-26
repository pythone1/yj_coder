from pathlib import Path
import math

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("D:/Users/Downloads/jinsiyuan_llm_diagnosis_flow_onepage.pptx")
PREVIEW = Path("D:/Users/Downloads/jinsiyuan_llm_diagnosis_flow_onepage_preview.png")

prs = Presentation()
prs.slide_width = Inches(13.333333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

BLUE = RGBColor(16, 76, 165)
DARK = RGBColor(13, 43, 94)
TEAL = RGBColor(14, 133, 124)
LIGHT_BLUE = RGBColor(238, 247, 255)
LIGHT_TEAL = RGBColor(237, 250, 247)
TEXT = RGBColor(24, 38, 64)
GRAY = RGBColor(92, 105, 125)
WHITE = RGBColor(255, 255, 255)

PW, PH = 1600, 900
SX, SY = PW / 13.333333, PH / 7.5
preview = Image.new("RGB", (PW, PH), (249, 252, 255))
draw = ImageDraw.Draw(preview)
FONT_REGULAR = "C:/Windows/Fonts/msyh.ttc"
FONT_BOLD = "C:/Windows/Fonts/msyhbd.ttc"


def cc(color):
    return (color[0], color[1], color[2])


def px(x):
    return int(round(x * SX))


def py(y):
    return int(round(y * SY))


def pbox(x, y, w, h):
    return (px(x), py(y), px(x + w), py(y + h))


def getfont(size, bold=False):
    try:
        font_path = FONT_BOLD if bold and Path(FONT_BOLD).exists() else FONT_REGULAR
        return ImageFont.truetype(font_path, max(7, int(size * 1.45)))
    except Exception:
        return ImageFont.load_default()


def set_line(shape, color=BLUE, width=1.2, dash=None):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    if dash:
        shape.line.dash_style = dash


def set_fill(shape, color=WHITE):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def preview_text(x, y, w, h, text, size=12, color=TEXT, bold=False, align="center"):
    fnt = getfont(size, bold)
    lines = text.split("\n")
    line_h = max(9, int(size * 1.55))
    yy = py(y) + max(0, (py(h) - line_h * len(lines)) // 2)
    for line in lines:
        bb = draw.textbbox((0, 0), line, font=fnt)
        tw = bb[2] - bb[0]
        if align == "left":
            xx = px(x) + 4
        elif align == "right":
            xx = px(x + w) - tw - 4
        else:
            xx = px(x) + (px(w) - tw) // 2
        draw.text((xx, yy), line, fill=cc(color), font=fnt)
        yy += line_h


def add_textbox(x, y, w, h, text, size=12, color=TEXT, bold=False, align="center"):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}[align]
    preview_text(x, y, w, h, text, size, color, bold, align)
    return tb


def add_round_rect(x, y, w, h, fill=WHITE, line=BLUE, lw=1.1, dash=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shp, fill)
    set_line(shp, line, lw, MSO_LINE_DASH_STYLE.DASH if dash else None)
    draw.rounded_rectangle(
        pbox(x, y, w, h),
        radius=10,
        fill=cc(fill),
        outline=cc(line),
        width=max(1, int(lw * 1.5)),
    )
    return shp


def add_badge(x, y, n, color):
    add_round_rect(x, y, 0.24, 0.24, color, color, 0.5)
    add_textbox(x, y + 0.005, 0.24, 0.225, str(n), 10, WHITE, True)


def add_arrow(x1, y1, x2, y2, color=BLUE, width=2.0, dashed=False, two=False):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    set_line(conn, color, width, MSO_LINE_DASH_STYLE.DASH if dashed else None)
    conn.line.end_arrowhead = True
    if two:
        conn.line.begin_arrowhead = True
    draw.line((px(x1), py(y1), px(x2), py(y2)), fill=cc(color), width=max(1, int(width * 1.4)))
    ang = math.atan2(py(y2) - py(y1), px(x2) - px(x1))
    ah = 9
    pts = [
        (px(x2), py(y2)),
        (int(px(x2) - ah * math.cos(ang - 0.45)), int(py(y2) - ah * math.sin(ang - 0.45))),
        (int(px(x2) - ah * math.cos(ang + 0.45)), int(py(y2) - ah * math.sin(ang + 0.45))),
    ]
    draw.polygon(pts, fill=cc(color))


def add_stage(x, y, w, h, idx, title, items, theme="blue"):
    color = BLUE if theme == "blue" else TEAL
    fill = LIGHT_BLUE if theme == "blue" else LIGHT_TEAL
    add_round_rect(x, y, w, h, fill, color, 1.1)
    add_badge(x + 0.10, y + 0.14, idx, color)
    add_textbox(x + 0.40, y + 0.08, w - 0.48, 0.42, title, 10.4, color, True, "left")
    item_h = (h - 0.95) / len(items)
    cy = y + 0.70
    for icon, label in items:
        add_round_rect(x + 0.10, cy, w - 0.20, item_h - 0.09, WHITE, color, 0.7)
        add_textbox(x + 0.18, cy + 0.03, 0.34, item_h - 0.15, icon, 14, color, True)
        add_textbox(x + 0.52, cy + 0.02, w - 0.66, item_h - 0.13, label, 8.4, TEXT, False, "left")
        cy += item_h


add_textbox(3.35, 0.26, 6.60, 0.50, "今世缘大模型智能诊断流程原型图", 25, DARK, True)
for x1, x2 in [(0.35, 1.75), (11.25, 12.95)]:
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(0.55), Inches(x2), Inches(0.55)
    )
    set_line(line, DARK, 1)
    draw.line((px(x1), py(0.55), px(x2), py(0.55)), fill=cc(DARK), width=2)

for x in [1.83, 11.18]:
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(0.515), Inches(0.065), Inches(0.065))
    set_fill(c, DARK)
    c.line.fill.background()
    draw.ellipse(pbox(x, 0.515, 0.065, 0.065), fill=cc(DARK))

stages = [
    ("用户输入", [("...", "文本问题"), ("◎", "业务数据"), ("▣", "文件/表格")], "blue"),
    ("今世缘大模型\n理解层", [("脑", "语义理解"), ("⌖", "意图识别"), ("文", "关键实体抽取"), ("☷", "上下文补全")], "teal"),
    ("任务分类与路由", [("诊", "诊断类"), ("↗", "预测类"), ("!", "风险识别类"), ("▥", "经营分析类"), ("?", "知识问答类")], "blue"),
    ("智能体调度中心", [("▧", "任务拆解"), ("⚒", "工具选择"), ("▣", "模型选择"), ("┬", "流程编排")], "teal"),
    ("垂类模型计算层", [("↗", "销售预测模型"), ("□", "库存分析模型"), ("⌕", "质量诊断模型"), ("⚙", "设备健康模型"), ("●", "用户画像模型"), ("盾", "风险评估模型")], "blue"),
    ("智能结果校验层", [("✓", "结果标准化"), ("盾", "一致性校验"), ("▽", "异常结果过滤"), ("◎", "结果整理汇总")], "teal"),
    ("今世缘大模型\n诊断分析", [("⌘", "多模型结果融合"), ("◎", "置信度评估"), ("△", "异常原因分析"), ("灯", "诊断建议生成"), ("▤", "用户可读报告")], "blue"),
]

x0, y0, w, h, gap = 0.30, 0.98, 1.27, 3.30, 0.32
centers = []
for i, (title, items, theme) in enumerate(stages, start=1):
    x = x0 + (i - 1) * (w + gap)
    add_stage(x, y0, w, h, i, title, items, theme)
    centers.append((x + w / 2, y0 + h))
    if i < len(stages):
        add_arrow(x + w + 0.02, y0 + 1.77, x + w + gap - 0.07, y0 + 1.77, BLUE, 2.1)

out_x = x0 + 7 * (w + gap) + 0.02
add_arrow(x0 + 6 * (w + gap) + w + 0.02, y0 + 1.77, out_x - 0.08, y0 + 1.77, BLUE, 2.1)
add_round_rect(out_x, y0 + 0.50, 0.98, 2.85, RGBColor(244, 249, 255), BLUE, 1.1)
add_round_rect(out_x + 0.17, y0 + 0.70, 0.64, 0.36, BLUE, BLUE, 0.5)
add_textbox(out_x + 0.17, y0 + 0.72, 0.64, 0.30, "输出", 13, WHITE, True)
add_textbox(out_x + 0.23, y0 + 1.18, 0.52, 0.52, "◔\n▥", 22, BLUE, True)
add_textbox(out_x + 0.16, y0 + 1.75, 0.68, 1.15, "详细诊断分析\n原因解释\n风险等级\n改进建议\n下一步行动", 9.3, TEXT, True)

support_y, support_h = 4.88, 1.36
add_round_rect(0.30, support_y, 12.72, support_h, RGBColor(252, 254, 255), BLUE, 0.8, True)
add_textbox(0.42, support_y + 0.28, 0.70, 0.70, "◎\n支撑层", 17, BLUE, True)
sep = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT,
    Inches(1.23),
    Inches(support_y + 0.18),
    Inches(1.23),
    Inches(support_y + support_h - 0.18),
)
set_line(sep, BLUE, 0.8)
draw.line((px(1.23), py(support_y + 0.18), px(1.23), py(support_y + support_h - 0.18)), fill=cc(BLUE), width=1)

supports = [
    ("知识库", ["行业知识", "产品知识", "技术文档", "常见问题"], "书"),
    ("业务规则库", ["业务规则", "计算规则", "风控规则", "策略规则"], "夹"),
    ("历史案例库", ["诊断案例", "处理记录", "解决方案", "经验沉淀"], "档"),
    ("模型监控", ["模型性能监控", "指标监控", "漂移监控", "告警通知"], "屏"),
    ("日志追踪", ["请求日志", "执行日志", "错误日志", "审计日志"], "文"),
    ("人工反馈闭环", ["结果反馈", "问题标注", "规则优化", "模型迭代"], "人"),
]

sx, sw, sg = 1.42, 1.84, 0.18
for i, (title, bullets, icon) in enumerate(supports):
    x = sx + i * (sw + sg)
    add_round_rect(x, support_y + 0.18, sw, 1.00, WHITE, RGBColor(137, 164, 210), 0.7)
    add_textbox(x + 0.10, support_y + 0.32, 0.45, 0.42, icon, 21, BLUE, True)
    add_textbox(x + 0.73, support_y + 0.20, 0.94, 0.25, title, 10.3, BLUE, True)
    add_textbox(x + 0.72, support_y + 0.48, 0.95, 0.58, "• " + "\n• ".join(bullets), 6.8, TEXT, False, "left")

base_y = support_y - 0.24
hline = slide.shapes.add_connector(
    MSO_CONNECTOR.STRAIGHT, Inches(0.98), Inches(base_y), Inches(12.35), Inches(base_y)
)
set_line(hline, BLUE, 0.8, MSO_LINE_DASH_STYLE.DASH)
draw.line((px(0.98), py(base_y), px(12.35), py(base_y)), fill=cc(BLUE), width=1)

for cx, _ in centers:
    v = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(cx), Inches(y0 + h + 0.17), Inches(cx), Inches(base_y)
    )
    set_line(v, BLUE, 0.8, MSO_LINE_DASH_STYLE.DASH)
    v.line.begin_arrowhead = True
    v.line.end_arrowhead = True
    draw.line((px(cx), py(y0 + h + 0.17), px(cx), py(base_y)), fill=cc(BLUE), width=1)

strip_y = 6.48
add_round_rect(0.45, strip_y, 12.45, 0.58, RGBColor(252, 254, 255), BLUE, 0.9)
add_textbox(0.65, strip_y + 0.13, 0.88, 0.30, "流程说明：", 10.5, BLUE, True, "left")
flow = [
    ("1", "用户输入", BLUE),
    ("2", "理解意图", TEAL),
    ("3", "分类路由", BLUE),
    ("4", "调度执行", TEAL),
    ("5", "模型计算", BLUE),
    ("6", "结果校验", TEAL),
    ("7", "诊断分析", BLUE),
]

fx = 1.75
for i, (n, label, color) in enumerate(flow):
    add_badge(fx, strip_y + 0.16, n, color)
    add_textbox(fx + 0.30, strip_y + 0.14, 0.70, 0.28, label, 9.5, TEXT, False, "left")
    fx += 1.30
    if i < len(flow) - 1:
        add_arrow(fx - 0.25, strip_y + 0.30, fx - 0.02, strip_y + 0.30, GRAY, 1.0)

prs.save(OUT)
preview.save(PREVIEW, quality=95)
print(OUT)
print(PREVIEW)
