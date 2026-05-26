from pathlib import Path

from pptx import Presentation


prs = Presentation(str(Path(r"E:\PY\research\0424\Jinshiyuan_AI_Production_Blueprint.pptx")))
for idx in [6, 9]:
    print("SLIDE", idx + 1, "size", prs.slide_width, prs.slide_height)
    for i, shape in enumerate(prs.slides[idx].shapes):
        text = ""
        if hasattr(shape, "text_frame") and shape.text_frame:
            text = shape.text_frame.text.replace("\n", " ")[:80]
        print(i, shape.shape_type, int(shape.left), int(shape.top), int(shape.width), int(shape.height), repr(text))
