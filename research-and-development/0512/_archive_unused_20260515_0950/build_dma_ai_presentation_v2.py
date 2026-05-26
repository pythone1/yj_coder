# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = max([p for p in ROOT.glob("*.pptx") if not p.name.startswith("~$")], key=lambda p: p.stat().st_size)
ASSET_ROOT = ROOT / "output" / "generated_ppt_assets_v6"
SPLIT = ASSET_ROOT / "split_plugins"
RAW = ASSET_ROOT / "raw_sheets"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
OUTPUT = OUT / "AI供水管网DMA漏损检测_模型应用教学汇报_模板重构版.pptx"


NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(246, 250, 253)
TEXT = RGBColor(28, 42, 58)
MUTED = RGBColor(91, 104, 120)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(35, 150, 120)
ORANGE = RGBColor(240, 143, 40)
RED = RGBColor(210, 78, 78)
GRAY_LINE = RGBColor(210, 226, 238)
FONT = "微软雅黑"


def A(folder, name):
    return str(SPLIT / folder / name)


def R(name):
    return str(RAW / name)


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def stroke(shape, color=GRAY_LINE, width=0.75):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def tb(slide, text, x, y, w, h, size=20, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=18, color=TEXT, gap=0.92):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.line_spacing = gap
        r = p.add_run()
        r.text = "· " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def pic_fit(slide, path, x, y, w, h, crop=False):
    pic = slide.shapes.add_picture(path, x, y)
    sx, sy = w / pic.width, h / pic.height
    scale = max(sx, sy) if crop else min(sx, sy)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int(x + (w - pic.width) / 2)
    pic.top = int(y + (h - pic.height) / 2)
    return pic


def rect(slide, x, y, w, h, color=PALE2, line=GRAY_LINE, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    fill(s, color)
    stroke(s, line, 0.8)
    return s


def arrow(slide, x1, y1, x2, y2, color=CYAN, width=1.5):
    a = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    a.line.color.rgb = color
    a.line.width = Pt(width)
    a.line.end_arrowhead = True
    return a


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.25), Inches(0.11), Inches(0.46))
    fill(bar, BLUE)
    tb(slide, title, Inches(0.62), Inches(0.18), Inches(7.8), Inches(0.36), 26, TEXT, True)
    tb(slide, sub, Inches(0.64), Inches(0.62), Inches(8.5), Inches(0.22), 11, MUTED)
    tb(slide, f"{no:02d}", Inches(12.15), Inches(0.22), Inches(0.6), Inches(0.25), 11, BLUE, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.88), Inches(11.7), Inches(0.015))
    fill(line, PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA系统漏损检测 · 教学汇报", Inches(0.62), Inches(7.08), Inches(5.2), Inches(0.18), 10, MUTED)


def content(prs, title, sub, no):
    s = blank(prs)
    header(s, title, sub, no)
    footer(s)
    return s


def cover(prs):
    s = blank(prs, NAVY)
    # Template-like cover: dark left panel + resolved image stage on the right, without cropping.
    left = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(6.35), Inches(7.5))
    fill(left, NAVY)
    accent = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(1.15), Inches(0.13), Inches(4.8))
    fill(accent, CYAN)
    pic_fit(s, R("00_algorithm_overall_architecture.png"), Inches(6.55), Inches(0.68), Inches(6.3), Inches(5.7), crop=False)
    white = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(6.25), Inches(0.45), Inches(7.05), Inches(6.35))
    fill(white, WHITE)
    s.shapes._spTree.remove(white._element)
    pic_fit(s, R("00_algorithm_overall_architecture.png"), Inches(6.55), Inches(0.78), Inches(6.2), Inches(5.55), crop=False)
    tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.34), Inches(5.0), Inches(1.45), 34, WHITE, True)
    tb(s, "模型机理 · 算法选型 · 场景落地 · 持续运营", Inches(1.0), Inches(3.15), Inches(5.0), Inches(0.32), 18, RGBColor(207, 232, 246))
    tb(s, "教学汇报版", Inches(1.0), Inches(5.65), Inches(2.1), Inches(0.28), 15, CYAN, True)
    tb(s, "基于DMA分区计量、时序预测、异常检测与水力模型融合", Inches(1.0), Inches(6.08), Inches(5.1), Inches(0.3), 12, RGBColor(220, 236, 247))


def toc(prs):
    s = blank(prs)
    header(s, "目录", "CONTENTS", 2)
    items = [
        ("PART 01", "DMA漏损检测的业务背景", "讲清DMA、漏损检测对象和数据基础"),
        ("PART 02", "AI模型应用的发展", "从阈值报警到动态基线、多源联动和模型组合"),
        ("PART 03", "核心算法体系", "LSTM、异常检测、监督排序、机理融合与图模型"),
        ("PART 04", "典型应用场景", "事前预警、事中定位、事后复盘、调度与规划"),
        ("PART 05", "模型建设与落地路径", "数据治理、训练验证、漂移监控和规模化推广"),
    ]
    for i, (part, title, desc) in enumerate(items):
        y = Inches(1.25 + i * 1.03)
        tb(s, part, Inches(1.0), y, Inches(1.2), Inches(0.25), 13, CYAN, True)
        tb(s, title, Inches(2.35), y - Inches(0.07), Inches(3.5), Inches(0.35), 22, TEXT, True)
        tb(s, desc, Inches(6.05), y, Inches(5.4), Inches(0.25), 15, MUTED)
        line = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), y + Inches(0.55), Inches(10.9), Inches(0.01))
        fill(line, PALE)
    pic_fit(s, A("04_small_infographic_plugins", "02_water_ai_droplet.png"), Inches(11.25), Inches(5.8), Inches(0.75), Inches(0.75))
    footer(s)


def section(prs, part, title, desc, icon, no):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.22), Inches(7.5)), CYAN)
    tb(s, part, Inches(0.85), Inches(1.30), Inches(1.5), Inches(0.3), 15, CYAN, True)
    tb(s, title, Inches(0.85), Inches(2.05), Inches(7.0), Inches(0.65), 38, WHITE, True)
    tb(s, desc, Inches(0.9), Inches(3.05), Inches(6.8), Inches(0.65), 18, RGBColor(207, 230, 244))
    c = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.25), Inches(1.55), Inches(3.15), Inches(3.15))
    fill(c, RGBColor(8, 63, 120), 12)
    stroke(c, CYAN, 1.0)
    pic_fit(s, icon, Inches(8.78), Inches(2.08), Inches(2.1), Inches(2.1))
    tb(s, f"{no:02d}", Inches(11.7), Inches(6.62), Inches(0.65), Inches(0.3), 14, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)


def two_col(slide, left_title, left_items, right_title, right_items, icon_left=None, icon_right=None):
    if icon_left:
        pic_fit(slide, icon_left, Inches(0.75), Inches(1.35), Inches(1.15), Inches(1.15))
        x1 = Inches(2.1)
    else:
        x1 = Inches(0.9)
    tb(slide, left_title, x1, Inches(1.33), Inches(4.2), Inches(0.35), 23, BLUE, True)
    bullets(slide, left_items, x1, Inches(1.92), Inches(4.75), Inches(2.4), 18, TEXT)
    if icon_right:
        pic_fit(slide, icon_right, Inches(6.65), Inches(1.35), Inches(1.15), Inches(1.15))
        x2 = Inches(7.95)
    else:
        x2 = Inches(6.75)
    tb(slide, right_title, x2, Inches(1.33), Inches(4.1), Inches(0.35), 23, BLUE, True)
    bullets(slide, right_items, x2, Inches(1.92), Inches(4.5), Inches(2.4), 18, TEXT)


def mini_card(slide, x, y, w, h, title, body, color=BLUE, img=None):
    rect(slide, x, y, w, h, WHITE, GRAY_LINE, False)
    if img:
        pic_fit(slide, img, x + Inches(0.12), y + Inches(0.15), Inches(0.62), Inches(0.62))
        tx = x + Inches(0.88)
    else:
        tx = x + Inches(0.18)
    tb(slide, title, tx, y + Inches(0.15), w - Inches(0.3), Inches(0.28), 17, color, True)
    tb(slide, body, x + Inches(0.18), y + Inches(0.58), w - Inches(0.36), h - Inches(0.75), 14, MUTED)


def timeline(slide, items, y):
    colors = [BLUE, CYAN, GREEN, ORANGE, RED]
    x0 = Inches(0.95)
    for i, (title, body) in enumerate(items):
        x = x0 + Inches(i * 2.28)
        c = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(0.82), Inches(0.82))
        fill(c, colors[i])
        tb(slide, str(i + 1), x, y + Inches(0.14), Inches(0.82), Inches(0.25), 17, WHITE, True, PP_ALIGN.CENTER)
        tb(slide, title, x - Inches(0.45), y + Inches(1.05), Inches(1.75), Inches(0.28), 17, TEXT, True, PP_ALIGN.CENTER)
        tb(slide, body, x - Inches(0.52), y + Inches(1.45), Inches(1.9), Inches(0.75), 13, MUTED, False, PP_ALIGN.CENTER)
        if i < 4:
            arrow(slide, x + Inches(0.95), y + Inches(0.42), x + Inches(1.75), y + Inches(0.42), RGBColor(150, 186, 210), 1.2)


def make_deck():
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    cover(prs)
    toc(prs)

    section(prs, "PART 01", "DMA漏损检测的业务背景", "先建立业务对象与数据基础，再讨论模型算法。", A("01_water_equipment_plugins", "04_dma_boundary_map.png"), 3)

    s = content(prs, "DMA系统在漏损管控中的位置", "PART 01 业务背景", 4)
    pic_fit(s, A("01_water_equipment_plugins", "04_dma_boundary_map.png"), Inches(0.75), Inches(1.25), Inches(2.35), Inches(2.35))
    tb(s, "DMA的核心作用", Inches(3.45), Inches(1.20), Inches(3.6), Inches(0.35), 24, BLUE, True)
    bullets(s, ["把大管网切分为可计量、可分析、可考核的分区", "通过入口流量、压力点和边界阀形成运行画像", "将漏损管理从全网粗判断变为分区精细化诊断"], Inches(3.45), Inches(1.75), Inches(4.7), Inches(2.0), 18)
    tb(s, "模型建设前提", Inches(8.45), Inches(1.20), Inches(2.8), Inches(0.35), 24, BLUE, True)
    bullets(s, ["边界清晰", "计量可靠", "压力点有效", "工单可回填"], Inches(8.45), Inches(1.75), Inches(3.0), Inches(1.75), 19)
    tb(s, "结论：DMA不是简单分区，而是AI模型学习正常状态、定位异常范围的基本业务单元。", Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.45), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "传统漏损检测方法与局限", "PART 01 业务背景", 5)
    two_col(
        s,
        "常用方法",
        ["夜间最小流量MNF判断", "固定流量/压力阈值报警", "人工巡检与听漏排查", "历史报修与投诉触发"],
        "主要局限",
        ["分区差异大，统一阈值难以适配", "节假日、天气、商业夜间用水易造成误报", "单点报警缺少证据链", "难以持续沉淀为模型训练样本"],
        A("04_small_infographic_plugins", "01_leak_warning_badge.png"),
        A("04_small_infographic_plugins", "07_anomaly_cluster_dots.png"),
    )
    tb(s, "AI模型的切入点：从固定规则转向动态基线，从单点报警转向多源证据链。", Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.42), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "DMA数据基础：模型能否有效取决于数据是否可用", "PART 01 业务背景", 6)
    cards = [
        ("流量", "入口流量、MNF、日内曲线、残差", A("01_water_equipment_plugins", "01_smart_inlet_flow_meter.png")),
        ("压力", "末梢压力、波动、压力响应", A("01_water_equipment_plugins", "02_pressure_sensor_node.png")),
        ("GIS资产", "管段、阀门、材质、口径、管龄", A("01_water_equipment_plugins", "07_gis_network_map.png")),
        ("工单", "报警、核查、维修、误报原因", A("01_water_equipment_plugins", "08_work_order_clipboard.png")),
        ("外部上下文", "天气、节假日、施工、用户结构", A("04_small_infographic_plugins", "15_cloud_edge_sync.png")),
    ]
    for i, (t, b, img) in enumerate(cards):
        mini_card(s, Inches(0.75 + i * 2.45), Inches(1.35), Inches(2.05), Inches(2.65), t, b, BLUE, img)
    bullets(s, ["建模前必须完成统一编码、时序对齐、缺失处理和工单结构化。", "数据治理不是附属工作，而是模型可信运行的前置条件。"], Inches(1.05), Inches(5.05), Inches(10.8), Inches(0.85), 20, BLUE)

    section(prs, "PART 02", "AI模型应用的发展", "从经验阈值走向动态基线、多源联动和模型组合。", A("04_small_infographic_plugins", "10_timeseries_mini_chart.png"), 7)

    s = content(prs, "从阈值报警到动态基线", "PART 02 发展逻辑", 8)
    timeline(s, [("经验阈值", "固定上下限\nMNF经验值"), ("统计基线", "移动平均\n同比环比"), ("动态预测", "按DMA学习\n预测区间"), ("残差预警", "持续偏离\n触发复核"), ("闭环更新", "工单回填\n迭代模型")], Inches(1.55))
    bullets(s, ["固定阈值适合启动阶段，但难以覆盖不同DMA的运行差异。", "动态基线通过历史流量、压力、日期、天气和用户结构学习正常状态。", "当实际曲线持续偏离预测区间，并与漏损特征吻合时，形成风险预警。"], Inches(1.05), Inches(5.15), Inches(10.8), Inches(1.05), 19)

    s = content(prs, "从单点报警到多源联动识别", "PART 02 发展逻辑", 9)
    icons = [
        (A("01_water_equipment_plugins", "01_smart_inlet_flow_meter.png"), "入口流量"),
        (A("01_water_equipment_plugins", "02_pressure_sensor_node.png"), "压力响应"),
        (A("01_water_equipment_plugins", "07_gis_network_map.png"), "空间资产"),
        (A("01_water_equipment_plugins", "08_work_order_clipboard.png"), "历史工单"),
        (A("01_water_equipment_plugins", "10_ai_model_cube.png"), "综合判断"),
    ]
    for i, (img, lab) in enumerate(icons):
        x = Inches(0.95 + i * 2.3)
        pic_fit(s, img, x, Inches(1.3), Inches(1.1), Inches(1.1))
        tb(s, lab, x - Inches(0.2), Inches(2.58), Inches(1.5), Inches(0.25), 15, TEXT, True, PP_ALIGN.CENTER)
        if i < 4:
            arrow(s, x + Inches(1.35), Inches(1.85), x + Inches(1.9), Inches(1.85), RGBColor(153, 190, 215), 1.2)
    bullets(s, ["漏损事件通常不是单点变化，而是流量抬升、夜间底流增加、末梢压力下降等组合特征。", "模型输出应包含异常时段、主要证据、疑似原因、候选管段和复核方式。", "汇报时应强调证据链和业务可解释性，而不是只展示异常分数。"], Inches(1.05), Inches(4.35), Inches(10.8), Inches(1.35), 19)

    s = content(prs, "从单一算法到模型组合", "PART 02 发展逻辑", 10)
    matrix = [
        ("预测", "LSTM / GRU", "学习正常曲线"),
        ("发现", "孤立森林 / DBSCAN", "识别候选异常"),
        ("排序", "随机森林 / GBDT", "评估管段风险"),
        ("定位", "水力仿真 + ML", "收敛候选管段"),
        ("解释", "知识图谱 / 智能体", "报告与派单辅助"),
    ]
    for i, (a, b, c) in enumerate(matrix):
        mini_card(s, Inches(0.78 + i * 2.43), Inches(1.35), Inches(2.05), Inches(2.75), a, b + "\n" + c, [BLUE, CYAN, GREEN, ORANGE, RED][i])
    tb(s, "核心观点：漏损检测是一条业务链，模型也应是组合体系。", Inches(1.05), Inches(5.35), Inches(10.9), Inches(0.42), 22, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "供水漏损AI模型技术演进路线图", "PART 02 发展逻辑", 11)
    rows = [
        ("经验阈值", "固定上下限、MNF经验阈值", "部署简单，适合基础监测"),
        ("统计分析", "移动平均、控制图、季节分解", "刻画历史偏离和周期变化"),
        ("机器学习", "孤立森林、DBSCAN、GBDT", "少标签异常识别与风险排序"),
        ("深度学习", "LSTM、GRU、自编码器、GNN", "高频时序与拓扑关系学习"),
        ("智能协同", "知识图谱、大模型智能体", "解释、问答、报告和派单辅助"),
    ]
    x0, y0 = Inches(0.7), Inches(1.18)
    widths = [1.6, 3.7, 5.1]
    for i, head in enumerate(["阶段", "主要方法", "应用特点"]):
        rect(s, x0 + sum(Inches(w) for w in widths[:i]), y0, Inches(widths[i]), Inches(0.48), BLUE)
        tb(s, head, x0 + sum(Inches(w) for w in widths[:i]), y0 + Inches(0.13), Inches(widths[i]), Inches(0.18), 14, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y0 + Inches(0.58 + r * 0.75)
        for c, val in enumerate(row):
            xx = x0 + sum(Inches(w) for w in widths[:c])
            rect(s, xx, yy, Inches(widths[c]), Inches(0.58), PALE2 if r % 2 == 0 else WHITE)
            tb(s, val, xx + Inches(0.1), yy + Inches(0.15), Inches(widths[c] - 0.2), Inches(0.20), 13.5, BLUE if c == 0 else TEXT, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)

    section(prs, "PART 03", "核心算法体系", "算法讲解围绕业务任务展开，而不是罗列模型名。", A("02_algorithm_plugins", "14_model_training_pipeline.png"), 12)

    s = content(prs, "LSTM / GRU：时序预测与动态基线", "PART 03 核心算法", 13)
    two_col(s, "适用场景", ["按DMA预测入口流量", "学习日周期、周周期和季节性", "输出预测值、预测区间和残差"], "讲解重点", ["LSTM适合长期依赖", "GRU结构更轻，便于快速迭代", "模型用于建立正常边界，不直接替代现场判定"], A("02_algorithm_plugins", "01_lstm_sequence_prediction.png"), A("02_algorithm_plugins", "02_gru_gate_motif.png"))

    s = content(prs, "CNN-LSTM：多传感器时序特征融合", "PART 03 核心算法", 14)
    pic_fit(s, A("02_algorithm_plugins", "03_cnn_lstm_hybrid.png"), Inches(0.9), Inches(1.35), Inches(2.1), Inches(2.1))
    bullets(s, ["一维卷积先提取多个压力点、流量点之间的局部变化特征。", "循环网络再学习时间依赖，适合高频监测数据。", "在DMA场景中可用于多压力点与入口流量的联合建模。"], Inches(3.35), Inches(1.35), Inches(8.0), Inches(1.6), 20)
    tb(s, "输入层 → 局部特征提取 → 时间依赖学习 → 残差判断 → 预警建议", Inches(1.05), Inches(5.2), Inches(10.8), Inches(0.45), 22, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "孤立森林 / DBSCAN：少标签条件下的异常发现", "PART 03 核心算法", 15)
    two_col(s, "孤立森林", ["适合多维特征快速筛查", "对少标签场景友好", "输出异常分数和候选时段"], "DBSCAN", ["识别低密度离群点", "可发现异常簇", "适合流量-压力特征空间分析"], A("02_algorithm_plugins", "05_isolation_forest_outlier.png"), A("02_algorithm_plugins", "06_dbscan_cluster_outlier.png"))
    tb(s, "注意：无监督模型输出是复核清单，不等同于真实漏损结论。", Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.4), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "自编码器：用重构误差识别异常状态", "PART 03 核心算法", 16)
    pic_fit(s, A("02_algorithm_plugins", "04_autoencoder_reconstruction.png"), Inches(0.95), Inches(1.25), Inches(2.6), Inches(2.6))
    bullets(s, ["模型学习正常状态下多变量之间的组合关系。", "当输入状态与正常模式差异较大时，重构误差升高。", "适合流量、压力、MNF、残差等多变量联合异常检测。", "复核结果可回填为监督模型训练标签。"], Inches(4.0), Inches(1.35), Inches(7.2), Inches(2.4), 20)
    tb(s, "正常样本学习能力越充分，异常识别边界越稳定。", Inches(1.05), Inches(5.45), Inches(10.8), Inches(0.4), 22, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "随机森林 / GBDT / HGB：管段风险排序", "PART 03 核心算法", 17)
    pic_fit(s, A("02_algorithm_plugins", "07_random_forest_ensemble.png"), Inches(0.8), Inches(1.35), Inches(1.35), Inches(1.35))
    pic_fit(s, A("02_algorithm_plugins", "08_gradient_boosting_trees.png"), Inches(0.8), Inches(3.15), Inches(1.35), Inches(1.35))
    bullets(s, ["输入：管龄、材质、口径、压力等级、历史维修、投诉频次、道路等级。", "输出：管段风险分、风险等级、检漏优先级和改造排序。", "优势：对结构化资产和工单数据适配性强，便于解释特征贡献。"], Inches(2.55), Inches(1.3), Inches(8.8), Inches(1.7), 20)
    tb(s, "适用前提：资产台账完整、历史工单结构化、漏损标签口径一致。", Inches(1.0), Inches(5.5), Inches(11.0), Inches(0.4), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "水力模型 + AI：机理与数据融合", "PART 03 核心算法", 18)
    icons = [
        (A("03_system_architecture_modules", "07_hydraulic_simulation_module.png"), "水力仿真"),
        (A("04_small_infographic_plugins", "03_pressure_wave_pulse.png"), "压力响应"),
        (A("03_system_architecture_modules", "08_leak_localization_module.png"), "候选定位"),
        (A("01_water_equipment_plugins", "13_field_inspection_tablet.png"), "现场复核"),
        (A("03_system_architecture_modules", "12_model_retraining_loop.png"), "模型更新"),
    ]
    for i, (img, lab) in enumerate(icons):
        x = Inches(0.95 + i * 2.3)
        pic_fit(s, img, x, Inches(1.35), Inches(1.1), Inches(1.1))
        tb(s, lab, x - Inches(0.2), Inches(2.63), Inches(1.5), Inches(0.25), 15, TEXT, True, PP_ALIGN.CENTER)
        if i < 4:
            arrow(s, x + Inches(1.35), Inches(1.9), x + Inches(1.9), Inches(1.9), RGBColor(153, 190, 215), 1.2)
    bullets(s, ["水力模型提供拓扑、供水路径、压力敏感性和仿真样本。", "AI模型负责从真实监测数据中识别异常模式和候选区域。", "融合后可将定位结果限制在水力学合理范围内。"], Inches(1.05), Inches(4.45), Inches(10.8), Inches(1.0), 19)

    s = content(prs, "GNN与知识图谱：拓扑关系与解释能力", "PART 03 核心算法", 19)
    two_col(s, "图神经网络GNN", ["管网天然是图结构", "节点为压力点、阀门、水池和用户", "边为管段，可学习拓扑邻接和压力传播"], "知识图谱", ["连接DMA、管段、设备、工单和案例", "支撑异常原因解释与历史案例检索", "可与智能体结合生成报告和问答"], A("02_algorithm_plugins", "10_gnn_pipe_topology.png"), A("02_algorithm_plugins", "11_knowledge_graph.png"))

    s = content(prs, "算法选型矩阵", "PART 03 核心算法", 20)
    rows = [
        ("动态基线", "LSTM / GRU / CNN-LSTM", "流量、压力、天气、日历", "预测区间、残差"),
        ("异常发现", "孤立森林 / DBSCAN / 自编码器", "MNF、残差、多维时序", "异常分数、异常时段"),
        ("风险排序", "随机森林 / GBDT / HGB", "管龄、材质、维修、投诉", "风险分、优先级"),
        ("定位收敛", "水力仿真 + ML / GNN", "拓扑、多点压力、仿真样本", "候选管段TopN"),
        ("业务协同", "知识图谱 / 智能体", "模型输出、GIS、工单", "报告、问答、派单建议"),
    ]
    x0, y0 = Inches(0.55), Inches(1.18)
    widths = [1.55, 3.3, 3.1, 2.65]
    for i, head in enumerate(["任务", "推荐模型", "输入数据", "输出结果"]):
        xx = x0 + sum(Inches(w) for w in widths[:i])
        rect(s, xx, y0, Inches(widths[i]), Inches(0.48), BLUE)
        tb(s, head, xx, y0 + Inches(0.13), Inches(widths[i]), Inches(0.18), 14, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y0 + Inches(0.58 + r * 0.78)
        for c, val in enumerate(row):
            xx = x0 + sum(Inches(w) for w in widths[:c])
            rect(s, xx, yy, Inches(widths[c]), Inches(0.6), PALE2 if r % 2 == 0 else WHITE)
            tb(s, val, xx + Inches(0.08), yy + Inches(0.15), Inches(widths[c] - 0.16), Inches(0.22), 13, BLUE if c == 0 else TEXT, c == 0, PP_ALIGN.CENTER if c == 0 else PP_ALIGN.LEFT)
    tb(s, "选型原则：先明确业务目标，再确定算法路线；复杂模型必须建立在可靠数据和稳定边界之上。", Inches(0.8), Inches(6.25), Inches(11.3), Inches(0.3), 16, BLUE, True, PP_ALIGN.CENTER)

    section(prs, "PART 04", "典型应用场景", "模型价值体现在预警、定位、复盘和调度等业务动作中。", A("01_water_equipment_plugins", "05_pipe_leak_alert.png"), 21)

    s = content(prs, "事前预警：DMA异常识别", "PART 04 应用场景", 22)
    timeline(s, [("建基线", "日内/周周期\n季节性"), ("算残差", "实际值-预测值\n持续跟踪"), ("看压力", "末梢响应\n波动变化"), ("定等级", "异常评分\n风险分级"), ("给动作", "复核方式\n派单建议")], Inches(1.45))
    bullets(s, ["预警结果应包含异常时间、持续时长、残差幅度、压力响应、相似历史事件和建议动作。", "目标是提前发现风险，而不是用模型直接替代现场核查。"], Inches(1.05), Inches(5.12), Inches(10.8), Inches(0.9), 20)

    s = content(prs, "事中定位：候选管段TopN收敛", "PART 04 应用场景", 23)
    pic_fit(s, A("01_water_equipment_plugins", "04_dma_boundary_map.png"), Inches(0.85), Inches(1.2), Inches(2.3), Inches(2.3))
    arrow(s, Inches(3.25), Inches(2.35), Inches(4.15), Inches(2.35), RGBColor(153, 190, 215), 1.3)
    pic_fit(s, A("03_system_architecture_modules", "08_leak_localization_module.png"), Inches(4.3), Inches(1.2), Inches(2.3), Inches(2.3))
    arrow(s, Inches(6.75), Inches(2.35), Inches(7.65), Inches(2.35), RGBColor(153, 190, 215), 1.3)
    pic_fit(s, A("01_water_equipment_plugins", "07_gis_network_map.png"), Inches(7.8), Inches(1.2), Inches(2.3), Inches(2.3))
    tb(s, "整个DMA", Inches(0.95), Inches(3.78), Inches(2.1), Inches(0.28), 18, TEXT, True, PP_ALIGN.CENTER)
    tb(s, "候选片区", Inches(4.4), Inches(3.78), Inches(2.1), Inches(0.28), 18, TEXT, True, PP_ALIGN.CENTER)
    tb(s, "候选管段TopN", Inches(7.8), Inches(3.78), Inches(2.4), Inches(0.28), 18, TEXT, True, PP_ALIGN.CENTER)
    bullets(s, ["候选项应包含：管段编号、空间位置、置信度、主要证据、影响用户数、建议复核方式。", "这种输出更符合现场检漏和派单实际。"], Inches(1.05), Inches(5.1), Inches(10.8), Inches(0.85), 20, BLUE)

    s = content(prs, "事后复盘：工单回填与持续学习", "PART 04 应用场景", 24)
    timeline(s, [("报警", "异常特征\n模型版本"), ("核查", "现场结果\n误报原因"), ("维修", "漏点位置\n漏损类型"), ("验证", "MNF回落\n压力恢复"), ("训练", "标签修正\n模型迭代")], Inches(1.45))
    bullets(s, ["每一次报警、核查、维修和误报都应形成结构化记录。", "复盘结果用于更新特征、调整阈值、修正标签和优化模型。"], Inches(1.05), Inches(5.12), Inches(10.8), Inches(0.9), 20)

    s = content(prs, "压力优化：降漏、稳压、节能平衡", "PART 04 应用场景", 25)
    two_col(s, "模型输入", ["需水预测", "压力监测", "水力模型", "泵阀状态"], "约束条件", ["最不利点服务压力", "消防保障", "二次供水补水", "用户体验"], A("01_water_equipment_plugins", "14_pump_station.png"), A("02_algorithm_plugins", "09_genetic_algorithm_optimization.png"))
    tb(s, "目标：在保障供水安全前提下，寻找降漏、节能和稳压之间的平衡。", Inches(1.05), Inches(5.55), Inches(10.8), Inches(0.4), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "DMA规划：分区与监测点优化", "PART 04 应用场景", 26)
    two_col(s, "分区方案评价", ["拓扑结构", "地形高差", "用户数量", "入口数量与边界阀数量", "施工成本"], "监测点优化", ["仿真不同漏点情景", "比较压力点组合贡献", "提升定位效果", "控制建设预算"], A("04_small_infographic_plugins", "05_dma_boundary_icon.png"), A("01_water_equipment_plugins", "02_pressure_sensor_node.png"))

    section(prs, "PART 05", "模型建设与落地路径", "把模型做成可运行、可验证、可迭代的工程系统。", A("03_system_architecture_modules", "03_data_cleaning_filter.png"), 27)

    s = content(prs, "数据治理：五张基础表", "PART 05 落地路径", 28)
    tables = [
        ("设备表", "流量计、压力计、RTU、边缘网关"),
        ("管网表", "管段、阀门、DMA归属、资产属性"),
        ("时序表", "采样时间、流量压力、缺失标记、清洗版本"),
        ("工单表", "报警、核查、维修、处置结果"),
        ("标签表", "真实漏损、误报原因、模型版本"),
    ]
    for i, (t, b) in enumerate(tables):
        mini_card(s, Inches(0.75 + i * 2.43), Inches(1.32), Inches(2.05), Inches(2.65), t, b, BLUE)
    tb(s, "治理重点：统一编码、时序对齐、缺失处理、异常毛刺剔除、工单结构化、标签口径统一。", Inches(1.0), Inches(5.45), Inches(11.0), Inches(0.45), 20, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "模型训练路径：基线 - 异常 - 定位 - 排序", "PART 05 落地路径", 29)
    timeline(s, [("MNF基线", "形成基础\n预警能力"), ("动态基线", "LSTM/GRU\n预测残差"), ("异常检测", "无监督模型\n交叉验证"), ("定位模型", "水力仿真\n候选管段"), ("风险排序", "监督学习\n检漏优先级")], Inches(1.45))
    bullets(s, ["先建立可解释、可复核的基础能力，再逐步引入复杂模型。", "训练过程应采用时间切分验证，避免未来信息泄漏。"], Inches(1.05), Inches(5.12), Inches(10.8), Inches(0.9), 20)

    s = content(prs, "模型验证：算法指标与业务指标并重", "PART 05 落地路径", 30)
    two_col(s, "算法指标", ["MAE / RMSE", "精确率 / 召回率", "F1 / AUC", "TopN命中率"], "业务指标", ["提前预警时间", "误报工单比例", "平均排查范围", "维修后MNF回落", "复盘回填率"], A("02_algorithm_plugins", "13_prediction_residual_chart.png"), A("03_system_architecture_modules", "16_executive_dashboard.png"))
    tb(s, "对外汇报中，业务指标通常比单一准确率更能体现应用价值。", Inches(1.05), Inches(5.55), Inches(10.8), Inches(0.4), 21, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "模型运营：漂移监控与再训练", "PART 05 落地路径", 31)
    pic_fit(s, A("02_algorithm_plugins", "15_model_drift_monitor.png"), Inches(0.85), Inches(1.25), Inches(2.25), Inches(2.25))
    bullets(s, ["季节变化、用户结构变化、管网改造和传感器更换都会改变数据分布。", "建议每月复盘误报警报，每季度检查模型漂移。", "每半年结合新增工单再训练，重大管网变更后重新评估基线。"], Inches(3.55), Inches(1.32), Inches(7.8), Inches(1.8), 20)
    tb(s, "模型上线不是结束，而是运营开始。", Inches(1.05), Inches(5.45), Inches(10.8), Inches(0.4), 23, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "试点到规模化推广路线", "PART 05 落地路径", 32)
    timeline(s, [("选点", "边界清晰\n数据稳定"), ("治理", "数据表\n标签口径"), ("试运行", "预警复核\n误报分析"), ("评估", "业务指标\n模型指标"), ("推广", "标准流程\n分区复制")], Inches(1.45))
    bullets(s, ["试点优先选择数据质量较好、漏损治理需求明确、现场反馈机制顺畅的DMA。", "规模化推广依赖标准化数据接口、模型版本管理和工单闭环机制。"], Inches(1.05), Inches(5.12), Inches(10.8), Inches(0.9), 20)

    s = content(prs, "总结：以业务闭环定义AI模型价值", "SUMMARY", 33)
    mini_card(s, Inches(0.9), Inches(1.35), Inches(3.2), Inches(2.6), "建设顺序", "先数据，再基线；先预警，再定位；先复核，再自动化。", BLUE, A("04_small_infographic_plugins", "16_closed_loop_arrows.png"))
    mini_card(s, Inches(5.0), Inches(1.35), Inches(3.2), Inches(2.6), "模型组合", "时序预测建立边界，异常检测发现候选，监督学习排序风险，水力模型约束定位。", GREEN, A("01_water_equipment_plugins", "10_ai_model_cube.png"))
    mini_card(s, Inches(9.1), Inches(1.35), Inches(3.2), Inches(2.6), "落地关键", "证据链可解释，工单可回填，模型可监控，流程可闭环。", ORANGE, A("02_algorithm_plugins", "16_human_in_loop_validation.png"))
    tb(s, "供水管网DMA漏损检测的AI应用，应从算法展示转向业务闭环：能发现、能解释、能派单、能复盘、能持续优化。", Inches(1.0), Inches(5.45), Inches(11.0), Inches(0.55), 22, BLUE, True, PP_ALIGN.CENTER)

    s = content(prs, "附：教学演讲中的素材使用建议", "APPENDIX", 34)
    bullets(s, ["封面和章节页使用大图，正文页以局部插件和结构图为主。", "算法页不使用抽象背景，优先展示模型输入、处理逻辑和业务输出。", "架构页采用分层方式讲解，不把大图整体塞满。", "每页控制一个主问题：是什么、为什么、怎么做、输出什么。"], Inches(0.95), Inches(1.35), Inches(7.1), Inches(2.4), 20)
    pic_fit(s, A("02_algorithm_plugins", "14_model_training_pipeline.png"), Inches(8.8), Inches(1.35), Inches(1.35), Inches(1.35))
    pic_fit(s, A("03_system_architecture_modules", "15_digital_twin_network.png"), Inches(10.2), Inches(1.35), Inches(1.35), Inches(1.35))
    pic_fit(s, A("04_small_infographic_plugins", "08_ai_chip_icon.png"), Inches(9.5), Inches(3.1), Inches(1.35), Inches(1.35))
    tb(s, "这份PPT应作为现场讲解底稿，不作为长篇阅读文档。", Inches(1.0), Inches(5.55), Inches(11.0), Inches(0.4), 21, BLUE, True, PP_ALIGN.CENTER)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(make_deck())
