# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme color definitions
    BG_COLOR = RGBColor(12, 17, 30)      # Dark blue #0C111E
    TEXT_TITLE = RGBColor(255, 255, 255) # White
    TEXT_BODY = RGBColor(203, 213, 225)  # Light slate #cbd5e1
    ACCENT_TEAL = RGBColor(13, 148, 136) # Teal #0d9488
    
    # ------------------ SLIDE 1: COVER ------------------
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # Add top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()
    
    # Main Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "中学物理电学与力学综合题解析报告"
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "第 15、16、17、18、19、25 题详细步骤分析与多方案解答"
    p2.font.name = 'Microsoft YaHei'
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_TEAL
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    
    # Footer Info
    p3 = tf.add_paragraph()
    p3.text = "\n报告生成形式：本地 Python-PPTX 导出 (100% 可编辑)"
    p3.font.name = 'Microsoft YaHei'
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_BODY
    p3.alignment = PP_ALIGN.CENTER
    
    # Helper to add standard slides
    def add_standard_slide(title, subtitle, bullets):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        # Left accent line decoration
        dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.6))
        dec.fill.solid()
        dec.fill.fore_color.rgb = ACCENT_TEAL
        dec.line.fill.background()
        
        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(12.0), Inches(0.5))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Microsoft YaHei'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_TITLE
        
        # Slide Subtitle
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Microsoft YaHei'
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = ACCENT_TEAL
        p_sub.font.bold = True
        
        # Bullets Body
        b_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.0))
        tf_b = b_box.text_frame
        tf_b.word_wrap = True
        
        for i, text in enumerate(bullets):
            p_b = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
            p_b.text = text
            p_b.font.name = 'Microsoft YaHei'
            p_b.font.size = Pt(14)
            p_b.font.color.rgb = TEXT_BODY
            p_b.level = 0
            p_b.space_after = Pt(10)
            
        return slide

    # ------------------ SLIDE 2: Q25 ------------------
    add_standard_slide(
        "第 25 题：液体密度杠杆秤（力矩平衡）",
        "力矩平衡方程与秤的灵敏度调精分析",
        [
            "1. 调平方向：图甲中杠杆左端下沉（左重右轻），应将平衡螺母向右调节，使杠杆在水平位置平衡。",
            "2. 力矩平衡方程计算：\n   由于两空容器挂在杠杆两侧等长 10cm 位置，空容器重力力矩在两侧相互抵消。\n   m_水 × 10 = m_液 × 10 + m_A × l_A  (已知 m_A = 10 g, l_A = 9 cm)\n   30 g × 10 = m_液 × 10 + 10 g × 9  =>  10 × m_液 = 210  =>  m_液 = 21 g\n   液体的密度为：ρ_液 = m_液 / V = 21 g / 30 cm³ = 0.7 g/cm³。",
            "3. 函数关系与图线选择：\n   ρ_液 = ρ_水 - [ m_A / (10·V) ] × l_A。该函数是一条斜率小于 0 的直线，因此在图丙中选择图线 b。",
            "4. 提高测量精确度：\n   密度变化 Δρ 对应的位移改变量 Δl = [ 10V / m_A ] × Δρ。为使刻度更稀疏（大位移），应减小重物 A 的质量 m_A。当重物变轻，需要移动更长的力臂来平衡微小的重量差，从而提高灵敏度。"
        ]
    )

    # ------------------ SLIDE 3: Q18 ------------------
    add_standard_slide(
        "第 18 题：多开关控制电路（并/串联）",
        "定值电阻阻值与灯泡产生的热量之比计算",
        [
            "1. 灯泡额定参数分析：\n   灯泡 L 标有 '6V 3W'，额定电流 I_L = 3W / 6V = 0.5 A，灯丝电阻 R_L = 6V / 0.5A = 12 Ω。",
            "2. 状态 1（S1, S2, S3 均闭合）：\n   L 被短路，电阻 R1 与 R2 并联在电源 U 两端。I1 = 1 A (流过 R1)，I_总 = 1.5 A。\n   流过 R2 的电流为：I2 = 1.5 A - 1 A = 0.5 A。\n   电源电压关系式：U = I1 × R1 = 1 × R1，R2 = U / I2 = U / 0.5 = 2U = 2R1。",
            "3. 状态 2（仅闭合 S1）：\n   R2 支路开路，灯泡 L 与 R1 串联。此时 L 正常发光，电路中电流为 I = 0.5 A，U_L = 6 V。\n   电源电压为：U = U_L + I × R1  =>  R1 = 6 + 0.5 × R1  =>  0.5 × R1 = 6  =>  R1 = 12 Ω。\n   代入可求：电源电压 U = 12 V，电阻 R2 = 24 Ω。",
            "4. 前后两次 R1 产生的热量之比：\n   Q1 : Q2 = I1² · R1 · t : I² · R1 · t = 1.0² : 0.5² = 4 : 1。"
        ]
    )

    # ------------------ SLIDE 4: Q19 ------------------
    add_standard_slide(
        "第 19 题：电饭锅双档功耗（串联电路阻值）",
        "高温与保温功率状态、电能消耗及功率提升分析",
        [
            "1. 保温电流与阻值分析：\n   U = 220 V，P_高 = 1100 W (接 1，仅 R1 工作)；P_保 = 22 W (接 2，R1 与 R2 串联)。\n   保温时电流：I_保 = P_保 / U = 22 W / 220 V = 0.1 A。",
            "2. 阻值计算：\n   高温档 R1 阻值：R1 = U² / P_高 = 220² / 1100 = 44 Ω。\n   保温档总电阻：R_总 = U² / P_保 = 220² / 22 = 2200 Ω。\n   定值电阻 R2 = R_总 - R1 = 2200 - 44 = 2156 Ω。",
            "3. 电能表闪烁次数计算：\n   高温档工作 10 min 消耗电能：W = P_高 × t = 1.1 kW × (10/60) h = 11/60 kWh。\n   电能表闪烁次数：N = (11/60 kWh) × 1200 imp/kWh = 220 次。",
            "4. 功率提升后的阻值之比：\n   为使两档功率均提升 10%，设新的阻值 R1' = R1 / 1.1，R2' = R2 / 1.1。\n   新阻值之比 R1' : R2' = (R1/1.1) : (R2/1.1) = R1 : R2，因此 R1' 和 R2' 的阻值之比保持不变。"
        ]
    )

    # ------------------ SLIDE 5: Q15 ------------------
    add_standard_slide(
        "第 15 题：恒流源驱动并联电路安全范围",
        "并联电路安全工作约束与总功率变化范围",
        [
            "1. 恒流源特点与并联电流：\n   恒流源输出恒定电流 I_总 = 1 A。定值电阻 R1 = 5 Ω，滑动变阻器 R2 最大阻值 25 Ω（允许最大电流 1 A）。\n   若 U = 2.6 V 时，R1 电流为 2.6V / 5Ω = 0.52 A，R2 电流为 1A - 0.52A = 0.48 A，说明电流表串联在 R2 支路上。",
            "2. 安全范围约束条件：\n   并联两端电压 U 受限：① 电压表量程限制：U_max = 3 V；\n   ② 电流表量程限制：流过 R2 的电流 I2 = 1 - U/R1 = 1 - U/5 ≤ 0.6 A  =>  U/5 ≥ 0.4  =>  U_min = 2 V。\n   因此，并联两端的电压安全范围为 2 V ~ 3 V。",
            "3. 总功率变化范围计算：\n   电路总功率：P_总 = U × I_总 = U × 1 A。\n   将电压范围代入，得到总功率的安全范围为：2 W ~ 3 W。"
        ]
    )

    # ------------------ SLIDE 6: Q16 ------------------
    add_standard_slide(
        "第 16 题：电路状态切换与灯泡参数推导",
        "串并联功率之比与灯泡额定电压多方案推导",
        [
            "1. 总功率之比 P1 : P2：\n   R1 = 2R2。S 断开时串联，R_串 = 3R2，P1 = U² / 3R2；S 闭合且电源对调时并联，R_联 = 2/3 R2，P2 = 1.5 U²/R2。\n   因此，P1 : P2 = (1/3) : (3/2) = 2 : 9。",
            "2. 灯正常发光求额定电压 U_L (U=12V, R1=8Ω, P_L=4W, R_L > 8Ω)：\n   【方案 1 (以电压为变量)】：R1 两端电压为 12 - U_L，电流 I = (12 - U_L) / 8。\n   P_L = U_L · (12 - U_L) / 8 = 4  =>  U_L² - 12U_L + 32 = 0  =>  U_L = 8V 或 4V。\n   若 U_L = 4V，R_L = 4²/4 = 4 Ω < 8 Ω (舍)；若 U_L = 8V，R_L = 8²/4 = 16 Ω > 8 Ω (合)。答案为 8 V。",
            "   【方案 2 (以电阻为变量)】：电路电流 I = 12 / (8 + R_L)，灯功率 P_L = I² · R_L = 144 R_L / (8+R_L)² = 4。\n   展开得：R_L² - 20 R_L + 64 = 0  =>  R_L = 16 Ω 或 4 Ω。由于 R_L > 8 Ω，选 R_L = 16 Ω。\n   额定电压：U_L = sqrt(P_L · R_L) = sqrt(4 × 16) = 8 V。",
            "   【方案 3 (估算特值法)】：由于 R_L > 8 Ω，在串联电路中分压超过一半，即 U_L > 6 V。尝试常见值 U_L = 8 V，则 R1 两端电压为 4 V，电流 I = 4V/8Ω = 0.5 A，实际功率 P = 8V × 0.5A = 4 W，与题设完全吻合，直接锁定答案。"
        ]
    )

    # ------------------ SLIDE 7: Q17 ------------------
    add_standard_slide(
        "第 17 题：电路功率随电流变化曲线分析",
        "电源、定值电阻及滑动变阻器功率变化与最值求解",
        [
            "1. 图线分析与电源电动势：\n   电源总功率 P_总 = E · I (正比例直线)，对应图乙的图线 a。代入点 (3A, 9W)，得 E = 9/3 = 3 V。\n   定值电阻功率 P0 = I² · R0 (开口向上抛物线)，对应图乙的图线 b。代入点 (3A, 9W)，得 R0 = 9/3² = 1 Ω。\n   滑动变阻器功率 P_R = I(E - I R0) = 3I - I² (开口向下抛物线)，对应图乙的图线 c。",
            "2. 滑动变阻器最大功率求解方法一 (配方法)：\n   P_R = 3I - I² = -(I - 1.5)² + 2.25。当 I = 1.5 A 时，P_R 最大值为 2.25 W。",
            "3. 滑动变阻器最大功率求解方法二 (等效内阻法)：\n   在串联电路中，当滑动变阻器的阻值 R 等于其余部分等效电阻（这里即为 R0 = 1 Ω）时，变阻器功率最大。\n   此时 R = R0 = 1 Ω，总电阻 R_总 = 2 Ω，电流 I = 3V / 2Ω = 1.5 A。\n   滑动变阻器的最大功率为：P_Rmax = I² · R = 1.5² × 1 = 2.25 W。"
        ]
    )
    
    # ------------------ SLIDE 8: END ------------------
    slide_end = prs.slides.add_slide(blank_layout)
    slide_end.background.fill.solid()
    slide_end.background.fill.fore_color.rgb = BG_COLOR
    
    # Bottom line
    bar_end = slide_end.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(13.333), Inches(0.2))
    bar_end.fill.solid()
    bar_end.fill.fore_color.rgb = ACCENT_TEAL
    bar_end.line.fill.background()
    
    # Main Title
    txBox_end = slide_end.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf_end = txBox_end.text_frame
    tf_end.word_wrap = True
    p_end = tf_end.paragraphs[0]
    p_end.text = "谢谢您的观看与审阅"
    p_end.font.name = 'Microsoft YaHei'
    p_end.font.size = Pt(38)
    p_end.font.bold = True
    p_end.font.color.rgb = ACCENT_TEAL
    p_end.alignment = PP_ALIGN.CENTER
    
    p_sub_end = tf_end.add_paragraph()
    p_sub_end.text = "如有疑问，欢迎随时讨论与指正！"
    p_sub_end.font.name = 'Microsoft YaHei'
    p_sub_end.font.size = Pt(18)
    p_sub_end.font.color.rgb = TEXT_TITLE
    p_sub_end.alignment = PP_ALIGN.CENTER
    
    prs.save("物理综合题解析报告.pptx")
    print("Successfully generated physics PPTX!")

if __name__ == '__main__':
    create_presentation()
