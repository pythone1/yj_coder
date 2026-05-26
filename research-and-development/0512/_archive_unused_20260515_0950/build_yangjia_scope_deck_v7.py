# -*- coding: utf-8 -*-
from pathlib import Path

import build_teaching_deck_v6 as b
from pptx import Presentation


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)

b.PPTX = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_v7.pptx"
b.AUDIT = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_v7_逐页检查.txt"
b.PREVIEW = OUT / "AI供水管网DMA漏损检测_杨佳负责部分_v7_预览联系表.png"
b.PREVIEW_DIR = OUT / "v7_page_previews"
b.PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
b.AUDIT_LINES.clear()


def cover(prs):
    s = b.blank(prs, b.NAVY)
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(0.62), b.emu(0.82), b.emu(0.12), b.emu(5.90)), b.CYAN)
    b.tb(s, "杨佳负责部分\nAI模型、核心算法\n与实施路径", 1.02, 1.08, 4.55, 1.78, 33, b.WHITE, True)
    b.tb(s, "供水管网DMA系统漏损检测", 1.05, 3.25, 3.60, 0.32, 18, b.CYAN, True)
    b.tb(s, "对应提纲：1.3、2.2、第四部分实施路径", 1.05, 5.74, 4.30, 0.30, 15, b.RGBColor(215, 236, 248))
    for i, (txt, col) in enumerate([("AI发展", b.BLUE), ("核心算法", b.TEAL), ("工程落地", b.ORANGE)]):
        b.chip(s, txt, 1.05 + i * 1.28, 6.18, 1.05, col)
    b.pic(s, "positioning", 5.72, 0.62, 7.20, 6.20, crop=False)
    border = s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, b.emu(5.62), b.emu(0.52), b.emu(7.40), b.emu(6.40))
    border.fill.background()
    b.stroke(border, b.RGBColor(8, 72, 140), 1.2)
    b.audit(s, 1, allow_short=True)


def scope(prs):
    s = b.content_slide(prs, 2, "本部分讲解范围：只覆盖杨佳负责内容", "按原始提纲拆成三组，不再扩展无关背景")
    items = [
        ("1.3 AI模型应用发展", "AI模型与智能体发展\n漏损检测适配算法谱系", b.BLUE),
        ("2.2 核心AI技术与价值", "时序异常检测\n水力模型融合\n大数据智能决策", b.TEAL),
        ("第四部分 实施路径", "规划边界\n数据治理\n模型建设\n工程落地\n长效运营", b.ORANGE),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 0.80 + i * 4.10
        b.rect(s, x, 1.45, 3.45, 4.45, b.WHITE, col, True)
        b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x), b.emu(1.45), b.emu(3.45), b.emu(0.48)), col)
        b.tb(s, h, x + 0.12, 1.59, 3.20, 0.16, 15, b.WHITE, True, b.PP_ALIGN.CENTER)
        b.tb(s, body, x + 0.35, 2.35, 2.75, 1.35, 20, b.TEXT, True, b.PP_ALIGN.CENTER)
    b.tb(s, "讲课主线：为什么需要AI模型 → 哪些算法适配漏损检测 → 项目如何从试点走向闭环运营。", 0.90, 6.42, 11.70, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 2)


def ai_evolution(prs):
    s = b.content_slide(prs, 4, "1.3 供水管网AI模型应用的发展", "从固定阈值到智能体协同，核心变化是模型从“报警器”变成“诊断助手”")
    stages = [
        ("规则阈值", "MNF固定阈值\n人工经验判断", b.BLUE),
        ("机器学习", "孤立森林\nDBSCAN\n随机森林", b.TEAL),
        ("深度学习", "LSTM/GRU\n自编码器\n时序残差", b.ORANGE),
        ("机理融合", "水力仿真\nGA校核\n候选管段排序", b.PURPLE),
        ("智能体协同", "知识检索\n工单复盘\n调度建议", b.RED),
    ]
    for i, (h, body, col) in enumerate(stages):
        x = 0.70 + i * 2.46
        b.oval(s, x + 0.64, 1.32, 0.70, 0.70, col)
        b.tb(s, str(i + 1), x + 0.64, 1.51, 0.70, 0.18, 15, b.WHITE, True, b.PP_ALIGN.CENTER)
        b.tb(s, h, x + 0.20, 2.28, 1.62, 0.28, 18, col, True, b.PP_ALIGN.CENTER)
        b.tb(s, body, x + 0.08, 2.76, 1.84, 0.78, 14.2, b.TEXT, False, b.PP_ALIGN.CENTER)
        if i < 4:
            b.line(s, x + 1.52, 1.66, x + 2.34, 1.66, b.RGBColor(160, 198, 222), 1.8)
    b.rect(s, 1.05, 4.75, 11.25, 1.10, b.PALE2, b.BLUE, True)
    b.tb(s, "讲解重点：AI不是替代水力模型，而是把多源数据、历史经验和现场处置连接成持续学习的闭环。", 1.45, 5.12, 10.45, 0.22, 18.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 4)


def algorithm_family(prs, no=7):
    s = b.content_slide(prs, no, "供水管网漏损检测适配算法谱系", "按任务归类，避免把算法名称简单堆叠")
    groups = [
        ("时序预测", "LSTM/GRU", "学习流量、压力的时间规律，形成动态基线。", b.BLUE),
        ("无监督异常", "孤立森林 / DBSCAN / 自编码器", "在少标签阶段识别离群、突变和重构误差。", b.TEAL),
        ("参数校核", "遗传算法", "反演粗糙度、节点需水量、背景漏失等不确定参数。", b.ORANGE),
        ("定位分类", "随机森林 / 梯度提升树", "基于仿真样本和实测响应输出候选漏点或管段。", b.PURPLE),
        ("智能决策", "知识图谱 / 边缘计算 / 云端AI", "连接资产、工单和模型，支撑调度与全生命周期管理。", b.GREEN),
    ]
    for i, (h, alg, body, col) in enumerate(groups):
        x = 0.75 + (i % 2) * 6.00
        y = 1.20 + (i // 2) * 1.72
        w = 5.25 if i < 4 else 11.25
        b.card(s, x, y, w, 1.28, h, f"{alg}\n{body}", col)
    b.audit(s, no)


def agent_development(prs, no=5):
    s = b.content_slide(prs, no, "1.3 AI模型与智能体发展：从识别异常到辅助决策", "智能体不是单一算法，而是“模型 + 知识 + 工具 + 流程”的组合")
    items = [
        ("感知模型", "接入流量、压力、水质、阀门、工单等数据，识别异常信号。", b.BLUE),
        ("知识库", "沉淀管网资产、DMA边界、维修案例、调度规则和专家经验。", b.TEAL),
        ("工具调用", "联动GIS、水力模型、SCADA、工单系统，完成查询和计算。", b.ORANGE),
        ("人机协同", "生成候选原因、核查路线和派单建议，由调度人员确认执行。", b.PURPLE),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 0.95 + (i % 2) * 5.80
        y = 1.32 + (i // 2) * 1.95
        b.card(s, x, y, 5.05, 1.45, h, body, col)
    b.rect(s, 1.15, 5.70, 11.00, 0.72, b.PALE2, b.BLUE, True)
    b.tb(s, "在DMA漏损场景中，智能体的价值是把“报警、分析、定位、派单、复盘”组织成可追踪的工作流。", 1.42, 5.96, 10.45, 0.18, 17, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def selection_rules(prs):
    s = b.content_slide(prs, 8, "算法选型原则：先匹配场景，再选择模型复杂度", "对外讲解时强调适配性，避免让听众感觉是在堆名词")
    rows = [
        ("数据标签少", "孤立森林、DBSCAN、自编码器", "先做异常筛查，减少人工盲查"),
        ("时序规律强", "LSTM/GRU", "建立动态基线，识别持续偏离"),
        ("拓扑影响明显", "水力模型 + 树模型", "把异常从片区收敛到候选管段"),
        ("需要解释与派单", "风险评分 + 规则引擎", "输出可复核的优先级和原因"),
        ("多系统联动", "知识图谱 + 智能体", "连接工单、资产、调度和复盘"),
    ]
    widths = [2.15, 3.35, 5.75]
    x0, y0 = 0.85, 1.25
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.25), b.emu(0.42)), b.BLUE)
    for i, h in enumerate(["适用条件", "推荐技术", "输出价值"]):
        b.tb(s, h, x0 + sum(widths[:i]) + 0.05, y0 + 0.13, widths[i] - 0.10, 0.13, 12.8, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.42 + r * 0.78
        b.rect(s, x0, y, 11.25, 0.72, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(224, 239, 248), False)
        for i, txt in enumerate(row):
            col = [b.BLUE, b.TEAL, b.TEXT][i]
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.10, y + 0.20, widths[i] - 0.20, 0.20, 13.3, col, i < 2, b.PP_ALIGN.CENTER)
    b.tb(s, "教学表达：模型选择不追求“最先进”，而要能解释、能部署、能回填、能持续优化。", 1.05, 6.22, 11.20, 0.24, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 8)


def lstm_process_detail(prs, no=12):
    s = b.content_slide(prs, no, "2.2.1 LSTM/GRU时序预测：动态基线如何形成", "适合用本地 LSTM 环境演示的核心页")
    steps = [
        ("数据窗口", "过去24小时/7天\n流量、压力、日期、天气", b.BLUE),
        ("特征构造", "节假日、小时、星期\nMNF时段标识", b.TEAL),
        ("模型训练", "学习正常用水规律\n输出预测区间", b.PURPLE),
        ("残差判断", "实际值-预测值\n持续偏离触发预警", b.ORANGE),
        ("业务输出", "异常等级\n疑似DMA\n复核建议", b.RED),
    ]
    for i, (h, body, col) in enumerate(steps):
        x = 0.58 + i * 2.52
        b.card(s, x, 1.35, 2.10, 3.50, h, body, col, i + 1)
        if i < 4:
            b.line(s, x + 2.13, 3.05, x + 2.42, 3.05, b.BLUE, 2)
    b.rect(s, 1.00, 5.70, 11.20, 0.72, b.PALE2, b.BLUE, True)
    b.tb(s, "业务价值：用动态基线替代固定阈值，区分正常用水波动和真实漏损异常，支撑事前预警。", 1.35, 5.95, 10.50, 0.18, 17.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def anomaly_compare_table(prs, no=14):
    s = b.content_slide(prs, no, "2.2.1 无监督异常检测：三类模型怎么配合", "少标签阶段先建立可复核的异常筛查能力")
    rows = [
        ("孤立森林", "少数异常点更容易被随机切分隔离", "高频流量、压力残差、波动特征", "异常分数、报警等级", "速度快，适合在线筛查"),
        ("DBSCAN", "正常状态形成密度簇，低密度点为异常", "压力-流量二维/多维状态点", "异常片段、离群簇", "不需预设类别数"),
        ("自编码器", "只学习正常模式，异常样本重构误差高", "多变量时序窗口", "重构误差、异常阈值", "适合识别隐性异常模式"),
    ]
    headers = ["算法", "识别逻辑", "输入", "输出", "适用特点"]
    widths = [1.55, 3.00, 2.65, 2.00, 2.30]
    x0, y0 = 0.58, 1.28
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(11.50), b.emu(0.42)), b.TEAL)
    for i, h in enumerate(headers):
        b.tb(s, h, x0 + sum(widths[:i]) + 0.05, y0 + 0.13, widths[i] - 0.10, 0.13, 12.5, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.42 + r * 1.16
        b.rect(s, x0, y, 11.50, 1.08, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(224, 239, 248), False)
        for i, txt in enumerate(row):
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.08, y + 0.25, widths[i] - 0.16, 0.36, 12.5, b.TEXT if i else b.BLUE, i == 0, b.PP_ALIGN.CENTER)
    b.tb(s, "组合方式：LSTM生成残差，孤立森林快速筛查，DBSCAN识别状态簇，自编码器补充隐性模式。", 0.95, 6.18, 11.45, 0.24, 17.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def unsupervised(prs, no=13):
    s = b.content_slide(prs, no, "2.2.1 异常检测：少标签阶段先把异常筛出来", "LSTM负责动态基线，无监督模型负责从残差和波动中识别异常模式")
    items = [
        ("孤立森林", "适合快速筛查少数异常点\n输出异常分数，便于分级预警", b.BLUE),
        ("DBSCAN", "按密度识别正常运行簇\n可发现突变点和离群片段", b.TEAL),
        ("自编码器", "只学习正常模式\n用重构误差识别隐性异常", b.ORANGE),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 0.85 + i * 4.10
        b.card(s, x, 1.45, 3.35, 2.65, h, body, col)
        # simple scatter illustration
        for j in range(18):
            px = x + 0.55 + (j % 6) * 0.32
            py = 3.08 + (j // 6) * 0.18
            b.oval(s, px, py, 0.055, 0.055, col)
        b.oval(s, x + 2.75, 3.33, 0.13, 0.13, b.RED)
    b.rect(s, 1.10, 5.05, 11.10, 0.95, b.PALE2, b.BLUE, True)
    b.tb(s, "业务价值：把固定上下限报警升级为“动态预测 + 残差识别 + 分级复核”，减少误报并提前发现漏损风险。", 1.42, 5.38, 10.45, 0.22, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def hydraulic_fusion_steps(prs, no=16):
    s = b.content_slide(prs, no, "2.2.2 水力模型融合：从异常信号到候选管段", "定位问题必须引入管网拓扑和水力响应，不能只看单点曲线")
    steps = [
        ("建立仿真样本", "在不同管段、漏量、时段下模拟压力和流量响应", b.BLUE),
        ("提取响应特征", "压力降幅、传播路径、相邻测点相关性、残差模式", b.TEAL),
        ("训练定位模型", "随机森林/梯度提升树学习“响应特征-漏点位置”关系", b.ORANGE),
        ("候选排序", "输出TopN管段、置信度、核查顺序和定位依据", b.PURPLE),
    ]
    for i, (h, body, col) in enumerate(steps):
        x = 0.82 + i * 3.00
        b.card(s, x, 1.38, 2.50, 3.25, h, body, col, i + 1)
        if i < 3:
            b.line(s, x + 2.55, 3.00, x + 2.90, 3.00, b.BLUE, 2)
    b.tb(s, "业务价值：把DMA片区级报警进一步收敛为候选管段，减少无效巡检和无效开挖。", 1.05, 6.05, 11.25, 0.28, 18.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def ga_rf_gbt(prs, no=17):
    s = b.content_slide(prs, no, "2.2.2 机器学习与水力模型融合", "核心不是替代水力学，而是提高参数校核、候选定位和排序效率")
    rows = [
        ("遗传算法 GA", "水力模型参数校核", "粗糙度、需水量、背景漏失", "降低仿真与实测偏差", b.BLUE),
        ("随机森林 RF", "漏点位置分类", "仿真样本、压力响应、管段属性", "输出候选管段类别", b.TEAL),
        ("梯度提升树 GBT", "风险评分与定位排序", "多维特征、残差、历史维修", "输出TopN优先级", b.ORANGE),
    ]
    for i, (a, task, inp, out, col) in enumerate(rows):
        y = 1.40 + i * 1.45
        b.rect(s, 0.80, y, 11.75, 1.08, b.WHITE, col, True)
        b.tb(s, a, 1.05, y + 0.34, 2.20, 0.18, 16, col, True)
        b.tb(s, task, 3.30, y + 0.28, 2.25, 0.30, 15, b.TEXT, True, b.PP_ALIGN.CENTER)
        b.tb(s, inp, 5.75, y + 0.24, 3.10, 0.34, 13.5, b.TEXT, False, b.PP_ALIGN.CENTER)
        b.tb(s, out, 9.10, y + 0.28, 2.85, 0.30, 15, col, True, b.PP_ALIGN.CENTER)
    b.tb(s, "业务价值：优化水力模型精度，实现漏点精准定位，并支持DMA分区和巡检优先级优化。", 1.05, 6.20, 11.25, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def dma_partition_ai(prs):
    s = b.content_slide(prs, 18, "2.2.2 DMA分区智能优化：模型也服务于分区规划", "AI不仅用于检漏，也可辅助DMA边界、监测点和阀门配置优化")
    items = [
        ("分区评价", "结合供水压力、管网拓扑、用户规模、历史漏损，评估现有DMA合理性。", b.BLUE),
        ("监测点优化", "在有限预算下选择最能提升定位能力的流量计、压力计布点。", b.TEAL),
        ("阀门边界优化", "分析边界阀状态和水力连通关系，减少分区串水和计量偏差。", b.ORANGE),
        ("巡检资源优化", "按风险评分安排巡检优先级，实现从平均巡检到重点巡检。", b.PURPLE),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 0.95 + (i % 2) * 5.80
        y = 1.36 + (i // 2) * 1.95
        b.card(s, x, y, 5.05, 1.42, h, body, col)
    b.tb(s, "落地表达：DMA优化不是一次性画边界，而是随管网改造、用户变化和漏损反馈持续调整。", 1.10, 6.15, 11.15, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 18)


def big_data_decision(prs, no=19):
    s = b.content_slide(prs, no, "2.2.3 大数据融合与智能决策技术", "把模型结果接入资产、工单和调度系统，形成可执行决策")
    items = [
        ("多源数据融合", "流量、压力、水质、GIS、SCADA、管网属性、运维工单统一治理", b.BLUE),
        ("知识图谱", "把设备、管段、阀门、工单、事件关联起来，支持原因追踪", b.TEAL),
        ("边缘计算", "在网关侧完成初筛、缓存、校时和轻量推理，降低延迟与带宽压力", b.ORANGE),
        ("云端AI协同", "云端完成模型训练、全局分析、策略优化和跨DMA比较", b.PURPLE),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 1.00 + (i % 2) * 5.85
        y = 1.35 + (i // 2) * 2.00
        b.card(s, x, y, 5.05, 1.42, h, body, col)
    b.tb(s, "业务价值：支撑漏损管控决策、巡检路径优化、改造优先级排序和全生命周期管理。", 1.10, 6.20, 11.25, 0.26, 18.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def edge_cloud_detail(prs):
    s = b.content_slide(prs, 20, "2.2.3 边缘计算 + 云端AI协同", "实时预警靠边缘端，深度分析和模型训练靠云端")
    layers = [
        ("传感器端", "流量、压力、水质、声学\n分钟级/秒级采样", b.BLUE),
        ("边缘端", "校时、缓存、异常初筛\n断网续传、轻量推理", b.TEAL),
        ("云端", "模型训练、全局分析\n跨DMA对比、策略优化", b.PURPLE),
        ("业务端", "预警看板、GIS定位\n工单派发、移动巡检", b.ORANGE),
    ]
    for i, (h, body, col) in enumerate(layers):
        x = 0.80 + i * 3.10
        b.card(s, x, 1.45, 2.55, 3.40, h, body, col)
        if i < 3:
            b.line(s, x + 2.62, 3.10, x + 3.00, 3.10, b.BLUE, 2)
    b.tb(s, "设计原则：边缘端负责快，云端负责准，业务端负责闭环，三者共同保证实时性和可运营性。", 1.05, 6.12, 11.25, 0.28, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, 20)


def planning(prs, no=23):
    s = b.content_slide(prs, no, "4.1 前期规划：明确建设目标与实施边界", "实施路径从现状诊断开始，避免先上算法再找场景")
    items = [
        ("现状诊断", "数据资产、硬件设备、信息化系统、漏损管控现状、DMA分区现状", b.BLUE),
        ("目标设定", "短期降本增效\n中期闭环管控\n长期全生命周期管理", b.TEAL),
        ("路径规划", "试点先行\n分步推广\n持续迭代", b.ORANGE),
    ]
    for i, (h, body, col) in enumerate(items):
        b.card(s, 0.95 + i * 4.05, 1.55, 3.35, 3.35, h, body, col, i + 1)
    b.tb(s, "输出物：试点DMA清单、数据接入清单、模型建设边界、阶段性验收指标。", 1.10, 6.18, 11.20, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def data_governance_detail(prs, no=25):
    s = b.content_slide(prs, no, "4.2 数据治理体系：AI模型的核心底座", "数据治理不是后台工作，而是决定模型能否上线的前置条件")
    items = [
        ("多源归集", "SCADA、GIS、DMA、工单、资产、天气、施工记录统一接入。", b.BLUE),
        ("统一编码", "设备ID、管段ID、DMA编号、阀门编号保持一致。", b.TEAL),
        ("时间对齐", "统一采样粒度和时间戳，处理延迟、断点和补传数据。", b.ORANGE),
        ("质量校验", "缺失值、异常值、漂移、重复点、仪表故障识别。", b.PURPLE),
        ("标签体系", "误报、漏损、施工、阀门调整、维修结果形成可训练标签。", b.GREEN),
    ]
    for i, (h, body, col) in enumerate(items):
        x = 0.70 + (i % 3) * 4.05
        y = 1.25 + (i // 3) * 1.85
        w = 3.35 if i < 3 else 5.35
        if i >= 3:
            x = 1.55 + (i - 3) * 5.70
        b.card(s, x, y, w, 1.35, h, body, col)
    b.tb(s, "输出物：设备表、管网表、时序表、工单表、标签表、规则表，支撑模型训练与追溯。", 1.05, 6.18, 11.25, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def model_construction(prs, no=26):
    s = b.content_slide(prs, no, "4.3 模型建设：选型、训练、验证与优化", "模型建设要贴合业务场景，优先选择轻量化、可解释、可验证的方案")
    steps = [
        ("场景拆解", "预警、定位、排序、复盘"),
        ("模型选型", "LSTM/GRU、无监督异常、树模型、水力融合"),
        ("训练验证", "历史数据、时间切分、指标评估、现场复核"),
        ("迭代机制", "工单回填、漂移监测、版本管理、灰度发布"),
    ]
    for i, (h, body) in enumerate(steps):
        x = 0.85 + i * 3.05
        col = [b.BLUE, b.TEAL, b.ORANGE, b.PURPLE][i]
        b.card(s, x, 1.55, 2.55, 3.35, h, body, col, i + 1)
        if i < 3:
            b.line(s, x + 2.60, 3.20, x + 2.95, 3.20, b.BLUE, 2)
    b.tb(s, "验收重点：不仅看模型准确率，还要看误报率、TopN命中率、派单有效率和闭环回填率。", 1.05, 6.18, 11.35, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def training_validation_detail(prs, no=27):
    s = b.content_slide(prs, no, "4.3 模型训练与验证：既看算法指标，也看业务指标", "训练过程必须避免未来信息泄漏，并保留现场复核链路")
    left = [
        ("训练数据", "历史正常运行数据、已确认漏损事件、仿真样本"),
        ("切分方式", "按时间切分训练/验证/测试，避免未来信息泄漏"),
        ("算法指标", "MAE、RMSE、召回率、F1、异常分数稳定性"),
    ]
    right = [
        ("业务指标", "TopN命中率、误报率、派单有效率、处置时长"),
        ("验证方式", "历史回放、在线灰度、现场复核、专家复盘"),
        ("版本管理", "记录数据版本、特征版本、模型版本和阈值版本"),
    ]
    for i, (h, body) in enumerate(left):
        b.card(s, 0.95, 1.35 + i * 1.30, 5.35, 0.98, h, body, [b.BLUE, b.TEAL, b.ORANGE][i])
    for i, (h, body) in enumerate(right):
        b.card(s, 7.05, 1.35 + i * 1.30, 5.35, 0.98, h, body, [b.PURPLE, b.GREEN, b.RED][i])
    b.tb(s, "对外表达：准确率只是起点，真正的验收标准是能否减少无效巡检、缩短定位时间、形成闭环数据。", 1.05, 6.20, 11.30, 0.26, 17.5, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def pilot(prs, no=29):
    s = b.content_slide(prs, no, "4.4 试点验证与优化：形成可复制经验", "工程落地需要先在典型DMA跑通闭环，再扩大到全域")
    stages = [
        ("选择试点", "典型管材\n高漏损率\n数据较完整", b.BLUE),
        ("联调运行", "硬件接入\n系统集成\n模型上线", b.TEAL),
        ("效果验证", "预警有效率\n定位命中率\n处置时长", b.ORANGE),
        ("优化推广", "规则沉淀\n模板复制\n多DMA推广", b.PURPLE),
    ]
    for i, (h, body, col) in enumerate(stages):
        x = 0.90 + i * 3.05
        b.card(s, x, 1.48, 2.55, 3.30, h, body, col, i + 1)
        if i < 3:
            b.line(s, x + 2.58, 3.10, x + 2.94, 3.10, b.BLUE, 2)
    b.tb(s, "试点不只证明算法能跑，还要证明业务人员愿意用、工单能闭环、指标能持续改善。", 1.10, 6.10, 11.20, 0.28, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def operation_roles(prs, no=31):
    s = b.content_slide(prs, no, "4.5 长效运营：组织能力与运维保障", "AI系统上线后需要制度、人员、指标和模型版本共同维护")
    rows = [
        ("管理层", "看投入产出、漏损率变化、试点推广节奏", "明确目标和资源"),
        ("调度人员", "看预警解释、定位证据、派单优先级", "确认业务动作"),
        ("运维人员", "看现场核查、维修记录、结果回填", "提供高质量标签"),
        ("算法/IT人员", "看数据漂移、模型性能、接口稳定性", "维护模型和系统"),
    ]
    widths = [1.80, 5.00, 4.10]
    x0, y0 = 0.85, 1.30
    b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(x0), b.emu(y0), b.emu(10.90), b.emu(0.42)), b.BLUE)
    for i, h in enumerate(["角色", "关注内容", "责任重点"]):
        b.tb(s, h, x0 + sum(widths[:i]) + 0.05, y0 + 0.13, widths[i] - 0.10, 0.13, 12.8, b.WHITE, True, b.PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = y0 + 0.42 + r * 0.90
        b.rect(s, x0, y, 10.90, 0.82, b.RGBColor(248, 252, 255) if r % 2 == 0 else b.WHITE, b.RGBColor(224, 239, 248), False)
        for i, txt in enumerate(row):
            b.tb(s, txt, x0 + sum(widths[:i]) + 0.08, y + 0.24, widths[i] - 0.16, 0.22, 13.3, b.BLUE if i == 0 else b.TEXT, i == 0, b.PP_ALIGN.CENTER)
    b.tb(s, "保障机制：传感器在线率监测、接口巡检、模型漂移监控、定期复盘会、版本发布审批。", 1.05, 6.18, 11.20, 0.26, 18, b.BLUE, True, b.PP_ALIGN.CENTER)
    b.audit(s, no)


def final_summary(prs, no=32):
    s = b.blank(prs, b.NAVY)
    b.tb(s, "杨佳部分最终落点", 0.95, 0.90, 5.30, 0.45, 34, b.WHITE, True)
    items = [
        ("讲发展", "AI模型从阈值报警走向动态基线、机理融合和智能体协同。"),
        ("讲算法", "按时序异常、水力融合、大数据决策三类技术栈组织。"),
        ("讲实施", "从规划、数据治理、模型建设、工程落地到长效运营闭环。"),
    ]
    for i, (h, body) in enumerate(items):
        y = 2.02 + i * 1.22
        col = [b.CYAN, b.TEAL, b.ORANGE][i]
        b.fill(s.shapes.add_shape(b.MSO_AUTO_SHAPE_TYPE.RECTANGLE, b.emu(0.95), b.emu(y), b.emu(0.12), b.emu(0.68)), col)
        b.tb(s, h, 1.25, y - 0.02, 2.10, 0.28, 22, col, True)
        b.tb(s, body, 3.20, y + 0.05, 5.15, 0.38, 16.5, b.RGBColor(220, 236, 248))
    b.pic(s, "overall", 8.15, 1.20, 4.55, 4.30, crop=False)
    b.tb(s, "AI供水管网DMA漏损检测｜杨佳负责部分", 0.95, 6.75, 5.0, 0.22, 11.5, b.RGBColor(166, 208, 235))
    b.audit(s, no, allow_short=True)


def build():
    missing = [str(p) for p in b.ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("\n".join(missing))

    prs = Presentation(str(b.TEMPLATE))
    b.clear_slides(prs)

    cover(prs)
    scope(prs)
    b.section(prs, 3, "PART 01", "1.3 AI模型应用发展", "介绍AI模型、智能体发展，以及漏损检测可适配算法。")
    ai_evolution(prs)
    agent_development(prs)
    b.full_visual(prs, 6, "positioning")
    algorithm_family(prs)
    selection_rules(prs)

    b.section(prs, 9, "PART 02", "2.2 核心AI技术及业务价值", "按时序异常、水力融合、大数据决策三类技术栈展开。", b.TEAL)
    b.full_visual(prs, 10, "baseline")
    b.lstm_slide(prs, 11)
    lstm_process_detail(prs)
    unsupervised(prs)
    anomaly_compare_table(prs)
    b.hydraulic_fusion(prs, 15)
    hydraulic_fusion_steps(prs)
    ga_rf_gbt(prs)
    dma_partition_ai(prs)
    big_data_decision(prs)
    edge_cloud_detail(prs)
    b.full_visual(prs, 21, "training")

    b.section(prs, 22, "PART 03", "第四部分 实施路径", "从规划、数据治理、模型建设、工程落地到长效运营。", b.ORANGE)
    planning(prs)
    b.full_visual(prs, 24, "training")
    data_governance_detail(prs)
    model_construction(prs)
    training_validation_detail(prs)
    b.deployment(prs, 28)
    pilot(prs)
    b.operations(prs, 30)
    operation_roles(prs)
    final_summary(prs, 32)

    prs.save(b.PPTX)
    b.AUDIT.write_text("\n".join(b.AUDIT_LINES), encoding="utf-8")
    b.build_previews(prs)
    b.append_repetition_check(prs)
    return b.PPTX


if __name__ == "__main__":
    print(build())
    print(b.AUDIT)
    print(b.PREVIEW)
