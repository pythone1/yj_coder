from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Pt


BASE = Path(r"E:\PY\research\0425")
INPUT = BASE / "生产流程.pptx"
OUTPUT = BASE / "生产流程_最终检查版.pptx"

REPLACEMENTS = {
    "迭线": "送线",
    "锡护": "锅炉",
    "异崇检到": "异常检测",
    "预湖": "预测",
    "遇终功能效益": "最终功能效益",
    "凤险": "风险",
    "橙惨计划": "检修计划",
    "摄矢": "损失",
    "赐标": "贴标",
    "波位": "液位",
    "思浮": "悬浮",
    "主霎": "主要",
    "分劃": "分割",
    "识刖": "识别",
    "复较": "复核",
    "缺船": "缺陷",
    "追潮": "追溯",
    "任努": "任务",
    "上瓶": "上甑",
    "昇常": "异常",
    "传燃": "传感",
    "温温度": "温度",
    "酒酷": "酒醅",
    "异案": "异常",
    "围": "图",
    "等缀": "等级",
    "掠入": "输入",
    "专案舰则": "专家规则",
    "检素": "检索",
    "规钏": "规则",
    "牧据": "数据",
    "措荐抉行": "推荐执行",
    "反懊": "反馈",
    "横型童训": "模型重训",
    "辖出": "输出",
    "璃参": "微调",
    "知玖沉渡": "知识沉淀",
    "产续": "产线",
    "发观": "发现",
    "该杀率": "误杀率",
    "络": "结论",
    "罄线": "基线",
    "横型": "模型",
    "尖时": "实时",
    "汇粲": "汇聚",
    "逮度": "速度",
    "裁重": "载重",
    "起焚": "起终",
    "紫行": "禁行",
    "A\"": "A*",
    "落径": "路径",
    "规刺": "规划",
    "重规刘": "重规划",
    "逼过": "通过",
    "交忖": "交付",
    "温装": "混装",
    "睾求": "要求",
    "关肤": "关联",
    "保分": "评分",
    "密令": "订单",
    "启寂丸": "启发式",
    "液灰计刻": "波次计划",
    "捺选": "拣选",
    "支排": "安排",
    "辅入": "输入",
    "装魉舰物": "装载规划",
    "约柬伏化": "约束优化",
    "籍出": "输出",
    "多合节赤": "多仓节点",
    "韩出": "输出",
    "肠逆": "配送",
    "庠区": "库区",
    "又车": "叉车",
    "拣 遮": "拣选",
    "疲儋": "疲倦",
    "淀粉辜": "淀粉率",
    "搠眼": "摊晾",
    "发   醇": "发酵",
    "出客": "出窖",
    "图期": "周期",
    "周图": "周期",
}


def fix_text(value):
    value = value.replace("ANotebookLM", "").replace("NotebookLM", "")
    for old, new in REPLACEMENTS.items():
        value = value.replace(old, new)
    return value


def tune_shape(shape):
    if hasattr(shape, "text_frame") and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = fix_text(run.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = fix_text(cell.text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            tune_shape(child)


def put_after_background(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(3, el)


def add_footer(slide, slide_width, slide_height):
    footer_h = Emu(213333)
    footer_y = Emu(int(slide_height) - int(footer_h))
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, footer_y, slide_width, footer_h)
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(71, 105, 158)
    footer.line.fill.background()
    put_after_background(slide, footer)


def add_body_card(slide, x, y, w, h, text):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(4, 28, 43)
    rect.line.fill.background()

    box = slide.shapes.add_textbox(x + 90000, y + 70000, w - 180000, h - 140000)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for pi, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        para.text = line
        para.space_after = Pt(0)
        for run in para.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(245, 250, 255)


prs = Presentation(str(INPUT))
for idx, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        tune_shape(shape)
    if idx == 1:
        add_body_card(
            slide,
            10230000,
            5600000,
            2350000,
            3350000,
            "核心工艺对象：厂区物流节点、AGV任务、车辆位置、充电、避障、任务优先级。\n"
            "主要算法任务：路径规划、任务分配、多车冲突消解、预防、仿真评估。\n"
            "最终功能效益：减少等待和绕路，提高准时率和设备利用率。",
        )
        add_body_card(
            slide,
            13330000,
            5600000,
            2350000,
            3350000,
            "核心工艺对象：库位、订单、库存、装车、车辆、客户路线、多仓发货。\n"
            "主要算法任务：库位优化、订单波次、装车排程、车辆路径优化。\n"
            "最终功能效益：缩短装车发货时间，提高库位利用率和配送效率。",
        )
    add_footer(slide, prs.slide_width, prs.slide_height)

prs.save(str(OUTPUT))
print(OUTPUT)
