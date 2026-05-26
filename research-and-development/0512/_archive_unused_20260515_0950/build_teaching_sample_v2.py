# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = max([p for p in ROOT.glob("*.pptx") if not p.name.startswith("~$")], key=lambda p: p.stat().st_size)
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_对外讲解样张_v2.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_对外讲解样张_v2_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_对外讲解样张_v2_预览联系表.png"

NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
FONT = "微软雅黑"

ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "dma": ASSET_DIR / "dma_clean.png",
    "data": ASSET_DIR / "data_clean.png",
    "ai": ASSET_DIR / "ai_clean.png",
    "locate": ASSET_DIR / "locate_clean.png",
    "loop": ASSET_DIR / "loop_clean.png",
    "algorithm": ASSET_DIR / "algorithm_clean.png",
}

AUDIT_LINES = []


def ensure_assets():
    missing = [str(p) for p in ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("\n".join(missing))


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, trans=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = trans
    shape.line.fill.background()


def stroke(shape, color=LINE, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(typ, x, y, w, h)
    fill(shp, color, trans)
    stroke(shp, line, 1.0)
    return shp


def tb(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pic(slide, path, x, y, w, h, crop=False):
    p = slide.shapes.add_picture(str(path), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def connector(slide, x1, y1, x2, y2, color=BLUE, width=1.6, arrow=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.25), Inches(0.12), Inches(0.48))
    fill(bar, BLUE)
    tb(slide, title, Inches(0.62), Inches(0.14), Inches(9.45), Inches(0.42), 25, TEXT, True)
    tb(slide, sub, Inches(0.64), Inches(0.62), Inches(8.2), Inches(0.22), 11.5, MUTED)
    tb(slide, f"{no:02d}", Inches(12.05), Inches(0.22), Inches(0.7), Inches(0.25), 12, BLUE, True, PP_ALIGN.RIGHT)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.90), Inches(11.70), Inches(0.016))
    fill(line, PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA漏损检测 · 对外讲解样张", Inches(0.62), Inches(7.08), Inches(5.8), Inches(0.2), 10, MUTED)


def content(prs, title, sub, no):
    slide = blank(prs)
    header(slide, title, sub, no)
    footer(slide)
    return slide


def pill(slide, text, x, y, w, color=PALE, text_color=BLUE, size=12):
    rect(slide, x, y, w, Inches(0.34), color, color, True)
    tb(slide, text, x + Inches(0.06), y + Inches(0.08), w - Inches(0.12), Inches(0.16), size, text_color, True, PP_ALIGN.CENTER)


def note(slide, text, y=6.10, color=BLUE):
    rect(slide, Inches(0.72), Inches(y), Inches(11.85), Inches(0.56), RGBColor(247, 252, 255), LINE, True)
    tb(slide, text, Inches(0.98), Inches(y + 0.155), Inches(11.3), Inches(0.22), 15.0, color, True, PP_ALIGN.CENTER)


def mini_text(slide, title, body, x, y, w, h, color=BLUE):
    rect(slide, x, y, w, h, WHITE, LINE, True)
    tb(slide, title, x + Inches(0.16), y + Inches(0.12), w - Inches(0.32), Inches(0.23), 15, color, True)
    tb(slide, body, x + Inches(0.16), y + Inches(0.50), w - Inches(0.32), h - Inches(0.62), 12.5, TEXT)


def audit_slide(slide, no, allow_short=False):
    prs_w, prs_h = 12192000, 6858000
    chars = 0
    pics = 0
    min_font = 99
    off = 0
    bottom = 0
    for sh in slide.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs_w + 50000 or sh.top + sh.height > prs_h + 50000:
            off += 1
        if sh.shape_type == 13:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off:
        flags.append(f"越界{off}")
    if min_font < 10:
        flags.append(f"小字{min_font}")
    if pics < 1:
        flags.append("缺少主视觉")
    min_chars = 120 if no == 1 else 150
    if not allow_short and chars < min_chars:
        flags.append(f"文字不足{chars}")
    if bottom < Inches(6.2):
        flags.append("下半页利用不足")
    status = "FAIL" if flags else "OK"
    AUDIT_LINES.append(f"{no:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise SystemExit(f"slide {no} failed: {'; '.join(flags)}")


def slide1_cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(5.75), Inches(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(0.92), Inches(0.13), Inches(5.75)), CYAN)
    tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.18), Inches(4.55), Inches(1.55), 33.5, WHITE, True)
    tb(s, "对外教学演讲样张", Inches(1.0), Inches(3.10), Inches(3.3), Inches(0.32), 18, CYAN, True)
    tb(s, "从“哪个区可能漏”讲到“哪段管优先查、谁去查、结果怎么回填”。", Inches(1.0), Inches(5.45), Inches(4.5), Inches(0.52), 13, RGBColor(218, 237, 248))
    tb(s, "面向智慧水务、管网运维、漏损治理和AI模型应用培训。", Inches(1.0), Inches(6.05), Inches(4.5), Inches(0.35), 12.2, RGBColor(218, 237, 248))
    tb(s, "覆盖DMA分区计量、动态基线、算法溯源、现场派单和模型复盘。", Inches(1.0), Inches(6.42), Inches(4.5), Inches(0.28), 11.5, RGBColor(218, 237, 248))
    rect(s, Inches(6.05), Inches(0.68), Inches(6.75), Inches(5.78), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS["cover"], Inches(6.25), Inches(0.90), Inches(6.35), Inches(5.30), crop=False)
    pill(s, "DMA锁定", Inches(6.68), Inches(6.47), Inches(1.35), PALE, BLUE, 12)
    pill(s, "AI溯源", Inches(8.35), Inches(6.47), Inches(1.35), PALE, GREEN, 12)
    pill(s, "现场闭环", Inches(10.02), Inches(6.47), Inches(1.45), PALE, ORANGE, 12)
    audit_slide(s, 1)


def slide2_goals(prs):
    s = content(prs, "课程目标：听完这部分应能回答什么", "对外讲解不是堆模型名，而是讲清工程问题如何被模型解决", 2)
    pic(s, ASSETS["cover"], Inches(0.78), Inches(1.22), Inches(3.15), Inches(2.15), crop=False)
    pic(s, ASSETS["algorithm"], Inches(4.98), Inches(1.22), Inches(3.15), Inches(2.15), crop=False)
    pic(s, ASSETS["loop"], Inches(9.15), Inches(1.22), Inches(3.15), Inches(2.15), crop=False)
    mini_text(s, "1. 业务对象", "DMA如何把大管网转成可计量、可分析、可考核的分区？", Inches(0.78), Inches(3.70), Inches(3.15), Inches(1.12), BLUE)
    mini_text(s, "2. 模型方法", "AI如何利用流量、压力、拓扑和工单，把异常片区继续收敛到候选管段？", Inches(4.98), Inches(3.70), Inches(3.15), Inches(1.12), GREEN)
    mini_text(s, "3. 落地闭环", "模型结果如何进入派单、维修、复盘和再训练，而不是停留在一张报警图上？", Inches(9.15), Inches(3.70), Inches(3.15), Inches(1.12), ORANGE)
    note(s, "讲课主线：先讲清问题，再讲模型怎么介入，最后讲如何在现场形成闭环。", 5.85)
    audit_slide(s, 2)


def slide3_mnf(prs):
    s = content(prs, "夜间最小流量MNF：AI预警的起点", "从夜间底流异常进入DMA漏损诊断", 3)
    rect(s, Inches(0.75), Inches(1.22), Inches(5.35), Inches(3.55), WHITE, LINE, True)
    tb(s, "MNF曲线示意", Inches(1.05), Inches(1.45), Inches(2.2), Inches(0.25), 17, BLUE, True)
    # editable chart
    x0, y0, w, h = Inches(1.15), Inches(2.00), Inches(4.45), Inches(2.05)
    connector(s, x0, y0 + h, x0 + w, y0 + h, MUTED, 1.0, False)
    connector(s, x0, y0 + h, x0, y0, MUTED, 1.0, False)
    pts = [(0.0, 0.55), (0.12, 0.60), (0.24, 0.58), (0.36, 0.52), (0.48, 0.50), (0.60, 0.38), (0.72, 0.30), (0.84, 0.28), (1.0, 0.25)]
    prev = None
    for px, py in pts:
        cx = x0 + int(w * px)
        cy = y0 + int(h * py)
        if prev:
            connector(s, prev[0], prev[1], cx, cy, BLUE, 2.2, False)
        prev = (cx, cy)
    # abnormal curve
    pts2 = [(0.0, 0.53), (0.12, 0.56), (0.24, 0.55), (0.36, 0.48), (0.48, 0.42), (0.60, 0.35), (0.72, 0.33), (0.84, 0.48), (1.0, 0.62)]
    prev = None
    for px, py in pts2:
        cx = x0 + int(w * px)
        cy = y0 + int(h * py)
        if prev:
            connector(s, prev[0], prev[1], cx, cy, RED, 2.2, False)
        prev = (cx, cy)
    pill(s, "正常夜间底流", Inches(1.2), Inches(4.22), Inches(1.45), PALE, BLUE, 10.5)
    pill(s, "异常抬升", Inches(2.85), Inches(4.22), Inches(1.25), RGBColor(255, 238, 236), RED, 10.5)
    pic(s, ASSETS["dma"], Inches(6.65), Inches(1.18), Inches(2.55), Inches(2.08), crop=False)
    mini_text(s, "讲解1：MNF不是漏损量", "夜间最小流量还包含合法夜间用水、计量误差和边界异常，需要结合上下文判断。", Inches(9.55), Inches(1.15), Inches(2.65), Inches(1.30), BLUE)
    mini_text(s, "讲解2：AI看持续偏离", "模型关注曲线是否持续偏离本DMA的正常基线，而不是单个时刻超过阈值。", Inches(9.55), Inches(2.75), Inches(2.65), Inches(1.30), GREEN)
    mini_text(s, "讲解3：预警要能复核", "输出应包括异常时段、残差幅度、压力响应和建议核查方式。", Inches(6.65), Inches(3.72), Inches(5.55), Inches(1.08), ORANGE)
    note(s, "这一页用于说明：DMA发现异常靠MNF，AI提高稳定性靠动态基线和多源证据。", 5.95)
    audit_slide(s, 3)


def slide4_architecture(prs):
    s = content(prs, "总体架构：DMA宏观锁定 + AI微观溯源", "从片区异常到候选管段的完整链路", 4)
    y = Inches(1.18)
    modules = [
        ("DMA锁定", ASSETS["dma"], "先圈定异常片区", BLUE),
        ("数据融合", ASSETS["data"], "汇集流量、压力、管网、工单", CYAN),
        ("AI溯源", ASSETS["ai"], "反推压力和流量变化来源", GREEN),
        ("现场闭环", ASSETS["locate"], "输出管段坐标并派单", ORANGE),
    ]
    for i, (title, img, desc, color) in enumerate(modules):
        x = Inches(0.58 + i * 3.05)
        rect(s, x, y, Inches(2.62), Inches(4.15), WHITE, RGBColor(205, 229, 243), True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(2.62), Inches(0.42)), color)
        tb(s, title, x, y + Inches(0.10), Inches(2.62), Inches(0.18), 14.5, WHITE, True, PP_ALIGN.CENTER)
        pic(s, img, x + Inches(0.16), y + Inches(0.58), Inches(2.30), Inches(2.42), crop=False)
        tb(s, desc, x + Inches(0.18), y + Inches(3.35), Inches(2.26), Inches(0.42), 12.2, TEXT, True, PP_ALIGN.CENTER)
        if i < 3:
            connector(s, x + Inches(2.67), y + Inches(2.35), x + Inches(2.95), y + Inches(2.35), BLUE, 2.0)
    note(s, "讲课主线：DMA锁片区，AI找原因，模型给管段，工单到现场闭环处置。", 5.90)
    audit_slide(s, 4)


def slide5_algorithm(prs):
    s = content(prs, "算法全景：不同模型解决不同任务", "不要讲成模型名词堆砌，要讲清每个模型的业务位置", 5)
    pic(s, ASSETS["algorithm"], Inches(0.78), Inches(1.18), Inches(5.00), Inches(4.55), crop=False)
    items = [
        ("预测正常曲线", "LSTM / GRU", "输出预测区间和残差", BLUE),
        ("发现候选异常", "孤立森林 / 自编码器", "输出异常时段和分数", CYAN),
        ("收敛候选管段", "水力仿真 + ML", "输出管段TopN", GREEN),
        ("排序检漏优先级", "GBDT / HGB", "输出风险分", ORANGE),
        ("解释和派单", "知识图谱 / 智能体", "输出报告和建议", PURPLE),
    ]
    for i, (a, b, c, color) in enumerate(items):
        y = Inches(1.18 + i * 0.82)
        rect(s, Inches(6.20), y, Inches(5.65), Inches(0.58), PALE2 if i % 2 == 0 else WHITE, LINE, True)
        pill(s, a, Inches(6.35), y + Inches(0.12), Inches(1.65), PALE, color, 10.5)
        tb(s, b, Inches(8.15), y + Inches(0.15), Inches(1.80), Inches(0.18), 12.5, TEXT, True)
        tb(s, c, Inches(10.05), y + Inches(0.15), Inches(1.60), Inches(0.18), 12.0, MUTED)
    note(s, "一句话讲清：预测模型负责正常边界，异常模型负责发现信号，机理模型负责缩小位置范围。", 5.90)
    audit_slide(s, 5)


def slide6_story(prs):
    s = content(prs, "业务闭环：从预警到复盘的一次现场故事", "用现场流程讲清模型价值", 6)
    pic(s, ASSETS["loop"], Inches(0.78), Inches(1.18), Inches(4.75), Inches(4.55), crop=False)
    steps = [
        ("1 事前预警", "模型发现MNF持续抬升，并给出异常评分。", BLUE),
        ("2 事中定位", "结合压力响应和拓扑，输出候选管段TopN。", CYAN),
        ("3 现场处置", "检漏人员按坐标核查，维修结果写回工单。", ORANGE),
        ("4 事后复盘", "真实漏损或误报原因回填，作为再训练样本。", GREEN),
    ]
    for i, (title, body, color) in enumerate(steps):
        x = Inches(6.05 + (i % 2) * 3.05)
        y = Inches(1.18 + (i // 2) * 2.02)
        mini_text(s, title, body, x, y, Inches(2.72), Inches(1.38), color)
    connector(s, Inches(8.78), Inches(1.85), Inches(9.08), Inches(1.85), BLUE, 1.6)
    connector(s, Inches(7.38), Inches(2.58), Inches(7.38), Inches(3.18), BLUE, 1.6)
    connector(s, Inches(8.78), Inches(3.88), Inches(9.08), Inches(3.88), BLUE, 1.6)
    note(s, "对外讲解时要落到人和流程：模型不是代替运维，而是让人员带着坐标和证据去现场。", 5.90)
    audit_slide(s, 6)


def build_preview(prs):
    W, H = 420, 236
    sx, sy = W / prs.slide_width, H / prs.slide_height
    try:
        font_title = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 9)
        font_body = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 7)
    except Exception:
        font_title = font_body = ImageFont.load_default()
    thumbs = []
    for slide in prs.slides:
        im = Image.new("RGB", (W, H), "white")
        d = ImageDraw.Draw(im)
        for sh in slide.shapes:
            x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
            if w <= 0 or h <= 0:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                    p.thumbnail((w, h))
                    im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
                except Exception:
                    pass
            else:
                fill_color = None
                line_color = None
                try:
                    rgb = sh.fill.fore_color.rgb
                    if rgb:
                        fill_color = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
                except Exception:
                    pass
                try:
                    rgb = sh.line.color.rgb
                    if rgb:
                        line_color = tuple(int(str(rgb)[i:i+2], 16) for i in (0, 2, 4))
                except Exception:
                    pass
                if fill_color:
                    d.rounded_rectangle([x, y, x+w, y+h], radius=5, fill=fill_color, outline=line_color)
                elif line_color:
                    d.rectangle([x, y, x+w, y+h], outline=line_color)
                if hasattr(sh, "text") and sh.text.strip() and w > 25 and h > 8:
                    font = font_title if any(r.font.bold for pgh in sh.text_frame.paragraphs for r in pgh.runs) else font_body
                    text = sh.text.strip().split("\n")[0][:28]
                    d.text((x+2, y+2), text, fill=(20, 40, 60), font=font)
        thumbs.append(im)
    sheet = Image.new("RGB", (W * 2, (H + 18) * 3), (230, 236, 242))
    d = ImageDraw.Draw(sheet)
    for i, t in enumerate(thumbs):
        x = (i % 2) * W
        y = (i // 2) * (H + 18)
        sheet.paste(t, (x, y))
        d.text((x + 4, y + H + 2), f"{i+1:02d}", fill=(0, 70, 148), font=font_title)
    sheet.save(PREVIEW)


def build():
    ensure_assets()
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    slide1_cover(prs)
    slide2_goals(prs)
    slide3_mnf(prs)
    slide4_architecture(prs)
    slide5_algorithm(prs)
    slide6_story(prs)
    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    build_preview(prs)
    return PPTX


if __name__ == "__main__":
    print(build())
    print(PREVIEW)
