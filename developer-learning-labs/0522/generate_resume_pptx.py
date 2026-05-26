# -*- coding: utf-8 -*-
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme color definitions (Deep Slate & Teal Accents)
    BG_COLOR = RGBColor(15, 23, 42)      # Deep slate #0F172A
    TEXT_TITLE = RGBColor(255, 255, 255) # White
    TEXT_BODY = RGBColor(203, 213, 225)  # Light slate #cbd5e1
    ACCENT_TEAL = RGBColor(13, 148, 136) # Teal #0d9488
    
    # ------------------ SLIDE 1: COVER ------------------
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_COLOR
    
    # Add top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()
    
    # Main Title
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "个人求职介绍与技术项目答辩汇报"
    p.font.name = 'Microsoft YaHei'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = TEXT_TITLE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "应聘岗位：资深 AI 算法与遥感系统工程师 | 空间大数据专家"
    p2.font.name = 'Microsoft YaHei'
    p2.font.size = Pt(18)
    p2.font.color.rgb = ACCENT_TEAL
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    
    # Contact Info Info
    p3 = tf.add_paragraph()
    p3.text = "\n汇报人：面试候选人 (Geospatial & AI Expert)   |   联系电话：+86 188-XXXX-XXXX\n电子邮箱：candidate@email.com   |   Github: github.com/geospatial-ai-unicorn"
    p3.font.name = 'Microsoft YaHei'
    p3.font.size = Pt(13)
    p3.font.color.rgb = TEXT_BODY
    p3.alignment = PP_ALIGN.CENTER
    
    # Helper to add standard slides
    def add_standard_slide(title, subtitle, bullets):
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = BG_COLOR
        
        # Left accent line decoration
        dec = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(0.08), Inches(0.6))
        dec.fill.solid()
        dec.fill.fore_color.rgb = ACCENT_TEAL
        dec.line.fill.background()
        
        # Slide Title
        t_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(12.0), Inches(0.5))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Microsoft YaHei'
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_TITLE
        
        # Slide Subtitle
        p_sub = tf_t.add_paragraph()
        p_sub.text = subtitle
        p_sub.font.name = 'Microsoft YaHei'
        p_sub.font.size = Pt(12)
        p_sub.font.color.rgb = ACCENT_TEAL
        p_sub.font.bold = True
        
        # Bullets Body
        b_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.13), Inches(5.0))
        tf_b = b_box.text_frame
        tf_b.word_wrap = True
        
        for i, text in enumerate(bullets):
            p_b = tf_b.paragraphs[0] if i == 0 else tf_b.add_paragraph()
            p_b.text = text
            p_b.font.name = 'Microsoft YaHei'
            p_b.font.size = Pt(14)
            p_b.font.color.rgb = TEXT_BODY
            p_b.level = 0
            p_b.space_after = Pt(12)
            
        return slide

    # ------------------ SLIDE 2: ADVANTAGES ------------------
    add_standard_slide(
        "个人简介与多维复合型优势总结",
        "兼具算法研发、空间大数据治理与无人机外业测绘全栈能力",
        [
            "• 多维复合背景：兼具遥感算法研发、大模型(GenAI/LLM)深度定制、高吞吐数据工程以及无人机(UAV)测绘外内业实力。",
            "• 前沿技术落地：擅长将前沿 CV (如 Segment Anything 1/2) 与 LLM (RAG、Agent) 技术落地于空间地理信息与复杂业务系统。",
            "• 生产项目经验：具备大型遥感数据湖建设及生产级 AI Agent 工作流管线部署经验，主导过多个千亿级空间数据处理项目。",
            "• 协同开发提速：熟练使用 Antigravity、Codex 等 AI 辅助编程工具进行结对编程，快速搭建原型并完成高质量交付。",
            "• 软硬件一体化：持有无人机测绘飞行执照，具备野外航线数据采集及 RTK 静态解算、控制点相似变换等软硬件集成能力。"
        ]
    )

    # ------------------ SLIDE 3: SKILL MATRIX ------------------
    add_standard_slide(
        "核心专业技能矩阵与技术广度",
        "覆盖大模型、计算机视觉、空间计算及无人机测绘四大技术支柱",
        [
            "• 【大模型与智能体】: 提示词工程 / 高级 RAG 检索管线 / 多模态对齐 / Nous Hermes & ReAct 自进化智能体框架 / LoRA 微调",
            "• 【计算机视觉与遥感】: SAM 1/2 模型微调 / 目标检测(YOLO) / 空三解算(SfM) / DSM与DOM生成 / 多光谱与高光谱数据解译",
            "• 【数据工程与空间计算】: Spark & Flink 分布式空间 ETL / Kafka 流处理 / 向量数据库(Milvus, Qdrant) / PostGIS / Delta Lake",
            "• 【无人机测绘外业】: 航线规划(DJI Terra, Pix4D) / 野外控制点刺点配准 / 坐标投影转换(CGCS2000与WGS84) / RTK高精度差分"
        ]
    )

    # ------------------ SLIDE 4: PROJECT 1 ------------------
    add_standard_slide(
        "项目一：空天地一体化生态精准监测与SAM视频微调分割",
        "基于 CV 基础模型与遥感影像的亚米级精准语义解译项目",
        [
            "• 航测数据采集与建模：主导了多光谱无人机外业航测，基于 SfM 解算生成了 5cm 高分辨率的正射影像(DOM)和数字表面模型(DSM)。",
            "• SAM 视频微调分割：利用轻量化多层级 Memory Bank 机制，引入了时序自注意力对齐。在视频帧上微调 SAM 2，解决了重叠遮挡和镜头移动下的水体、绿植持续追踪分割难题，使追踪 IoU 精度提升至 91.2%。",
            "• 旋转目标检测(OBB)：使用 YOLO-OBB 对密集排布的温室大棚和地块实施旋转框检测，解决了传统水平边框(HBB)导致的重叠误杀及地物背景干扰，在测试集上 mAP@0.5 达到了 94.6%。",
            "• 地物分类模型微调：基于 MAE (Masked Autoencoders) 预训练骨干网络微调，结合残差跳跃连接(Skip Connection)，大幅提升了细小地物（如田垄、水渠、细小道路）的分割精确率。"
        ]
    )

    # ------------------ SLIDE 5: PROJECT 2 ------------------
    add_standard_slide(
        "项目二：大模型多模态智能体检索RAG与Nous Hermes自进化平台",
        "提升企业级专有地理与航测规范检索精度并具备自主代码执行修正能力的 Agent 平台",
        [
            "• 高可靠混合检索管线：针对专业航测与遥感标准文档，设计了语义切割(Semantic Chunking)保持切片语义连续。结合密集向量与稀疏 BM25 检索，利用 RRF (Reciprocal Rank Fusion) 倒数排名无参融合，Top-5 检索召回率提升至 93.8%。",
            "• 交叉重排(Rerank)防幻觉：串接 Cross-Encoder 重排模型，充分捕获 Query 与切片的精细交互特征。输入大模型前将 Chunk 数量压缩 60% 以上，结合系统提示词将生成结果中的“幻觉率”大幅削减 40% 左右。",
            "• 智能体代码自进化系统：基于 Nous Hermes 3 指令遵循能力，设计了自进化 Skill Engine 闭环。当智能体接收到复杂地球物理模型或坐标系三维转换计算任务时，可自主编写 Python 脚本，并在本地 Sandbox 中调试运行。",
            "• 异常自我纠偏与沉淀：若报错，智能体可捕获 Traceback 并重试修正，验证通过后将该代码封装为标准 Method 存盘，构建了“越用技能越多、计算越精准”的 Agent 自我进化环路。"
        ]
    )

    # ------------------ SLIDE 6: PROJECT 3 ------------------
    add_standard_slide(
        "项目三：基于 Spark/Flink 的大规模遥感与航测空间数据高吞吐处理湖仓",
        "解决海量空间数据读写 IO 瓶颈与高并发流处理计算架构",
        [
            "• 大规模地理数据治湖：基于 Delta Lake 搭建了海量遥感影像和 GPS 数据湖仓。优化数据物理存储分区规则，合并微小文件，克服了数百万张航测小文件照片的频繁随机读写 I/O 瓶颈，每天高稳定承载新增 10TB+ 空间影像。",
            "• 分布式空间 ETL：利用 Spark 配合 Apache Sedona 对传统单机 GDAL 数据清洗流程进行重构，实现了大尺寸 DOM 分布式重采样、波段拼接和空间重投影。在 20 节点集群上，处理时间由 18 小时大幅降至 42 分钟。",
            "• 实时 NDVI 流计算：利用 Flink 引擎消费来自无人机物联网、地面传感器网络的 Kafka 数据流，利用滑窗对植被指数进行实时的异常监控；采用异步 Async I/O 技术与 Milvus 向量集群连接，入库吞吐突破 50,000 ops/sec，并实现了 Exactly-Once 语义。",
            "• 空间查询优化：基于 PostGIS 建立二级 R-Tree 空间索引，优化 Bounding Box 范围查询与 ST_DWithin 空间邻近检索。实现了千亿级空间数据联合拓扑检索，响应时间控制在 100ms 以内。"
        ]
    )
    
    # ------------------ SLIDE 7: END ------------------
    slide_end = prs.slides.add_slide(blank_layout)
    slide_end.background.fill.solid()
    slide_end.background.fill.fore_color.rgb = BG_COLOR
    
    # Bottom line
    bar_end = slide_end.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.4), Inches(13.333), Inches(0.2))
    bar_end.fill.solid()
    bar_end.fill.fore_color.rgb = ACCENT_TEAL
    bar_end.line.fill.background()
    
    # Main Title
    txBox_end = slide_end.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
    tf_end = txBox_end.text_frame
    tf_end.word_wrap = True
    p_end = tf_end.paragraphs[0]
    p_end.text = "谢谢您的观看与审阅"
    p_end.font.name = 'Microsoft YaHei'
    p_end.font.size = Pt(38)
    p_end.font.bold = True
    p_end.font.color.rgb = ACCENT_TEAL
    p_end.alignment = PP_ALIGN.CENTER
    
    p_sub_end = tf_end.add_paragraph()
    p_sub_end.text = "期待有机会与贵公司共同打造前沿行业应用！"
    p_sub_end.font.name = 'Microsoft YaHei'
    p_sub_end.font.size = Pt(18)
    p_sub_end.font.color.rgb = TEXT_TITLE
    p_sub_end.alignment = PP_ALIGN.CENTER
    
    prs.save("面试汇报答辩_可编辑.pptx")
    print("Successfully generated resume PPTX!")

if __name__ == '__main__':
    create_presentation()
