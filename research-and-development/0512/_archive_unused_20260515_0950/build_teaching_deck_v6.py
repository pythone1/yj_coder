# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO
import hashlib
from collections import defaultdict

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PREVIEW_DIR = OUT / "v6_page_previews"
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

TEMPLATE = ROOT / "机理、算法与实践：人工智能行业应用实证分析.pptx"
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
ARCH = ASSET_DIR / "v4_arch"
FULL = ASSET_DIR / "v6_full"

PPTX = OUT / "AI供水管网DMA漏损检测_教学大图可编辑版_v6.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_教学大图可编辑版_v6_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_教学大图可编辑版_v6_预览联系表.png"

NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 74, 153)
CYAN = RGBColor(0, 176, 240)
TEAL = RGBColor(0, 151, 143)
GREEN = RGBColor(28, 148, 92)
ORANGE = RGBColor(245, 132, 32)
RED = RGBColor(212, 57, 57)
PURPLE = RGBColor(92, 95, 214)
TEXT = RGBColor(18, 38, 62)
MUTED = RGBColor(86, 105, 124)
LINE = RGBColor(185, 220, 240)
PALE = RGBColor(239, 249, 255)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
FONT = "微软雅黑"

ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "positioning": FULL / "dma_ai_positioning_full.png",
    "baseline": FULL / "dynamic_baseline_full.png",
    "training": FULL / "data_training_full.png",
    "overall": ARCH / "overall_architecture.png",
    "pain": ARCH / "pain_ai_path.png",
    "alg": ARCH / "algorithm_combo.png",
    "impl": ARCH / "implementation_path.png",
    "workorder": ARCH / "workorder_no_person.png",
}

AUDIT_LINES = []


def emu(v):
    return Inches(v)


def clear_slides(prs):
    ids = prs.slides._sldIdLst
    for sid in list(ids):
        prs.part.drop_rel(sid.rId)
        ids.remove(sid)


def blank(prs, bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def fill(shp, color, trans=0):
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.fill.transparency = trans
    shp.line.fill.background()


def stroke(shp, color=LINE, width=1):
    shp.line.color.rgb = color
    shp.line.width = Pt(width)


def rect(s, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = s.shapes.add_shape(typ, emu(x), emu(y), emu(w), emu(h))
    fill(shp, color, trans)
    stroke(shp, line, 1)
    return shp


def oval(s, x, y, w, h, color, line=None):
    shp = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(x), emu(y), emu(w), emu(h))
    fill(shp, color)
    stroke(shp, line or color, 1)
    return shp


def tb(s, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.02)
    tf.margin_right = emu(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pic(s, key, x, y, w, h, crop=False):
    p = s.shapes.add_picture(str(ASSETS[key]), emu(x), emu(y))
    sx, sy = emu(w) / p.width, emu(h) / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(emu(x) + (emu(w) - p.width) / 2)
    p.top = int(emu(y) + (emu(h) - p.height) / 2)
    return p


def line(s, x1, y1, x2, y2, color=BLUE, width=2.0, arrow=True):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(s, no, title, kicker=None):
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.42), emu(0.28), emu(0.10), emu(0.48)), BLUE)
    tb(s, title, 0.65, 0.18, 10.8, 0.42, 25, TEXT, True)
    if kicker:
        tb(s, kicker, 0.67, 0.64, 9.5, 0.25, 12.5, MUTED)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.65), emu(0.90), emu(11.55), emu(0.02)), RGBColor(220, 238, 248))
    tb(s, f"{no:02d}", 12.0, 0.23, 0.55, 0.20, 13, BLUE, True, PP_ALIGN.RIGHT)


def footer(s):
    tb(s, "AI供水管网DMA漏损检测｜教学演讲版", 0.65, 7.18, 5.0, 0.18, 10.5, MUTED)


def content_slide(prs, no, title, kicker=None):
    s = blank(prs)
    header(s, no, title, kicker)
    footer(s)
    return s


def chip(s, text, x, y, w, color=BLUE):
    rect(s, x, y, w, 0.34, PALE2, color, True)
    tb(s, text, x + 0.06, y + 0.08, w - 0.12, 0.14, 12.5, color, True, PP_ALIGN.CENTER)


def card(s, x, y, w, h, title, body, color=BLUE, num=None):
    rect(s, x, y, w, h, WHITE, LINE, True)
    if h < 0.95:
        tb(s, title, x + 0.18, y + 0.17, 1.00, 0.16, 14.5, color, True)
        tb(s, body, x + 1.10, y + 0.17, w - 1.28, 0.18, 13.8, TEXT, False)
        return
    if num:
        oval(s, x + 0.18, y + 0.18, 0.42, 0.42, color)
        tb(s, str(num), x + 0.18, y + 0.28, 0.42, 0.14, 12.5, WHITE, True, PP_ALIGN.CENTER)
        tx = x + 0.72
        tw = w - 0.90
    else:
        tx = x + 0.20
        tw = w - 0.40
    tb(s, title, tx, y + 0.17, tw, 0.26, 16.5, color, True)
    tb(s, body, x + 0.22, y + 0.62, w - 0.44, h - 0.82, 14.5, TEXT)


def mini_chart(s, x, y, w, h, color=BLUE):
    rect(s, x, y, w, h, RGBColor(250, 254, 255), RGBColor(202, 225, 240), True)
    points = [(x + 0.25, y + h - 0.35), (x + 0.70, y + h - 0.70), (x + 1.15, y + h - 0.55),
              (x + 1.65, y + h - 1.05), (x + 2.10, y + h - 0.80), (x + 2.55, y + h - 1.35)]
    for i in range(len(points) - 1):
        line(s, points[i][0], points[i][1], points[i + 1][0], points[i + 1][1], color, 2, False)
    line(s, x + 0.20, y + h - 0.20, x + w - 0.20, y + h - 0.20, RGBColor(190, 210, 224), 1, False)
    line(s, x + 0.20, y + 0.20, x + 0.20, y + h - 0.20, RGBColor(190, 210, 224), 1, False)


def map_schematic(s, x, y, w, h):
    rect(s, x, y, w, h, RGBColor(248, 252, 255), RGBColor(190, 220, 238), True)
    for gx in [x + w * i / 5 for i in range(1, 5)]:
        line(s, gx, y + 0.20, gx, y + h - 0.20, RGBColor(220, 234, 244), 0.7, False)
    for gy in [y + h * i / 4 for i in range(1, 4)]:
        line(s, x + 0.20, gy, x + w - 0.20, gy, RGBColor(220, 234, 244), 0.7, False)
    pts = [(x + 0.55, y + 0.70), (x + 1.30, y + 1.10), (x + 2.00, y + 0.75), (x + 2.65, y + 1.45),
           (x + 3.35, y + 1.05), (x + 3.85, y + 1.80), (x + 3.05, y + 2.35), (x + 2.25, y + 2.05),
           (x + 1.55, y + 2.55), (x + 0.82, y + 2.05)]
    for i in range(len(pts) - 1):
        line(s, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], BLUE, 1.7, False)
    line(s, pts[0][0], pts[0][1], pts[-1][0], pts[-1][1], BLUE, 1.7, False)
    for px, py in pts:
        oval(s, px - 0.05, py - 0.05, 0.10, 0.10, BLUE)
    # DMA boundary
    line(s, x + 0.35, y + 0.35, x + w - 0.35, y + 0.35, PURPLE, 1.4, False)
    line(s, x + w - 0.35, y + 0.35, x + w - 0.35, y + h - 0.35, PURPLE, 1.4, False)
    line(s, x + w - 0.35, y + h - 0.35, x + 0.35, y + h - 0.35, PURPLE, 1.4, False)
    line(s, x + 0.35, y + h - 0.35, x + 0.35, y + 0.35, PURPLE, 1.4, False)
    oval(s, x + w - 1.05, y + h - 1.10, 0.34, 0.34, RED)
    tb(s, "疑似漏点", x + w - 1.38, y + h - 0.67, 1.10, 0.18, 10.8, RED, True, PP_ALIGN.CENTER)


def cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.62), emu(0.82), emu(0.12), emu(5.90)), CYAN)
    tb(s, "AI模型在供水管网\nDMA系统漏损检测\n中的应用", 1.02, 1.08, 4.15, 1.78, 32, WHITE, True)
    tb(s, "教学演讲版", 1.05, 3.13, 2.50, 0.32, 18, CYAN, True)
    tb(s, "从DMA锁区、动态基线，到AI溯源定位与工单闭环", 1.05, 5.65, 4.35, 0.30, 14.5, RGBColor(215, 236, 248))
    for i, (txt, col) in enumerate([("DMA宏观锁区", BLUE), ("AI微观定位", TEAL), ("工单闭环", ORANGE)]):
        chip(s, txt, 1.05 + i * 1.42, 6.18, 1.18, col)
    pic(s, "positioning", 5.72, 0.62, 7.20, 6.20, crop=False)
    border = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, emu(5.62), emu(0.52), emu(7.40), emu(6.40))
    border.fill.background()
    stroke(border, RGBColor(8, 72, 140), 1.2)
    audit(s, 1, allow_short=True)


def agenda(prs):
    s = content_slide(prs, 2, "课程结构：从问题到落地", "对外讲解按“场景-架构-算法-闭环-实施”展开")
    items = [
        ("1 背景", "行业痛点\nDMA边界\nAI补位价值", BLUE),
        ("2 架构", "数据层\n模型层\n业务闭环", TEAL),
        ("3 算法", "动态基线\n异常检测\n定位排序", ORANGE),
        ("4 应用", "预警复核\n候选管段\n派单回填", GREEN),
        ("5 实施", "数据治理\n训练验证\n系统集成", PURPLE),
        ("6 总结", "建设路线\n课堂演示\n推广边界", RED),
    ]
    for i, (h, b, col) in enumerate(items):
        x = 0.72 + (i % 3) * 4.05
        y = 1.35 + (i // 3) * 2.30
        rect(s, x, y, 3.50, 1.70, WHITE, col, True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(3.50), emu(0.36)), col)
        tb(s, h, x + 0.15, y + 0.11, 3.20, 0.14, 14.5, WHITE, True, PP_ALIGN.CENTER)
        tb(s, b, x + 0.35, y + 0.58, 2.80, 0.72, 15, TEXT, True, PP_ALIGN.CENTER)
    line(s, 2.55, 3.27, 10.85, 3.27, RGBColor(200, 224, 238), 2, False)
    tb(s, "主线：DMA先缩小范围，AI再形成定位证据，工单结果继续反哺模型。", 1.08, 6.28, 11.20, 0.26, 18, BLUE, True, PP_ALIGN.CENTER)
    audit(s, 2)


def section(prs, no, part, title, subtitle, color=CYAN):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(0.24), emu(7.5)), color)
    tb(s, part, 0.92, 1.25, 1.55, 0.30, 17, color, True)
    tb(s, title, 0.92, 1.88, 7.60, 0.68, 36, WHITE, True)
    tb(s, subtitle, 0.94, 3.04, 7.20, 0.55, 18, RGBColor(210, 232, 245))
    topic_items = {
        3: [("DMA边界", "先锁定异常片区"), ("AI证据", "再收敛候选管段"), ("闭环回填", "让模型持续变准")],
        8: [("动态基线", "识别正常范围"), ("异常检测", "发现持续偏离"), ("定位排序", "输出核查优先级")],
        14: [("预警", "提前发现异常"), ("派单", "现场核查处置"), ("复盘", "结果回填再训练")],
        19: [("演示", "LSTM环境跑通"), ("验证", "算法指标+业务指标"), ("案例", "从曲线到工单")],
    }.get(no, [("结构", "讲清逻辑"), ("方法", "讲清边界"), ("输出", "讲清价值")])
    for i, (h, b) in enumerate(topic_items):
        y = 1.35 + i * 1.20
        rect(s, 8.35, y, 3.65, 0.82, RGBColor(255, 255, 255), RGBColor(126, 201, 238), True, 16)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(8.35), emu(y), emu(0.10), emu(0.82)), color)
        tb(s, h, 8.62, y + 0.16, 1.30, 0.20, 17, color, True)
        tb(s, b, 9.92, y + 0.18, 1.85, 0.18, 14.5, TEXT)
    tb(s, "本章讲解路径", 8.35, 5.22, 2.20, 0.24, 16, color, True)
    line(s, 8.35, 5.78, 11.65, 5.78, color, 2)
    tb(s, f"{no:02d}", 11.80, 6.58, 0.62, 0.22, 15, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)
    audit(s, no, allow_short=True)


def dma_ai_editable(prs, no=4):
    s = content_slide(prs, no, "DMA宏观锁区 + AI微观定位：可编辑架构图", "所有标题、模块、箭头、表格均为PPT原生对象")
    colors = [BLUE, CYAN, TEAL, PURPLE, ORANGE, RED]
    heads = ["DMA分区计量", "夜间最小流量", "多源数据融合", "AI溯源分析", "候选管段排序", "现场核查闭环"]
    bodies = ["关闭边界阀\n形成独立计量单元", "识别MNF抬升\n判断疑似漏损", "流量 压力 GIS\n资产 工单 天气", "模型反推异常来源\n融合水力约束", "输出TopN管段\n给出风险评分", "检漏维修回填\n形成训练标签"]
    for i, (h, b, col) in enumerate(zip(heads, bodies, colors)):
        x = 0.42 + i * 2.08
        rect(s, x, 1.18, 1.86, 1.16, WHITE, col, True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(1.18), emu(1.86), emu(0.32)), col)
        tb(s, f"{i+1} {h}", x + 0.05, 1.26, 1.76, 0.12, 10.8, WHITE, True, PP_ALIGN.CENTER)
        tb(s, b, x + 0.12, 1.62, 1.62, 0.44, 12.5, TEXT, False, PP_ALIGN.CENTER)
        if i < 5:
            line(s, x + 1.88, 1.76, x + 2.04, 1.76, BLUE, 1.8)
    map_schematic(s, 0.58, 2.82, 4.32, 3.00)
    rect(s, 5.28, 2.80, 2.25, 3.05, PALE2, TEAL, True)
    tb(s, "融合特征", 5.55, 3.02, 1.72, 0.24, 17, TEAL, True, PP_ALIGN.CENTER)
    for i, t in enumerate(["流量序列", "压力序列", "GIS拓扑", "管龄管材", "历史工单", "天气施工"]):
        chip(s, t, 5.55, 3.45 + i * 0.34, 1.70, TEAL)
    oval(s, 8.18, 3.16, 1.42, 1.42, BLUE)
    tb(s, "AI", 8.46, 3.55, 0.84, 0.32, 30, WHITE, True, PP_ALIGN.CENTER)
    tb(s, "残差分析\n水力约束\n多模型集成", 7.78, 4.80, 2.22, 0.55, 16, BLUE, True, PP_ALIGN.CENTER)
    line(s, 7.55, 4.10, 8.12, 4.10, BLUE, 2)
    line(s, 9.68, 4.10, 10.15, 4.10, BLUE, 2)
    rect(s, 10.18, 2.82, 2.18, 2.90, WHITE, ORANGE, True)
    tb(s, "候选管段 Top 3", 10.38, 3.04, 1.78, 0.22, 15, ORANGE, True, PP_ALIGN.CENTER)
    rows = [("P-1256", "0.92", "优先"), ("P-1289", "0.71", "复核"), ("P-1345", "0.56", "观察")]
    for i, row in enumerate(rows):
        y = 3.48 + i * 0.55
        rect(s, 10.38, y, 1.76, 0.36, RGBColor(252, 254, 255), RGBColor(225, 235, 242), False)
        tb(s, f"{i+1}. {row[0]}   {row[1]}   {row[2]}", 10.45, y + 0.08, 1.62, 0.12, 11.5, TEXT if i else RED, True)
    tb(s, "讲解逻辑：DMA负责“缩小范围”，AI负责“形成证据链”，工单负责“验证并回填”。", 1.05, 6.43, 11.0, 0.24, 18, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def full_visual(prs, no, key, overlay=None):
    s = blank(prs)
    pic(s, key, 0, 0, 13.333, 7.5, crop=False)
    if overlay:
        rect(s, 0.45, 6.86, 12.45, 0.36, RGBColor(255, 255, 255), RGBColor(220, 235, 246), True, 8)
        tb(s, overlay, 0.70, 6.94, 11.95, 0.15, 14, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no, allow_short=True)


def data_architecture(prs, no=7):
    s = content_slide(prs, no, "数据架构：模型不是单点算法，而是数据链路", "把管网运行、空间拓扑和业务处置放到同一张表述中")
    cols = [("感知数据", ["流量计", "压力计", "水质仪", "阀门状态"], BLUE),
            ("空间资产", ["GIS拓扑", "管段属性", "管龄管材", "DMA边界"], TEAL),
            ("业务事件", ["工单记录", "投诉报修", "施工信息", "天气节假日"], ORANGE)]
    for i, (h, items, col) in enumerate(cols):
        x = 0.72 + i * 4.05
        rect(s, x, 1.25, 3.35, 3.25, WHITE, col, True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(1.25), emu(3.35), emu(0.42)), col)
        tb(s, h, x + 0.12, 1.37, 3.10, 0.18, 15, WHITE, True, PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            chip(s, item, x + 0.42, 1.95 + j * 0.52, 2.50, col)
    line(s, 2.40, 4.66, 2.40, 5.15, BLUE, 2)
    line(s, 6.45, 4.66, 6.45, 5.15, TEAL, 2)
    line(s, 10.50, 4.66, 10.50, 5.15, ORANGE, 2)
    rect(s, 1.18, 5.15, 11.05, 0.82, PALE, BLUE, True)
    tb(s, "统一编码、时间对齐、缺失补全、异常剔除、标签回填", 1.55, 5.38, 10.30, 0.20, 19, BLUE, True, PP_ALIGN.CENTER)
    tb(s, "输出给模型的不是原始台账，而是可追踪、可复核、可持续更新的训练样本。", 1.10, 6.40, 11.20, 0.24, 18, TEXT, True, PP_ALIGN.CENTER)
    audit(s, no)


def lstm_slide(prs, no=9):
    s = content_slide(prs, no, "LSTM动态基线：用历史序列预测“正常范围”", "本地演示环境使用 conda 的 LSTM 环境")
    mini_chart(s, 0.75, 1.36, 3.00, 1.70, BLUE)
    tb(s, "输入窗口", 1.25, 3.20, 2.00, 0.22, 17, BLUE, True, PP_ALIGN.CENTER)
    tb(s, "过去24h/7d流量、压力、节假日、天气", 0.72, 3.58, 3.10, 0.45, 14.5, TEXT, False, PP_ALIGN.CENTER)
    for i in range(3):
        rect(s, 4.25 + i * 0.64, 1.52 + i * 0.16, 0.48, 0.88, RGBColor(238, 248, 255), BLUE, True)
        tb(s, "LSTM", 4.20 + i * 0.64, 1.84 + i * 0.16, 0.60, 0.15, 10.5, BLUE, True, PP_ALIGN.CENTER)
    line(s, 3.95, 2.18, 4.22, 2.18, BLUE, 2)
    line(s, 6.10, 2.22, 6.70, 2.22, BLUE, 2)
    mini_chart(s, 6.75, 1.36, 2.75, 1.70, GREEN)
    tb(s, "预测基线", 7.18, 3.20, 1.90, 0.22, 17, GREEN, True, PP_ALIGN.CENTER)
    tb(s, "输出下一时段正常区间，随用水规律自动更新", 6.48, 3.58, 3.30, 0.45, 14.5, TEXT, False, PP_ALIGN.CENTER)
    line(s, 9.68, 2.22, 10.08, 2.22, BLUE, 2)
    rect(s, 10.18, 1.36, 2.15, 1.70, RGBColor(255, 248, 246), RED, True)
    tb(s, "残差阈值", 10.48, 1.68, 1.55, 0.22, 18, RED, True, PP_ALIGN.CENTER)
    tb(s, "实际 - 预测\n持续偏离才报警", 10.43, 2.10, 1.60, 0.42, 15, TEXT, True, PP_ALIGN.CENTER)
    rows = [("训练集", "历史正常运行数据"), ("验证集", "按时间切分，避免未来泄漏"), ("输出", "动态基线、残差、异常等级"), ("使用", "预警复核，不替代现场确认")]
    for i, (a, b) in enumerate(rows):
        card(s, 0.95 + (i % 2) * 5.95, 4.82 + (i // 2) * 0.76, 5.35, 0.55, a, b, [BLUE, TEAL, ORANGE, RED][i])
    audit(s, no)


def algorithm_matrix(prs, no=10):
    s = content_slide(prs, no, "算法选型：按任务分工，而不是堆算法名称", "教学重点放在输入、输出和业务动作")
    headers = ["任务", "推荐模型", "输入特征", "输出结果", "讲解重点"]
    widths = [1.55, 2.25, 3.45, 2.25, 2.20]
    x0, y0 = 0.55, 1.25
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x0), emu(y0), emu(12.25), emu(0.46)), BLUE)
    x = x0
    for h, w in zip(headers, widths):
        tb(s, h, x + 0.05, y0 + 0.14, w - 0.10, 0.14, 12.8, WHITE, True, PP_ALIGN.CENTER)
        x += w
    rows = [
        ("动态基线", "LSTM/GRU", "历史流量、压力、日历、天气", "预测区间、残差", "判断“是否异常”"),
        ("异常识别", "孤立森林/自编码器", "残差、波动、压力响应", "异常等级", "减少固定阈值误报"),
        ("空间定位", "水力模型+树模型", "拓扑、仿真样本、传感器响应", "候选管段TopN", "把片区收敛到管段"),
        ("风险排序", "规则+机器学习", "管龄、材质、维修、压力", "巡检优先级", "指导先查哪里"),
        ("闭环优化", "再训练/漂移监测", "工单结果、误报原因", "版本更新", "越用越准"),
    ]
    y = y0 + 0.46
    for i, row in enumerate(rows):
        bg = RGBColor(248, 252, 255) if i % 2 == 0 else WHITE
        rect(s, x0, y, 12.25, 0.72, bg, RGBColor(225, 238, 246), False)
        x = x0
        for j, (txt, w) in enumerate(zip(row, widths)):
            col = [BLUE, TEAL, TEXT, ORANGE, MUTED][min(j, 4)]
            tb(s, txt, x + 0.08, y + 0.18, w - 0.16, 0.25, 12.8 if j != 2 else 12.0, col, j in [0, 1], PP_ALIGN.CENTER)
            x += w
        y += 0.72
    tb(s, "落地建议：课堂演示只用 LSTM 环境完成动态基线与残差报警，其余算法作为系统扩展讲清楚边界。", 0.85, 6.20, 11.55, 0.28, 17, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def hydraulic_fusion(prs, no=11):
    s = content_slide(prs, no, "水力模型 + AI：定位可信度的关键", "AI发现异常，水力约束判断异常是否符合管网物理规律")
    map_schematic(s, 0.70, 1.28, 4.42, 3.25)
    card(s, 5.55, 1.20, 2.70, 1.10, "AI证据", "残差升高、压力响应、相邻点协同变化", BLUE)
    card(s, 5.55, 2.75, 2.70, 1.10, "水力约束", "流向、阀门状态、节点压力、候选漏量", TEAL)
    card(s, 5.55, 4.30, 2.70, 1.10, "综合评分", "候选管段置信度与核查顺序", ORANGE)
    line(s, 4.95, 2.85, 5.45, 1.75, BLUE, 2)
    line(s, 4.95, 2.85, 5.45, 3.30, TEAL, 2)
    line(s, 8.35, 1.75, 9.08, 3.38, BLUE, 2)
    line(s, 8.35, 3.30, 9.08, 3.38, TEAL, 2)
    rect(s, 9.18, 2.20, 3.10, 2.40, RGBColor(252, 254, 255), ORANGE, True)
    tb(s, "定位输出", 9.56, 2.55, 2.34, 0.26, 20, ORANGE, True, PP_ALIGN.CENTER)
    tb(s, "TopN候选管段\n节点坐标\n置信度\n现场核查建议", 9.58, 3.10, 2.30, 0.88, 17, TEXT, True, PP_ALIGN.CENTER)
    tb(s, "核心讲法：模型不是直接“猜点位”，而是把多个证据按管网物理关系进行收敛。", 1.15, 6.32, 10.95, 0.24, 18, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def localization_workflow(prs, no=12):
    s = content_slide(prs, no, "事中定位：从片区收敛到候选管段", "适合课堂逐步讲解的可编辑流程页")
    steps = [("1 发现", "DMA夜间最小流量异常"), ("2 复核", "排除施工、计量、用水扰动"), ("3 反推", "水力模型与AI残差融合"),
             ("4 排序", "候选管段TopN与置信度"), ("5 派单", "检漏、维修、结果回填")]
    for i, (h, b) in enumerate(steps):
        x = 0.70 + i * 2.48
        color = [BLUE, TEAL, PURPLE, ORANGE, RED][i]
        oval(s, x + 0.62, 1.30, 0.68, 0.68, color)
        tb(s, str(i + 1), x + 0.62, 1.48, 0.68, 0.18, 15, WHITE, True, PP_ALIGN.CENTER)
        tb(s, h, x + 0.30, 2.20, 1.35, 0.26, 18, color, True, PP_ALIGN.CENTER)
        tb(s, b, x + 0.05, 2.62, 1.86, 0.50, 14.0, TEXT, False, PP_ALIGN.CENTER)
        if i < 4:
            line(s, x + 1.42, 1.64, x + 2.38, 1.64, RGBColor(160, 198, 222), 1.6)
    rect(s, 1.00, 4.20, 11.25, 1.35, PALE2, BLUE, True)
    tb(s, "输出物", 1.35, 4.53, 1.12, 0.24, 19, BLUE, True)
    tb(s, "异常类型、可疑DMA、候选管段、核查优先级、派单建议、回填字段", 2.40, 4.55, 9.15, 0.22, 18, TEXT, True)
    tb(s, "教学提示：这页讲“流程”，不要讲复杂公式；公式和训练细节放到算法页。", 1.30, 6.28, 10.75, 0.24, 18, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def deployment(prs, no=16):
    s = content_slide(prs, no, "工程部署：端边云协同与业务系统集成", "把实时性、成本和可靠性放在同一张架构里")
    parts = [("现场感知层", "流量计 压力计 水质仪\n阀门状态 高频采样", BLUE),
             ("边缘网关", "数据缓存 初筛 校时\n断网续传 异常压缩", TEAL),
             ("云端模型服务", "LSTM基线 异常检测\n定位排序 模型监控", PURPLE),
             ("业务应用层", "SCADA GIS 工单\n移动端 巡检看板", ORANGE)]
    for i, (h, b, col) in enumerate(parts):
        x = 0.68 + i * 3.13
        card(s, x, 1.42, 2.55, 2.55, h, b, col)
        # Native pictograms keep the slide editable while avoiding blank cards.
        if i == 0:
            line(s, x + 0.35, 3.22, x + 2.12, 3.22, col, 5, False)
            oval(s, x + 0.58, 2.86, 0.38, 0.38, col)
            oval(s, x + 1.50, 2.78, 0.48, 0.48, col)
            line(s, x + 0.77, 2.86, x + 0.77, 2.42, col, 2, False)
            line(s, x + 1.74, 2.78, x + 1.74, 2.32, col, 2, False)
        elif i == 1:
            for k in range(3):
                rect(s, x + 0.58, 2.72 + k * 0.24, 1.35, 0.16, RGBColor(235, 249, 247), col, False)
            line(s, x + 1.25, 2.32, x + 1.25, 3.42, col, 1.6, False)
            oval(s, x + 0.88, 2.18, 0.18, 0.18, col)
            oval(s, x + 1.44, 2.18, 0.18, 0.18, col)
        elif i == 2:
            oval(s, x + 0.62, 2.62, 1.25, 0.64, RGBColor(238, 241, 255), col)
            tb(s, "AI", x + 0.86, 2.80, 0.72, 0.18, 17, col, True, PP_ALIGN.CENTER)
            line(s, x + 0.72, 3.48, x + 1.90, 3.48, col, 2, False)
            line(s, x + 0.72, 3.62, x + 1.55, 3.62, col, 2, False)
        else:
            rect(s, x + 0.48, 2.55, 1.15, 0.70, RGBColor(255, 248, 238), col, True)
            rect(s, x + 1.76, 2.48, 0.46, 0.86, RGBColor(255, 248, 238), col, True)
            line(s, x + 0.65, 3.07, x + 1.45, 3.07, col, 2, False)
            line(s, x + 1.88, 3.12, x + 2.10, 3.12, col, 2, False)
        if i < 3:
            line(s, x + 2.62, 2.70, x + 3.05, 2.70, BLUE, 2)
    rect(s, 1.00, 4.75, 11.25, 1.10, RGBColor(248, 252, 255), RGBColor(198, 225, 241), True)
    tb(s, "关键约束：传感器在线率、时间同步、数据质量、模型版本、工单回填字段。", 1.25, 5.12, 10.75, 0.24, 18, TEXT, True, PP_ALIGN.CENTER)
    audit(s, no)


def validation(prs, no=18):
    s = content_slide(prs, no, "验证指标：算法准确，不等于业务可用", "验证要同时看模型指标和现场处置指标")
    left = [("MAE/RMSE", 0.78, BLUE), ("召回率", 0.70, TEAL), ("误报率控制", 0.62, ORANGE), ("TopN命中", 0.82, GREEN)]
    right = [("发现时长", 0.72, BLUE), ("平均定位半径", 0.68, TEAL), ("无效派单率", 0.55, ORANGE), ("闭环回填率", 0.86, GREEN)]
    for title, data, x0 in [("模型指标", left, 1.00), ("业务指标", right, 7.00)]:
        tb(s, title, x0, 1.30, 4.90, 0.30, 22, BLUE if x0 < 2 else TEAL, True, PP_ALIGN.CENTER)
        for i, (name, val, col) in enumerate(data):
            y = 1.95 + i * 0.78
            tb(s, name, x0, y + 0.05, 1.55, 0.18, 13.5, TEXT, True)
            rect(s, x0 + 1.72, y, 3.05, 0.28, RGBColor(232, 241, 248), RGBColor(232, 241, 248), False)
            rect(s, x0 + 1.72, y, 3.05 * val, 0.28, col, col, False)
    tb(s, "课堂结论：模型评价必须回到“少漏报、少误派、快定位、能复盘”。", 1.20, 6.10, 10.90, 0.28, 21, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def case_exercise(prs, no=19):
    s = content_slide(prs, no, "课堂案例：夜间最小流量异常如何判断", "用一张图讲清楚从报警到派单的判断链")
    mini_chart(s, 0.78, 1.28, 4.20, 2.55, BLUE)
    tb(s, "现象：凌晨2:00-4:00最小流量连续三天抬升，压力传感器P-07同步下降。", 0.90, 4.05, 3.95, 0.55, 16, TEXT, True, PP_ALIGN.CENTER)
    questions = [
        ("先排除什么？", "施工、阀门调整、计量故障、异常用水"),
        ("需要融合什么？", "DMA边界、压力响应、管网拓扑、历史工单"),
        ("输出什么？", "候选管段、置信度、核查路线、回填字段"),
    ]
    for i, (q, a) in enumerate(questions):
        card(s, 5.55, 1.30 + i * 1.35, 6.25, 1.02, q, a, [BLUE, TEAL, ORANGE][i])
    tb(s, "讲解目标：让听众理解AI不是替代检漏人员，而是把核查范围从“片区”缩小到“优先管段”。", 1.25, 6.22, 10.80, 0.27, 17.5, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def operations(prs, no=20):
    s = content_slide(prs, no, "长效运营：上线后如何保持有效", "漏损AI系统的价值来自持续回填和版本管理")
    items = [("监测", "数据漂移\n设备离线\n报警分布", BLUE),
             ("复盘", "误报原因\n漏报原因\n处置时长", ORANGE),
             ("标注", "工单结果\n现场确认\n维修记录", TEAL),
             ("再训练", "样本更新\n模型评估\n灰度发布", PURPLE),
             ("推广", "试点复核\n规则沉淀\n多DMA复制", GREEN)]
    for i, (h, b, col) in enumerate(items):
        x = 0.75 + i * 2.48
        card(s, x, 1.68, 2.05, 2.65, h, b, col, i + 1)
        if i < 4:
            line(s, x + 2.07, 3.00, x + 2.40, 3.00, BLUE, 2)
    rect(s, 1.00, 5.22, 11.20, 0.92, PALE2, BLUE, True)
    tb(s, "没有工单回填，就没有高质量标签；没有版本管理，就无法证明模型变好。", 1.40, 5.55, 10.35, 0.22, 18.5, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def future(prs, no=21):
    s = content_slide(prs, no, "发展方向：从漏损检测到管网全生命周期管理", "面向未来的讲解保持克制，落在可解释、可验证、可运营")
    items = [
        ("数字孪生", "把水力模型、实时数据和业务状态统一到同一对象。", BLUE),
        ("图模型", "利用管网拓扑表达压力响应和异常传播关系。", TEAL),
        ("智能体辅助", "用自然语言检索工单、生成复盘、辅助调度。", PURPLE),
        ("资产决策", "从漏点定位延伸到改造优先级和风险治理。", ORANGE),
    ]
    for i, (h, b, col) in enumerate(items):
        x = 1.00 + (i % 2) * 5.85
        y = 1.45 + (i // 2) * 2.12
        card(s, x, y, 5.05, 1.55, h, b, col)
    tb(s, "底线：任何新技术都必须能被数据验证、被现场复核、被工单闭环。", 1.18, 6.25, 10.95, 0.28, 20, BLUE, True, PP_ALIGN.CENTER)
    audit(s, no)


def summary(prs, no=22):
    s = blank(prs, NAVY)
    tb(s, "三句话总结", 0.92, 0.88, 4.5, 0.48, 34, WHITE, True)
    items = [("DMA先锁区", "把复杂管网划分为可计量、可考核、可预警的片区。"),
             ("AI再定位", "融合流量、压力、拓扑、资产和工单，收敛到候选管段。"),
             ("闭环才变准", "现场核查和维修结果回填，形成模型持续优化的数据基础。")]
    for i, (h, b) in enumerate(items):
        y = 1.90 + i * 1.38
        color = [CYAN, TEAL, ORANGE][i]
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.95), emu(y), emu(0.12), emu(0.72)), color)
        tb(s, h, 1.25, y - 0.02, 2.80, 0.30, 22, color, True)
        tb(s, b, 1.25, y + 0.42, 6.15, 0.30, 16.5, RGBColor(220, 236, 248))
    pic(s, "overall", 7.70, 1.10, 5.05, 4.70, crop=False)
    tb(s, "AI供水管网DMA漏损检测｜教学演讲版", 0.95, 6.72, 5.0, 0.22, 11.5, RGBColor(166, 208, 235))
    audit(s, no, allow_short=True)


def audit(s, no, allow_short=False):
    prs_w, prs_h = 12192000, 6858000
    chars = pics = off = 0
    min_font = 99
    bottom = 0
    for sh in s.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -60000 or sh.top < -60000 or sh.left + sh.width > prs_w + 60000 or sh.top + sh.height > prs_h + 60000:
            off += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off:
        flags.append(f"越界{off}")
    if min_font < 10:
        flags.append(f"小字{min_font}")
    if not allow_short and chars < 55:
        flags.append(f"可见内容偏少{chars}")
    if bottom < emu(5.75):
        flags.append("页面下半部偏空")
    AUDIT_LINES.append(f"{no:02d}\t{'FAIL' if flags else 'OK'}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise RuntimeError(f"slide {no}: {'; '.join(flags)}")


def build():
    missing = [str(p) for p in ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("\n".join(missing))
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    cover(prs)
    agenda(prs)
    section(prs, 3, "PART 01", "问题与总体逻辑", "先讲清DMA能解决什么，再讲AI补上哪一段。")
    dma_ai_editable(prs, 4)
    full_visual(prs, 5, "positioning")
    full_visual(prs, 6, "overall")
    data_architecture(prs, 7)
    section(prs, 8, "PART 02", "模型与算法", "算法按任务组织，课堂演示以LSTM动态基线为主。", TEAL)
    full_visual(prs, 9, "baseline")
    lstm_slide(prs, 10)
    algorithm_matrix(prs, 11)
    hydraulic_fusion(prs, 12)
    localization_workflow(prs, 13)
    section(prs, 14, "PART 03", "应用闭环", "从预警、定位、派单到复盘回填。", ORANGE)
    full_visual(prs, 15, "workorder")
    full_visual(prs, 16, "training")
    full_visual(prs, 17, "impl")
    deployment(prs, 18)
    section(prs, 19, "PART 04", "教学演示与验证", "把模型效果转化成现场可理解的判断链。", PURPLE)
    validation(prs, 20)
    case_exercise(prs, 21)
    operations(prs, 22)
    future(prs, 23)
    summary(prs, 24)

    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    build_previews(prs)
    append_repetition_check(prs)
    return PPTX


def rgb_tuple(rgb):
    if rgb is None:
        return None
    s = str(rgb)
    return tuple(int(s[i:i + 2], 16) for i in (0, 2, 4))


def draw_text_wrapped(d, xy, wh, text, font, fill, align="left"):
    x, y = xy
    w, h = wh
    lines = []
    for raw in text.splitlines():
        cur = ""
        for ch in raw:
            trial = cur + ch
            if d.textbbox((0, 0), trial, font=font)[2] <= w or not cur:
                cur = trial
            else:
                lines.append(cur)
                cur = ch
        lines.append(cur)
    yy = y
    step = max(12, int(font.size * 1.28))
    for ln in lines[: max(1, h // step)]:
        tw = d.textbbox((0, 0), ln, font=font)[2]
        xx = x + (w - tw) // 2 if align == "center" else x
        d.text((xx, yy), ln, font=font, fill=fill)
        yy += step


def render_slide(slide, path, W=1600, H=900):
    sx, sy = W / 12192000, H / 6858000
    im = Image.new("RGB", (W, H), "white")
    d = ImageDraw.Draw(im)
    try:
        bg = rgb_tuple(slide.background.fill.fore_color.rgb)
        if bg:
            im.paste(bg, [0, 0, W, H])
    except Exception:
        pass
    font_cache = {}

    def get_font(size, bold=False):
        key = (round(size), bold)
        if key not in font_cache:
            fp = "C:/Windows/Fonts/msyhbd.ttc" if bold else "C:/Windows/Fonts/msyh.ttc"
            try:
                font_cache[key] = ImageFont.truetype(fp, max(8, int(size)))
            except Exception:
                font_cache[key] = ImageFont.load_default()
        return font_cache[key]

    for sh in slide.shapes:
        x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
        if w <= 0 or h <= 0:
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                p.thumbnail((w, h), Image.LANCZOS)
                im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
            except Exception:
                pass
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.LINE:
            try:
                lc = rgb_tuple(sh.line.color.rgb) or (0, 74, 153)
            except Exception:
                lc = (0, 74, 153)
            d.line((x, y, x + w, y + h), fill=lc, width=max(1, int(sh.line.width.pt if sh.line.width else 1)))
            continue
        fc = None
        lc = None
        try:
            fc = rgb_tuple(sh.fill.fore_color.rgb)
        except Exception:
            pass
        try:
            lc = rgb_tuple(sh.line.color.rgb)
        except Exception:
            pass
        if fc:
            d.rounded_rectangle([x, y, x + w, y + h], radius=10, fill=fc, outline=lc)
        elif lc:
            d.rectangle([x, y, x + w, y + h], outline=lc)
        txt = getattr(sh, "text", "").strip()
        if txt:
            p = sh.text_frame.paragraphs[0]
            r = p.runs[0] if p.runs else None
            size = (r.font.size.pt if r and r.font.size else 14) * W / 1600
            bold = bool(r.font.bold) if r else False
            try:
                color = rgb_tuple(r.font.color.rgb) if r else (18, 38, 62)
            except Exception:
                color = (18, 38, 62)
            align = "center" if p.alignment == PP_ALIGN.CENTER else "left"
            draw_text_wrapped(d, (x + 4, y + 3), (max(1, w - 8), max(1, h - 6)), txt, get_font(size, bold), color or (18, 38, 62), align)
    im.save(path)
    return path


def build_previews(prs):
    page_paths = []
    for idx, slide in enumerate(prs.slides, 1):
        p = PREVIEW_DIR / f"slide_{idx:02d}.png"
        render_slide(slide, p)
        page_paths.append(p)
    W, H = 400, 225
    cols = 4
    rows = (len(page_paths) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * W, rows * (H + 26)), (230, 236, 242))
    d = ImageDraw.Draw(sheet)
    try:
        f = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 12)
    except Exception:
        f = ImageFont.load_default()
    for idx, path in enumerate(page_paths):
        im = Image.open(path).convert("RGB")
        im.thumbnail((W, H), Image.LANCZOS)
        x = (idx % cols) * W
        y = (idx // cols) * (H + 26)
        sheet.paste(im, (x, y))
        d.text((x + 6, y + H + 5), f"{idx + 1:02d}", fill=(0, 74, 153), font=f)
    sheet.save(PREVIEW)


def append_repetition_check(prs):
    by = defaultdict(list)
    for si, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                by[hashlib.md5(sh.image.blob).hexdigest()[:10]].append(si)
    with AUDIT.open("a", encoding="utf-8") as f:
        f.write("\n\n图片重复检查：\n")
        f.write(f"唯一图片={len(by)}，总图片={sum(len(v) for v in by.values())}\n")
        for h, slides in sorted(by.items(), key=lambda kv: len(kv[1]), reverse=True):
            if len(slides) > 3:
                f.write(f"{h}\t{len(slides)}\t{slides}\n")


if __name__ == "__main__":
    print(build())
    print(AUDIT)
    print(PREVIEW)
