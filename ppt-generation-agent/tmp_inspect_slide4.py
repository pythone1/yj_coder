from pathlib import Path

from pptx import Presentation

p = Path(r"E:\PY\research\0424\Jinshiyuan_AI_Production_Blueprint.pptx")
prs = Presentation(str(p))
slide = prs.slides[3]
for i, shape in enumerate(slide.shapes, 1):
    text = ""
    if hasattr(shape, "text_frame") and shape.text_frame:
        text = "\n".join(p.text for p in shape.text_frame.paragraphs).strip()
    print(i, shape.shape_type, int(shape.left), int(shape.top), int(shape.width), int(shape.height), repr(text[:200]))
