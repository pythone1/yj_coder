from pptx import Presentation
from pathlib import Path
p=Path(r'E:\PY\research\output\ppt\Jinshiyuan_AI_Production_Blueprint_可编辑版.pptx')
prs=Presentation(str(p))
print('file', p)
print('slides', len(prs.slides))
print('size', p.stat().st_size)
for i,s in enumerate(prs.slides,1):
    pics=0; texts=0; shapes=0; nonpic=0
    for sh in s.shapes:
        shapes += 1
        if sh.shape_type == 13:
            pics += 1
        else:
            nonpic += 1
        if getattr(sh,'has_text_frame',False) and sh.text.strip():
            texts += 1
    print(i, 'shapes', shapes, 'pictures', pics, 'editable_non_picture', nonpic, 'text_shapes', texts)
