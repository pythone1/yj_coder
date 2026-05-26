from pptx import Presentation
from pathlib import Path
p=Path(r'E:\PY\research\output\ppt\Jinshiyuan_AI_Production_Blueprint_可编辑版.pptx')
if not p.exists():
    matches=list(Path(r'E:\PY\research\output\ppt').glob('Jinshiyuan_AI_Production_Blueprint*.pptx'))
    print('matches', [str(x) for x in matches])
    p=matches[-1]
prs=Presentation(str(p))
print('file', str(p))
for i,s in enumerate(prs.slides,1):
    for sh in s.shapes:
        if sh.shape_type == 13:
            print(i, 'picture_inches', round(sh.left/914400,2), round(sh.top/914400,2), round(sh.width/914400,2), round(sh.height/914400,2))
