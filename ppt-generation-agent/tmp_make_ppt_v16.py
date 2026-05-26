from io import BytesIO
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.util import Emu, Pt


BASE = Path(r"E:\PY\research\0424")
SOURCE = BASE / "Jinshiyuan_AI_Production_Blueprint.pptx"
TEMPLATE = BASE / "PPT模版.pptx"
OUTPUT = BASE / "Jinshiyuan_AI_Production_Blueprint_套模板_可读性修正_v16.pptx"

TITLES = [
    "今世缘酒业生产模块AI工艺分析与实施路径",
    "生产模块总体逻辑",
    "匠心传承与数字未来的交汇点",
    "酿酒指挥中心",
    "包装智能质检",
    "设备预测性维护",
    "AGV路径优化",
    "仓储物流优化",
    "分阶段推进建议与实施路径",
    "标杆案例与前沿技术背书（参考资料）",
]

TEXT_FIXES = {
    "设备墓础": "设备基础",
    "历皮故障": "历史故障",
    "送行时序": "运行时序",
    "负乾": "负载",
    "启俘": "启停",
    "产绫": "产线",
    "艳出": "输出",
    "造垮春板": "运维看板",
    "识剔": "识别",
    "建汉": "建议",
    "限饥森林": "随机森林",
    "故  标签": "故障标签",
}
MARKERS = {"南大五维", "NotebookLM", "ANotebookLM"}


def clean_text(value):
    value = value.replace("ANotebookLM", "").replace("NotebookLM", "")
    for old, new in TEXT_FIXES.items():
        value = value.replace(old, new)
    value = re.sub(r"(?<=[\u4e00-\u9fff])\s+(?=[\u4e00-\u9fffA-Za-z0-9])", "", value)
    value = re.sub(r"(?<=[A-Za-z])\s{2,}(?=[A-Za-z])", " ", value)
    value = re.sub(r"\s+([，。；：、])", r"\1", value)
    return value


def text_of(shape):
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def remove(shape):
    shape._element.getparent().remove(shape._element)


def tune_text(shape, slide_idx):
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    if para.text.strip() in MARKERS:
                        for run in para.runs:
                            run.text = ""
                        continue
                    for run in para.runs:
                        run.text = clean_text(run.text)
                        if not run.text.strip():
                            continue
                        run.font.color.rgb = RGBColor(245, 250, 255)
                        floor = 13 if slide_idx == 1 else 12
                        run.font.size = Pt(max(run.font.size.pt if run.font.size else floor, floor))
    if hasattr(shape, "text_frame") and shape.text_frame is not None:
        for para in shape.text_frame.paragraphs:
            if para.text.strip() in MARKERS:
                for run in para.runs:
                    run.text = ""
                continue
            for run in para.runs:
                run.text = clean_text(run.text)
                if not run.text.strip():
                    continue
                if slide_idx >= 1:
                    run.font.color.rgb = RGBColor(245, 250, 255)
                    floor = 13 if slide_idx == 1 else 12
                    run.font.size = Pt(max(run.font.size.pt if run.font.size else floor, floor))
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            tune_text(child, slide_idx)


def is_old_title(shape, title):
    text = text_of(shape).replace("\n", "").replace(" ", "")
    target = title.replace(" ", "")
    if not text:
        return False
    if text == target:
        return True
    return text in target and len(text) >= 8 and shape.top < 3000000


def is_old_logo(shape, width, height):
    text = text_of(shape).strip()
    if text == "南大五维":
        return True
    return shape.left > width * 0.84 and shape.top < 900000 and shape.width < 1400000 and shape.height < 1200000


def is_bottom_mark_picture(shape, width, height):
    return shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.top > height * 0.9 and shape.left > width * 0.75


def put_after_background(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(3, el)


def add_dark_logo_mask(slide):
    mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 14760000, 1170000, 1380000, 430000)
    mask.fill.solid()
    mask.fill.fore_color.rgb = RGBColor(4, 32, 48)
    mask.line.fill.background()


def add_old_frame_mask(slide, slide_idx):
    if slide_idx == 3:
        # Cover the residual edge of the original title frame without touching body text.
        mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 3600000, 1195000, 3200000, 220000)
        mask.fill.solid()
        mask.fill.fore_color.rgb = RGBColor(5, 31, 47)
        mask.line.fill.background()
    if slide_idx in (4, 5):
        # Old logo text is baked into the image in these slides.
        mask = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 14700000, 1180000, 1400000, 430000)
        mask.fill.solid()
        mask.fill.fore_color.rgb = RGBColor(5, 31, 47)
        mask.line.fill.background()


def add_frame(slide, logo_blob, title, scale, width, slide_idx):
    def s(value):
        return Emu(int(value * scale))

    if slide_idx in (3, 5):
        add_dark_logo_mask(slide)
    add_old_frame_mask(slide, slide_idx)

    header_base = 1380000 if slide_idx == 9 else 900000
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, width, Emu(int(header_base * scale)))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(255, 255, 255)
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(s(262255), s(245745), s(8756650), s(521970))
    tf = title_box.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    run = tf.paragraphs[0].add_run()
    run.text = title
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(25)
    run.font.color.rgb = RGBColor(31, 78, 121)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, s(349759), s(810548), s(11492483), s(16000))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(178, 198, 220)
    line.line.fill.background()

    slide.shapes.add_picture(BytesIO(logo_blob), s(11188065), s(114300), width=s(640080), height=s(640080))

    footer_y = 6680000 if slide_idx == 3 else 6550000
    footer_h = 160000 if slide_idx == 3 else 300000
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, s(footer_y), width, s(footer_h))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(71, 105, 158)
    footer.line.fill.background()
    put_after_background(slide, footer)


source = Presentation(str(SOURCE))
template = Presentation(str(TEMPLATE))
scale = source.slide_width / template.slide_width

logo_blob = None
for shape in template.slides[3].shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        logo_blob = shape.image.blob
        break
if logo_blob is None:
    raise RuntimeError("template logo not found")

for idx, slide in enumerate(source.slides):
    title = TITLES[idx] if idx < len(TITLES) else ""
    for shape in list(slide.shapes):
        tune_text(shape, idx)
        if (
            is_bottom_mark_picture(shape, source.slide_width, source.slide_height)
            or is_old_logo(shape, source.slide_width, source.slide_height)
            or (title and is_old_title(shape, title))
        ):
            remove(shape)
    add_frame(slide, logo_blob, title, scale, source.slide_width, idx)

source.save(str(OUTPUT))
print(OUTPUT)
