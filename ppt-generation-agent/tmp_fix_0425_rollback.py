from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu


BASE = Path(r"E:\PY\research\0425")
INPUT = BASE / "生产流程.pptx"
OUTPUT = BASE / "生产流程_回退干净版.pptx"

REPLACEMENTS = {
    "输 迭线": "输送线",
    "迭线": "送线",
    "锡护": "锅炉",
    "异崇检到": "异常检测",
    "异崇": "异常",
    "预湖": "预测",
    "遇终": "最终",
    "凤险": "风险",
    "橙惨": "检修",
    "摄矢": "损失",
    "赐标": "贴标",
    "波位": "液位",
    "思浮": "悬浮",
    "主霎": "主要",
    "分劃": "分割",
    "识刖": "识别",
    "复较": "复核",
    "缺船": "缺陷",
    "追潮": "追溯",
    "任努": "任务",
    "上瓶": "上甑",
    "昇常": "异常",
    "传燃": "传感",
    "温温度": "温度",
    "酒酷": "酒醅",
    "异案": "异常",
    "等缀": "等级",
    "掠入": "输入",
    "专案舰则": "专家规则",
    "检素": "检索",
    "规钏": "规则",
    "牧据": "数据",
    "措荐抉行": "推荐执行",
    "反懊": "反馈",
    "横型童训": "模型重训",
    "辖出": "输出",
    "璃参": "微调",
    "知玖沉渡": "知识沉淀",
    "产续": "产线",
    "发观": "发现",
    "该杀率": "误杀率",
    "罄线": "基线",
    "横型": "模型",
    "尖时": "实时",
    "汇粲": "汇聚",
    "逮度": "速度",
    "裁重": "载重",
    "起焚": "起终",
    "紫行": "禁行",
    "A\"": "A*",
    "落径": "路径",
    "规刺": "规划",
    "重规刘": "重规划",
    "逼过": "通过",
    "交忖": "交付",
    "温装": "混装",
    "睾求": "要求",
    "关肤": "关联",
    "保分": "评分",
    "密令": "订单",
    "启寂丸": "启发式",
    "液灰计刻": "波次计划",
    "捺选": "拣选",
    "支排": "安排",
    "辅入": "输入",
    "装魉舰物": "装载规划",
    "约柬伏化": "约束优化",
    "籍出": "输出",
    "多合节赤": "多仓节点",
    "韩出": "输出",
    "肠逆": "配送",
    "庠区": "库区",
    "又车": "叉车",
    "拣 遮": "拣选",
    "疲儋": "疲倦",
    "淀粉辜": "淀粉率",
    "搠眼": "摊晾",
    "发   醇": "发酵",
    "出客": "出窖",
    "图期": "周期",
}


def fix_text(text):
    text = text.replace("ANotebookLM", "").replace("NotebookLM", "")
    for old, new in REPLACEMENTS.items():
        text = text.replace(old, new)
    return text


def tune_shape(shape):
    if hasattr(shape, "text_frame") and shape.text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = fix_text(run.text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                cell.text = fix_text(cell.text)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            tune_shape(child)


def put_after_background(slide, shape):
    sp_tree = slide.shapes._spTree
    el = shape._element
    sp_tree.remove(el)
    sp_tree.insert(3, el)


def has_footer(slide, slide_height):
    for shape in slide.shapes:
        if shape.top > slide_height * 0.94 and shape.height < 350000:
            fill = getattr(shape, "fill", None)
            if fill and fill.type:
                return True
    return False


def add_footer(slide, slide_width, slide_height):
    if has_footer(slide, slide_height):
        return
    footer_h = Emu(213333)
    footer_y = Emu(int(slide_height) - int(footer_h))
    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, footer_y, slide_width, footer_h)
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(71, 105, 158)
    footer.line.fill.background()
    put_after_background(slide, footer)


prs = Presentation(str(INPUT))
for slide in prs.slides:
    for shape in slide.shapes:
        tune_shape(shape)
    add_footer(slide, prs.slide_width, prs.slide_height)

prs.save(str(OUTPUT))
print(OUTPUT)
