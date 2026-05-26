from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


TEMPLATE = Path(r"D:\Users\Downloads\Gemini_Generated_Image_j53320j53320j533.png")
OUT_PPTX = Path(r"D:\Users\Downloads\一厂一网_智慧水务系统全景图_文字版_editable.pptx")
OUT_PREVIEW = Path(r"D:\Users\Downloads\一厂一网_智慧水务系统全景图_文字版_editable_preview.png")

W, H = 2626, 1600
SLIDE_W = Emu(13.13 * 914400)
SLIDE_H = Emu(8.0 * 914400)
SX = SLIDE_W / W
SY = SLIDE_H / H

FONT_FACE = "Microsoft YaHei"
FONT_FILE = r"C:\Windows\Fonts\msyh.ttc"
FONT_BOLD_FILE = r"C:\Windows\Fonts\msyhbd.ttc"

WHITE = "F4FAFF"
MUTED = "C9D8E8"
CYAN = "8DDCFF"
GREEN = "A9F28A"
YELLOW = "FFF06A"
PEACH = "F7C99B"


def emu_x(v):
    return Emu(int(v * SX))


def emu_y(v):
    return Emu(int(v * SY))


def rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[:2], 16), int(h[2:4], 16), int(h[4:], 16))


def ppt_text(slide, x, y, w, h, text, size, color=WHITE, bold=False, align="left", valign="mid"):
    shape = slide.shapes.add_textbox(emu_x(x), emu_y(y), emu_x(w), emu_y(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE if valign == "mid" else MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = FONT_FACE
        # The design coordinates and preview font sizes are in source-image pixels.
        # PowerPoint font sizes are points; 1600 px maps to an 8 in / 576 pt slide.
        p.font.size = Pt(size * 0.36)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]


def preview_font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD_FILE if bold and Path(FONT_BOLD_FILE).exists() else FONT_FILE, size)


def preview_text(draw, x, y, w, h, text, size, color=WHITE, bold=False, align="left", valign="mid"):
    f = preview_font(size, bold)
    lines = text.split("\n")
    line_h = max(12, int(size * 1.25))
    total_h = line_h * len(lines)
    yy = y + (h - total_h) / 2 if valign == "mid" else y
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=f)
        tw = bbox[2] - bbox[0]
        if align == "center":
            xx = x + (w - tw) / 2
        elif align == "right":
            xx = x + w - tw
        else:
            xx = x
        draw.text((xx, yy), line, font=f, fill="#" + color, stroke_width=2, stroke_fill="#03101C")
        yy += line_h


TEXTS = []


def add(x, y, w, h, text, size, color=WHITE, bold=False, align="left", valign="mid"):
    TEXTS.append((x, y, w, h, text, size, color, bold, align, valign))


# Top header, matched to the second image.
add(720, 24, 1185, 72, "一厂一网 · 智慧水务系统全景图", 56, WHITE, True, "center")
add(850, 103, 930, 40, "全面感知 · 智能决策 · 精准控制 · 低碳运行 · 韧性保障", 27, WHITE, False, "center")

# Left title block.
add(345, 160, 500, 48, "污水管网（感知神经系统）", 34, WHITE, True, "center")
add(402, 215, 365, 30, "全域感知 · 智能诊断 · 精准养护", 21, WHITE, False, "center")

# Left rail: align every text row to the template bullet dots.
add(55, 252, 100, 36, "感知层", 25, WHITE, True)
for y, text in zip(
    [307, 349, 391, 433, 475, 517, 559],
    ["液位监测", "流量监测", "水质监测", "气体监测", "视频监控", "井盖监测", "管道结构监测"],
):
    add(80, y, 150, 28, text, 20, WHITE)

add(55, 655, 100, 36, "应用层", 25, WHITE, True)
for y, text in zip(
    [710, 752, 794, 836, 878],
    ["管网健康评估", "溢流污染监测", "内涝风险预警", "管道病害诊断", "养护计划优化"],
):
    add(80, y, 160, 28, text, 19, WHITE)

add(55, 1040, 100, 36, "运维层", 25, WHITE, True)
for y, text in zip([1097, 1139, 1181, 1223], ["巡检调度", "工事管理", "养护作业", "绩效考核"]):
    add(80, y, 140, 28, text, 20, WHITE)

# Left diagram labels.
add(405, 275, 95, 28, "居民区", 18, PEACH, True, "center")
for x, y, w, text in [
    (290, 404, 100, "流量监测"),
    (272, 499, 100, "液位监测"),
    (640, 456, 105, "水质监测"),
    (705, 544, 105, "气体监测"),
    (648, 674, 105, "井盖监测"),
    (525, 842, 165, "污水提升泵站"),
    (406, 1008, 110, "干线管网"),
    (654, 1010, 150, "管道结构监测"),
    (535, 1164, 95, "排放口"),
]:
    add(x, y, w, 30, text, 18 if w > 120 else 19, WHITE, False, "center")
add(302, 803, 95, 32, "商业区", 22, WHITE, True, "center")
add(742, 824, 125, 32, "工业园区", 22, WHITE, True, "center")
add(845, 1198, 100, 28, "自然水体", 15, MUTED, False, "center")

# Indicators.
add(875, 309, 105, 32, "关键指标", 20, GREEN, True, "center")
for y, text in zip(
    [365, 398, 431, 464, 497, 530],
    ["流量：12,580 m³/h", "液位：2.45 m", "COD：28 mg/L", "氨氮：3.2 mg/L", "温度：18.6℃", "井盖状态：正常"],
):
    add(840, y, 165, 26, text, 15, WHITE)
add(888, 880, 105, 32, "风险预警", 22, YELLOW, True, "center")
for y, text in zip([940, 977, 1014, 1051], ["溢流风险：低", "内涝风险：中", "管道缺陷：2处", "设备告警：1处"]):
    add(865, y, 135, 26, text, 16, WHITE)

# Center cloud and AI.
add(1160, 300, 305, 70, "智慧水务云平台\n（AI中枢大脑）", 27, WHITE, True, "center")
for y, text in zip([438, 507, 576], ["数据汇聚", "模型训练", "智能诊断"]):
    add(1088, y, 105, 28, text, 17, WHITE)
for y, text in zip([438, 507, 576], ["决策优化", "数字孪生", "可视化呈现"]):
    add(1400, y, 110, 28, text, 17, WHITE)
add(1225, 742, 190, 34, "数据中台", 24, CYAN, True, "center")
for x, text in zip([1110, 1215, 1320, 1425], ["数据接入", "数据治理", "数据存储", "数据服务"]):
    add(x, 866, 90, 24, text, 15, WHITE, False, "center")
add(1160, 990, 295, 34, "统一标准与安全体系", 22, WHITE, True, "center")
for x, text in zip([1110, 1275, 1430], ["物联感知标准化", "数据安全", "权限管理"]):
    add(x, 1082, 120, 24, text, 15, WHITE, False, "center")

# Right title and process sections.
add(1765, 160, 545, 48, "污水处理厂（智慧中枢系统）", 34, WHITE, True, "center")
add(1840, 215, 390, 30, "精准控制 · 提质增效 · 低碳运行", 21, WHITE, False, "center")
add(1815, 300, 145, 28, "进水监测", 18, CYAN, True, "center")
for x, text in zip([1612, 1732, 1850, 1965, 2083, 2200], ["流量", "水质", "水温", "pH", "氨氮", "重金属"]):
    add(x, 374, 70, 24, text, 16, WHITE, False, "center")
add(1810, 406, 145, 28, "预处理系统", 19, CYAN, True, "center")
for x, text in zip([1685, 1845, 2050], ["粗格栅", "细格栅", "曝气沉砂池"]):
    add(x, 486, 110, 24, text, 16, WHITE, False, "center")
add(1810, 560, 145, 28, "生化处理系统", 19, CYAN, True, "center")
for x, text in zip([1715, 1878, 2040], ["厌氧池", "缺氧池", "好氧池"]):
    add(x, 660, 95, 24, text, 16, WHITE, False, "center")
add(1795, 716, 170, 28, "深度处理系统", 19, CYAN, True, "center")
for x, text in zip([1685, 1845, 2015, 2160], ["沉淀池", "高效沉淀池", "滤布滤池", "消毒池"]):
    add(x, 800, 110, 24, text, 16, WHITE, False, "center")
add(1810, 895, 170, 28, "污泥处理系统", 19, CYAN, True, "center")
for x, text in zip([1685, 1847, 2020, 2165], ["浓缩池", "厌氧消化", "脱水机房", "干化/焚烧"]):
    add(x, 980, 110, 24, text, 16, WHITE, False, "center")
add(1810, 1095, 145, 28, "出水监测", 19, CYAN, True, "center")
for x, text in zip([1597, 1710, 1824, 1935, 2050, 2165], ["COD", "氨氮", "总磷", "总氮", "浊度", "余氯"]):
    add(x, 1168, 70, 24, text, 16, WHITE, False, "center")
add(2380, 1180, 225, 30, "达标排放 / 回用利用", 17, WHITE, False, "center")

# Right side rail: each line aligned to its bullet dot.
add(2415, 272, 130, 30, "AI优化控制", 20, CYAN, True, "center")
for y, text in zip([333, 375, 417, 459], ["曝气智能调控", "加药智能控制", "回流比优化", "能耗优化"]):
    add(2425, y, 140, 24, text, 17, WHITE)
add(2415, 560, 130, 30, "能效管理", 21, CYAN, True, "center")
for y, text in zip([625, 667, 709, 751], ["碳排监测", "能效分析", "碳绩效管理", "设备优化"]):
    add(2425, y, 140, 24, text, 17, WHITE)
add(2415, 850, 130, 30, "生产运营", 21, CYAN, True, "center")
for y, text in zip([914, 956, 998, 1040], ["运行监控", "异常预警", "报表管理", "绩效分析"]):
    add(2425, y, 140, 24, text, 17, WHITE)

# Bottom rails.
add(1120, 1258, 390, 40, "一厂一网协同联动", 30, WHITE, True, "center")
for x, text in zip([700, 1030, 1350, 1700], ["水量水质协同调度", "风险预警联动", "应急响应联动", "调度策略优化"]):
    add(x, 1335, 230, 32, text, 23, WHITE, False, "center")
add(1220, 1425, 190, 38, "价值体系", 28, WHITE, True, "center")
for x, title, subtitle in [
    (390, "环境效益", "减少污染排放"),
    (725, "社会效益", "提升履约保障"),
    (1060, "经济效益", "降低运营成本"),
    (1395, "管理效益", "提升决策效率"),
    (1735, "安全效益", "保障系统安全"),
    (2085, "可持续发展", "推动绿色低碳"),
]:
    add(x, 1482, 160, 27, title, 21, WHITE, True)
    add(x, 1515, 170, 24, subtitle, 15, WHITE)


def build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(TEMPLATE), 0, 0, width=SLIDE_W, height=SLIDE_H)
    for item in TEXTS:
        ppt_text(slide, *item)
    prs.save(OUT_PPTX)


def build_preview():
    img = Image.open(TEMPLATE).convert("RGB")
    draw = ImageDraw.Draw(img)
    for item in TEXTS:
        preview_text(draw, *item)
    img.save(OUT_PREVIEW, quality=95)


if __name__ == "__main__":
    build_preview()
    build_pptx()
    print(OUT_PPTX)
    print(OUT_PREVIEW)
