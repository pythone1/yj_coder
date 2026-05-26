# -*- coding: utf-8 -*-
from pathlib import Path
from io import BytesIO

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "机理、算法与实践：人工智能行业应用实证分析.pptx"
ASSET_DIR = ROOT / "output" / "mixed_editable_assets"
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_对外讲解完整稿_v3.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_对外讲解完整稿_v3_逐页检查.txt"
PREVIEW = OUT / "AI供水管网DMA漏损检测_对外讲解完整稿_v3_预览联系表.png"

NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(247, 252, 255)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
FONT = "微软雅黑"

ASSETS = {
    "cover": ASSET_DIR / "cover_clean.png",
    "dma": ASSET_DIR / "dma_clean.png",
    "data": ASSET_DIR / "data_clean.png",
    "ai": ASSET_DIR / "ai_clean.png",
    "locate": ASSET_DIR / "locate_clean.png",
    "loop": ASSET_DIR / "loop_clean.png",
    "algorithm": ASSET_DIR / "algorithm_clean.png",
    "mnf": ASSET_DIR / "mnf_clean.png",
    "pressure": ASSET_DIR / "pressure_clean.png",
    "governance": ASSET_DIR / "governance_clean.png",
    "ops": ASSET_DIR / "ops_clean.png",
}

AUDIT_LINES = []


def emu(v):
    return Inches(v)


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, trans=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = trans
    shape.line.fill.background()


def stroke(shape, color=LINE, width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(typ, x, y, w, h)
    fill(shp, color, trans)
    stroke(shp, line, 1)
    return shp


def tb(slide, text, x, y, w, h, size=16, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = emu(0.03)
    tf.margin_right = emu(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pic(slide, path, x, y, w, h, crop=False):
    p = slide.shapes.add_picture(str(path), x, y)
    sx, sy = w / p.width, h / p.height
    scale = max(sx, sy) if crop else min(sx, sy)
    p.width = int(p.width * scale)
    p.height = int(p.height * scale)
    p.left = int(x + (w - p.width) / 2)
    p.top = int(y + (h - p.height) / 2)
    return p


def connector(slide, x1, y1, x2, y2, color=BLUE, width=1.8, arrow=True):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.38), emu(0.23), emu(0.12), emu(0.50))
    fill(bar, BLUE)
    tb(slide, title, emu(0.62), emu(0.12), emu(9.45), emu(0.45), 24, TEXT, True)
    tb(slide, sub, emu(0.64), emu(0.62), emu(8.7), emu(0.25), 11.3, MUTED)
    tb(slide, f"{no:02d}", emu(12.05), emu(0.20), emu(0.7), emu(0.24), 12, BLUE, True, PP_ALIGN.RIGHT)
    fill(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.62), emu(0.90), emu(11.7), emu(0.016)), PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA漏损检测 · 对外讲解完整稿", emu(0.62), emu(7.08), emu(6.0), emu(0.18), 9.8, MUTED)


def content(prs, title, sub, no):
    s = blank(prs)
    header(s, title, sub, no)
    footer(s)
    return s


def pill(slide, text, x, y, w, color=PALE, text_color=BLUE, size=11):
    rect(slide, x, y, w, emu(0.34), color, color, True)
    tb(slide, text, x + emu(0.06), y + emu(0.085), w - emu(0.12), emu(0.16), size, text_color, True, PP_ALIGN.CENTER)


def note(slide, text, y=6.05, color=BLUE):
    rect(slide, emu(0.72), emu(y), emu(11.85), emu(0.56), PALE2, LINE, True)
    tb(slide, text, emu(0.98), emu(y + 0.15), emu(11.30), emu(0.25), 14.5, color, True, PP_ALIGN.CENTER)


def label(slide, title, body, x, y, w, h, color=BLUE, fill_color=PALE2):
    rect(slide, x, y, w, h, fill_color, LINE, True)
    tb(slide, title, x + emu(0.14), y + emu(0.11), w - emu(0.28), emu(0.23), 14, color, True)
    tb(slide, body, x + emu(0.14), y + emu(0.46), w - emu(0.28), h - emu(0.56), 11.6, TEXT)


def overlay_label(slide, title, body, x, y, w, h, color=BLUE):
    rect(slide, x, y, w, h, WHITE, RGBColor(210, 232, 246), True)
    tb(slide, title, x + emu(0.12), y + emu(0.10), w - emu(0.24), emu(0.20), 13.0, color, True)
    tb(slide, body, x + emu(0.12), y + emu(0.34), w - emu(0.24), h - emu(0.38), 10.8, TEXT)


def module(slide, title, img_key, desc, x, y, w, h, color):
    rect(slide, x, y, w, h, WHITE, RGBColor(201, 228, 242), True)
    fill(slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, emu(0.42)), color)
    tb(slide, title, x, y + emu(0.10), w, emu(0.18), 13.8, WHITE, True, PP_ALIGN.CENTER)
    pic(slide, ASSETS[img_key], x + emu(0.12), y + emu(0.55), w - emu(0.24), h - emu(1.25), crop=False)
    tb(slide, desc, x + emu(0.14), y + h - emu(0.46), w - emu(0.28), emu(0.32), 11.2, TEXT, False, PP_ALIGN.CENTER)


def audit_slide(slide, no, allow_short=False):
    prs_w, prs_h = 12192000, 6858000
    chars, pics, off, bottom, min_font = 0, 0, 0, 0, 99
    for sh in slide.shapes:
        bottom = max(bottom, sh.top + sh.height)
        if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs_w + 50000 or sh.top + sh.height > prs_h + 50000:
            off += 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
        if hasattr(sh, "text_frame") and sh.text.strip():
            chars += len(sh.text.strip())
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        min_font = min(min_font, r.font.size.pt)
    flags = []
    if off:
        flags.append(f"越界{off}")
    if min_font < 9.8:
        flags.append(f"小字{min_font}")
    if pics < 1:
        flags.append("缺少主视觉")
    if not allow_short and chars < 135:
        flags.append(f"文字不足{chars}")
    if bottom < emu(6.15):
        flags.append("下半页利用不足")
    status = "FAIL" if flags else "OK"
    AUDIT_LINES.append(f"{no:02d}\t{status}\tchars={chars}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\tbottom={round(bottom/914400,2)}\t{'；'.join(flags)}")
    if flags:
        raise RuntimeError(f"slide {no} failed: {'; '.join(flags)}")


def cover(prs):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(5.85), emu(7.5)), NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0.62), emu(0.95), emu(0.13), emu(5.78)), CYAN)
    tb(s, "AI模型在供水管网\nDMA系统漏损检测\n中的应用", emu(1.02), emu(1.18), emu(4.4), emu(1.75), 31, WHITE, True)
    tb(s, "对外教学演讲版", emu(1.05), emu(3.18), emu(2.7), emu(0.28), 17, CYAN, True)
    tb(s, "从DMA分区异常识别，讲到AI溯源定位、候选管段收敛与工单闭环处置。", emu(1.05), emu(5.35), emu(4.55), emu(0.44), 12.8, RGBColor(220, 238, 248))
    tb(s, "讲解重点：模型机理 · 算法选型 · 场景落地 · 持续运营", emu(1.05), emu(5.95), emu(4.55), emu(0.28), 12.0, RGBColor(220, 238, 248))
    rect(s, emu(6.16), emu(0.62), emu(6.48), emu(5.70), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS["cover"], emu(6.34), emu(0.85), emu(6.12), emu(5.25), crop=False)
    overlay_label(s, "DMA宏观锁定", "先把全网问题压缩到可核查片区", emu(6.70), emu(6.34), emu(1.72), emu(0.64), BLUE)
    overlay_label(s, "AI微观溯源", "再用多源特征收敛候选管段", emu(8.62), emu(6.34), emu(1.72), emu(0.64), GREEN)
    overlay_label(s, "工单闭环", "把模型结果变成现场行动", emu(10.54), emu(6.34), emu(1.72), emu(0.64), ORANGE)
    audit_slide(s, 1)


def toc(prs):
    s = content(prs, "课程结构：从业务问题到模型落地", "按照“为什么、怎么做、如何验证、怎样运营”的顺序展开", 2)
    chapters = [
        ("01", "DMA漏损检测的业务底座", "分区计量、MNF、传统检漏边界", "dma", BLUE),
        ("02", "AI如何介入漏损识别", "动态基线、多源融合、溯源定位", "ai", GREEN),
        ("03", "核心算法与模型组合", "LSTM、无监督、监督排序、水力约束", "algorithm", ORANGE),
        ("04", "场景落地与系统运营", "预警、定位、派单、复盘、再训练", "loop", PURPLE),
    ]
    for i, (num, title, body, img, color) in enumerate(chapters):
        x = emu(0.68 + i * 3.05)
        y = emu(1.28)
        module(s, f"{num} {title}", img, body, x, y, emu(2.62), emu(4.35), color)
        if i < 3:
            connector(s, x + emu(2.68), y + emu(2.2), x + emu(2.92), y + emu(2.2), BLUE, 1.8)
    note(s, "对外讲解不从模型名开始，而是从管网运维人员真正要解决的问题开始。", 6.03)
    audit_slide(s, 2)


def section(prs, no, part, title, subtitle, img):
    s = blank(prs, NAVY)
    fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(0), emu(0), emu(0.24), emu(7.5)), CYAN)
    tb(s, part, emu(0.88), emu(1.16), emu(1.5), emu(0.30), 16, CYAN, True)
    tb(s, title, emu(0.88), emu(1.86), emu(6.65), emu(0.72), 34, WHITE, True)
    tb(s, subtitle, emu(0.92), emu(3.00), emu(6.55), emu(0.74), 16.5, RGBColor(210, 232, 245))
    rect(s, emu(8.05), emu(1.20), emu(4.25), emu(4.25), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS[img], emu(8.28), emu(1.48), emu(3.78), emu(3.68), crop=False)
    tb(s, f"{no:02d}", emu(11.75), emu(6.60), emu(0.65), emu(0.25), 14, RGBColor(166, 208, 235), True, PP_ALIGN.RIGHT)
    audit_slide(s, no, allow_short=True)


def stage_slide(prs, no, title, sub, img, labels, bottom):
    s = content(prs, title, sub, no)
    rect(s, emu(0.72), emu(1.15), emu(11.85), emu(4.82), WHITE, RGBColor(205, 229, 243), True)
    pic(s, ASSETS[img], emu(0.98), emu(1.36), emu(11.35), emu(4.35), crop=False)
    for item in labels:
        overlay_label(s, item["t"], item["b"], emu(item["x"]), emu(item["y"]), emu(item["w"]), emu(item["h"]), item.get("c", BLUE))
    note(s, bottom, 6.15)
    audit_slide(s, no)


def compare_slide(prs, no, title, sub, left, right, img_left, img_right, bottom):
    s = content(prs, title, sub, no)
    pic(s, ASSETS[img_left], emu(0.78), emu(1.18), emu(3.2), emu(2.35), crop=False)
    pic(s, ASSETS[img_right], emu(9.22), emu(1.18), emu(3.2), emu(2.35), crop=False)
    label(s, left[0], left[1], emu(0.78), emu(3.70), emu(4.20), emu(1.58), BLUE)
    label(s, right[0], right[1], emu(8.20), emu(3.70), emu(4.20), emu(1.58), GREEN)
    rect(s, emu(5.30), emu(1.38), emu(2.70), emu(3.60), PALE2, LINE, True)
    tb(s, "关键变化", emu(5.55), emu(1.62), emu(2.20), emu(0.24), 17, ORANGE, True, PP_ALIGN.CENTER)
    for i, txt in enumerate(["从经验巡检", "到数据预警", "从片区判断", "到管段排序"]):
        pill(s, txt, emu(5.70), emu(2.12 + i * 0.55), emu(1.90), PALE, ORANGE if i % 2 else BLUE, 11)
    connector(s, emu(4.20), emu(2.80), emu(5.25), emu(2.80), BLUE, 1.8)
    connector(s, emu(8.05), emu(2.80), emu(9.02), emu(2.80), BLUE, 1.8)
    note(s, bottom, 6.06)
    audit_slide(s, no)


def flow_slide(prs, no, title, sub, steps, img, bottom):
    s = content(prs, title, sub, no)
    pic(s, ASSETS[img], emu(0.82), emu(1.12), emu(4.0), emu(4.30), crop=False)
    x0 = 5.05
    for i, (head, body, color) in enumerate(steps):
        y = 1.17 + i * 1.08
        rect(s, emu(x0), emu(y), emu(6.95), emu(0.78), WHITE, LINE, True)
        fill(s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(x0 + 0.17), emu(y + 0.18), emu(0.36), emu(0.36)), color)
        tb(s, str(i + 1), emu(x0 + 0.17), emu(y + 0.255), emu(0.36), emu(0.12), 10.5, WHITE, True, PP_ALIGN.CENTER)
        tb(s, head, emu(x0 + 0.68), emu(y + 0.13), emu(1.65), emu(0.22), 13.3, color, True)
        tb(s, body, emu(x0 + 2.30), emu(y + 0.14), emu(3.55), emu(0.38), 11.3, TEXT)
        if i < len(steps) - 1:
            connector(s, emu(x0 + 0.35), emu(y + 0.79), emu(x0 + 0.35), emu(y + 1.03), BLUE, 1.4)
    note(s, bottom, 6.05)
    audit_slide(s, no)


def algorithm_slide(prs, no, title, sub, img, items, bottom):
    s = content(prs, title, sub, no)
    pic(s, ASSETS[img], emu(0.82), emu(1.18), emu(4.55), emu(4.65), crop=False)
    for i, (name, task, output, color) in enumerate(items):
        x = 5.70 + (i % 2) * 3.35
        y = 1.18 + (i // 2) * 1.42
        label(s, name, f"{task}\n输出：{output}", emu(x), emu(y), emu(3.05), emu(1.02), color)
    note(s, bottom, 6.06)
    audit_slide(s, no)


def matrix_slide(prs, no, title, sub, img, rows, bottom):
    s = content(prs, title, sub, no)
    pic(s, ASSETS[img], emu(0.75), emu(1.18), emu(3.8), emu(4.75), crop=False)
    headers = ["对象", "模型输入", "判断方式", "业务输出"]
    widths = [1.42, 2.05, 2.18, 2.15]
    x = emu(4.75)
    y = emu(1.20)
    for i, head in enumerate(headers):
        xx = x + sum(emu(w) for w in widths[:i])
        rect(s, xx, y, emu(widths[i]), emu(0.42), BLUE, BLUE, False)
        tb(s, head, xx, y + emu(0.11), emu(widths[i]), emu(0.15), 11.5, WHITE, True, PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        yy = y + emu(0.54 + r * 0.74)
        for c, val in enumerate(row):
            xx = x + sum(emu(w) for w in widths[:c])
            rect(s, xx, yy, emu(widths[c]), emu(0.62), PALE2 if r % 2 == 0 else WHITE, LINE, False)
            tb(s, val, xx + emu(0.06), yy + emu(0.12), emu(widths[c] - 0.12), emu(0.32), 10.6, BLUE if c == 0 else TEXT, c == 0)
    note(s, bottom, 6.08)
    audit_slide(s, no)


def build():
    missing = [p.name for p in ASSETS.values() if not p.exists()]
    if missing:
        raise FileNotFoundError("缺少素材：" + "、".join(missing))
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    cover(prs)
    toc(prs)
    section(prs, 3, "PART 01", "DMA漏损检测的业务底座", "先讲清楚DMA能解决什么，以及为什么仍然需要AI。", "dma")
    stage_slide(prs, 4, "DMA是什么：把复杂管网转成可计量单元", "关闭边界阀门、布设进出水计量点，形成可核算的分区", "dma", [
        {"t": "边界清晰", "b": "边界阀门定义片区范围，避免水量混入。", "x": 0.98, "y": 1.42, "w": 2.25, "h": 0.78, "c": BLUE},
        {"t": "进出水计量", "b": "用流量计持续记录进入与离开DMA的水量。", "x": 9.78, "y": 1.50, "w": 2.18, "h": 0.78, "c": GREEN},
        {"t": "异常入口", "b": "夜间低用水时段更容易暴露持续漏损。", "x": 5.25, "y": 4.84, "w": 2.38, "h": 0.70, "c": ORANGE},
    ], "这一页用于建立共识：DMA先把全网问题缩小到一个可管理、可核算、可复盘的片区。")
    compare_slide(prs, 5, "传统检漏的局限：经验强、成本高、响应慢", "人工巡检、听漏、分段排查仍然重要，但很难支撑全域连续监测", ("传统方式", "依赖巡检经验和现场听漏，适合确认问题，但难以及时发现隐蔽漏损。"), ("AI辅助方式", "先用模型筛出高风险片区和管段，再安排人员带着坐标和证据核查。"), "locate", "ai", "AI不是替代现场人员，而是把“先查哪里”变成有依据的排序问题。")
    stage_slide(prs, 6, "DMA能回答“哪个区异常”，但不能直接回答“哪根管漏”", "这正是AI溯源分析进入业务链条的核心位置", "locate", [
        {"t": "DMA输出", "b": "片区流量异常、MNF抬升、压力趋势变化。", "x": 1.10, "y": 4.82, "w": 2.60, "h": 0.78, "c": BLUE},
        {"t": "定位缺口", "b": "片区内部可能有多条候选管段，需要进一步收敛。", "x": 5.05, "y": 1.42, "w": 2.75, "h": 0.82, "c": ORANGE},
        {"t": "AI任务", "b": "融合拓扑、压力响应和历史记录，输出候选管段排序。", "x": 9.28, "y": 4.78, "w": 2.70, "h": 0.82, "c": GREEN},
    ], "讲解重点：DMA完成宏观锁定，AI负责微观溯源，两者不是替代关系而是前后衔接。")
    stage_slide(prs, 7, "夜间最小流量MNF：漏损预警的关键观测量", "低用水时段的持续底流更能反映背景漏损变化", "mnf", [
        {"t": "正常基线", "b": "每个DMA都有自己的季节、工作日和节假日用水规律。", "x": 1.08, "y": 1.42, "w": 2.70, "h": 0.78, "c": BLUE},
        {"t": "持续偏离", "b": "AI重点识别连续多日偏离，而不是单点波动。", "x": 8.95, "y": 1.42, "w": 2.70, "h": 0.78, "c": RED},
        {"t": "复核条件", "b": "需同步排查边界阀、计量误差、临时用水和施工扰动。", "x": 5.10, "y": 4.78, "w": 3.25, "h": 0.82, "c": ORANGE},
    ], "MNF不是漏损量本身，而是触发进一步诊断的信号；AI提升的是稳定识别和证据整合能力。")
    flow_slide(prs, 8, "从片区异常到漏点定位：问题如何一步步收敛", "把一个模糊报警转成可执行工单，需要连续四步", [
        ("片区识别", "发现DMA进出水量、MNF或压力曲线出现持续异常。", BLUE),
        ("原因排查", "排除边界阀、计量设备、计划用水和施工扰动。", ORANGE),
        ("候选收敛", "结合拓扑和水力响应，形成候选管段TopN。", GREEN),
        ("现场闭环", "派单核查、维修记录回写，形成再训练样本。", PURPLE),
    ], "loop", "这一页是后续所有模型章节的业务主线：先预警，再解释，再定位，最后闭环。")

    section(prs, 9, "PART 02", "AI如何介入DMA漏损检测", "从数据融合、动态基线到AI溯源，建立一套可讲清的技术框架。", "ai")
    stage_slide(prs, 10, "总体架构：DMA宏观锁定 + AI微观溯源", "数据、模型和工单不是三套系统，而是一条业务链", "ai", [
        {"t": "多源输入", "b": "流量、压力、拓扑、管龄、工单、施工和天气。", "x": 1.02, "y": 1.42, "w": 2.58, "h": 0.80, "c": BLUE},
        {"t": "模型分析", "b": "动态基线识别异常，水力约束收敛候选管段。", "x": 5.25, "y": 1.40, "w": 2.72, "h": 0.82, "c": GREEN},
        {"t": "业务输出", "b": "形成片区报警、管段排序、核查建议和派单依据。", "x": 9.30, "y": 4.74, "w": 2.68, "h": 0.86, "c": ORANGE},
    ], "对外讲解可用一句话概括：DMA负责发现异常区域，AI负责解释异常来源并指导核查顺序。")
    matrix_slide(prs, 11, "多源数据如何进入模型", "数据层决定模型上限，算法层决定模型效率", "data", [
        ["流量压力", "SCADA、远传水表、压力记录", "时序建模、残差识别", "MNF异常、压力响应"],
        ["管网拓扑", "管线、节点、阀门、泵站", "连通关系与影响范围", "候选管段收敛"],
        ["资产属性", "管材、管龄、口径、埋深", "风险特征编码", "管段优先级排序"],
        ["外部扰动", "施工、天气、节假日、工商业用水", "异常解释与误报过滤", "复核建议"],
        ["工单记录", "巡检、维修、漏点、处置结果", "监督标签与反馈学习", "模型再训练"],
    ], "讲数据时要落到业务字段，避免只说“多源融合”而不说明具体融合什么。")
    stage_slide(prs, 12, "动态基线：每个DMA都应有自己的正常曲线", "同一条阈值很难适配不同片区、季节和用户结构", "mnf", [
        {"t": "时间规律", "b": "日周期、周周期、节假日和季节性用水差异。", "x": 1.05, "y": 4.72, "w": 2.88, "h": 0.80, "c": BLUE},
        {"t": "模型预测", "b": "LSTM/GRU根据历史窗口预测下一时段正常范围。", "x": 5.02, "y": 1.34, "w": 3.05, "h": 0.82, "c": GREEN},
        {"t": "残差判断", "b": "实际值持续高于预测区间，才进入异常复核。", "x": 8.96, "y": 4.68, "w": 2.95, "h": 0.86, "c": RED},
    ], "动态基线的价值是减少固定阈值带来的误报，让每个DMA按自己的历史规律被判断。")
    flow_slide(prs, 13, "异常识别：从单点超限转向持续证据", "实际项目中，稳定的异常证据比一次报警更重要", [
        ("预测区间", "模型给出某时段流量或压力的合理范围。", BLUE),
        ("残差序列", "计算实际值与预测值之间的偏差并形成连续序列。", CYAN),
        ("异常评分", "结合偏差幅度、持续时间和空间一致性给出分数。", RED),
        ("复核建议", "自动提示是否需要检查计量、边界阀或现场漏损。", ORANGE),
    ], "algorithm", "异常识别要讲清楚“为什么报”和“先排除什么”，否则现场很难信任模型。")
    stage_slide(prs, 14, "压力响应：漏点会改变局部水力状态", "AI定位需要利用压力波动和流量重分布，而不只看总流量", "pressure", [
        {"t": "压力下降", "b": "漏点附近压力可能出现持续性或突发性下降。", "x": 1.05, "y": 1.42, "w": 2.72, "h": 0.80, "c": BLUE},
        {"t": "流量重分布", "b": "相邻管段流量方向和大小可能随漏点发生变化。", "x": 5.18, "y": 4.75, "w": 2.98, "h": 0.82, "c": GREEN},
        {"t": "拓扑约束", "b": "只有与异常水力响应相符的管段才保留为候选。", "x": 8.95, "y": 1.42, "w": 2.86, "h": 0.82, "c": ORANGE},
    ], "这一页用于解释为什么仅靠MNF不够：定位需要把水力响应和管网结构一起看。")
    stage_slide(prs, 15, "AI溯源定位：从整个DMA收敛到候选管段", "定位模型输出的是优先核查顺序，而不是替代现场确认", "locate", [
        {"t": "候选管段TopN", "b": "按异常解释度、历史风险和可达性排序。", "x": 1.05, "y": 4.74, "w": 2.70, "h": 0.82, "c": GREEN},
        {"t": "现场坐标", "b": "输出节点坐标、管段编号和周边阀门信息。", "x": 5.20, "y": 1.42, "w": 2.80, "h": 0.82, "c": BLUE},
        {"t": "核查路径", "b": "把检漏路线、人员和工单结果纳入闭环管理。", "x": 9.00, "y": 4.74, "w": 2.80, "h": 0.82, "c": ORANGE},
    ], "对外表达要准确：模型提高“先查哪里”的命中率，最终漏点仍由现场核查确认。")
    compare_slide(prs, 16, "AI应用价值：从报警系统升级为决策辅助系统", "价值不只在发现异常，更在缩短定位和处置时间", ("仅报警", "告诉运维人员某个DMA异常，但需要人工再判断原因、范围和核查路径。"), ("辅助决策", "同时输出异常解释、候选管段、核查优先级和处置闭环要求。"), "mnf", "loop", "讲价值时不要停在“检测准确率”，还要讲“缩短漏损持续时间”和“减少无效排查”。")

    section(prs, 17, "PART 03", "核心算法与模型组合", "用工程语言讲清算法分工，避免变成模型名称罗列。", "algorithm")
    algorithm_slide(prs, 18, "算法全景：不同模型承担不同环节", "把算法放回业务链路中讲，听众更容易理解", "algorithm", [
        ("LSTM / GRU", "学习DMA流量、压力的时间规律", "正常曲线与残差", BLUE),
        ("无监督模型", "识别少标签场景下的异常模式", "异常时段和分数", RED),
        ("监督学习", "利用历史漏点和工单训练风险排序", "管段风险等级", ORANGE),
        ("水力模型", "模拟不同漏点对压力和流量的影响", "候选位置收敛", GREEN),
        ("GNN / 知识图谱", "表达管网拓扑和资产关系", "解释路径与影响范围", PURPLE),
        ("规则引擎", "纳入边界阀、计量、施工等业务规则", "复核清单", CYAN),
    ], "算法不是越复杂越好，关键是与数据条件和业务输出匹配。")
    flow_slide(prs, 19, "LSTM/GRU：预测正常曲线，再用残差发现异常", "适合连续流量、压力时序，且可在本地LSTM环境中演示", [
        ("输入窗口", "取过去24小时或7天的流量、压力、时间特征。", BLUE),
        ("序列学习", "LSTM/GRU学习日周期、周周期和趋势变化。", GREEN),
        ("预测区间", "输出下一时段的正常值或置信范围。", ORANGE),
        ("残差报警", "实际值持续偏离预测范围时触发复核。", RED),
    ], "mnf", "这一页可连接你的本地LSTM环境：用小样例演示“预测正常曲线—残差报警”的核心逻辑。")
    algorithm_slide(prs, 20, "CNN-LSTM：多压力点联合判断", "当DMA内有多个压力监测点时，可同时看空间和时间变化", "pressure", [
        ("空间局部特征", "CNN提取多监测点之间的同步变化", "压力响应图谱", BLUE),
        ("时间演化特征", "LSTM跟踪异常从出现到持续的过程", "异常持续性", GREEN),
        ("联合判断", "同时满足局部压力变化和全局流量异常", "更稳的定位证据", ORANGE),
        ("适用条件", "需要较密的压力点和稳定的数据采样", "试点DMA优先", PURPLE),
    ], "CNN-LSTM适合讲“多点联动”，但不要把它讲成所有项目的默认模型。")
    matrix_slide(prs, 21, "无监督异常检测：少标签阶段的起步方案", "当历史漏点标签不足时，先做异常发现和人工复核", "algorithm", [
        ["孤立森林", "流量、压力残差、MNF变化", "孤立异常点更容易被分离", "异常时段候选"],
        ["DBSCAN", "多维特征聚类", "识别偏离正常簇的样本", "异常片区聚类"],
        ["自编码器", "重构正常模式", "重构误差高说明异常", "异常分数"],
        ["统计控制图", "均值、方差、连续偏离", "规则透明、易解释", "报警阈值"],
        ["人工复核", "工单和现场结果", "把结果回填为标签", "进入监督训练"],
    ], "无监督模型适合早期建设，但必须配合复核闭环，否则只能发现“异常”，不能形成可靠定位能力。")
    algorithm_slide(prs, 22, "监督学习：把历史漏损经验转成管段风险排序", "当工单与漏点记录质量提升后，可训练管段级风险模型", "governance", [
        ("资产特征", "管材、管龄、口径、埋深、维修历史", "基础风险", BLUE),
        ("水力特征", "压力波动、流速、供水方向变化", "运行风险", GREEN),
        ("环境特征", "道路施工、土壤、温度、地面荷载", "外部风险", ORANGE),
        ("标签来源", "真实漏点、误报、维修闭环记录", "训练样本", PURPLE),
        ("模型输出", "管段风险分、解释因子、核查优先级", "派单依据", RED),
        ("常用模型", "随机森林、GBDT、HistGradientBoosting", "可解释排序", CYAN),
    ], "监督学习的关键不是换算法，而是把工单、资产和现场结果整理成可信标签。")
    stage_slide(prs, 23, "水力模型 + AI：用机理约束提升可信度", "AI发现异常，水力模型判断候选漏点是否符合管网响应", "pressure", [
        {"t": "仿真候选", "b": "在不同管段设置假想漏点，观察压力和流量响应。", "x": 1.05, "y": 1.40, "w": 2.82, "h": 0.86, "c": BLUE},
        {"t": "匹配观测", "b": "将仿真结果与实际监测曲线对比，剔除不符合的候选。", "x": 5.10, "y": 4.72, "w": 3.18, "h": 0.88, "c": GREEN},
        {"t": "增强解释", "b": "把模型结果讲成可验证的水力证据，而不是黑箱分数。", "x": 8.85, "y": 1.40, "w": 3.02, "h": 0.86, "c": ORANGE},
    ], "这一页适合对工程听众解释：AI不是脱离机理，而是用机理提高定位可信度。")
    algorithm_slide(prs, 24, "GNN与知识图谱：表达管网关系和解释路径", "适用于拓扑复杂、资产关系丰富、需要解释的场景", "ai", [
        ("节点表示", "监测点、阀门、泵站、水表、用户区", "图节点", BLUE),
        ("边表示", "管段连接、阀门状态、供水方向", "图边关系", GREEN),
        ("图传播", "异常影响沿拓扑关系扩散与衰减", "影响范围", ORANGE),
        ("知识约束", "管材、管龄、历史维修与规则库", "解释依据", PURPLE),
        ("输出形式", "候选管段、关联证据、原因链条", "讲解路径", RED),
        ("建设前提", "GIS拓扑准确、阀门状态及时维护", "数据底座", CYAN),
    ], "GNN不是必须一开始就上，适合在拓扑数据质量较好、需要解释复杂关系时引入。")

    section(prs, 25, "PART 04", "场景落地与系统运营", "把模型结果放进预警、定位、派单、复盘和再训练闭环。", "loop")
    flow_slide(prs, 26, "事前预警：从监测曲线进入处置准备", "预警阶段的目标是尽早发现并判断是否值得核查", [
        ("连续监测", "实时采集DMA流量、压力和设备状态。", BLUE),
        ("异常评分", "模型综合MNF、残差、压力响应形成评分。", RED),
        ("复核清单", "提示边界阀、计量设备、临时用水等排查项。", ORANGE),
        ("预警升级", "满足持续性和影响范围条件后进入定位阶段。", GREEN),
    ], "mnf", "事前预警要强调“减少误报”和“提前准备”，不是一报警就派人开挖。")
    stage_slide(prs, 27, "事中定位：候选管段、节点坐标与核查路线", "定位阶段的核心输出必须能够指导现场行动", "locate", [
        {"t": "候选排序", "b": "给出TopN管段和每段风险解释。", "x": 1.08, "y": 1.42, "w": 2.55, "h": 0.78, "c": BLUE},
        {"t": "核查路线", "b": "结合道路、阀门和人员位置安排路线。", "x": 5.08, "y": 4.76, "w": 2.85, "h": 0.80, "c": ORANGE},
        {"t": "证据包", "b": "附带异常曲线、压力响应和历史维修记录。", "x": 9.18, "y": 1.42, "w": 2.65, "h": 0.78, "c": GREEN},
    ], "定位页面要让听众看到：AI输出不是一句“可能漏”，而是一套可执行的现场证据。")
    flow_slide(prs, 28, "现场处置：模型结果如何进入工单", "现场闭环决定模型能否持续变准", [
        ("派单", "系统按候选管段生成核查任务和位置说明。", BLUE),
        ("核查", "人员使用听漏、相关仪、阀门操作等手段确认。", ORANGE),
        ("维修", "记录漏点位置、漏损原因、修复时间和照片。", GREEN),
        ("回写", "结果进入样本库，用于模型评估和再训练。", PURPLE),
    ], "loop", "没有工单回写，AI只能做一次性展示；有闭环，模型才会随城市管网经验持续成长。")
    matrix_slide(prs, 29, "事后复盘：把每次处置变成训练资产", "复盘不是写总结，而是补齐模型需要的标签和证据", "ops", [
        ["确认漏点", "管段编号、节点、位置坐标", "作为正样本", "提升定位模型"],
        ["误报原因", "计量异常、阀门状态、临时用水", "作为排除规则", "减少后续误报"],
        ["维修记录", "破损类型、管材、修复方式", "更新资产风险", "优化排序"],
        ["时间指标", "发现、派单、到场、修复时间", "评估业务收益", "缩短持续漏损"],
        ["现场反馈", "路线可达性、核查难度、遗漏点", "校正模型输出", "改进工单模板"],
    ], "复盘页要让管理层看到：每一次处置都在形成可复用的数据资产。")
    compare_slide(prs, 30, "压力优化：从漏损检测延伸到控漏运行", "定位只是第一步，长期降漏还需要压力管理和运行优化", ("被动检漏", "发现异常后再现场排查，漏损持续时间受发现速度影响。"), ("主动控漏", "利用压力分区、泵阀调度和风险预测，降低漏损发生和扩大概率。"), "pressure", "ops", "这一页用于扩展演讲宽度：AI漏损检测可以与压力优化、能耗管理共同形成智慧水务应用。")
    stage_slide(prs, 31, "数据治理：五张基础表支撑模型可用", "模型建设不是只训练算法，还要建立可维护的数据底座", "governance", [
        {"t": "资产表", "b": "管段、节点、阀门、泵站、水表的编码与属性。", "x": 1.05, "y": 1.40, "w": 2.75, "h": 0.80, "c": BLUE},
        {"t": "监测表", "b": "流量、压力、水质、设备状态等时序数据。", "x": 5.05, "y": 1.40, "w": 2.85, "h": 0.80, "c": GREEN},
        {"t": "工单表", "b": "报警、巡检、维修、复核和现场照片。", "x": 9.05, "y": 1.40, "w": 2.75, "h": 0.80, "c": ORANGE},
        {"t": "标签表", "b": "真实漏点、误报原因、处置结果和影响范围。", "x": 3.08, "y": 4.78, "w": 2.85, "h": 0.80, "c": PURPLE},
        {"t": "规则表", "b": "边界阀状态、施工计划、临时供水和阈值策略。", "x": 7.38, "y": 4.78, "w": 2.85, "h": 0.80, "c": RED},
    ], "数据治理讲清这五张表，听众才能理解为什么模型落地需要业务、运维和信息化共同参与。")
    flow_slide(prs, 32, "模型运营：漂移监控与再训练闭环", "管网状态会变化，模型也需要被持续运营", [
        ("监控漂移", "观察输入分布、残差分布和报警命中率是否变化。", BLUE),
        ("复核样本", "把近期待查、确认漏点和误报原因整理入库。", ORANGE),
        ("再训练", "按月或按季度更新模型，保留版本和评估结果。", GREEN),
        ("灰度发布", "先在试点DMA验证，再逐步扩大到更多片区。", PURPLE),
    ], "ops", "模型运营的重点是版本、指标、反馈和责任分工，避免模型上线后无人维护。")
    stage_slide(prs, 33, "总结：DMA提供范围，AI提供排序，工单形成闭环", "对外演讲最后要留下三句话", "cover", [
        {"t": "第一句话", "b": "DMA把全网漏损问题变成可计量、可管理的片区问题。", "x": 1.05, "y": 1.42, "w": 3.05, "h": 0.86, "c": BLUE},
        {"t": "第二句话", "b": "AI把片区异常继续收敛为候选管段和现场核查优先级。", "x": 5.00, "y": 4.72, "w": 3.35, "h": 0.88, "c": GREEN},
        {"t": "第三句话", "b": "工单闭环把一次检漏转化为下一轮模型优化的数据资产。", "x": 8.88, "y": 1.42, "w": 3.05, "h": 0.86, "c": ORANGE},
    ], "最终观点：高质量漏损检测不是单个算法，而是DMA架构、模型分析和运维闭环的系统协同。")

    prs.save(PPTX)
    AUDIT.write_text("\n".join(AUDIT_LINES), encoding="utf-8")
    build_preview(prs)
    return PPTX


def build_preview(prs):
    thumb_w, thumb_h = 360, 203
    cols = 4
    gap = 22
    rows = (len(prs.slides) + cols - 1) // cols
    sx, sy = thumb_w / prs.slide_width, thumb_h / prs.slide_height
    try:
        font_bold = ImageFont.truetype("C:/Windows/Fonts/msyhbd.ttc", 8)
        font = ImageFont.truetype("C:/Windows/Fonts/msyh.ttc", 7)
    except Exception:
        font_bold = font = ImageFont.load_default()
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + gap)), (230, 236, 242))
    draw_sheet = ImageDraw.Draw(sheet)
    for idx, slide in enumerate(prs.slides):
        im = Image.new("RGB", (thumb_w, thumb_h), "white")
        d = ImageDraw.Draw(im)
        bg = None
        try:
            bg = slide.background.fill.fore_color.rgb
        except Exception:
            pass
        if bg:
            rgb = tuple(int(str(bg)[i:i + 2], 16) for i in (0, 2, 4))
            im.paste(rgb, [0, 0, thumb_w, thumb_h])
        for sh in slide.shapes:
            x, y, w, h = int(sh.left * sx), int(sh.top * sy), int(sh.width * sx), int(sh.height * sy)
            if w <= 0 or h <= 0:
                continue
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    p = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                    p.thumbnail((w, h))
                    im.paste(p, (x + (w - p.width) // 2, y + (h - p.height) // 2))
                except Exception:
                    pass
                continue
            fill_color = None
            line_color = None
            try:
                rgb = sh.fill.fore_color.rgb
                if rgb:
                    fill_color = tuple(int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            try:
                rgb = sh.line.color.rgb
                if rgb:
                    line_color = tuple(int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            if fill_color:
                d.rounded_rectangle([x, y, x + w, y + h], radius=5, fill=fill_color, outline=line_color)
            elif line_color:
                d.rectangle([x, y, x + w, y + h], outline=line_color)
            if hasattr(sh, "text") and sh.text.strip() and w > 28 and h > 8:
                bold = False
                try:
                    bold = any(r.font.bold for pgh in sh.text_frame.paragraphs for r in pgh.runs)
                except Exception:
                    pass
                text = sh.text.strip().split("\n")[0][:24]
                d.text((x + 2, y + 2), text, fill=(20, 40, 60), font=font_bold if bold else font)
        cx = (idx % cols) * thumb_w
        cy = (idx // cols) * (thumb_h + gap)
        sheet.paste(im, (cx, cy))
        draw_sheet.text((cx + 5, cy + thumb_h + 3), f"{idx + 1:02d}", fill=(0, 70, 148), font=font_bold)
    sheet.save(PREVIEW)


if __name__ == "__main__":
    print(build())
    print(AUDIT)
    print(PREVIEW)
