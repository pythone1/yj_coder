from pathlib import Path
import shutil
from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


ROOT = Path(__file__).resolve().parent
GEN = Path(r"C:\Users\Administrator\.codex\generated_images\019e1a29-183a-70c1-bc80-166e84fb9582")
OUT = ROOT / "output" / "generated_ppt_assets_v6"
RAW = OUT / "raw_sheets"
ICONS = OUT / "split_plugins"
PPT = ROOT / "output" / "ppt" / "AI供水管网漏损检测_生图素材插件包_v6.pptx"
ZIP = ROOT / "output" / "AI供水管网漏损检测_生图素材插件包_v6.zip"
for p in [RAW, ICONS, PPT.parent]:
    p.mkdir(parents=True, exist_ok=True)


files = sorted(GEN.glob("*.png"), key=lambda p: p.stat().st_mtime)
selected = {
    "00_algorithm_overall_architecture.png": files[-5],
    "01_water_equipment_plugins_sheet.png": files[-4],
    "02_algorithm_plugins_sheet.png": files[-3],
    "03_system_architecture_modules_sheet.png": files[-2],
    "04_small_infographic_plugins_sheet.png": files[-1],
}


names = {
    "01_water_equipment_plugins_sheet.png": [
        "smart_inlet_flow_meter", "pressure_sensor_node", "boundary_valve_chamber", "dma_boundary_map",
        "pipe_leak_alert", "scada_control_screen", "gis_network_map", "work_order_clipboard",
        "data_lake_cylinder", "ai_model_cube", "cloud_platform", "edge_gateway",
        "field_inspection_tablet", "pump_station", "water_quality_sensor", "maintenance_feedback_loop",
    ],
    "02_algorithm_plugins_sheet.png": [
        "lstm_sequence_prediction", "gru_gate_motif", "cnn_lstm_hybrid", "autoencoder_reconstruction",
        "isolation_forest_outlier", "dbscan_cluster_outlier", "random_forest_ensemble", "gradient_boosting_trees",
        "genetic_algorithm_optimization", "gnn_pipe_topology", "knowledge_graph", "anomaly_score_gauge",
        "prediction_residual_chart", "model_training_pipeline", "model_drift_monitor", "human_in_loop_validation",
    ],
    "03_system_architecture_modules_sheet.png": [
        "sensing_layer_module", "data_ingestion_module", "data_cleaning_filter", "feature_engineering_module",
        "timeseries_prediction_module", "anomaly_detection_module", "hydraulic_simulation_module", "leak_localization_module",
        "risk_ranking_module", "work_order_dispatch_module", "repair_verification_module", "model_retraining_loop",
        "edge_computing_gateway", "cloud_training_module", "digital_twin_network", "executive_dashboard",
    ],
    "04_small_infographic_plugins_sheet.png": [
        "leak_warning_badge", "water_ai_droplet", "pressure_wave_pulse", "flow_direction_pipe",
        "dma_boundary_icon", "sensor_wireless_signal", "anomaly_cluster_dots", "ai_chip_icon",
        "graph_topology_icon", "timeseries_mini_chart", "prediction_confidence_band", "risk_heatmap_tile",
        "work_order_checkmark", "repair_wrench_pipe", "cloud_edge_sync", "closed_loop_arrows",
    ],
}


def copy_raw():
    copied = {}
    for out_name, src in selected.items():
        dst = RAW / out_name
        shutil.copy2(src, dst)
        copied[out_name] = dst
    return copied


def split_sheet(sheet_path, out_dir, tile_names):
    out_dir.mkdir(parents=True, exist_ok=True)
    im = Image.open(sheet_path).convert("RGB")
    w, h = im.size
    paths = []
    for r in range(4):
        for c in range(4):
            idx = r * 4 + c
            x1 = round(c * w / 4) + 4
            x2 = round((c + 1) * w / 4) - 4
            y1 = round(r * h / 4) + 4
            y2 = round((r + 1) * h / 4) - 4
            tile = im.crop((x1, y1, x2, y2))
            tile = ImageOps.contain(tile, (512, 512), Image.Resampling.LANCZOS)
            canvas = Image.new("RGB", (512, 512), "white")
            canvas.paste(tile, ((512 - tile.width) // 2, (512 - tile.height) // 2))
            path = out_dir / f"{idx+1:02d}_{tile_names[idx]}.png"
            canvas.save(path, quality=95)
            paths.append(path)
    return paths


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(12.6), Inches(0.45))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(23, 32, 51)


def add_caption(slide, x, y, w, text, size=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.23))
    tf = box.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(80, 92, 110)


def make_ppt(raw_paths, split):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    add_title(slide, "算法整体架构图（生图素材）")
    slide.shapes.add_picture(str(raw_paths["00_algorithm_overall_architecture.png"]), Inches(0.25), Inches(0.8), width=Inches(12.85))

    for raw_name, raw_path in raw_paths.items():
        if raw_name.startswith("00_"):
            continue
        slide = prs.slides.add_slide(blank)
        add_title(slide, raw_name.replace(".png", ""))
        slide.shapes.add_picture(str(raw_path), Inches(0.75), Inches(0.75), height=Inches(6.35))

    for category, paths in split.items():
        slide = prs.slides.add_slide(blank)
        add_title(slide, category)
        size = 1.25
        gap_x = 0.48
        gap_y = 0.36
        start_x = 0.75
        start_y = 0.85
        for i, path in enumerate(paths):
            r, c = divmod(i, 8)
            x = start_x + c * (size + gap_x)
            y = start_y + r * (size + gap_y + 0.22)
            slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(size), height=Inches(size))
            label = path.stem[3:].replace("_", " ")
            add_caption(slide, x - 0.06, y + size + 0.03, size + 0.12, label[:24], 6.5)

    prs.save(PPT)


def make_index(raw_paths, split):
    md = OUT / "素材清单_v6.md"
    lines = ["# AI供水管网漏损检测 生图素材插件包 v6", ""]
    lines += ["## 原始组件板", ""]
    for k, v in raw_paths.items():
        lines.append(f"- {k}: `{v}`")
    lines += ["", "## 拆分插件", ""]
    for cat, paths in split.items():
        lines.append(f"### {cat}")
        for p in paths:
            lines.append(f"- `{p.name}`")
        lines.append("")
    md.write_text("\n".join(lines), encoding="utf-8")
    return md


def main():
    raw_paths = copy_raw()
    split = {}
    for raw_name in list(raw_paths.keys()):
        if raw_name not in names:
            continue
        cat = raw_name.replace("_sheet.png", "")
        split[cat] = split_sheet(raw_paths[raw_name], ICONS / cat, names[raw_name])
    make_ppt(raw_paths, split)
    index = make_index(raw_paths, split)
    print(RAW)
    print(ICONS)
    print(PPT)
    print(index)


if __name__ == "__main__":
    main()
