from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt


BASE = Path(r"E:\PY\research\0425")
INPUT = BASE / "生产流程_最终版_v2.pptx"
OUTPUT = BASE / "生产流程_最终交付版_v2.pptx"


def add_clean_block(slide, x, y, w, h, text, font_size=10):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(4, 28, 43)
    rect.line.fill.background()

    box = slide.shapes.add_textbox(x + 120000, y + 100000, w - 240000, h - 200000)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(text.split("\n")):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.text = line
        para.space_after = Pt(0)
        for run in para.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(245, 250, 255)


prs = Presentation(str(INPUT))

# Slide 2: cover old dark text in the two right cards and restore clean editable copy.
slide = prs.slides[1]
add_clean_block(
    slide,
    9850000,
    5750000,
    3050000,
    3350000,
    "核心工艺对象：厂区物流节点、AGV任务、车辆位置、充电、避障、任务优先级。\n"
    "主要算法任务：路径规划、任务分配、多车冲突消解、预防、仿真评估。\n"
    "最终功能效益：减少等待和绕路，提高厂内物流准时率和设备利用率。",
    9,
)
add_clean_block(
    slide,
    12850000,
    5750000,
    3350000,
    3350000,
    "核心工艺对象：库位、订单、库存、装车、车辆、客户路线、多仓发货。\n"
    "主要算法任务：库位优化、订单波次、装车排程、车辆路径优化。\n"
    "最终功能效益：缩短装车发货时间，提高库位利用率、满载率和配送效率。",
    9,
)

# Slide 5: cover residual low-contrast old text in the酒体异物 row.
slide = prs.slides[4]
add_clean_block(
    slide,
    5050000,
    6100000,
    8200000,
    3000000,
    "酒体异物检测\n"
    "采集：光学工艺与高速视频，覆盖气泡、发丝、悬浮点等疑似异物。\n"
    "算法：颗粒轨迹提取、时序筛选、目标检测、异常分类模型。\n"
    "输出：异物检出、过滤复核、疑似样本入库。",
    10,
)

prs.save(str(OUTPUT))
print(OUTPUT)
