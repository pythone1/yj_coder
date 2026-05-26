from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.text import PP_ALIGN


PPTX = Path(r"D:\Users\Downloads\一厂一网_智慧水务系统全景图.pptx")
OUT = Path(r"D:\Users\Downloads\一厂一网_智慧水务系统全景图_高清.png")
EXPORT_W, EXPORT_H = 5252, 3200
BASE_W, BASE_H = 2626, 1600
FONT_FILE = r"C:\Windows\Fonts\msyh.ttc"
FONT_BOLD_FILE = r"C:\Windows\Fonts\msyhbd.ttc"


def rels_for_slide(z: ZipFile, slide_index: int = 1):
    path = f"ppt/slides/_rels/slide{slide_index}.xml.rels"
    rels = {}
    if path not in z.namelist():
        return rels
    import xml.etree.ElementTree as ET

    root = ET.fromstring(z.read(path))
    for rel in root:
        rid = rel.attrib.get("Id")
        target = rel.attrib.get("Target", "")
        if rid and target.startswith("../media/"):
            rels[rid] = "ppt/media/" + target.split("../media/", 1)[1]
    return rels


def extract_background(prs: Presentation) -> Image.Image:
    slide = prs.slides[0]
    with ZipFile(PPTX) as z:
        rels = rels_for_slide(z, 1)
        for shape in slide.shapes:
            if shape.shape_type == 13:  # picture
                rid = shape._element.blip_rId
                media = rels.get(rid)
                if media:
                    return Image.open(BytesIO(z.read(media))).convert("RGB")
    raise RuntimeError("No slide background picture found")


def font_for(size_px: int, bold: bool):
    return ImageFont.truetype(FONT_BOLD_FILE if bold and Path(FONT_BOLD_FILE).exists() else FONT_FILE, max(1, size_px))


def draw_text(draw: ImageDraw.ImageDraw, x, y, w, h, text, size_px, bold=False, align="left"):
    f = font_for(size_px, bold)
    lines = text.split("\n")
    line_h = int(size_px * 1.22)
    yy = y + (h - line_h * len(lines)) / 2
    for line in lines:
        bbox = draw.textbbox((0, 0), line, font=f)
        tw = bbox[2] - bbox[0]
        if align == "center":
            xx = x + (w - tw) / 2
        elif align == "right":
            xx = x + w - tw
        else:
            xx = x
        draw.text((xx, yy), line, font=f, fill=(245, 250, 255), stroke_width=max(2, size_px // 12), stroke_fill=(3, 12, 24))
        yy += line_h


def ppt_color_to_rgb(run):
    try:
        c = run.font.color
        if c.type == MSO_COLOR_TYPE.RGB and c.rgb:
            return tuple(int(str(c.rgb)[i : i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return (245, 250, 255)


def run_font_px(run, fallback_pt: float = 10.0) -> int:
    try:
        if run.font.size:
            pt = run.font.size.pt
        else:
            pt = fallback_pt
    except Exception:
        pt = fallback_pt
    return max(8, int(pt / 72 * 400))


def paragraph_alignment(paragraph):
    if paragraph.alignment == PP_ALIGN.CENTER:
        return "center"
    if paragraph.alignment == PP_ALIGN.RIGHT:
        return "right"
    return "left"


def export():
    prs = Presentation(PPTX)
    slide = prs.slides[0]
    bg = extract_background(prs)
    canvas = bg.resize((EXPORT_W, EXPORT_H), Image.Resampling.LANCZOS)
    draw = ImageDraw.Draw(canvas)

    sx = EXPORT_W / prs.slide_width
    sy = EXPORT_H / prs.slide_height

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
            continue
        text = shape.text.strip()
        if not text:
            continue
        x = int(shape.left * sx)
        y = int(shape.top * sy)
        w = int(shape.width * sx)
        h = int(shape.height * sy)
        paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
        if not paragraphs:
            continue

        # Estimate paragraph heights with their largest run size.
        para_infos = []
        for p in paragraphs:
            runs = [r for r in p.runs if r.text]
            if not runs:
                runs = []
            fallback_pt = 10.0
            if runs and runs[0].font.size:
                fallback_pt = runs[0].font.size.pt
            max_size = max([run_font_px(r, fallback_pt) for r in runs] or [int(fallback_pt / 72 * 400)])
            line_h = int(max_size * 1.22)
            para_infos.append((p, runs, max_size, line_h))
        total_h = sum(info[3] for info in para_infos)
        yy = y + (h - total_h) / 2

        for p, runs, max_size, line_h in para_infos:
            if not runs:
                yy += line_h
                continue
            align = paragraph_alignment(p)
            pieces = []
            total_w = 0
            for r in runs:
                size_px = run_font_px(r, max_size / 400 * 72)
                bold = bool(r.font.bold)
                f = font_for(size_px, bold)
                bbox = draw.textbbox((0, 0), r.text, font=f)
                rw = bbox[2] - bbox[0]
                pieces.append((r.text, f, ppt_color_to_rgb(r), size_px, rw))
                total_w += rw
            if align == "center":
                xx = x + (w - total_w) / 2
            elif align == "right":
                xx = x + w - total_w
            else:
                xx = x
            for text_piece, f, color, size_px, rw in pieces:
                draw.text(
                    (xx, yy),
                    text_piece,
                    font=f,
                    fill=color,
                    stroke_width=max(2, size_px // 12),
                    stroke_fill=(3, 12, 24),
                )
                xx += rw
            yy += line_h

    canvas.save(OUT, quality=95)
    print(OUT)


if __name__ == "__main__":
    export()
