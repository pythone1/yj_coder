# -*- coding: utf-8 -*-
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parent
TEMPLATE = max([p for p in ROOT.glob("*.pptx") if not p.name.startswith("~$")], key=lambda p: p.stat().st_size)
OUT = ROOT / "output" / "ppt"
OUT.mkdir(parents=True, exist_ok=True)
PPTX = OUT / "AI供水管网DMA漏损检测_可编辑视觉样张版.pptx"
AUDIT = OUT / "AI供水管网DMA漏损检测_可编辑视觉样张版_逐页检查.txt"


NAVY = RGBColor(0, 38, 84)
BLUE = RGBColor(0, 70, 148)
CYAN = RGBColor(0, 176, 240)
SKY = RGBColor(72, 191, 235)
PALE = RGBColor(231, 244, 252)
PALE2 = RGBColor(246, 251, 254)
WHITE = RGBColor(255, 255, 255)
TEXT = RGBColor(22, 38, 58)
MUTED = RGBColor(83, 102, 120)
LINE = RGBColor(184, 218, 237)
GREEN = RGBColor(31, 151, 122)
ORANGE = RGBColor(241, 142, 42)
RED = RGBColor(214, 72, 72)
PURPLE = RGBColor(96, 103, 220)
FONT = "微软雅黑"


def clear_slides(prs):
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)


def blank(prs, bg=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = bg
    return slide


def fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.fill.transparency = transparency
    shape.line.fill.background()


def stroke(shape, color=LINE, width=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def tb(slide, text, x, y, w, h, size=18, color=TEXT, bold=False,
       align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def bullets(slide, items, x, y, w, h, size=13.5, color=TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ""
        p.line_spacing = 0.9
        r = p.add_run()
        r.text = "· " + item
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def rect(slide, x, y, w, h, color=WHITE, line=LINE, radius=True, trans=0):
    typ = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    s = slide.shapes.add_shape(typ, x, y, w, h)
    fill(s, color, trans)
    stroke(s, line, 1.0)
    return s


def line(slide, x1, y1, x2, y2, color=LINE, width=1.2, arrow=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    c.line.color.rgb = color
    c.line.width = Pt(width)
    if arrow:
        c.line.end_arrowhead = True
    return c


def header(slide, title, sub, no):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.25), Inches(0.12), Inches(0.48))
    fill(bar, BLUE)
    tb(slide, title, Inches(0.62), Inches(0.16), Inches(8.6), Inches(0.38), 26, TEXT, True)
    tb(slide, sub, Inches(0.64), Inches(0.62), Inches(7.6), Inches(0.22), 11, MUTED)
    tb(slide, f"{no:02d}", Inches(12.05), Inches(0.22), Inches(0.7), Inches(0.25), 12, BLUE, True, PP_ALIGN.RIGHT)
    divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(0.90), Inches(11.7), Inches(0.016))
    fill(divider, PALE)


def footer(slide):
    tb(slide, "AI供水管网DMA漏损检测 · 可编辑视觉样张", Inches(0.62), Inches(7.08), Inches(5.4), Inches(0.2), 10, MUTED)


def content(prs, title, sub, no):
    s = blank(prs)
    header(s, title, sub, no)
    footer(s)
    return s


def pill(slide, text, x, y, w, h=0.32, color=PALE, text_color=BLUE, size=12, bold=True):
    shape = rect(slide, x, y, w, Inches(h), color, color, True)
    shape.adjustments[0] = 0.2
    tb(slide, text, x + Inches(0.06), y + Inches(0.07), w - Inches(0.12), Inches(h - 0.08), size, text_color, bold, PP_ALIGN.CENTER)
    return shape


def node(slide, x, y, label, color=BLUE, size=0.38):
    o = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, y, Inches(size), Inches(size))
    fill(o, color)
    stroke(o, WHITE, 1)
    tb(slide, label, x - Inches(0.22), y + Inches(size + 0.05), Inches(size + 0.44), Inches(0.18), 8.5, MUTED, False, PP_ALIGN.CENTER)
    return o


def draw_pipe_network(slide, x, y, w, h, alert=False):
    # Editable pipe network illustration.
    coords = [
        (0.10, 0.55), (0.30, 0.55), (0.50, 0.55), (0.70, 0.55), (0.90, 0.55),
        (0.30, 0.25), (0.30, 0.85), (0.70, 0.25), (0.70, 0.85), (0.50, 0.30), (0.50, 0.80),
    ]
    pts = [(x + Inches(w * a), y + Inches(h * b)) for a, b in coords]
    pipe_color = RGBColor(57, 143, 197)
    for a, b in [(0, 1), (1, 2), (2, 3), (3, 4), (1, 5), (1, 6), (3, 7), (3, 8), (2, 9), (2, 10)]:
        line(slide, pts[a][0], pts[a][1], pts[b][0], pts[b][1], pipe_color, 4)
    for i, (px, py) in enumerate(pts):
        col = RED if alert and i == 3 else CYAN if i in [0, 4, 5, 8] else BLUE
        slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, px - Inches(0.065), py - Inches(0.065), Inches(0.13), Inches(0.13)).fill.solid()
        shp = slide.shapes[-1]
        fill(shp, col)
        stroke(shp, WHITE, 0.7)
    if alert:
        ring = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, pts[3][0] - Inches(0.20), pts[3][1] - Inches(0.20), Inches(0.40), Inches(0.40))
        fill(ring, RGBColor(255, 235, 232), 30)
        stroke(ring, RED, 1.5)
        tb(slide, "MNF异常", pts[3][0] - Inches(0.35), pts[3][1] + Inches(0.22), Inches(0.7), Inches(0.18), 8.5, RED, True, PP_ALIGN.CENTER)


def draw_data_stack(slide, x, y, w, h):
    # Editable data platform with incoming feature chips.
    base = rect(slide, x + Inches(0.45), y + Inches(0.92), Inches(w - 0.9), Inches(0.82), WHITE, LINE, True)
    fill(base, RGBColor(248, 253, 255))
    tb(slide, "数据中台", x + Inches(0.65), y + Inches(1.18), Inches(w - 1.3), Inches(0.24), 15, BLUE, True, PP_ALIGN.CENTER)
    chips = [("流量压力", BLUE), ("管网拓扑", CYAN), ("管材管龄", GREEN), ("土壤环境", ORANGE), ("历史工单", PURPLE)]
    for i, (t, c) in enumerate(chips):
        cx = x + Inches(0.05 + (i % 3) * (w / 3.1))
        cy = y + Inches(0.05 + (i // 3) * 0.47)
        pill(slide, t, cx, cy, Inches(0.95), 0.28, RGBColor(235, 247, 253), c, 9.5)
        line(slide, cx + Inches(0.48), cy + Inches(0.31), x + Inches(w/2), y + Inches(0.94), c, 0.8, True)


def draw_ai_engine(slide, x, y, w, h, dark=False):
    c_fill = RGBColor(229, 246, 255) if not dark else RGBColor(5, 72, 136)
    chip = rect(slide, x + Inches(w * 0.33), y + Inches(0.55), Inches(w * 0.34), Inches(0.9), c_fill, CYAN, True)
    tb(slide, "AI", x + Inches(w * 0.33), y + Inches(0.78), Inches(w * 0.34), Inches(0.28), 23, BLUE if not dark else WHITE, True, PP_ALIGN.CENTER)
    modules = [("水力模型", BLUE), ("深度学习", GREEN), ("反向推演", ORANGE), ("压力波动", PURPLE)]
    locs = [(0.02, 0.10), (0.68, 0.10), (0.02, 1.55), (0.68, 1.55)]
    for (t, c), (a, b0) in zip(modules, locs):
        pill(slide, t, x + Inches(w * a), y + Inches(b0), Inches(w * 0.30), 0.30, PALE, c, 10)
        line(slide, x + Inches(w * (a + 0.15)), y + Inches(b0 + 0.34), x + Inches(w * 0.50), y + Inches(1.00), c, 0.9, True)


def draw_funnel(slide, x, y, w, h):
    labels = [("DMA片区", BLUE), ("候选管段TopN", CYAN), ("节点坐标", ORANGE), ("现场派单", GREEN)]
    for i, (t, c) in enumerate(labels):
        yy = y + Inches(0.08 + i * 0.43)
        ww = Inches(w - i * 0.35)
        xx = x + Inches(i * 0.175)
        shape = rect(slide, xx, yy, ww, Inches(0.30), RGBColor(244, 250, 253), c, True)
        tb(slide, t, xx, yy + Inches(0.07), ww, Inches(0.15), 9.5, c, True, PP_ALIGN.CENTER)
    pin = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(w * 0.44), y + Inches(1.88), Inches(0.25), Inches(0.25))
    fill(pin, RED)
    stroke(pin, WHITE, 1)
    tb(slide, "维修闭环", x, y + Inches(2.20), Inches(w), Inches(0.22), 10, GREEN, True, PP_ALIGN.CENTER)


def module_card(slide, x, y, w, h, title, subtitle, color, draw_func=None):
    rect(slide, x, y, w, h, WHITE, RGBColor(190, 222, 240), True)
    fill(slide.shapes[-1], WHITE)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, Inches(0.44))
    fill(top, color)
    tb(slide, title, x + Inches(0.14), y + Inches(0.11), w - Inches(0.28), Inches(0.18), 13.5, WHITE, True, PP_ALIGN.CENTER)
    if draw_func:
        draw_func(slide, x + Inches(0.20), y + Inches(0.62), w.inches - 0.40, h.inches - 1.35)
    tb(slide, subtitle, x + Inches(0.18), y + h - Inches(0.48), w - Inches(0.36), Inches(0.34), 10.5, MUTED, False, PP_ALIGN.CENTER)


def slide_cover(prs):
    s = blank(prs, NAVY)
    # left template-style dark block
    accent = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.58), Inches(1.05), Inches(0.13), Inches(5.05))
    fill(accent, CYAN)
    tb(s, "AI模型在供水管网DMA系统\n漏损检测中的应用", Inches(0.95), Inches(1.28), Inches(5.25), Inches(1.55), 35, WHITE, True)
    tb(s, "可编辑视觉样张 · 风格确认版", Inches(1.0), Inches(3.20), Inches(4.8), Inches(0.32), 18, CYAN, True)
    tb(s, "DMA宏观锁定 + AI微观溯源 + 工单闭环处置", Inches(1.0), Inches(5.86), Inches(5.1), Inches(0.28), 12.5, RGBColor(210, 232, 245))
    # right editable architecture illustration
    stage = rect(s, Inches(6.55), Inches(0.65), Inches(6.2), Inches(5.9), WHITE, RGBColor(213, 235, 247), True)
    fill(stage, RGBColor(249, 253, 255))
    draw_pipe_network(s, Inches(6.95), Inches(4.65), 5.35, 1.25, True)
    for i, name in enumerate(["动态基线", "异常检测", "溯源定位"]):
        x = Inches(7.0 + i * 1.65)
        rect(s, x, Inches(1.2), Inches(1.25), Inches(0.85), RGBColor(246, 251, 254), RGBColor(206, 230, 243), True)
        tb(s, name, x + Inches(0.08), Inches(1.38), Inches(1.09), Inches(0.18), 9.5, BLUE, True, PP_ALIGN.CENTER)
        if i == 0:
            line(s, x + Inches(0.25), Inches(1.78), x + Inches(1.0), Inches(1.48), SKY, 1)
        elif i == 1:
            for n in range(10):
                node(s, x + Inches(0.15 + (n % 5)*0.18), Inches(1.42 + (n // 5)*0.18), "", CYAN if n % 3 else BLUE, 0.08)
        else:
            line(s, x + Inches(0.25), Inches(1.45), x + Inches(0.95), Inches(1.75), BLUE, 1)
            line(s, x + Inches(0.95), Inches(1.45), x + Inches(0.25), Inches(1.75), BLUE, 1)
    rect(s, Inches(11.05), Inches(1.15), Inches(1.2), Inches(1.1), RGBColor(247, 251, 254), RGBColor(206, 230, 243), True)
    draw_funnel(s, Inches(11.18), Inches(1.30), 0.95, 0.95)
    tb(s, "教学汇报版", Inches(1.0), Inches(5.42), Inches(1.8), Inches(0.25), 16, CYAN, True)


def slide_arch(prs):
    s = content(prs, "DMA宏观锁定 + AI微观溯源总体架构", "可编辑架构图样张", 2)
    y = Inches(1.20)
    w = Inches(2.65)
    module_card(s, Inches(0.62), y, w, Inches(3.95), "1 DMA分区计量", "MNF异常升高 → 区域锁定", BLUE,
                lambda sl, x, y, ww, hh: draw_pipe_network(sl, x, y, ww, hh, True))
    module_card(s, Inches(3.55), y, w, Inches(3.95), "2 多源数据融合", "工程 + 空间 + 历史信息", CYAN,
                lambda sl, x, y, ww, hh: draw_data_stack(sl, x, y, ww, hh))
    module_card(s, Inches(6.48), y, w, Inches(3.95), "3 AI溯源分析", "反向推演水力影响", GREEN,
                lambda sl, x, y, ww, hh: draw_ai_engine(sl, x, y, ww, hh))
    module_card(s, Inches(9.41), y, w, Inches(3.95), "4 精准定位处置", "管段TopN → 节点坐标", ORANGE,
                lambda sl, x, y, ww, hh: draw_funnel(sl, x, y, ww, hh))
    for x in [Inches(3.28), Inches(6.21), Inches(9.14)]:
        line(s, x, Inches(3.02), x + Inches(0.23), Inches(3.02), BLUE, 2.0, True)
    # information-rich lower belt
    rect(s, Inches(0.62), Inches(5.55), Inches(11.44), Inches(0.78), RGBColor(247, 252, 255), RGBColor(195, 224, 240), True)
    tb(s, "DMA锁定区域  →  AI融合分析  →  漏点溯源定位  →  工单闭环处置", Inches(0.9), Inches(5.78), Inches(10.9), Inches(0.26), 18, BLUE, True, PP_ALIGN.CENTER)
    footer(s)


def alg_card(slide, x, y, title, lines, color):
    rect(slide, x, y, Inches(3.35), Inches(1.18), WHITE, RGBColor(205, 228, 242), True)
    fill(slide.shapes[-1], WHITE)
    circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.16), y + Inches(0.22), Inches(0.38), Inches(0.38))
    fill(circ, color)
    tb(slide, title, x + Inches(0.65), y + Inches(0.16), Inches(2.55), Inches(0.23), 14.5, TEXT, True)
    tb(slide, "\n".join(lines), x + Inches(0.65), y + Inches(0.47), Inches(2.55), Inches(0.50), 10.5, MUTED)


def slide_algorithm(prs):
    s = content(prs, "AI算法体系：按业务任务组合模型", "可编辑算法架构样张", 3)
    # central engine
    center = rect(s, Inches(4.80), Inches(2.40), Inches(3.35), Inches(1.38), RGBColor(236, 249, 255), CYAN, True)
    tb(s, "AI模型引擎", Inches(4.80), Inches(2.70), Inches(3.35), Inches(0.30), 24, BLUE, True, PP_ALIGN.CENTER)
    tb(s, "预测 · 识别 · 定位 · 排序 · 解释", Inches(4.80), Inches(3.15), Inches(3.35), Inches(0.22), 12.5, MUTED, False, PP_ALIGN.CENTER)
    alg_card(s, Inches(0.75), Inches(1.30), "时序预测", ["LSTM / GRU", "动态基线、预测区间"], BLUE)
    alg_card(s, Inches(0.75), Inches(3.05), "异常检测", ["孤立森林 / DBSCAN", "自编码器、异常分数"], CYAN)
    alg_card(s, Inches(4.90), Inches(4.75), "机理融合", ["水力模型 + ML", "候选管段TopN"], GREEN)
    alg_card(s, Inches(8.95), Inches(1.30), "风险排序", ["随机森林 / GBDT", "风险分、优先级"], ORANGE)
    alg_card(s, Inches(8.95), Inches(3.05), "解释协同", ["GNN / 知识图谱", "报告、问答、派单"], PURPLE)
    for x1, y1 in [(Inches(4.10), Inches(1.90)), (Inches(4.10), Inches(3.65)), (Inches(6.55), Inches(4.75)), (Inches(8.95), Inches(1.90)), (Inches(8.95), Inches(3.65))]:
        line(s, x1, y1, Inches(6.48), Inches(3.08), RGBColor(154, 196, 220), 1.2, True)
    rect(s, Inches(0.75), Inches(6.15), Inches(11.55), Inches(0.48), RGBColor(247, 252, 255), LINE, True)
    tb(s, "讲课逻辑：先用时序模型建立正常边界，再用异常检测发现候选，随后通过水力机理和风险排序把结果转成可执行工单。", Inches(1.0), Inches(6.29), Inches(11.0), Inches(0.20), 13.2, BLUE, True, PP_ALIGN.CENTER)


def slide_loop(prs):
    s = content(prs, "典型应用闭环：事前预警、事中定位、事后复盘", "可编辑业务闭环样张", 4)
    stages = [
        ("事前预警", "动态基线\n残差持续偏离\n压力联动", BLUE),
        ("事中定位", "候选片区\n管段TopN\n节点坐标", CYAN),
        ("现场处置", "派单核查\n维修记录\n结果确认", ORANGE),
        ("事后复盘", "误报归因\n标签回填\n模型再训练", GREEN),
    ]
    for i, (title, body, color) in enumerate(stages):
        x = Inches(0.85 + i * 3.0)
        rect(s, x, Inches(1.35), Inches(2.35), Inches(2.65), WHITE, RGBColor(201, 227, 241), True)
        top = s.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, Inches(1.35), Inches(2.35), Inches(0.44))
        fill(top, color)
        tb(s, title, x, Inches(1.46), Inches(2.35), Inches(0.18), 14.5, WHITE, True, PP_ALIGN.CENTER)
        tb(s, body, x + Inches(0.22), Inches(2.05), Inches(1.9), Inches(0.98), 16, TEXT, True, PP_ALIGN.CENTER)
        pill(s, "业务输出", x + Inches(0.48), Inches(3.35), Inches(1.40), 0.28, PALE, color, 9.5)
        if i < 3:
            line(s, x + Inches(2.40), Inches(2.65), x + Inches(2.86), Inches(2.65), BLUE, 1.8, True)
    # closed-loop return
    line(s, Inches(10.25), Inches(4.35), Inches(2.05), Inches(4.35), RGBColor(150, 190, 214), 1.2, True)
    tb(s, "工单结果回填，驱动模型持续学习", Inches(4.2), Inches(4.47), Inches(4.0), Inches(0.22), 13.5, GREEN, True, PP_ALIGN.CENTER)
    rect(s, Inches(0.88), Inches(5.55), Inches(11.2), Inches(0.72), RGBColor(247, 252, 255), LINE, True)
    tb(s, "核心价值：把“哪个区可能漏”转化为“哪段管优先查、谁去查、结果如何回填”。", Inches(1.1), Inches(5.76), Inches(10.75), Inches(0.25), 17, BLUE, True, PP_ALIGN.CENTER)


def slide_library(prs):
    s = content(prs, "可编辑组件库：后续整套PPT按此风格扩展", "组件库样张", 5)
    # module examples
    module_card(s, Inches(0.68), Inches(1.18), Inches(2.25), Inches(2.15), "DMA模块", "分区/阀门/传感器", BLUE,
                lambda sl, x, y, ww, hh: draw_pipe_network(sl, x, y, ww, hh, True))
    module_card(s, Inches(3.18), Inches(1.18), Inches(2.25), Inches(2.15), "数据模块", "多源特征汇聚", CYAN,
                lambda sl, x, y, ww, hh: draw_data_stack(sl, x, y, ww, hh))
    module_card(s, Inches(5.68), Inches(1.18), Inches(2.25), Inches(2.15), "AI模块", "模型引擎/推演", GREEN,
                lambda sl, x, y, ww, hh: draw_ai_engine(sl, x, y, ww, hh))
    module_card(s, Inches(8.18), Inches(1.18), Inches(2.25), Inches(2.15), "定位模块", "TopN/坐标/派单", ORANGE,
                lambda sl, x, y, ww, hh: draw_funnel(sl, x, y, ww, hh))
    # primitives
    tb(s, "标题条 / 箭头 / 标签 / 讲解带", Inches(0.75), Inches(3.86), Inches(3.6), Inches(0.28), 18, BLUE, True)
    pill(s, "动态基线", Inches(0.78), Inches(4.35), Inches(1.15), 0.32, PALE, BLUE, 11)
    pill(s, "异常检测", Inches(2.05), Inches(4.35), Inches(1.15), 0.32, PALE, CYAN, 11)
    pill(s, "工单闭环", Inches(3.32), Inches(4.35), Inches(1.15), 0.32, PALE, GREEN, 11)
    line(s, Inches(5.1), Inches(4.52), Inches(6.4), Inches(4.52), BLUE, 2, True)
    rect(s, Inches(6.75), Inches(4.12), Inches(4.95), Inches(0.85), RGBColor(247, 252, 255), LINE, True)
    tb(s, "每页保留：主观点 + 可编辑图示 + 讲课要点 + 结论", Inches(7.0), Inches(4.38), Inches(4.45), Inches(0.22), 13.5, TEXT, True, PP_ALIGN.CENTER)
    rect(s, Inches(0.75), Inches(5.55), Inches(11.35), Inches(0.72), RGBColor(247, 252, 255), LINE, True)
    tb(s, "确认这5页风格后，可以把原34页整套PPT按同一组件语言重做。", Inches(1.0), Inches(5.78), Inches(10.85), Inches(0.24), 17, BLUE, True, PP_ALIGN.CENTER)


def audit(prs):
    lines = []
    fail = False
    for idx, slide in enumerate(prs.slides, 1):
        chars = 0
        pics = 0
        min_font = 99
        off = 0
        shape_count = len(slide.shapes)
        for sh in slide.shapes:
            if sh.left < -50000 or sh.top < -50000 or sh.left + sh.width > prs.slide_width + 50000 or sh.top + sh.height > prs.slide_height + 50000:
                off += 1
            if sh.shape_type == 13:
                pics += 1
            if hasattr(sh, "text_frame") and sh.text.strip():
                chars += len(sh.text.strip())
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        if r.font.size:
                            min_font = min(min_font, r.font.size.pt)
        flags = []
        if off:
            flags.append(f"越界{off}")
        if min_font < 8.5:
            flags.append(f"小字{min_font}")
        if shape_count < 30:
            flags.append(f"组件不足{shape_count}")
        if idx > 1 and chars < 120:
            flags.append(f"文字不足{chars}")
        status = "FAIL" if flags else "OK"
        if flags:
            fail = True
        lines.append(f"{idx:02d}\t{status}\tchars={chars}\tshapes={shape_count}\tpics={pics}\tmin_font={min_font if min_font != 99 else '-'}\t{'；'.join(flags)}")
    AUDIT.write_text("\n".join(lines), encoding="utf-8")
    if fail:
        print("\n".join(lines))
        raise SystemExit("audit failed")


def build():
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)
    slide_cover(prs)
    slide_arch(prs)
    slide_algorithm(prs)
    slide_loop(prs)
    slide_library(prs)
    prs.save(PPTX)
    audit(prs)
    return PPTX


if __name__ == "__main__":
    print(build())
