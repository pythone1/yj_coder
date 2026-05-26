from pathlib import Path
import math
import random

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PRGBColor


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "output"
DOC_DIR = OUT / "doc"
ASSET_DIR = OUT / "assets_v2"
PPT_DIR = OUT / "ppt"
for p in [DOC_DIR, ASSET_DIR, PPT_DIR]:
    p.mkdir(parents=True, exist_ok=True)

FONT_PATH = Path(r"C:\Windows\Fonts\simhei.ttf")
FONT = ImageFont.truetype(str(FONT_PATH), 30)
FONT_M = ImageFont.truetype(str(FONT_PATH), 26)
FONT_S = ImageFont.truetype(str(FONT_PATH), 22)
FONT_XS = ImageFont.truetype(str(FONT_PATH), 18)
FONT_B = ImageFont.truetype(str(FONT_PATH), 38)
FONT_XL = ImageFont.truetype(str(FONT_PATH), 50)

RED = RGBColor(192, 0, 0)
BLACK = RGBColor(0, 0, 0)
PINK = RGBColor(237, 28, 36)

INK = (26, 37, 55)
MUTED = (88, 101, 119)
BLUE = (36, 99, 168)
CYAN = (19, 145, 167)
GREEN = (59, 143, 103)
ORANGE = (224, 142, 38)
RED_P = (206, 75, 64)
PURPLE = (101, 90, 190)
LINE = (203, 213, 225)
PALE = (246, 248, 251)


def wrap(draw, text, font, max_width):
    lines, cur = [], ""
    for ch in text:
        test = cur + ch
        if draw.textlength(test, font=font) <= max_width:
            cur = test
        else:
            if cur:
                lines.append(cur)
            cur = ch
    if cur:
        lines.append(cur)
    return lines


def text_center(draw, box, text, font, fill=INK, spacing=6):
    x1, y1, x2, y2 = box
    lines = []
    for raw in text.split("\n"):
        lines += wrap(draw, raw, font, max(10, x2 - x1 - 20))
    heights = [draw.textbbox((0, 0), line, font=font)[3] - draw.textbbox((0, 0), line, font=font)[1] for line in lines]
    total = sum(heights) + spacing * (len(lines) - 1)
    y = y1 + (y2 - y1 - total) / 2
    for line, h in zip(lines, heights):
        bb = draw.textbbox((0, 0), line, font=font)
        draw.text((x1 + (x2 - x1 - (bb[2] - bb[0])) / 2, y), line, font=font, fill=fill)
        y += h + spacing


def text_left(draw, xy, text, font, fill=INK, max_width=480, leading=6):
    x, y = xy
    for line in wrap(draw, text, font, max_width):
        draw.text((x, y), line, font=font, fill=fill)
        y += (draw.textbbox((0, 0), line, font=font)[3] + leading)
    return y


def arrow(draw, start, end, color=MUTED, width=4):
    draw.line((start, end), fill=color, width=width)
    ang = math.atan2(end[1] - start[1], end[0] - start[0])
    size = 15
    pts = [
        end,
        (end[0] - size * math.cos(ang - math.pi / 6), end[1] - size * math.sin(ang - math.pi / 6)),
        (end[0] - size * math.cos(ang + math.pi / 6), end[1] - size * math.sin(ang + math.pi / 6)),
    ]
    draw.polygon(pts, fill=color)


def bg():
    return Image.new("RGBA", (1600, 900), (255, 255, 255, 0))


def save(img, name):
    p = ASSET_DIR / name
    img.save(p)
    return p


def pill(draw, box, label, color, font=FONT_S, fill_alpha=34):
    x1, y1, x2, y2 = box
    draw.rounded_rectangle(box, radius=20, fill=color + (fill_alpha,), outline=color + (230,), width=3)
    text_center(draw, box, label, font, INK)


def fig_dma_reference_architecture():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "DMA漏损管控参考架构", font=FONT_XL, fill=INK)
    d.text((82, 112), "分区边界 + 计量闭合 + 压力监测 + 模型识别 + 工单闭环", font=FONT_M, fill=MUTED)
    cols = [
        ("规划层", ["边界划分", "入口/出口识别", "阀门封闭校验", "压力分区"], BLUE),
        ("感知层", ["入口流量计", "关键压力点", "远传水表", "RTU/边缘网关"], CYAN),
        ("数据层", ["SCADA", "GIS资产", "营收抄表", "抢维修工单"], GREEN),
        ("模型层", ["MNF基线", "LSTM预测", "异常检测", "水力-数据融合"], ORANGE),
        ("业务层", ["预警分级", "候选管段", "巡检派单", "复盘再训练"], RED_P),
    ]
    x0, y0, w, h, gap = 82, 220, 250, 430, 48
    for i, (title, items, color) in enumerate(cols):
        x = x0 + i * (w + gap)
        d.rounded_rectangle((x, y0, x + w, y0 + h), radius=24, fill=(255, 255, 255, 245), outline=color + (255,), width=4)
        d.ellipse((x + 24, y0 + 26, x + 72, y0 + 74), fill=color + (255,))
        d.text((x + 88, y0 + 34), title, font=FONT_M, fill=INK)
        yy = y0 + 120
        for item in items:
            pill(d, (x + 26, yy, x + w - 26, yy + 54), item, color, FONT_S, 24)
            yy += 78
        if i < len(cols) - 1:
            arrow(d, (x + w + 10, y0 + h / 2), (x + w + gap - 12, y0 + h / 2), width=5)
    d.rounded_rectangle((280, 730, 1320, 802), radius=24, fill=(245, 248, 252, 230), outline=LINE + (255,), width=2)
    text_center(d, (300, 730, 1300, 802), "关键判断：DMA不是单个软件模块，而是“工程分区、仪表数据、模型算法、业务处置”组合成的管理单元", FONT_S, INK)
    return save(img, "v2_01_dma_reference_architecture.png")


def fig_dma_field_topology():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 45), "DMA现场拓扑与监测布点", font=FONT_XL, fill=INK)
    d.text((82, 108), "入口计量用于水量闭合，压力点用于响应定位，边界阀用于分区可靠性", font=FONT_M, fill=MUTED)
    # zone outline
    d.rounded_rectangle((170, 180, 1240, 760), radius=55, fill=(242, 247, 253, 190), outline=BLUE + (230,), width=5)
    d.text((220, 205), "DMA-A", font=FONT_B, fill=BLUE)
    nodes = {
        "J1": (300, 350), "J2": (520, 300), "J3": (760, 360), "J4": (1020, 315),
        "J5": (1180, 460), "J6": (420, 580), "J7": (665, 620), "J8": (900, 565), "J9": (1120, 650),
    }
    edges = [("J1","J2"),("J2","J3"),("J3","J4"),("J4","J5"),("J1","J6"),("J2","J6"),("J3","J7"),("J3","J8"),("J4","J8"),("J5","J9"),("J6","J7"),("J7","J8"),("J8","J9")]
    for a, b in edges:
        d.line((nodes[a], nodes[b]), fill=(121, 161, 205, 255), width=18)
        d.line((nodes[a], nodes[b]), fill=(222, 236, 252, 255), width=7)
    for k, (x, y) in nodes.items():
        d.ellipse((x-23, y-23, x+23, y+23), fill=(255,255,255,255), outline=BLUE + (255,), width=4)
        text_center(d, (x-24, y-20, x+24, y+20), k, FONT_XS, INK)
    # devices
    d.rounded_rectangle((70, 308, 198, 392), radius=18, fill=CYAN + (255,))
    text_center(d, (70, 308, 198, 392), "入口\n流量计", FONT_S, (255,255,255))
    arrow(d, (198, 350), (270, 350), CYAN, 6)
    for x, y, t in [(520,238,"P1高压点"), (905,500,"P2中部点"), (1120,708,"P3末梢点")]:
        d.rounded_rectangle((x-70,y-34,x+70,y+34), radius=16, fill=ORANGE+(255,))
        text_center(d, (x-70,y-34,x+70,y+34), t, FONT_XS, (255,255,255))
    d.rounded_rectangle((1010, 246, 1126, 294), radius=14, fill=GREEN+(255,))
    text_center(d, (1010,246,1126,294), "边界阀", FONT_XS, (255,255,255))
    d.ellipse((820, 430, 900, 510), fill=RED_P+(255,))
    text_center(d, (820,430,900,510), "疑似\n漏点", FONT_XS, (255,255,255))
    # labels
    d.rounded_rectangle((1280, 230, 1510, 690), radius=26, fill=(255,255,255,245), outline=LINE+(255,), width=2)
    yy = 260
    for title, desc, color in [
        ("流量", "MNF、日总量、突增突降", CYAN),
        ("压力", "多点响应、压降传播", ORANGE),
        ("边界", "阀门状态、误连通校验", GREEN),
        ("工单", "定位结果回填为标签", RED_P),
    ]:
        d.ellipse((1310, yy+6, 1332, yy+28), fill=color+(255,))
        d.text((1348, yy), title, font=FONT_S, fill=INK)
        d.text((1348, yy+34), desc, font=FONT_XS, fill=MUTED)
        yy += 100
    return save(img, "v2_02_dma_field_topology.png")


def fig_mnf_logic():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "MNF夜间最小流量分析逻辑", font=FONT_XL, fill=INK)
    d.text((82, 112), "先分解夜间合法用水，再识别背景漏损和突发漏损", font=FONT_M, fill=MUTED)
    left, top, right, bottom = 140, 230, 1050, 705
    d.rectangle((left, top, right, bottom), fill=(255,255,255,230), outline=LINE+(255,), width=2)
    for i in range(6):
        y = top + i*(bottom-top)/5
        d.line((left, y, right, y), fill=(226,232,240,255), width=1)
    random.seed(2)
    pts = []
    n = 96
    for i in range(n):
        hour = i/4
        x = left + i*(right-left)/(n-1)
        v = 0.42 + 0.23*math.sin((hour-7)/24*2*math.pi) + 0.14*math.sin((hour-18)/24*2*math.pi)
        v += 0.03*random.random()
        if 2 <= hour <= 4:
            v = 0.23 + 0.02*math.sin(i)
        y = bottom - (v*0.92)*(bottom-top)
        pts.append((x,y))
    d.line(pts, fill=BLUE+(255,), width=5)
    mnf_x1 = left + 8/96*(right-left)
    mnf_x2 = left + 16/96*(right-left)
    d.rectangle((mnf_x1, top, mnf_x2, bottom), fill=(255, 231, 186, 100))
    d.text((mnf_x1+10, top+20), "2:00-4:00\nMNF窗口", font=FONT_XS, fill=ORANGE)
    d.line((left, bottom-0.23*0.92*(bottom-top), right, bottom-0.23*0.92*(bottom-top)), fill=RED_P+(220,), width=3)
    d.text((right+18, bottom-0.23*0.92*(bottom-top)-18), "MNF", font=FONT_S, fill=RED_P)
    # decomposition stack
    x0, y0, w, h = 1135, 250, 310, 410
    parts = [
        ("背景漏损", 0.42, RED_P),
        ("突发漏损", 0.23, ORANGE),
        ("合法夜间用水", 0.20, CYAN),
        ("仪表/边界误差", 0.15, MUTED),
    ]
    y = y0 + h
    for label, frac, color in parts:
        hh = h*frac
        d.rectangle((x0, y-hh, x0+w, y), fill=color+(210,), outline=(255,255,255,255), width=2)
        text_center(d, (x0, y-hh, x0+w, y), label, FONT_S, (255,255,255) if color!=MUTED else INK)
        y -= hh
    d.text((1135, 690), "MNF ≠ 漏损量\n需要扣除合法夜间用水并校验边界", font=FONT_S, fill=INK)
    return save(img, "v2_03_mnf_logic.png")


def fig_ai_pipeline():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "DMA漏损AI模型流水线", font=FONT_XL, fill=INK)
    d.text((82, 112), "从原始时序到可派单结论，中间必须经过治理、预测、识别、复核", font=FONT_M, fill=MUTED)
    steps = [
        ("数据接入", "流量/压力/水表/工单/GIS", BLUE),
        ("时序治理", "对齐、补缺、去毛刺、节假日特征", CYAN),
        ("正常基线", "MNF基线 + LSTM/GRU预测", GREEN),
        ("异常识别", "残差、孤立森林、自编码器", ORANGE),
        ("定位收敛", "拓扑 + 水力响应 + 历史案例", PURPLE),
        ("业务闭环", "预警、派单、复盘、再训练", RED_P),
    ]
    x, y, w, h, gap = 70, 270, 220, 250, 35
    for i, (title, desc, color) in enumerate(steps):
        xx = x + i*(w+gap)
        d.rounded_rectangle((xx, y, xx+w, y+h), radius=26, fill=(255,255,255,245), outline=color+(255,), width=4)
        d.ellipse((xx+76, y+28, xx+144, y+96), fill=color+(255,))
        text_center(d, (xx+76, y+28, xx+144, y+96), str(i+1), FONT_B, (255,255,255))
        text_center(d, (xx+18, y+118, xx+w-18, y+158), title, FONT_M, INK)
        text_center(d, (xx+18, y+168, xx+w-18, y+230), desc, FONT_XS, MUTED)
        if i < len(steps)-1:
            arrow(d, (xx+w+4, y+125), (xx+w+gap-8, y+125), MUTED, 4)
    d.arc((430, 575, 1170, 830), 0, 180, fill=GREEN+(255,), width=5)
    arrow(d, (430, 703), (300, 520), GREEN, 5)
    d.text((560, 740), "工单结果回填，形成标签库，模型定期再训练", font=FONT_S, fill=GREEN)
    return save(img, "v2_04_ai_pipeline.png")


def fig_model_selection_tree():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "漏损检测模型选型决策树", font=FONT_XL, fill=INK)
    d.text((82, 112), "先判断数据和标签，再决定模型复杂度", font=FONT_M, fill=MUTED)
    # node helper
    def node(box, text, color):
        d.rounded_rectangle(box, radius=22, fill=(255,255,255,245), outline=color+(255,), width=4)
        text_center(d, box, text, FONT_S, INK)
    node((650, 190, 950, 275), "是否有稳定DMA边界\n和入口流量？", BLUE)
    node((260, 355, 560, 445), "否：先做工程治理\n边界校验/计量补齐", RED_P)
    node((1040, 355, 1340, 445), "是：进入模型建设", GREEN)
    arrow(d, (650, 235), (560, 400), MUTED, 4)
    arrow(d, (950, 235), (1040, 400), MUTED, 4)
    d.text((575, 328), "否", font=FONT_S, fill=RED_P)
    d.text((994, 328), "是", font=FONT_S, fill=GREEN)
    node((900, 525, 1200, 615), "是否有可靠漏损标签？", BLUE)
    arrow(d, (1190, 445), (1050, 525), MUTED, 4)
    node((540, 700, 850, 800), "少标签：MNF基线\nLSTM残差 + 孤立森林", ORANGE)
    node((1110, 700, 1420, 800), "有标签：GBDT/随机森林\nLSTM分类/定位模型", PURPLE)
    arrow(d, (900, 570), (850, 750), MUTED, 4)
    arrow(d, (1200, 570), (1110, 750), MUTED, 4)
    d.text((850, 655), "少/无", font=FONT_S, fill=ORANGE)
    d.text((1215, 655), "有", font=FONT_S, fill=PURPLE)
    node((95, 560, 420, 690), "最低可交付成果：\n分区台账、数据质量报告、\nMNF人工分析模板", CYAN)
    arrow(d, (410, 445), (280, 560), MUTED, 4)
    return save(img, "v2_05_model_selection_tree.png")


def fig_hybrid_method():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "机理模型与AI模型融合路线", font=FONT_XL, fill=INK)
    d.text((82, 112), "机理模型给约束和仿真样本，AI模型给识别和排序能力", font=FONT_M, fill=MUTED)
    left = (120, 250, 560, 710)
    right = (1040, 250, 1480, 710)
    mid = (650, 250, 950, 710)
    d.rounded_rectangle(left, radius=30, fill=(255,255,255,245), outline=BLUE+(255,), width=4)
    d.rounded_rectangle(right, radius=30, fill=(255,255,255,245), outline=ORANGE+(255,), width=4)
    d.rounded_rectangle(mid, radius=30, fill=(245,248,252,240), outline=GREEN+(255,), width=4)
    text_center(d, (120,260,560,320), "水力机理模型", FONT_M, BLUE)
    text_center(d, (1040,260,1480,320), "AI数据模型", FONT_M, ORANGE)
    text_center(d, (650,260,950,320), "融合接口", FONT_M, GREEN)
    left_items = ["拓扑/管径/高程", "泵阀工况", "压力-流量响应", "漏点情景仿真"]
    right_items = ["时序预测", "异常检测", "风险排序", "候选管段识别"]
    mid_items = ["仿真样本库", "特征工程", "参数校准", "置信度输出"]
    for arr, x1, y1, color in [(left_items, 165, 370, BLUE), (mid_items, 690, 370, GREEN), (right_items, 1085, 370, ORANGE)]:
        for item in arr:
            pill(d, (x1, y1, x1+350 if x1!=690 else x1+220, y1+52), item, color, FONT_S, 28)
            y1 += 76
    arrow(d, (560, 480), (650, 480), MUTED, 5)
    arrow(d, (950, 480), (1040, 480), MUTED, 5)
    arrow(d, (1040, 600), (950, 600), MUTED, 5)
    d.text((700, 755), "输出：异常等级 + 漏损原因解释 + 候选管段TopN + 巡检建议", font=FONT_S, fill=INK)
    return save(img, "v2_06_hydraulic_ai_hybrid.png")


def fig_edge_cloud():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "端-边-云-业务协同部署", font=FONT_XL, fill=INK)
    d.text((82, 112), "实时性放在边缘，训练和知识沉淀放在云端，处置闭环回到业务系统", font=FONT_M, fill=MUTED)
    levels = [
        ("端侧", "流量计、压力计、RTU、水表\n采集、缓存、校时", BLUE),
        ("边缘侧", "数据质检、轻量异常识别\n断网告警、局部联动", CYAN),
        ("云平台", "全量数据湖、模型训练\n跨DMA对比、知识库", ORANGE),
        ("业务侧", "GIS定位、巡检派单\n维修复盘、指标考核", GREEN),
    ]
    xs = [120, 485, 850, 1215]
    for (title, desc, color), x in zip(levels, xs):
        d.rounded_rectangle((x, 260, x+270, 620), radius=28, fill=(255,255,255,245), outline=color+(255,), width=4)
        d.ellipse((x+95, 300, x+175, 380), fill=color+(255,))
        text_center(d, (x+95,300,x+175,380), title, FONT_M, (255,255,255))
        text_center(d, (x+28,430,x+242,555), desc, FONT_S, INK)
    for x in [390,755,1120]:
        arrow(d, (x, 440), (x+85, 440), MUTED, 5)
    d.arc((390, 600, 1210, 835), 0, 180, fill=GREEN+(255,), width=5)
    arrow(d, (390, 718), (275, 620), GREEN, 5)
    d.text((575, 760), "工单标签、设备状态、管网改造信息反向更新模型", font=FONT_S, fill=GREEN)
    return save(img, "v2_07_edge_cloud_business.png")


def fig_teaching_storyboard():
    img = bg()
    d = ImageDraw.Draw(img)
    d.text((80, 50), "教学演讲PPT推荐叙事结构", font=FONT_XL, fill=INK)
    d.text((82, 112), "从行业问题讲到技术方法，再落到工程实施和案例复盘", font=FONT_M, fill=MUTED)
    steps = [
        ("1", "为什么要做", "政策目标、漏损成本、DMA价值", BLUE),
        ("2", "DMA怎么建", "边界、计量、压力、MNF", CYAN),
        ("3", "AI怎么判", "LSTM基线、异常检测、融合定位", ORANGE),
        ("4", "系统怎么落", "端边云、数据治理、工单闭环", GREEN),
        ("5", "案例怎么讲", "异常曲线、候选区域、维修复盘", RED_P),
    ]
    for i, (num, title, desc, color) in enumerate(steps):
        x = 150 + i*285
        y = 360 + (i % 2)*80
        d.ellipse((x-48, y-48, x+48, y+48), fill=color+(255,))
        text_center(d, (x-48, y-48, x+48, y+48), num, FONT_B, (255,255,255))
        d.text((x-88, y+72), title, font=FONT_M, fill=INK)
        text_left(d, (x-108, y+115), desc, FONT_XS, MUTED, 220)
        if i < len(steps)-1:
            arrow(d, (x+55, y), (x+225, 360 + ((i+1)%2)*80), MUTED, 4)
    return save(img, "v2_08_teaching_storyboard.png")


def make_figs():
    return [
        ("DMA漏损管控参考架构", fig_dma_reference_architecture(), "总览页：讲清DMA不是单点算法，而是工程、数据、模型、业务闭环。"),
        ("DMA现场拓扑与监测布点", fig_dma_field_topology(), "用于说明入口流量、压力点、边界阀、疑似漏点和工单回填关系。"),
        ("MNF夜间最小流量分析逻辑", fig_mnf_logic(), "用于讲夜间最小流量、合法夜间用水、背景漏损和突发漏损分解。"),
        ("DMA漏损AI模型流水线", fig_ai_pipeline(), "用于讲LSTM/GRU、异常检测、水力融合和工单闭环。"),
        ("漏损检测模型选型决策树", fig_model_selection_tree(), "用于讲模型选型：先数据质量，再算法复杂度。"),
        ("机理模型与AI模型融合路线", fig_hybrid_method(), "用于讲传统水力模型与AI模型的关系。"),
        ("端-边-云-业务协同部署", fig_edge_cloud(), "用于讲工程部署和系统集成。"),
        ("教学演讲PPT推荐叙事结构", fig_teaching_storyboard(), "用于规划最终PPT章节。"),
    ]


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_doc_style(doc):
    sec = doc.sections[0]
    sec.top_margin = Inches(0.75)
    sec.bottom_margin = Inches(0.75)
    sec.left_margin = Inches(0.85)
    sec.right_margin = Inches(0.85)
    styles = doc.styles
    for name in ["Normal", "List Bullet", "List Number"]:
        styles[name].font.name = "宋体"
        styles[name]._element.rPr.rFonts.set(qn("w:eastAsia"), "宋体")
        styles[name].font.size = Pt(10.5)
    for name in ["Heading 1", "Heading 2", "Heading 3"]:
        styles[name].font.name = "黑体"
        styles[name]._element.rPr.rFonts.set(qn("w:eastAsia"), "黑体")
        styles[name].font.color.rgb = BLACK


def run_style(run, color=BLACK, bold=False, size=10.5, font="宋体"):
    run.font.name = font
    run._element.rPr.rFonts.set(qn("w:eastAsia"), font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def title(doc, text, sub):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    run_style(r, BLACK, True, 22, "黑体")
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(sub)
    run_style(r, RED, False, 11, "宋体")
    doc.add_paragraph()


def h(doc, text, level):
    p = doc.add_heading(text, level=level)
    for r in p.runs:
        run_style(r, BLACK, True, 16 if level == 1 else 13 if level == 2 else 11.5, "黑体")


def p(doc, text, red=True, style=None):
    para = doc.add_paragraph(style=style)
    para.paragraph_format.first_line_indent = Pt(21) if not style else None
    para.paragraph_format.line_spacing = 1.25
    r = para.add_run(text)
    run_style(r, RED if red else BLACK, False, 10.5, "宋体")
    return para


def bullets(doc, items, red=True):
    for item in items:
        para = doc.add_paragraph(style="List Bullet")
        para.paragraph_format.line_spacing = 1.2
        r = para.add_run(item)
        run_style(r, RED if red else BLACK, False, 10.5, "宋体")


def table(doc, headers, rows):
    t = doc.add_table(rows=1, cols=len(headers))
    t.alignment = WD_TABLE_ALIGNMENT.CENTER
    t.style = "Table Grid"
    for i, head in enumerate(headers):
        cell = t.rows[0].cells[i]
        set_cell_shading(cell, "D9EAF7")
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        pp = cell.paragraphs[0]
        pp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = pp.add_run(head)
        run_style(r, BLACK, True, 9.5, "黑体")
    for row in rows:
        cells = t.add_row().cells
        for i, val in enumerate(row):
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
            pp = cells[i].paragraphs[0]
            pp.paragraph_format.line_spacing = 1.1
            r = pp.add_run(val)
            run_style(r, RED, False, 9, "宋体")
    doc.add_paragraph()


def pic(doc, path, cap):
    pp = doc.add_paragraph()
    pp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    pp.add_run().add_picture(str(path), width=Inches(6.6))
    pp = doc.add_paragraph()
    pp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = pp.add_run(cap)
    run_style(r, RED, False, 9, "宋体")


def make_main_doc(figs):
    doc = Document()
    set_doc_style(doc)
    title(doc, "AI模型在供水管网DMA系统漏损检测中的应用", "详细扩充版：红色文字为补充写入内容；面向后续对外教学演讲和PPT制作")
    p(doc, "说明：原始文件未检测到蓝色字体层面的可机读标记。本版按“杨佳负责的模型相关内容”处理，重点扩充1.3、2.1、2.2、3、4中与DMA架构、AI模型、数据治理、工程落地相关的内容。", True)
    p(doc, "本版刻意采用“讲稿型”写法，内容比PPT页面更展开，便于后续拆分为演讲稿、PPT正文和备注页。涉及准确率、降漏率、投资收益等项目量化指标，当前没有真实项目数据，文中不虚构数值，后续应替换为项目实测结果。", True)

    h(doc, "目录", 1)
    for item in ["1. 模型应用背景", "2. DMA系统与AI技术基础", "3. 核心应用场景", "4. 工程实施路径", "5. 典型案例包装建议", "6. 教学演讲PPT转化建议", "7. 总结与展望", "8. 参考资料"]:
        p(doc, item, False)

    h(doc, "第一部分 模型应用背景", 1)
    h(doc, "1.1 供水管网漏损管控的国家战略与行业刚需", 2)
    p(doc, "公共供水管网漏损控制已经从传统运维问题，变成水资源节约、城市韧性、基础设施精细化治理和企业经营降本的共同问题。住建部、国家发展改革委发布的建办城〔2022〕2号文明确提出推进分区计量、压力调控、智能化建设，并提出到2025年全国城市公共供水管网漏损率力争控制在9%以内的目标。这个目标意味着漏损控制不能只依靠阶段性检漏行动，而要建立可持续的数据化、模型化和闭环化机制。", True)
    p(doc, "从水司经营角度看，漏损造成的不是单一水量损失，而是一组连锁成本：制水电耗与药耗被无效消耗，管网高压和爆管风险增加，抢修与投诉压力上升，新增供水能力被无收益水量挤占。对于老城区和管龄较长区域，漏损治理还与道路开挖、交通组织、营商环境和居民满意度直接相关。", True)
    p(doc, "DMA系统是漏损治理从粗放走向精细的基础设施。它通过相对封闭的分区边界，把大管网拆成可计量、可比较、可定位的管理单元。AI模型则是在DMA基础上进一步提升识别效率：它能够学习每个DMA的正常用水模式，发现传统阈值难以及时识别的缓慢渗漏、夜间底流抬升和多点压力响应异常。", True)

    h(doc, "1.2 供水管网机理模型的发展与边界", 2)
    p(doc, "传统供水管网模型以水力机理为基础，核心是连续性方程、能量方程、管道阻力、泵阀边界条件和节点需水量分配。早期工程计算常用哈代-克罗斯平差方法，随后大量软件采用牛顿-拉夫逊等数值迭代方法求解稳态和延时段水力状态；瞬变分析则常使用特征线法处理水锤等快速动态过程。", True)
    p(doc, "EPANET是水力建模的重要开源工具。根据美国EPA说明，EPANET可模拟有压供水管网在多时段内的水力和水质行为，跟踪每条管道流量、各节点压力、水池水位、水龄和水质组分等结果。商业软件在此基础上增强了GIS、模型校准、工况管理、成果制图和企业协同能力。", True)
    p(doc, "机理模型在方案校核、压力分区、泵阀调度、水龄分析和供水安全评估中仍不可替代。但在DMA漏损应用中，它的短板也很明显：基础数据要求高，阀门实际状态难核验，用户需水模式随时间变化，粗糙系数和局部阻力难持续校准。若用一个未及时维护的机理模型直接判断漏点，模型可能给出看似精确、实际偏离现场的结果。", True)
    p(doc, "因此，合理定位不是“AI替代水力模型”，而是“机理模型给物理约束，AI模型吸收高频数据和历史经验”。对外演讲时建议强调这一点：AI不是魔法，也不是脱离工程的黑箱；它必须落在DMA边界、仪表数据、工单反馈和水力常识之上。", True)

    h(doc, "1.3 供水管网AI模型应用的发展", 2)
    pic(doc, figs[0][1], "图1 DMA漏损管控参考架构")
    p(doc, "人工智能的发展大致经历了符号智能、统计学习、深度学习和大模型/智能体四个阶段。符号智能强调专家规则和推理；统计学习强调从历史样本中学习分类、回归和聚类规律；深度学习通过神经网络学习复杂非线性特征，LSTM、GRU等循环神经网络尤其适合时间序列；Transformer和大模型则进一步提升了文本、知识和工具协同能力。", True)
    p(doc, "供水管网AI应用并不是从大模型直接开始的。第一阶段是经验阈值，例如夜间最小流量超过某个固定值即判断异常；第二阶段是统计分析，例如同比、环比、移动平均、控制图和季节性分解；第三阶段是机器学习和深度学习，例如LSTM预测正常流量、孤立森林识别离群异常、随机森林评估管段风险；第四阶段是机理-数据融合和智能体协同，把水力仿真、拓扑结构、实时监测、工单知识和自然语言交互统一到一个业务流程中。", True)
    p(doc, "在DMA漏损场景中，AI模型最核心的任务不是炫技，而是把“正常”和“异常”的边界讲清楚。每个DMA都有自己的用户结构、地形高差、压力制度、工业商业用水比例、夜间用水习惯和历史漏损特征。固定阈值难以同时适配所有分区，而AI模型可以为每个DMA建立动态基线：工作日与周末不同，夏季与冬季不同，节假日与普通日不同，突发天气与正常天气也不同。", True)
    p(doc, "LSTM/GRU模型适合处理流量、压力这类具有周期性、趋势性和滞后性的时序数据。典型做法是用过去一段时间的入口流量、关键压力点、天气、日期、节假日和历史用水特征预测未来短时窗口的正常流量或压力，再用实际值与预测值的残差判断异常。如果残差持续扩大，并且出现在夜间低用水窗口，模型即可触发疑似漏损预警。", True)
    p(doc, "孤立森林、DBSCAN和自编码器适合漏损标签不足的项目。许多水司历史工单并不完全结构化，漏点发生时间、修复时间、漏量估算和坐标信息不完整，直接训练监督模型风险较高。无监督模型可以先学习多数正常样本的分布，再把离群样本、异常簇或高重构误差时段推送给运维人员复核。", True)
    p(doc, "随机森林、梯度提升树和逻辑回归等模型更适合管段风险评估。输入可以包括管龄、材质、口径、压力等级、历史爆管、道路等级、周边施工、土壤环境、维修频次和用户投诉等结构化变量。输出不是实时报警，而是中长期风险排序，用于指导检漏计划、换管计划和改造优先级。", True)
    p(doc, "图神经网络和机理融合模型是更进一步的方向。供水管网天然是图结构，节点代表水表、阀门、压力点、用户和水池，边代表管段。GNN能够把上下游拓扑关系纳入异常识别，但其落地前提是拓扑质量、传感器布点和训练样本足够可靠。当前工程项目中，更稳妥的路线通常是先用MNF、LSTM残差和树模型跑通闭环，再逐步引入拓扑模型。", True)
    p(doc, "大模型和智能体更适合作为业务协同层。它可以读取模型输出、GIS信息、工单记录和知识库，生成日报、周报、异常解释、巡检建议和复盘报告；也可以帮助管理人员用自然语言查询“本周哪些DMA风险最高”“某DMA异常是否与历史事件相似”“建议先查哪些管段”。但底层异常识别仍应由专业时序模型、异常检测模型和水力模型支撑。", True)

    h(doc, "第二部分 DMA系统与AI技术基础", 1)
    h(doc, "2.1 DMA架构：从工程分区到数据闭环", 2)
    pic(doc, figs[1][1], "图2 DMA现场拓扑与监测布点")
    p(doc, "DMA即District Metered Area，通常译为独立计量分区或分区计量区域。公开研究中对DMA的基本定义是：通过关闭边界阀或设置边界控制，将供水管网划分为边界明确的水力区域，并对进入和离开该区域的水量进行计量。它的本质不是“画一个区域”，而是建立一个可以进行水量平衡和异常归因的管理单元。", True)
    p(doc, "一个合格的DMA至少包含六个要素。第一是边界清晰，区域内外连通关系必须可解释，边界阀门状态必须可核验；第二是计量闭合，入口和必要出口均应计量，避免漏计和重复计量；第三是压力可观测，关键高点、低点、末梢和代表性节点应布置压力监测；第四是用户结构可识别，居民、商业、工业、学校、医院等用水规律差异较大；第五是工单能回填，维修结果必须回到DMA台账；第六是模型能迭代，新的异常和处置结果应进入训练样本库。", True)
    p(doc, "DMA建设通常不是一次性完成。第一步是基础资料梳理，包括管网GIS、阀门台账、供水边界、用户数量、历史爆管和表计情况；第二步是水力边界划分，兼顾地形、压力制度、道路条件、管网连通性和施工可行性；第三步是边界封闭和计量安装；第四步是夜间流量和压力测试，确认分区是否真实封闭；第五步是建立MNF分析和报警机制；第六步是接入AI模型和工单闭环。", True)
    p(doc, "DMA架构中最容易被忽视的是边界可靠性。若边界阀未完全关闭、存在临时连通、旁通管未入库，入口流量就无法准确代表区域供水量，后续MNF和AI模型都会受到污染。因此在模型建设前，应先做边界校验：夜间压降试验、阀门状态核查、流量守恒检查、异常压力传播分析和现场复核。", True)
    p(doc, "压力点布设也直接影响漏损识别能力。入口流量能告诉我们“这个DMA是否异常”，但很难告诉我们“异常在哪一段”。多点压力能够提供空间响应信息：漏点附近和下游压力可能出现持续性下降或波动模式改变；若结合水力模型和拓扑关系，就可以把排查范围从整个DMA收敛到若干候选管段。", True)
    table(doc, ["DMA架构层级", "建设内容", "对AI模型的作用", "常见风险"], [
        ["工程边界", "边界阀、入口出口、分区规模、压力分区", "决定水量平衡是否成立", "误连通、阀门状态不明、边界频繁变化"],
        ["感知计量", "流量计、压力计、RTU、远传水表", "提供模型输入和实时监测信号", "时间不同步、漂移、断点、仪表量程不匹配"],
        ["数据治理", "统一编码、清洗补缺、时序对齐、标签口径", "决定训练数据是否可用", "工单无坐标、漏点时间不准、设备编码不一致"],
        ["模型算法", "MNF、LSTM、异常检测、树模型、水力融合", "识别异常、排序风险、收敛定位", "过拟合、黑箱输出、缺少业务解释"],
        ["业务闭环", "预警、派单、巡检、维修、复盘、再训练", "把模型结果转化为降漏效果", "只报警不处置、复盘不回填、模型不更新"],
    ])

    h(doc, "2.2 MNF夜间最小流量：DMA漏损分析的基础指标", 2)
    pic(doc, figs[2][1], "图3 MNF夜间最小流量分析逻辑")
    p(doc, "MNF（Minimum Night Flow）是DMA漏损分析中最经典、最实用的指标之一。其基本逻辑是：在凌晨用户活动最低的时段，DMA入口流量达到日内较低水平，此时合法用水较少，漏损在总流量中的占比更高，因此更容易观察漏损变化。公开文献通常指出MNF需要结合合法夜间用水、背景漏损和压力条件进行解释，不能简单把MNF等同于漏损量。", True)
    p(doc, "MNF分析一般包括四步。第一步，确定夜间分析窗口，常见为凌晨2点到4点，但应结合本地居民作息、工业用户和二供补水规律调整。第二步，估算合法夜间用水，包括居民夜间用水、商业夜间营业、连续生产企业、医院学校、消防和市政用水等。第三步，结合压力水平和历史基线判断背景漏损。第四步，观察MNF趋势是否出现持续抬升、阶跃变化或维修后回落。", True)
    p(doc, "MNF的局限性也必须讲清楚。第一，它对用户结构敏感，夜间商业和工业用水会抬升底流；第二，它对压力敏感，压力升高会增加背景漏损；第三，它对边界可靠性敏感，边界未封闭会导致水量误差；第四，它对仪表精度敏感，低流量段仪表误差可能被放大。AI模型的价值不是取消MNF，而是把MNF与时间序列、压力响应、节假日特征和历史工单结合起来。", True)
    p(doc, "在教学演讲中，MNF可以作为听众理解AI漏损模型的入口。先讲传统MNF如何发现异常，再讲AI如何扩展MNF：从固定夜间窗口扩展到全天动态基线，从单一入口流量扩展到多点压力联动，从人工阈值扩展到模型残差和异常概率，从单次报警扩展到工单复盘和持续学习。", True)

    h(doc, "2.3 AI模型流水线：从数据到派单", 2)
    pic(doc, figs[3][1], "图4 DMA漏损AI模型流水线")
    p(doc, "一个可落地的DMA漏损AI系统，应当被设计成流水线，而不是单个模型脚本。第一段是数据接入，包括SCADA流量压力、远传水表、GIS资产、阀门状态、泵站运行、工单系统和外部天气日历。第二段是数据治理，包括时间对齐、缺失值补齐、异常毛刺处理、设备漂移识别、采样粒度统一和特征生成。第三段是基线建模，包括MNF基线、日周期基线、周周期基线和LSTM/GRU预测基线。", True)
    p(doc, "第四段是异常识别。模型不应只看某个时刻是否超过阈值，而要判断异常是否持续、是否与压力响应一致、是否与历史同类时段显著不同、是否可能由仪表故障或合法用水解释。第五段是定位收敛，结合管网拓扑、压力点响应、水力仿真样本、历史维修记录和用户投诉，把异常从DMA级收敛到片区或候选管段。第六段是业务闭环，将异常推送为工单，并把现场核查结果回填模型。", True)
    p(doc, "LSTM模型在流水线中通常承担“正常曲线预测器”的角色。它不是直接替现场人员判断漏点，而是根据历史窗口输出未来流量或压力的合理范围。一旦实际值长期偏离预测区间，就形成残差信号。残差再与MNF、压力联动、工单历史和水力拓扑结合，形成更可靠的漏损风险判断。", True)
    p(doc, "模型上线前应设置人工复核机制。早期模型一定会出现误报和漏报，尤其是在节假日、设备维修、阀门调整、消防用水、二供补水等特殊场景中。工程上不应追求“一上线就全自动”，而应设置灰度运行期：模型报警先进入观察清单，由运维人员标注原因，积累足够样本后再进入自动派单。", True)

    h(doc, "2.4 模型选型原则", 2)
    pic(doc, figs[4][1], "图5 漏损检测模型选型决策树")
    p(doc, "模型选型必须从数据条件出发，而不是从算法名词出发。如果DMA边界不可靠、入口流量不完整、压力点缺失、工单标签混乱，则最先进的深度学习模型也无法给出可靠结果。此时优先工作应是工程治理和数据治理，而不是直接训练复杂模型。", True)
    p(doc, "当只有稳定入口流量和少量压力数据时，可采用MNF、移动基线、LSTM/GRU预测残差和孤立森林等轻量方法。它们对标签依赖较低，能够较快形成预警能力。当具备较完整工单标签和管网资产特征时，可引入随机森林、GBDT等监督模型进行管段风险排序。当具备高质量拓扑、多个压力点和水力模型时，再考虑机理-数据融合定位模型。", True)
    table(doc, ["应用目标", "优先模型", "适用条件", "输出结果"], [
        ["短时预警", "LSTM/GRU、Prophet、移动基线", "连续流量压力数据较完整", "预测值、残差、预警等级"],
        ["少标签异常识别", "孤立森林、DBSCAN、自编码器", "缺少准确漏损标签，但正常样本较多", "异常分数、异常时段、异常簇"],
        ["管段风险排序", "随机森林、GBDT、逻辑回归", "有资产属性和历史维修记录", "风险分、优先级、主要影响因素"],
        ["漏点定位收敛", "水力仿真+分类模型/GNN", "拓扑可靠、多点压力、仿真样本可生成", "候选管段TopN、定位置信度"],
        ["业务协同", "知识库+大模型智能体", "模型结果、GIS、工单接口可用", "日报周报、问答、派单建议、复盘材料"],
    ])

    h(doc, "2.5 机理模型与AI模型融合", 2)
    pic(doc, figs[5][1], "图6 机理模型与AI模型融合路线")
    p(doc, "机理模型与AI模型融合有三种常见方式。第一种是特征融合，把水力模型输出的节点压力、管段流量、压力敏感性、供水路径等结果作为AI模型输入特征。第二种是样本融合，用水力模型模拟不同漏点、漏量和阀门状态下的压力响应，构建训练样本库。第三种是约束融合，在AI模型输出候选管段后，用水力常识和现场可达性进行过滤和排序。", True)
    p(doc, "融合路线的优点是能够缓解真实漏损标签不足的问题。很多水司真实漏点样本数量有限，而且标签质量参差不齐。通过仿真生成多种漏损情景，再用少量真实事件校准，可以让模型先具备基本定位能力，再通过实际工单逐步修正。", True)
    p(doc, "融合模型的风险在于“仿真世界”和“真实世界”之间存在差距。如果水力模型拓扑、粗糙系数、阀门状态、需水模式不准确，仿真样本也会把错误规律带给AI模型。因此，融合模型上线前必须做现场校验：选取已知事件或压力试验数据，对比仿真响应与实际响应的一致性。", True)

    h(doc, "第三部分 核心应用场景", 1)
    h(doc, "3.1 事前防控：DMA异常预警", 2)
    p(doc, "事前防控的目标是提前发现漏损风险，而不是等到漏损量显著扩大后再组织排查。AI模型可对每个DMA建立动态基线，持续计算实际流量与预测流量的偏差、夜间底流变化、压力异常响应和历史同类日期差异，并输出风险等级。", True)
    p(doc, "预警分级建议采用业务可理解的表达：蓝色为关注，表示数据轻微偏离但可能由正常用水解释；黄色为核查，表示异常持续且超过历史波动范围；橙色为派单建议，表示流量和压力同时出现异常；红色为重点处置，表示异常持续多日或与历史漏损事件高度相似。", True)
    p(doc, "预警结果不应只显示一个分数，还应附带解释项，例如：异常发生时段、相对基线抬升量、是否位于MNF窗口、压力点是否同步下降、近期是否有阀门操作、是否存在同类历史事件。这样运维人员才能判断是漏损、用水扰动还是设备故障。", True)

    h(doc, "3.2 事中处置：漏点定位与巡检派单", 2)
    p(doc, "事中处置的核心是缩小排查范围。传统做法可能在一个DMA内进行大范围人工听漏或开挖验证，成本高且效率不稳定。AI模型可以先根据入口流量和压力响应锁定异常片区，再结合管网拓扑、历史爆管、材质管龄和道路条件输出候选管段清单。", True)
    p(doc, "建议模型输出采用TopN候选清单，而不是单一坐标。每个候选项包括管段编号、位置、置信度、主要证据、影响用户数、历史维修次数、建议复核方式和现场注意事项。这种输出方式更符合工程实际，也能降低模型误判导致的决策风险。", True)
    p(doc, "派单系统应允许现场人员回填“确认漏损、非漏损、设备故障、边界异常、合法用水、其他原因”等结果。只有当现场结果能够结构化回填，模型才能从误报中学习，逐步提升精度。", True)

    h(doc, "3.3 事后复盘：模型持续学习", 2)
    p(doc, "事后复盘是AI系统能否长期有效的关键。每一次报警、核查、维修和误报都应变成训练样本。复盘字段至少包括DMA编号、报警时间、异常特征、派单时间、现场发现、漏点位置、漏损类型、维修完成时间、估算漏量、误报原因和模型版本。", True)
    p(doc, "复盘不是为了写报告，而是为了更新模型。若某类误报频繁出现，例如夜间二供补水、商户集中用水或仪表漂移，就应把这些场景转化为特征或规则，减少下次误报。若某类漏损未被模型发现，则应检查数据采样、压力点布设、阈值设置和模型训练样本是否存在缺陷。", True)

    h(doc, "3.4 调度优化与压力管理", 2)
    p(doc, "压力管理是漏损控制的重要手段。漏损流量与压力通常存在正相关关系，高压会增加背景漏损并提高爆管概率。AI模型可以预测不同日期和时段的用水需求，再结合压力监测和水力模型，给出分时压力控制、泵站启停、阀门调节和压力区优化建议。", True)
    p(doc, "需要注意的是，压力优化不是简单降压。供水系统必须满足最不利点服务压力、消防保障、二供补水和高峰需求。模型应在“服务压力合格”的约束下寻找降漏和节能空间，避免为了降低漏损而影响用户供水体验。", True)

    h(doc, "3.5 DMA规划与监测点优化", 2)
    p(doc, "AI还可用于DMA规划和监测点优化。规划阶段可以综合管网拓扑、用户数量、地形高差、道路条件、施工成本、压力制度和管理半径，生成候选分区方案。评价指标包括分区规模是否合理、入口数量是否可控、边界阀数量是否过多、压力是否稳定、未来扩展是否方便。", True)
    p(doc, "监测点优化则关注“少量压力点如何提供最大定位信息”。模型可通过水力仿真模拟不同漏点对压力点的响应，比较不同布点方案对漏点定位的区分能力。对预算有限的项目，这比盲目增加传感器更实用。", True)

    h(doc, "第四部分 工程实施路径", 1)
    h(doc, "4.1 总体路线：先工程可靠，再模型智能", 2)
    pic(doc, figs[6][1], "图7 端-边-云-业务协同部署")
    p(doc, "实施路径建议遵循“先工程可靠，再数据可用，再模型有效，再业务闭环”的顺序。第一阶段完成DMA边界核查、仪表校验和基础台账整理；第二阶段完成数据接入、编码统一、时间同步和质量评估；第三阶段建设MNF和LSTM等基础模型；第四阶段接入工单和GIS，形成预警、派单、复盘闭环；第五阶段进行跨DMA推广和模型运营。", True)
    p(doc, "端侧负责数据采集和设备状态，边缘侧负责本地缓存、简单质检和轻量异常识别，云端负责全量数据存储、模型训练、跨DMA分析和知识库建设，业务侧负责GIS定位、巡检派单、维修闭环和指标考核。这个架构能同时兼顾实时性、算力效率和管理闭环。", True)
    p(doc, "试点选择很重要。建议优先选择边界清晰、仪表稳定、历史漏损事件较多、用户结构相对典型、运维队伍配合度高的DMA。不要一开始追求全市覆盖，否则主要精力会消耗在数据接入、设备缺陷和边界问题上，模型效果反而难以展示。", True)

    h(doc, "4.2 数据治理清单", 2)
    table(doc, ["治理对象", "具体要求", "模型影响"], [
        ["时间", "统一时区、统一采样粒度、处理延迟和乱序数据", "时间错位会直接破坏LSTM训练和残差判断"],
        ["设备", "统一设备编码、记录量程、精度、安装位置、维护状态", "设备漂移会被模型误判为漏损"],
        ["空间", "GIS坐标、管段编号、阀门状态、DMA归属一致", "空间错误会导致定位结果不可用"],
        ["业务", "工单原因、漏点位置、维修时间、处置结果结构化", "标签质量决定监督模型上限"],
        ["异常", "区分真实漏损、合法用水、设备故障、边界异常", "异常类别不清会造成模型学习混乱"],
    ])
    p(doc, "数据治理是项目成败的隐性关键。很多AI项目失败并不是算法不先进，而是数据无法支撑模型：同一个流量计在SCADA、GIS和运维系统中有不同名称；压力数据和流量数据时间戳不一致；工单只记录“某路漏水”，没有管段编号和修复时间；阀门状态长期未更新。这些问题必须在建模前或建模早期解决。", True)

    h(doc, "4.3 模型训练与验证", 2)
    p(doc, "训练数据应按时间切分，而不是随机切分。供水时序数据具有明显周期性，如果随机切分训练集和测试集，模型可能提前看到相似时段，从而高估效果。更稳妥的做法是用前一段时间训练，后一段时间验证，并单独保留若干真实漏损事件做事件级评估。", True)
    p(doc, "评价指标应同时包含算法指标和业务指标。算法指标包括MAE、RMSE、召回率、精确率、F1、AUC等；业务指标包括提前预警时间、误报工单比例、平均排查范围、定位TopN命中率、平均处置时长和复盘回填率。对外演讲时，建议少讲抽象准确率，多讲“模型如何减少无效排查、如何缩短发现时间、如何让经验沉淀”。", True)
    p(doc, "模型上线后还要关注漂移。季节变化、用户结构变化、管网改造、传感器更换、压力制度调整都会改变数据分布。建议建立月度误报复盘、季度模型评估、半年再训练和重大管网变更后模型重校准机制。", True)

    h(doc, "第五部分 典型案例包装建议", 1)
    p(doc, "当前没有真实案例素材，因此不应虚构具体城市、具体降漏率和准确率。可先包装为教学型案例：某老城区DMA在连续三天凌晨2点到4点出现MNF抬升，入口流量较历史同类日期显著增加，P2和P3压力点出现轻微下降。LSTM预测残差持续为正，孤立森林给出异常分数升高，水力模型匹配结果提示J3-J8和H区附近管段为候选区域。", True)
    p(doc, "教学案例的演讲顺序建议为：第一页展示原始曲线，说明传统阈值看到什么；第二页展示LSTM预测值与实际值的残差，说明AI如何发现持续偏离；第三页展示压力点响应和拓扑图，说明如何从DMA级异常收敛到候选管段；第四页展示派单和现场复核，说明如何形成业务闭环；第五页展示维修后MNF回落和标签回填，说明模型如何持续学习。", True)

    h(doc, "第六部分 教学演讲PPT转化建议", 1)
    pic(doc, figs[7][1], "图8 教学演讲PPT推荐叙事结构")
    p(doc, "这份Word不建议逐字搬到PPT。PPT应采用“一页一个观点”的方式拆分：背景页讲为什么做，架构页讲DMA是什么，方法页讲AI怎么判，应用页讲如何派单，实施页讲如何落地，案例页讲如何复盘。Word中的长段落可作为演讲备注，PPT正文只保留标题、关键句和图。", True)
    p(doc, "建议最终PPT结构为18到25页：封面1页，政策与痛点2页，DMA基础架构3页，MNF与压力分析2页，AI模型方法4页，机理融合2页，工程实施3页，案例演示3到5页，总结展望1到2页。素材图应占据页面主体，正文只放少量提示性文字。", True)

    h(doc, "第七部分 总结与展望", 1)
    p(doc, "AI模型在供水管网DMA系统中的价值，可以概括为四句话：DMA让漏损从全网问题变成分区问题；时序模型让异常从经验判断变成动态基线判断；机理-数据融合让定位从人工盲查变成候选管段排序；业务闭环让一次次处置变成可沉淀的组织经验。", True)
    p(doc, "未来的发展方向不是单一大模型替代所有算法，而是专业模型、机理模型、知识库和智能体协同。底层用LSTM、异常检测、树模型和水力仿真处理专业判断；中间用数据平台和工单系统完成闭环；上层用智能体完成解释、问答、报告和协同。只有这样，AI漏损管控才能从演示系统变成长期有效的运营系统。", True)

    h(doc, "第八部分 参考资料", 1)
    refs = [
        "[1] 住房和城乡建设部办公厅、国家发展改革委办公厅. 关于加强公共供水管网漏损控制的通知（建办城〔2022〕2号）. https://www.mohurd.gov.cn/gongkai/zc/wjk/art/2022/art_17339_764316.html",
        "[2] US EPA. EPANET. https://www.epa.gov/water-research/epanet",
        "[3] Di Nardo A. et al. Water Network Partitioning into District Metered Areas: A State-Of-The-Art Review. Water, 2020. https://www.mdpi.com/2073-4441/12/4/1002",
        "[4] Farley M., Trow S. Losses in Water Distribution Networks. IWA Publishing, 2003.",
        "[5] Adedeji K. B. et al. Leak detection in water distribution networks: an introductory overview. Smart Water, 2019. https://link.springer.com/article/10.1186/s40713-019-0017-x",
        "[6] Hochreiter S., Schmidhuber J. Long Short-Term Memory. Neural Computation, 1997. https://direct.mit.edu/neco/article/9/8/1735/6109/Long-Short-Term-Memory",
        "[7] Vaswani A. et al. Attention Is All You Need. arXiv:1706.03762, 2017. https://arxiv.org/abs/1706.03762",
        "[8] Zhang Z., Fink O. Algorithm-Informed Graph Neural Networks for Leakage Detection and Localization in Water Distribution Networks. arXiv:2408.02797, 2024. https://arxiv.org/abs/2408.02797",
    ]
    for ref in refs:
        p(doc, ref, True)
    path = DOC_DIR / "AI模型在供水管网DMA系统漏损检测中的应用_详细扩充红色标注版_v2.docx"
    doc.save(path)
    return path


def make_asset_doc(figs):
    doc = Document()
    set_doc_style(doc)
    title(doc, "DMA漏损AI模型PPT插图素材说明文档", "v2：去网页化，按正式PPT素材风格重绘；PNG另存于 assets_v2")
    p(doc, "本素材文档用于后续制作正式PPT。图片均为16:9透明底PNG，适合放在白底、浅灰底或深色底上；另提供可编辑PPT素材包，便于二次调整颜色、字体和布局。", True)
    for i, (name, path, usage) in enumerate(figs, 1):
        h(doc, f"{i}. {name}", 2)
        pic(doc, path, name)
        p(doc, f"素材文件：{path}", False)
        p(doc, f"建议用途：{usage}", False)
    out = DOC_DIR / "DMA漏损AI模型_PPT插图素材说明文档_v2.docx"
    doc.save(out)
    return out


def add_ppt_textbox(slide, x, y, w, h, text, size=20, color=(26,37,55), bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = box.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = PPt(size)
    run.font.bold = bold
    run.font.color.rgb = PRGBColor(*color)
    return box


def add_ppt_pill(slide, x, y, w, h, text, fill, size=16):
    shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PRGBColor(255,255,255)
    shp.line.color.rgb = PRGBColor(*fill)
    shp.line.width = PPt(2)
    tf = shp.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = PPt(size)
    run.font.color.rgb = PRGBColor(26,37,55)
    return shp


def make_ppt_assets(figs):
    prs = Presentation()
    prs.slide_width = PInches(13.333333)
    prs.slide_height = PInches(7.5)
    blank = prs.slide_layouts[6]
    for name, path, usage in figs:
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PRGBColor(255,255,255)
        add_ppt_textbox(slide, 0.55, 0.30, 7.5, 0.45, name, 23, bold=True)
        slide.shapes.add_picture(str(path), PInches(0.45), PInches(0.95), width=PInches(12.45))
        add_ppt_textbox(slide, 0.7, 7.05, 11.8, 0.3, usage, 10, color=(88,101,119))

    # one editable example slide built with native shapes
    slide = prs.slides.add_slide(blank)
    add_ppt_textbox(slide, 0.55, 0.35, 6.0, 0.5, "可编辑示例：DMA AI模型链路", 24, bold=True)
    labels = [("数据接入", BLUE), ("治理对齐", CYAN), ("LSTM基线", GREEN), ("异常识别", ORANGE), ("派单复盘", RED_P)]
    for i, (lab, col) in enumerate(labels):
        add_ppt_pill(slide, 0.75 + i*2.45, 2.9, 1.55, 0.55, lab, col, 15)
        if i < len(labels)-1:
            add_ppt_textbox(slide, 2.34 + i*2.45, 3.02, 0.5, 0.3, "→", 20, color=(88,101,119), bold=True, align=PP_ALIGN.CENTER)
    add_ppt_textbox(slide, 0.75, 4.1, 11.2, 0.7, "这一页为PPT原生形状，后续可直接改字、改色、拆分重排。前面8页为PNG素材预览页。", 16, color=(88,101,119))

    out = PPT_DIR / "DMA漏损AI模型_PPT素材包_v2.pptx"
    prs.save(out)
    return out


def main():
    figs = make_figs()
    main_doc = make_main_doc(figs)
    asset_doc = make_asset_doc(figs)
    ppt = make_ppt_assets(figs)
    print(main_doc)
    print(asset_doc)
    print(ppt)
    for _, path, _ in figs:
        print(path)


if __name__ == "__main__":
    main()
