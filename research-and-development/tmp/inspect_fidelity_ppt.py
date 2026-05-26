from pptx import Presentation
from pathlib import Path
p=Path(r'E:\PY\research\output\ppt\Jinshiyuan_AI_Production_Blueprint_可编辑保真版.pptx')
prs=Presentation(str(p))
print('file', p)
print('slides', len(prs.slides), 'size', p.stat().st_size)
for i,s in enumerate(prs.slides,1):
    pics=texts=shapes=0
    full=[]
    for sh in s.shapes:
        shapes += 1
        if sh.shape_type == 13:
            pics += 1
            full.append((round(sh.left/914400,2), round(sh.top/914400,2), round(sh.width/914400,2), round(sh.height/914400,2)))
        if getattr(sh,'has_text_frame',False) and sh.text.strip():
            texts += 1
    print(i, 'shapes', shapes, 'pictures', pics, 'pic_boxes', full[:2], 'text_shapes', texts)
